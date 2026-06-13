"""
python-pptx data-access layer for PowerPoint MCP Phase 1 — thin facade.

This module owns:
  - Error taxonomy (PPTMCPError hierarchy)
  - Governance helpers (_check_path, _check_image_path, _check_write, _check_confirm)
  - Audit logging (_audit_log)
  - Cache constants (_PRS_CACHE_MAX, _prs_cache) — re-exports from _presentation_cache
  - Shape / text helpers (_validate_slide_index, _find_shape_by_id, etc.)
  - Re-exports from _presentation_cache, _presentation_read, _presentation_write
  - Backward-compatible __getattr__ for names moved to shapes_pptx / content_pptx / _pptx_caps

No FastMCP / MCP imports — pure python-pptx + stdlib only.
"""
from __future__ import annotations

import json
import logging
import os
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation  # noqa: F401  (re-exported so tests can patch presentation_pptx.Presentation)
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Error taxonomy (sourced from _presentation_governance to allow zero-dependency
# import by sub-modules; re-exported here for backward compatibility)
# ---------------------------------------------------------------------------

from pptmcp._presentation_governance import (  # noqa: E402
    OfficeCOMError,  # noqa: F401
    NotAllowedError,  # noqa: F401
    PPTMCPError,  # noqa: F401
    ValidationError,  # noqa: F401
)

# ---------------------------------------------------------------------------
# Governance helpers — defined here so tests can patch presentation_pptx.*
# ---------------------------------------------------------------------------

_MAX_PPTX_BYTES = int(os.getenv("PPT_MAX_FILE_MB", "256")) * 1024 * 1024

_ALLOWED_IMAGE_EXTENSIONS = frozenset({
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".svg"
})


def _check_path(path: str) -> Path:
    """Validate path against PPT_ALLOWLIST_ROOTS and extension rules."""
    roots_env = os.getenv("PPT_ALLOWLIST_ROOTS", "").strip()
    if not roots_env:
        raise ValidationError("No allowlist configured")

    if "\x00" in str(path):
        raise ValidationError("Invalid path: null byte")

    resolved = Path(path).resolve()

    try:
        if resolved.stat().st_size > _MAX_PPTX_BYTES:
            raise ValidationError(
                f"File exceeds size limit ({_MAX_PPTX_BYTES // (1024 * 1024)} MB)"
            )
    except FileNotFoundError:
        pass

    ext = resolved.suffix.lower()
    if ext != ".pptx":
        raise ValidationError(f"Unsupported extension: {ext!r}")

    roots = [r.strip() for r in roots_env.split(",") if r.strip()]
    for root in roots:
        try:
            if resolved.is_relative_to(Path(root).resolve()):
                return resolved
        except (ValueError, TypeError):
            continue

    raise ValidationError("Path not in allowlist")


def _check_image_path(path: str) -> Path:
    """Validate an image path against PPT_ALLOWLIST_ROOTS; accepts common image extensions."""
    roots_env = os.getenv("PPT_ALLOWLIST_ROOTS", "").strip()
    if not roots_env:
        raise ValidationError("No allowlist configured")

    if "\x00" in str(path):
        raise ValidationError("Invalid path: null byte detected")

    resolved = Path(path).resolve()

    try:
        if resolved.stat().st_size > _MAX_PPTX_BYTES:
            raise ValidationError(
                f"File exceeds size limit ({_MAX_PPTX_BYTES // (1024 * 1024)} MB)"
            )
    except FileNotFoundError:
        pass

    ext = resolved.suffix.lower()
    if ext not in _ALLOWED_IMAGE_EXTENSIONS:
        raise ValidationError(f"Unsupported image extension: {ext!r}")

    roots = [r.strip() for r in roots_env.split(",") if r.strip()]
    for root in roots:
        try:
            if resolved.is_relative_to(Path(root).resolve()):
                return resolved
        except (ValueError, TypeError):
            continue

    raise ValidationError("Path not in allowlist")


def _check_write() -> None:
    """Raise NotAllowedError unless PPT_ENABLE_WRITE=true."""
    if os.getenv("PPT_ENABLE_WRITE", "").strip().lower() != "true":
        raise NotAllowedError("Write operations require PPT_ENABLE_WRITE=true")


def _check_confirm(confirm: bool) -> None:
    """Raise ValidationError if confirm is not True."""
    if not confirm:
        raise ValidationError("confirm=True required for mutating operations")


def _audit_log(
    op: str,
    path: Path | str,
    slide_idx: int | None = None,
    extra: dict | None = None,
) -> None:
    """Write a JSON audit record to stderr."""
    record = {
        "ts": datetime.now(timezone.utc).isoformat(),
        "op": op,
        "path": str(path),
        "slide_index": slide_idx,
        **(extra or {}),
    }
    logger.info(json.dumps(record))


# ---------------------------------------------------------------------------
# Cache helpers — re-exported from _presentation_cache
# Cache constants also exposed here so tests can patch presentation_pptx.*
# ---------------------------------------------------------------------------

from pptmcp._presentation_cache import (  # noqa: E402
    _atomic_save,  # noqa: F401
    _evict_prs,  # noqa: F401
    _load_prs,  # noqa: F401
    _prs_cache,  # noqa: F401
    _PRS_CACHE_MAX,  # noqa: F401
)

# ---------------------------------------------------------------------------
# Read functions — re-exported from _presentation_read
# ---------------------------------------------------------------------------

from pptmcp._presentation_read import (  # noqa: E402
    _slide_title,  # noqa: F401
    get_presentation_metadata,  # noqa: F401
    get_shape,  # noqa: F401
    list_layouts,  # noqa: F401
    list_shapes,  # noqa: F401
    list_slides,  # noqa: F401
    read_presentation,  # noqa: F401
    read_slide,  # noqa: F401
    read_speaker_notes,  # noqa: F401
)

# ---------------------------------------------------------------------------
# Write functions — re-exported from _presentation_write
# ---------------------------------------------------------------------------

from pptmcp._presentation_write import (  # noqa: E402
    _copy_slide_background,  # noqa: F401
    _find_blank_layout_idx,  # noqa: F401
    add_slide,  # noqa: F401
    apply_slide_layout,  # noqa: F401
    clear_slide_content,  # noqa: F401
    copy_slide,  # noqa: F401
    create_presentation,  # noqa: F401
    delete_slide,  # noqa: F401
    reorder_slides,  # noqa: F401
    save,  # noqa: F401
)

# ---------------------------------------------------------------------------
# Shape / text helpers (used by shapes_pptx and other modules)
# ---------------------------------------------------------------------------


def _validate_slide_index(prs: Any, slide_index: int, context: str = "slide") -> Any:
    """Validate slide_index within bounds; return the slide or raise ValidationError."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValidationError(
            f"{context}_index {slide_index} out of range (0-{len(prs.slides) - 1})"
        )
    return prs.slides[slide_index]


def _find_shape_by_id(slide: Any, shape_id: int) -> Any:
    """Find shape by shape_id on a slide; raise ValidationError if not found."""
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValidationError(f"shape_id {shape_id} not found on slide")


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convert hex color (e.g., '#FF0000') to (R, G, B) tuple."""
    if not re.match(r"^#[0-9A-Fa-f]{6}$", hex_color):
        raise ValidationError(f"Invalid hex color: {hex_color!r}")
    return tuple(int(hex_color[i : i + 2], 16) for i in (1, 3, 5))  # type: ignore[return-value]


def _resolve_shape_type(shape_type_str: str) -> Any:
    """Convert string (case-insensitive) to MSO_SHAPE enum member."""
    upper = shape_type_str.upper()
    aliases = {
        "CIRCLE": "OVAL",
        "RECT": "RECTANGLE",
        "SQUARE": "RECTANGLE",
    }
    upper = aliases.get(upper, upper)
    try:
        return getattr(MSO_SHAPE, upper)
    except AttributeError:
        valid = sorted([name for name in dir(MSO_SHAPE) if not name.startswith("_")])
        raise ValidationError(
            f"Unknown shape type: {shape_type_str!r}. "
            f"Valid shapes: {', '.join(valid[:15])}... "
            f"({len(valid)} available total). "
            f"See PowerPoint shape gallery for complete list."
        )


def _emu_to_inches(emu: int | None) -> float | None:
    """Convert EMU to inches (4 dp); returns None if emu is None."""
    return None if emu is None else round(float(Emu(emu).inches), 4)


_ALIGNMENT_MAP: dict[str, Any] = {
    "CENTER": PP_ALIGN.CENTER,
    "LEFT": PP_ALIGN.LEFT,
    "RIGHT": PP_ALIGN.RIGHT,
    "JUSTIFY": PP_ALIGN.JUSTIFY,
    "DISTRIBUTE": PP_ALIGN.DISTRIBUTE,
}

# ---------------------------------------------------------------------------
# Backward-compatible re-exports (implementations moved to shapes_pptx /
# _pptx_caps / content_pptx).
# Must remain at end of file; lazy import avoids circular-import at load time.
# ---------------------------------------------------------------------------


def __getattr__(name: str) -> Any:  # noqa: ANN201
    """PEP 562 module __getattr__: re-export functions moved to shapes_pptx / content_pptx."""
    _SHAPES = frozenset({
        "set_speaker_notes", "edit_text_placeholder", "replace_slide_text",
        "add_textbox", "add_shape", "add_table_to_slide",
        "set_table_style",
        "set_text_format", "set_paragraph_format",
        "add_hyperlink", "manage_hyperlinks", "insert_image",
    })
    _CONTENT = frozenset({
        "extract_tables", "extract_images",
        "export_slide_as_text", "extract_presentation_text",
        "produce_evidence_bundle",
    })
    _CAPS = frozenset({"capabilities"})
    if name in _SHAPES:
        from pptmcp import shapes_pptx  # noqa: PLC0415
        return getattr(shapes_pptx, name)
    if name in _CONTENT:
        from pptmcp import content_pptx  # noqa: PLC0415
        return getattr(content_pptx, name)
    if name in _CAPS:
        from pptmcp import _pptx_caps  # noqa: PLC0415
        return getattr(_pptx_caps, name)
    raise AttributeError(f"module {__name__!r} has no attribute {name!r}")
