"""_shapes_geometry — geometry, fill, table, and image functions for pptmcp.

detect_overlapping_shapes lives in _shapes_overlap; re-exported here for backwards compat.
"""
from __future__ import annotations

import logging
from typing import Any

from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from pptmcp._shapes_helpers import (
    _apply_cell_fill, _apply_cell_text_style, _calculate_overflow_risk,
    _first_overlap_for_shape, _normalize_color_hex,
    _remove_empty_overlapping_content_placeholders,
)
from pptmcp._shapes_overlap import detect_overlapping_shapes  # noqa: F401
from pptmcp.presentation_pptx import (
    ValidationError, _audit_log, _check_confirm, _check_image_path, _check_path,
    _check_write, _find_shape_by_id, _load_prs, _resolve_shape_type, _validate_slide_index,
)

_log = logging.getLogger(__name__)


def _save_and_evict(prs: Any, resolved: Any) -> None:
    from pptmcp._presentation_cache import _atomic_save, _evict_prs
    _atomic_save(prs, resolved)
    _evict_prs(resolved)


def _style_table_shape(
    table_shape: Any, *, font_size_pt: float | None = None, bold: bool | None = None,
    color_hex: str | None = None, header_font_size_pt: float | None = None,
    header_bold: bool | None = None, header_color_hex: str | None = None,
    fill_color_hex: str | None = None, header_fill_color_hex: str | None = None,
    row_height_pt: float | None = None, header_row_height_pt: float | None = None,
    row_heights_pt: list[float] | None = None,
) -> dict:
    if not table_shape.has_table:
        raise ValidationError(f"shape_id {table_shape.shape_id} is not a table")
    if row_height_pt is not None and row_height_pt <= 0:
        raise ValidationError("row_height_pt must be greater than 0")
    if header_row_height_pt is not None and header_row_height_pt <= 0:
        raise ValidationError("header_row_height_pt must be greater than 0")
    if row_heights_pt is not None:
        for i, h in enumerate(row_heights_pt):
            if h <= 0:
                raise ValidationError(f"row_heights_pt[{i}] must be > 0, got {h}")
    table = table_shape.table
    runs_updated = cells_updated = rows_resized = 0
    for row_idx, row in enumerate(table.rows):
        is_header = row_idx == 0
        effective_row_height = (
            header_row_height_pt if is_header and header_row_height_pt is not None else row_height_pt
        )
        if effective_row_height is not None:
            row.height = Pt(effective_row_height)
            rows_resized += 1
        if row_heights_pt is not None and row_idx < len(row_heights_pt):
            row.height = Pt(row_heights_pt[row_idx])
            rows_resized += 1
        for cell in row.cells:
            effective_font_size = (
                header_font_size_pt if is_header and header_font_size_pt is not None else font_size_pt
            )
            effective_bold = header_bold if is_header and header_bold is not None else bold
            effective_color = header_color_hex if is_header and header_color_hex is not None else color_hex
            effective_fill = (
                header_fill_color_hex if is_header and header_fill_color_hex is not None else fill_color_hex
            )
            runs_updated += _apply_cell_text_style(
                cell, font_size_pt=effective_font_size, bold=effective_bold, color_hex=effective_color,
            )
            _apply_cell_fill(cell, effective_fill)
            if any(v is not None for v in (effective_font_size, effective_bold, effective_color, effective_fill)):
                cells_updated += 1
    return {"shape_id": table_shape.shape_id, "cells_updated": cells_updated,
            "runs_updated": runs_updated, "rows_resized": rows_resized}


def add_textbox(
    path: str, slide_index: int, left: float, top: float,
    width: float, height: float, text: str = "", confirm: bool = False,
) -> dict:
    """Add a text box to a slide (inches). Requires PPT_ENABLE_WRITE=true and confirm=True."""
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    if width <= 0:
        raise ValidationError("width and height must be positive")
    if height <= 0:
        raise ValidationError("width and height must be positive")
    prs = _load_prs(resolved)
    slide = _validate_slide_index(prs, slide_index)
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    shape.text_frame.word_wrap = True
    if text:
        shape.text_frame.text = text
    _save_and_evict(prs, resolved)
    _audit_log("add_textbox", resolved, slide_index)
    _risk = _calculate_overflow_risk(text, width, height)
    _overlap = _first_overlap_for_shape(slide, shape.shape_id)
    return {
        "path": str(resolved), "slide_index": slide_index, "shape_id": shape.shape_id,
        "width_inches": width, "height_inches": height,
        "overflow_risk": _risk["overflow_risk"], "overflow_detail": _risk["overflow_detail"],
        "overlap_warning": _overlap,
    }


def add_shape(
    path: str, slide_index: int, shape_type: str,
    left: float, top: float, width: float, height: float,
    text: str = "", confirm: bool = False,
) -> dict:
    """Add an autoshape (RECTANGLE, OVAL, etc.) to a slide. Requires PPT_ENABLE_WRITE=true and confirm=True."""
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    if width <= 0:
        raise ValidationError("width and height must be positive")
    if height <= 0:
        raise ValidationError("width and height must be positive")
    try:
        mso_shape = _resolve_shape_type(shape_type)
    except ValidationError:
        raise
    prs = _load_prs(resolved)
    slide = _validate_slide_index(prs, slide_index)
    shape = slide.shapes.add_shape(mso_shape, Inches(left), Inches(top), Inches(width), Inches(height))
    if shape.has_text_frame:
        shape.text_frame.word_wrap = True
    if text and shape.has_text_frame:
        shape.text_frame.text = text
    _save_and_evict(prs, resolved)
    _audit_log("add_shape", resolved, slide_index, {"shape_type": shape_type})
    _risk = _calculate_overflow_risk(text, width, height)
    _overlap = _first_overlap_for_shape(slide, shape.shape_id)
    return {
        "path": str(resolved), "slide_index": slide_index, "shape_id": shape.shape_id,
        "shape_type": shape_type,
        "overflow_risk": _risk["overflow_risk"], "overflow_detail": _risk["overflow_detail"],
        "overlap_warning": _overlap,
    }


def add_table_to_slide(
    path: str, slide_index: int, rows: int, cols: int,
    left: float, top: float, width: float, height: float,
    data: list[list[str]] | None = None, confirm: bool = False,
    font_size_pt: float | None = None, header_font_size_pt: float | None = None,
    header_bold: bool | None = None, row_height_pt: float | None = None,
    suppress_content_placeholder: bool = True,
) -> dict:
    """Add a table to a slide. Requires PPT_ENABLE_WRITE=true and confirm=True."""
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    if rows < 1 or cols < 1:
        raise ValidationError(f"rows and cols must each be >= 1 (got rows={rows}, cols={cols})")
    if data is not None:
        if len(data) != rows:
            raise ValidationError(f"data has {len(data)} rows but rows={rows} was specified")
        for r, row in enumerate(data):
            if len(row) != cols:
                raise ValidationError(f"data[{r}] has {len(row)} columns but cols={cols} was specified")
    prs = _load_prs(resolved)
    slide = _validate_slide_index(prs, slide_index)
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    if data is not None:
        for r in range(rows):
            for c in range(cols):
                table_shape.table.cell(r, c).text = data[r][c]
    runs_updated = 0
    rows_resized = 0
    placeholders_removed = 0
    if any(v is not None for v in (font_size_pt, header_font_size_pt, header_bold, row_height_pt)):
        style_result = _style_table_shape(
            table_shape, font_size_pt=font_size_pt, header_font_size_pt=header_font_size_pt,
            header_bold=header_bold, row_height_pt=row_height_pt,
        )
        runs_updated = style_result["runs_updated"]
        rows_resized = style_result["rows_resized"]
    if suppress_content_placeholder:
        placeholders_removed = _remove_empty_overlapping_content_placeholders(slide, table_shape)
    _save_and_evict(prs, resolved)
    _audit_log("add_table_to_slide", resolved, slide_index, {"rows": rows, "cols": cols})
    return {
        "path": str(resolved), "slide_index": slide_index, "shape_id": table_shape.shape_id,
        "rows": rows, "cols": cols, "runs_updated": runs_updated,
        "rows_resized": rows_resized, "placeholders_removed": placeholders_removed,
    }


def set_table_style(
    path: str, slide_index: int, shape_id: int,
    font_size_pt: float | None = None, bold: bool | None = None, color_hex: str | None = None,
    header_font_size_pt: float | None = None, header_bold: bool | None = None,
    header_color_hex: str | None = None, fill_color_hex: str | None = None,
    header_fill_color_hex: str | None = None, row_height_pt: float | None = None,
    header_row_height_pt: float | None = None, row_heights_pt: list[float] | None = None,
    confirm: bool = False,
) -> dict:
    """Apply font, colour, header, fill, and row-height styling to a table. Requires PPT_ENABLE_WRITE=true."""
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    slide = _validate_slide_index(prs, slide_index)
    table_shape = _find_shape_by_id(slide, shape_id)
    result = _style_table_shape(
        table_shape, font_size_pt=font_size_pt, bold=bold, color_hex=color_hex,
        header_font_size_pt=header_font_size_pt, header_bold=header_bold,
        header_color_hex=header_color_hex, fill_color_hex=fill_color_hex,
        header_fill_color_hex=header_fill_color_hex, row_height_pt=row_height_pt,
        header_row_height_pt=header_row_height_pt, row_heights_pt=row_heights_pt,
    )
    _save_and_evict(prs, resolved)
    _audit_log("set_table_style", resolved, slide_index, {"shape_id": shape_id})
    return {"path": str(resolved), "slide_index": slide_index, **result}


def set_shape_fill(
    path: str, slide_index: int, shape_id: int, fill_color_hex: str, confirm: bool = False,
) -> dict:
    """Apply a solid fill colour to a shape. Requires PPT_ENABLE_WRITE=true and confirm=True."""
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    slide = _validate_slide_index(prs, slide_index)
    shape = _find_shape_by_id(slide, shape_id)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*_normalize_color_hex(fill_color_hex))
    normalized = fill_color_hex.strip().lstrip("#").upper()
    _save_and_evict(prs, resolved)
    _audit_log("set_shape_fill", resolved, slide_index, {"shape_id": shape_id})
    return {
        "path": str(resolved), "slide_index": slide_index, "shape_id": shape_id,
        "fill_color_hex": normalized, "fill_updated": True,
    }


def set_shape_geometry(
    path: str, slide_index: int, shape_id: int,
    left_in: float | None = None, top_in: float | None = None,
    width_in: float | None = None, height_in: float | None = None,
    confirm: bool = False,
) -> dict:
    """Move and/or resize a shape by shape_id (inches). Requires PPT_ENABLE_WRITE=true and confirm=True."""
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    if all(v is None for v in (left_in, top_in, width_in, height_in)):
        raise ValidationError(
            "set_shape_geometry requires at least one of left_in/top_in/width_in/height_in"
        )
    prs = _load_prs(resolved)
    slide = _validate_slide_index(prs, slide_index)
    shape = _find_shape_by_id(slide, shape_id)
    if left_in is not None:
        shape.left = Inches(left_in)
    if top_in is not None:
        shape.top = Inches(top_in)
    if width_in is not None:
        shape.width = Inches(width_in)
    if height_in is not None:
        shape.height = Inches(height_in)
    EMU_PER_INCH = 914400
    final_left_in = shape.left / EMU_PER_INCH if shape.left is not None else 0.0
    final_top_in = shape.top / EMU_PER_INCH if shape.top is not None else 0.0
    final_w_in = shape.width / EMU_PER_INCH if shape.width is not None else 0.0
    final_h_in = shape.height / EMU_PER_INCH if shape.height is not None else 0.0
    slide_w_in = prs.slide_width / EMU_PER_INCH if prs.slide_width is not None else 0.0
    slide_h_in = prs.slide_height / EMU_PER_INCH if prs.slide_height is not None else 0.0
    warnings: list[str] = []
    if final_left_in < 0:
        warnings.append(f"left {final_left_in:.3f}in is before slide left edge")
    if final_top_in < 0:
        warnings.append(f"top {final_top_in:.3f}in is above slide top edge")
    if slide_w_in > 0 and final_left_in + final_w_in > slide_w_in:
        warnings.append(f"right edge {final_left_in + final_w_in:.3f}in exceeds slide width {slide_w_in:.3f}in")
    if slide_h_in > 0 and final_top_in + final_h_in > slide_h_in:
        warnings.append(f"bottom edge {final_top_in + final_h_in:.3f}in exceeds slide height {slide_h_in:.3f}in")
    _save_and_evict(prs, resolved)
    _audit_log("set_shape_geometry", resolved, slide_index, {"shape_id": shape_id})
    return {
        "updated": True, "path": str(resolved), "slide_index": slide_index, "shape_id": shape_id,
        "geometry_in": {
            "left": round(final_left_in, 4), "top": round(final_top_in, 4),
            "width": round(final_w_in, 4), "height": round(final_h_in, 4),
        },
        "warnings": warnings,
    }


def delete_shape(path: str, slide_index: int, shape_id: int, confirm: bool = False) -> dict:
    """Remove a shape from a slide by shape_id. Requires PPT_ENABLE_WRITE=true and confirm=True."""
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    slide = _validate_slide_index(prs, slide_index)
    shape = _find_shape_by_id(slide, shape_id)
    sp = shape._element
    sp.getparent().remove(sp)
    _save_and_evict(prs, resolved)
    _audit_log("delete_shape", resolved, slide_index, {"shape_id": shape_id})
    return {"deleted": True, "shape_id": shape_id, "slide_index": slide_index}


def insert_image(
    path: str, slide_index: int, image_path: str,
    left: float = 1.0, top: float = 1.0,
    width: float | None = None, height: float | None = None,
    confirm: bool = False,
) -> dict:
    """Insert an allowlisted image onto a slide. Requires PPT_ENABLE_WRITE=true and confirm=True."""
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    if width is not None and width <= 0:
        raise ValidationError("width must be greater than 0")
    if height is not None and height <= 0:
        raise ValidationError("height must be greater than 0")
    image_resolved = _check_image_path(image_path)
    prs = _load_prs(resolved)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValidationError(f"slide_index {slide_index} out of range")
    slide = prs.slides[slide_index]
    pic = slide.shapes.add_picture(
        str(image_resolved), Inches(left), Inches(top),
        Inches(width) if width is not None else None,
        Inches(height) if height is not None else None,
    )
    _save_and_evict(prs, resolved)
    _audit_log("insert_image", resolved, slide_index, {"image_path": str(image_resolved)})
    return {"slide_index": slide_index, "shape_name": pic.name}


