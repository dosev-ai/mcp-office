"""_shapes_text — text manipulation functions for pptmcp shape operations."""
from __future__ import annotations

import copy
import logging
import re

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn as _qn
from pptx.util import Pt

from pptmcp.presentation_pptx import (
    _ALIGNMENT_MAP,
    ValidationError,
    _audit_log,
    _check_confirm,
    _check_path,
    _check_write,
    _find_shape_by_id,
    _hex_to_rgb,
    _load_prs,
    _validate_slide_index,
)

_log = logging.getLogger(__name__)

_A_R = _qn('a:r')
_A_BR = _qn('a:br')
_A_T = _qn('a:t')
_A_RPR = _qn('a:rPr')
_A_PPR = _qn('a:pPr')
_XML_SPACE = '{http://www.w3.org/XML/1998/namespace}space'


def _collect_para_text(para_p) -> str:
    """Return paragraph text treating <a:br> elements as '\\n'."""
    parts: list[str] = []
    for child in para_p:
        if child.tag == _A_R:
            t = child.find(_A_T)
            parts.append(t.text if (t is not None and t.text) else '')
        elif child.tag == _A_BR:
            parts.append('\n')
    return ''.join(parts)


def _rebuild_para_text(para_p, new_text: str) -> None:
    """Replace all <a:r>/<a:br> children of para_p with runs for new_text.

    Preserves the run-properties (<a:rPr>) of the first existing run.
    Newlines in new_text become <a:br> elements.
    """
    first_r = para_p.find(_A_R)
    saved_rpr = None
    if first_r is not None:
        rpr = first_r.find(_A_RPR)
        if rpr is not None:
            saved_rpr = copy.deepcopy(rpr)

    for elem in list(para_p):
        if elem.tag in (_A_R, _A_BR):
            para_p.remove(elem)

    insert_idx = 0
    for i, child in enumerate(para_p):
        if child.tag == _A_PPR:
            insert_idx = i + 1
            break

    segments = new_text.split('\n')
    for seg_num, segment in enumerate(segments):
        if seg_num > 0:
            br = etree.Element(_A_BR)
            if saved_rpr is not None:
                br.append(copy.deepcopy(saved_rpr))
            para_p.insert(insert_idx, br)
            insert_idx += 1
        if segment:
            r = etree.Element(_A_R)
            if saved_rpr is not None:
                r.append(copy.deepcopy(saved_rpr))
            t = etree.SubElement(r, _A_T)
            t.text = segment
            if segment != segment.strip():
                t.set(_XML_SPACE, 'preserve')
            para_p.insert(insert_idx, r)
            insert_idx += 1


def set_speaker_notes(
    path: str,
    slide_index: int,
    notes_text: str | None = None,
    confirm: bool = False,
    notes: str | None = None,
) -> dict:
    """Write speaker notes for a slide. Requires PPT_ENABLE_WRITE=true and confirm=True.

    Accept either notes_text (primary, kept for backward compat) or notes (preferred alias).
    If both are provided, notes_text takes precedence.
    Changes are held in memory — call save(path, confirm=True) to persist to disk.
    """
    effective_notes = notes_text if notes_text is not None else notes
    if effective_notes is None:
        raise ValidationError("Either notes_text or notes must be provided")
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValidationError(f"slide_index {slide_index} out of range")
    slide = prs.slides[slide_index]
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = effective_notes
    _audit_log("set_speaker_notes", resolved, slide_index)
    return {"slide_index": slide_index, "notes_text": effective_notes}


def edit_text_placeholder(
    path: str,
    slide_index: int,
    placeholder_idx: int = 0,
    text: str = "",
    confirm: bool = False,
    shape_id: int | None = None,
) -> dict:
    """Replace text in a placeholder by index, or in any text-containing shape by shape_id.

    If shape_id is provided: find shape by shape.shape_id, check it has text_frame, replace text.
    If shape_id is NOT provided: use placeholder_idx to find placeholder by
    ph.placeholder_format.idx (not list position).
    Requires PPT_ENABLE_WRITE=true and confirm=True.
    """
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValidationError(f"slide_index {slide_index} out of range")
    slide = prs.slides[slide_index]

    if shape_id is not None:
        shape = _find_shape_by_id(slide, shape_id)
        if not shape.has_text_frame:
            raise ValidationError(f"shape_id {shape_id} has no text frame")
        shape.text_frame.text = text
        _audit_log(
            "edit_text_placeholder",
            resolved,
            slide_index,
            {"shape_id": shape_id},
        )
        return {"slide_index": slide_index, "shape_id": shape_id, "text": text}

    valid_idxs = [ph.placeholder_format.idx for ph in slide.placeholders]
    try:
        ph = slide.placeholders[placeholder_idx]
        ph.text = text
    except KeyError:
        raise ValidationError(
            f"placeholder_idx {placeholder_idx} not found on slide {slide_index}. "
            f"Valid indices: {valid_idxs}"
        )
    _audit_log(
        "edit_text_placeholder",
        resolved,
        slide_index,
        {"placeholder_idx": placeholder_idx},
    )
    return {"slide_index": slide_index, "placeholder_idx": placeholder_idx, "text": text}


def replace_slide_text(
    path: str,
    find: str,
    replace: str,
    confirm: bool = False,
) -> dict:
    """Bulk find-and-replace text across all slides. Handles <a:br> line-break elements.

    Requires PPT_ENABLE_WRITE=true and confirm=True.
    """
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    replacements_made = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full_text = _collect_para_text(para._p)
                    if find not in full_text:
                        continue
                    count = full_text.count(find)
                    replacements_made += count
                    new_text = full_text.replace(find, replace)
                    _rebuild_para_text(para._p, new_text)
    _audit_log(
        "replace_slide_text",
        resolved,
        None,
        {"replacements_made": replacements_made},
    )
    return {"replacements_made": replacements_made}


def set_text_format(
    path: str,
    slide_index: int,
    shape_id: int,
    bold: bool | None = None,
    italic: bool | None = None,
    font_size_pt: float | None = None,
    font_name: str | None = None,
    color_hex: str | None = None,
    confirm: bool = False,
    paragraph_range: list[int] | None = None,
    run_index: int | None = None,
) -> dict:
    """Apply text formatting to runs in a shape's text frame.

    paragraph_range: optional [start, end] inclusive; None = all paragraphs.
    run_index: optional 0-based run index per matched paragraph; None = all runs.
    color_hex: 6-char hex without '#', e.g. 'FF0000'. Requires PPT_ENABLE_WRITE=true + confirm=True.
    Changes are held in memory — call save() to persist.
    """
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    if color_hex is not None and not re.match(r'^[0-9A-Fa-f]{6}$', color_hex):
        raise ValidationError(f"Invalid hex color {color_hex!r}. Must be exactly 6 hex chars, e.g. 'FF0000'.")
    if paragraph_range is not None and len(paragraph_range) != 2:
        raise ValidationError("paragraph_range must be a list of exactly 2 ints: [start, end] (inclusive).")
    prs = _load_prs(resolved)
    slide = _validate_slide_index(prs, slide_index)
    shape = _find_shape_by_id(slide, shape_id)
    if not shape.has_text_frame:
        raise ValidationError(f"Shape {shape_id} has no text frame")
    paragraphs = shape.text_frame.paragraphs
    if paragraph_range is not None:
        para_start, para_end = paragraph_range[0], paragraph_range[1]
        para_indices = range(para_start, para_end + 1)
    else:
        para_indices = range(len(paragraphs))
    runs_updated = 0
    for para_idx in para_indices:
        if para_idx < 0 or para_idx >= len(paragraphs):
            continue
        para = paragraphs[para_idx]
        runs = para.runs
        if run_index is not None:
            run_indices = [run_index] if 0 <= run_index < len(runs) else []
        else:
            run_indices = range(len(runs))
        for r_idx in run_indices:
            run = runs[r_idx]
            if bold is not None:
                run.font.bold = bold
            if italic is not None:
                run.font.italic = italic
            if font_size_pt is not None:
                run.font.size = Pt(font_size_pt)
            if font_name is not None:
                run.font.name = font_name
            if color_hex is not None:
                rgb = _hex_to_rgb("#" + color_hex)
                run.font.color.rgb = RGBColor(*rgb)
            runs_updated += 1
    _audit_log("set_text_format", resolved, slide_index, {"shape_id": shape_id})
    return {
        "path": str(resolved),
        "slide_index": slide_index,
        "shape_id": shape_id,
        "runs_updated": runs_updated,
    }


def set_paragraph_format(
    path: str,
    slide_index: int,
    shape_id: int,
    paragraph_index: int = 0,
    alignment: str | None = None,
    line_spacing: float | None = None,
    space_before_pt: float | None = None,
    space_after_pt: float | None = None,
    confirm: bool = False,
    paragraph_range: list[int] | None = None,
) -> dict:
    """Set paragraph formatting (alignment, spacing) for paragraphs in a shape.

    paragraph_range: optional [start, end] inclusive; overrides paragraph_index when set.
    alignment: CENTER, LEFT, RIGHT, JUSTIFY, DISTRIBUTE.
    Requires PPT_ENABLE_WRITE=true and confirm=True. Changes held in memory — call save() to persist.
    """
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    if paragraph_range is not None and len(paragraph_range) != 2:
        raise ValidationError("paragraph_range must be a list of exactly 2 ints: [start, end] (inclusive).")
    prs = _load_prs(resolved)
    slide = _validate_slide_index(prs, slide_index)
    shape = _find_shape_by_id(slide, shape_id)
    if not shape.has_text_frame:
        raise ValidationError(f"Shape {shape_id} has no text frame")
    paragraphs = shape.text_frame.paragraphs
    n_para = len(paragraphs)
    if paragraph_range is not None:
        para_start, para_end = paragraph_range[0], paragraph_range[1]
        para_indices = [i for i in range(para_start, para_end + 1) if 0 <= i < n_para]
    else:
        if paragraph_index < 0 or paragraph_index >= n_para:
            raise ValidationError(
                f"paragraph_index {paragraph_index} out of range (0-{n_para - 1})"
            )
        para_indices = [paragraph_index]
    if alignment is not None:
        upper = alignment.upper()
        if upper not in _ALIGNMENT_MAP:
            raise ValidationError(
                f"Invalid alignment {alignment!r}. Valid: {sorted(_ALIGNMENT_MAP.keys())}"
            )
    for idx in para_indices:
        para = paragraphs[idx]
        if alignment is not None:
            para.alignment = _ALIGNMENT_MAP[upper]
        if line_spacing is not None:
            para.line_spacing = line_spacing
        if space_before_pt is not None:
            para.space_before = Pt(space_before_pt)
        if space_after_pt is not None:
            para.space_after = Pt(space_after_pt)
    _audit_log(
        "set_paragraph_format", resolved, slide_index,
        {"shape_id": shape_id, "paragraph_index": paragraph_index},
    )
    return {
        "path": str(resolved),
        "slide_index": slide_index,
        "shape_id": shape_id,
        "paragraph_index": paragraph_index,
    }
