"""
Write operations for pptmcp: create_presentation, add_slide, delete_slide,
reorder_slides, save, copy_slide, clear_slide_content, apply_slide_layout.

No FastMCP / MCP imports — pure python-pptx + stdlib only.
_atomic_save, _evict_prs, and _load_prs are top-level imports from
_presentation_cache; tests that need to patch _load_prs must patch
pptmcp._presentation_write._load_prs (the name bound in this module's
namespace), not pptmcp._presentation_cache._load_prs.
Governance helpers (_check_write, _check_confirm, _check_path, _audit_log)
are imported from presentation_pptx at call time so tests patching
pptmcp.presentation_pptx.* take effect without restarting the server.
"""
from __future__ import annotations

import copy
import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from pptmcp._presentation_cache import _atomic_save, _evict_prs, _load_prs
from pptmcp._presentation_governance import ValidationError

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Internal write helpers
# ---------------------------------------------------------------------------


def _find_blank_layout_idx(prs_obj: Any) -> int:
    """Find the index of the 'Blank' slide layout by name, falling back to 0."""
    for i, layout in enumerate(prs_obj.slide_layouts):
        if "blank" in layout.name.lower():
            return i
    return 0


def _copy_slide_background(source_slide: Any, dest_slide: Any) -> None:
    """Deep-copy the <p:bg> element from source to dest slide, if present."""
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    source_bg = source_slide._element.find(".//p:cSld/p:bg", nsmap)
    if source_bg is not None:
        dest_cSld = dest_slide._element.find(".//p:cSld", nsmap)
        if dest_cSld is None:
            return
        existing_bg = dest_cSld.find("p:bg", nsmap)
        if existing_bg is not None:
            dest_cSld.remove(existing_bg)
        dest_spTree = dest_cSld.find("p:spTree", nsmap)
        if dest_spTree is not None:
            dest_cSld.insert(list(dest_cSld).index(dest_spTree), copy.deepcopy(source_bg))
        else:
            dest_cSld.insert(0, copy.deepcopy(source_bg))


def _validate_slide_index_w(prs: Any, slide_index: int, context: str = "slide") -> Any:
    """Validate slide_index within bounds; return the slide or raise ValidationError."""
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValidationError(
            f"{context}_index {slide_index} out of range (0-{len(prs.slides) - 1})"
        )
    return prs.slides[slide_index]


_EMPTY_PLACEHOLDER_PATTERNS: frozenset[str] = frozenset({
    "click to add title",
    "click to add text",
    "click to add subtitle",
    "click to edit master title style",
    "click to edit master text styles",
    "click to edit master subtitle style",
    "click to add notes",
})


def _is_title_placeholder(shape: Any) -> bool:
    try:
        return shape.placeholder_format.type in {
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
            PP_PLACEHOLDER.VERTICAL_TITLE,
        }
    except Exception:
        return False


def _remove_empty_content_placeholders(slide: Any) -> int:
    placeholders = []
    for shape in slide.shapes:
        if not shape.is_placeholder or _is_title_placeholder(shape):
            continue
        if not shape.has_text_frame:
            placeholders.append(shape)
            continue
        text = shape.text_frame.text.strip().lower()
        if text == "" or text in _EMPTY_PLACEHOLDER_PATTERNS:
            placeholders.append(shape)
    for placeholder in placeholders:
        placeholder._element.getparent().remove(placeholder._element)
    return len(placeholders)


# ---------------------------------------------------------------------------
# Public write functions
# Each function imports governance helpers from presentation_pptx at call time
# so that tests patching pptmcp.presentation_pptx.* take effect.
# ---------------------------------------------------------------------------


def create_presentation(
    path: str,
    title: str | None = None,
    template_path: str | None = None,
    confirm: bool = False,
) -> dict:
    """Create a brand-new .pptx file at path. Optionally add a title slide.

    Optionally initialise with slide masters/theme from a template .pptx.
    Requires PPT_ENABLE_WRITE=true and confirm=True.
    Fails if file already exists — use a new destination path.
    """
    from pptmcp.presentation_pptx import (  # noqa: PLC0415
        _audit_log,
        _check_confirm,
        _check_path,
        _check_write,
    )

    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    if resolved.exists():
        raise ValidationError(f"File already exists: {resolved}")

    resolved_template: Path | None = None
    if template_path is not None:
        resolved_template = _check_path(template_path)
        if not resolved_template.exists():
            raise ValidationError(f"Template file not found: {resolved_template}")
        new_prs = Presentation(str(resolved_template))
        xml_slides = new_prs.slides._sldIdLst
        for slide_elem in list(xml_slides):
            rId = slide_elem.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            if rId:
                try:
                    new_prs.part.drop_rel(rId)
                except Exception as exc:
                    logger.warning(
                        "drop_rel failed for rId=%r during template cleanup: %s — "
                        "continuing; template PPTX may contain orphaned relationships",
                        rId,
                        exc,
                    )
            xml_slides.remove(slide_elem)
    else:
        new_prs = Presentation()

    if title is not None:
        if not new_prs.slide_layouts:
            raise ValidationError(
                "Template has no slide layouts — cannot add a title slide"
            )
        layout = new_prs.slide_layouts[0]
        slide = new_prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
    _atomic_save(new_prs, resolved)
    _audit_log("create_presentation", resolved)
    return {
        "path": str(resolved),
        "slide_count": len(new_prs.slides),
        "title": title,
        "template_path": str(resolved_template) if resolved_template is not None else None,
    }


def add_slide(
    path: str,
    layout_index: int = 1,
    title: str | None = None,
    confirm: bool = False,
    layout_name: str | None = None,
    suppress_content_placeholder: bool = False,
) -> dict:
    """Append a new slide using layout_index or layout_name.

    Requires PPT_ENABLE_WRITE=true and confirm=True.
    """
    from pptmcp.presentation_pptx import (  # noqa: PLC0415
        _audit_log,
        _check_confirm,
        _check_path,
        _check_write,
    )

    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    n_layouts = len(prs.slide_layouts)
    if layout_name is not None:
        names = [sl.name for sl in prs.slide_layouts]
        if layout_name not in names:
            raise ValidationError(
                f"layout_name {layout_name!r} not found; available: {names}"
            )
        layout_index = names.index(layout_name)
    if layout_index < 0 or layout_index >= n_layouts:
        raise ValidationError(
            f"layout_index {layout_index} out of range (0-{n_layouts - 1})"
        )
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)
    if title is not None and slide.shapes.title is not None:
        slide.shapes.title.text = title
    placeholders_removed = 0
    if suppress_content_placeholder:
        placeholders_removed = _remove_empty_content_placeholders(slide)
    new_index = len(prs.slides) - 1
    _atomic_save(prs, resolved)
    _evict_prs(resolved)
    _audit_log("add_slide", resolved, new_index)
    return {
        "slide_index": new_index,
        "title": title,
        "placeholders_removed": placeholders_removed,
    }


def reorder_slides(path: str, new_order: list[int], confirm: bool = False) -> dict:
    """Reorder slides to new_order permutation.

    Requires PPT_ENABLE_WRITE=true and confirm=True.
    """
    from pptmcp.presentation_pptx import (  # noqa: PLC0415
        _audit_log,
        _check_confirm,
        _check_path,
        _check_write,
    )

    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    n = len(prs.slides)
    if sorted(new_order) != list(range(n)):
        raise ValidationError(
            f"new_order must be a permutation of 0..{n - 1}, got {new_order}"
        )
    xml_slides = list(prs.slides._sldIdLst)
    for elem in xml_slides:
        prs.slides._sldIdLst.remove(elem)
    for i in new_order:
        prs.slides._sldIdLst.append(xml_slides[i])
    _atomic_save(prs, resolved)
    _evict_prs(resolved)
    _audit_log("reorder_slides", resolved, None, {"new_order": new_order})
    return {"new_order": new_order}


def delete_slide(path: str, slide_index: int, confirm: bool = False) -> dict:
    """Remove a slide by index. Requires PPT_ENABLE_WRITE=true and confirm=True."""
    from pptmcp.presentation_pptx import (  # noqa: PLC0415
        _audit_log,
        _check_confirm,
        _check_path,
        _check_write,
    )

    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValidationError(f"slide_index {slide_index} out of range")
    xml_slides = prs.slides._sldIdLst
    slide_elem = xml_slides[slide_index]
    rId = slide_elem.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception as exc:
            logger.warning(
                "drop_rel failed for rId=%r during delete_slide: %s — "
                "continuing; PPTX may contain an orphaned relationship",
                rId,
                exc,
            )
    xml_slides.remove(slide_elem)
    _atomic_save(prs, resolved)
    _evict_prs(resolved)
    _audit_log("delete_slide", resolved, slide_index)
    return {"deleted_index": slide_index}


def save(path: str, confirm: bool = False) -> dict:
    """Persist the in-memory presentation to disk.

    Requires PPT_ENABLE_WRITE=true and confirm=True.
    """
    from pptmcp.presentation_pptx import (  # noqa: PLC0415
        _audit_log,
        _check_confirm,
        _check_path,
        _check_write,
    )

    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    _atomic_save(prs, resolved)
    _evict_prs(resolved)
    _audit_log("save", resolved)
    return {"saved": True, "path": str(resolved)}


def copy_slide(
    source_path: str | None = None,
    source_slide_index: int | None = None,
    target_path: str | None = None,
    target_slide_index: int | None = None,
    confirm: bool = False,
    path: str | None = None,
    slide_index: int | None = None,
) -> dict:
    """Copy a slide from source_path[source_slide_index] into target_path.

    If source_path == target_path, copies within the same presentation.
    Param aliases: path= for source_path, slide_index= for source_slide_index.
    target_path defaults to source_path when omitted (within-deck copy).
    target_slide_index: optional 0-based position for the copied slide (default: append).
    Requires PPT_ENABLE_WRITE=true and confirm=True.
    """
    from pptmcp.presentation_pptx import (  # noqa: PLC0415
        _audit_log,
        _check_confirm,
        _check_path,
        _check_write,
    )

    source_path = source_path if source_path is not None else path
    source_slide_index = source_slide_index if source_slide_index is not None else slide_index
    target_path = target_path if target_path is not None else source_path
    if source_path is None:
        raise ValidationError("source_path (or path) is required")
    if source_slide_index is None:
        raise ValidationError("source_slide_index (or slide_index) is required")
    _check_write()
    _check_confirm(confirm)
    source_resolved = _check_path(source_path)
    target_resolved = _check_path(target_path)

    source_prs = _load_prs(source_resolved)
    source_slide = _validate_slide_index_w(
        source_prs, source_slide_index, context="source_slide"
    )
    blank_layout_idx = _find_blank_layout_idx(source_prs)

    if source_resolved == target_resolved:
        new_slide = source_prs.slides.add_slide(source_prs.slide_layouts[blank_layout_idx])
        sp_tree = new_slide.shapes._spTree
        for child in list(sp_tree):
            sp_tree.remove(child)
        for element in source_slide.shapes._spTree:
            sp_tree.append(copy.deepcopy(element))
        _copy_slide_background(source_slide, new_slide)
        new_index = len(source_prs.slides) - 1
        if target_slide_index is not None:
            n = len(source_prs.slides)
            if target_slide_index < 0 or target_slide_index >= n:
                raise ValidationError(
                    f"target_slide_index {target_slide_index} out of range (0-{n - 1})"
                )
            new_slide_elem = list(source_prs.slides._sldIdLst)[-1]
            source_prs.slides._sldIdLst.remove(new_slide_elem)
            source_prs.slides._sldIdLst.insert(target_slide_index, new_slide_elem)
            new_index = target_slide_index
        _atomic_save(source_prs, source_resolved)
        _evict_prs(source_resolved)
        _audit_log(
            "copy_slide", source_resolved, source_slide_index, {"new_slide_index": new_index}
        )
        return {
            "source_path": str(source_resolved),
            "target_path": str(target_resolved),
            "source_slide_index": source_slide_index,
            "new_slide_index": new_index,
        }

    if not target_resolved.exists():
        raise ValidationError(f"Target file does not exist: {target_resolved}")
    target_prs = _load_prs(target_resolved)
    tgt_blank_idx = _find_blank_layout_idx(target_prs)
    new_slide = target_prs.slides.add_slide(target_prs.slide_layouts[tgt_blank_idx])
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        sp_tree.remove(child)
    for element in source_slide.shapes._spTree:
        sp_tree.append(copy.deepcopy(element))
    _copy_slide_background(source_slide, new_slide)
    new_index = len(target_prs.slides) - 1
    if target_slide_index is not None:
        n = len(target_prs.slides)
        if target_slide_index < 0 or target_slide_index >= n:
            raise ValidationError(
                f"target_slide_index {target_slide_index} out of range (0-{n - 1})"
            )
        new_slide_elem = list(target_prs.slides._sldIdLst)[-1]
        target_prs.slides._sldIdLst.remove(new_slide_elem)
        target_prs.slides._sldIdLst.insert(target_slide_index, new_slide_elem)
        new_index = target_slide_index
    _atomic_save(target_prs, target_resolved)
    _evict_prs(target_resolved)
    _audit_log(
        "copy_slide",
        target_resolved,
        source_slide_index,
        {"source_path": str(source_resolved), "new_slide_index": new_index},
    )
    return {
        "source_path": str(source_resolved),
        "target_path": str(target_resolved),
        "source_slide_index": source_slide_index,
        "new_slide_index": new_index,
    }


def clear_slide_content(path: str, slide_index: int, confirm: bool = False) -> dict:
    """Remove ALL shapes from a slide, leaving it blank.

    Uses a snapshot before the deletion loop to avoid index corruption.
    Requires PPT_ENABLE_WRITE=true and confirm=True.

    Returns: {"cleared_count": int, "slide_index": int}
    """
    from pptmcp.presentation_pptx import (  # noqa: PLC0415
        _audit_log,
        _check_confirm,
        _check_path,
        _check_write,
    )

    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    slide = _validate_slide_index_w(prs, slide_index)
    shapes_snapshot = list(slide.shapes)
    for sp in shapes_snapshot:
        sp._element.getparent().remove(sp._element)
    _atomic_save(prs, resolved)
    _evict_prs(resolved)
    cleared = len(shapes_snapshot)
    _audit_log("clear_slide_content", resolved, slide_index, {"cleared_count": cleared})
    return {"cleared_count": cleared, "slide_index": slide_index}


def apply_slide_layout(
    path: str,
    slide_index: int,
    layout_index: int,
    remove_placeholders: bool = True,
    confirm: bool = False,
) -> dict:
    """Apply a different slide layout to an existing slide via OPC relationship swap.

    python-pptx has no ``slide.slide_layout`` setter; this function swaps the
    underlying OPC relationship directly (verified to survive save/reload).

    If ``remove_placeholders=True`` (default), all placeholder shapes are removed
    after the layout swap so the agent can re-populate using add_content.

    Requires PPT_ENABLE_WRITE=true and confirm=True.

    Returns: {"layout_applied": str, "placeholders_removed": int, "slide_index": int}
    """
    from pptmcp.presentation_pptx import (  # noqa: PLC0415
        _audit_log,
        _check_confirm,
        _check_path,
        _check_write,
    )

    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    slide = _validate_slide_index_w(prs, slide_index)
    if layout_index < 0 or layout_index >= len(prs.slide_layouts):
        raise ValidationError(
            f"layout_index {layout_index} out of range (0-{len(prs.slide_layouts) - 1})"
        )
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT  # noqa: PLC0415

    new_layout_part = prs.slide_layouts[layout_index].part
    slide_part = slide.part
    existing_layout_rels = [
        r for r in slide_part.rels.values() if r.reltype == RT.SLIDE_LAYOUT
    ]
    for rel in existing_layout_rels:
        slide_part.drop_rel(rel.rId)
    slide_part.relate_to(new_layout_part, RT.SLIDE_LAYOUT)
    placeholders_removed = 0
    if remove_placeholders:
        ph_snapshot = [s for s in slide.shapes if s.is_placeholder]
        for ph in ph_snapshot:
            ph._element.getparent().remove(ph._element)
        placeholders_removed = len(ph_snapshot)
    layout_name = prs.slide_layouts[layout_index].name
    _atomic_save(prs, resolved)
    _evict_prs(resolved)
    _audit_log(
        "apply_slide_layout",
        resolved,
        slide_index,
        {
            "layout_index": layout_index,
            "layout_name": layout_name,
            "placeholders_removed": placeholders_removed,
        },
    )
    return {
        "layout_applied": layout_name,
        "placeholders_removed": placeholders_removed,
        "slide_index": slide_index,
    }
