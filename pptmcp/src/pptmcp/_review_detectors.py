"""_review_detectors — slide quality detector functions for pptmcp review pipeline."""
from __future__ import annotations

import logging
from typing import Any

from pptx.exc import InvalidXmlError

_log = logging.getLogger(__name__)

_EMU_PER_INCH: int = 914_400

DEFAULT_TEXT: frozenset[str] = frozenset({
    "click to add title",
    "click to add text",
    "click to add subtitle",
    "click to edit master title style",
    "click to edit master text styles",
    "title",
})
_PT_HIGH: int = 10 * 12700      # 127_000 EMU  (10pt threshold)
_PT_MEDIUM: int = 14 * 12700    # 177_800 EMU  (14pt threshold)
_SPARSE_FILL_THRESHOLD: float = 0.15
_SPARSE_COUNT_THRESHOLD: int = 4
_TEXT_TRUNCATION_LIMIT: int = 80


def _detect_png_issues(
    png_path: Any,
    expected_resolution: tuple[int, int] | None = None,
    blank_variance_threshold: float = 50.0,
) -> list[dict]:
    """Pixel-read PNG checks: PIXEL_DIMENSION_MISMATCH and BLANK_SLIDE."""
    findings: list[dict] = []
    try:
        from PIL import Image  # noqa: PLC0415
    except Exception as exc:  # noqa: BLE001
        _log.warning("Pillow not available; skipping pixel-read PNG checks: %s", exc)
        return findings
    try:
        with Image.open(str(png_path)) as img:
            img.load()
            actual_w, actual_h = img.size
            mode = img.mode
            thumb = img.convert("RGB").resize((64, 64))
    except Exception as exc:  # noqa: BLE001
        _log.warning("failed to read PNG %s: %s", png_path, exc)
        return findings
    _MIN_RENDER_DIM = 100
    is_render_sized = (actual_w >= _MIN_RENDER_DIM and actual_h >= _MIN_RENDER_DIM)

    if expected_resolution is not None:
        exp_w, exp_h = expected_resolution
        if (actual_w, actual_h) != (exp_w, exp_h):
            findings.append({
                "criterion": "PIXEL_DIMENSION_MISMATCH",
                "result": "FAIL",
                "detail": (
                    f"PNG dimensions {actual_w}x{actual_h} differ from expected "
                    f"{exp_w}x{exp_h}"
                ),
                "severity": "high",
                "actual_resolution": [actual_w, actual_h],
                "expected_resolution": [exp_w, exp_h],
            })

    try:
        raw = thumb.tobytes()
        n_channels = len(raw)
        if n_channels == 0:
            return findings
        r_bytes, g_bytes, b_bytes = raw[0::3], raw[1::3], raw[2::3]
        n = len(r_bytes)
        if n == 0:
            return findings
        mean_r, mean_g, mean_b = sum(r_bytes) / n, sum(g_bytes) / n, sum(b_bytes) / n
        var_total = 0.0
        for v in r_bytes:
            var_total += (v - mean_r) ** 2
        for v in g_bytes:
            var_total += (v - mean_g) ** 2
        for v in b_bytes:
            var_total += (v - mean_b) ** 2
        variance = var_total / (3 * n)
    except Exception as exc:  # noqa: BLE001
        _log.warning("variance computation failed for %s: %s", png_path, exc)
        return findings
    if is_render_sized and variance < blank_variance_threshold:
        findings.append({
            "criterion": "BLANK_SLIDE",
            "result": "FAIL",
            "detail": (
                f"PNG pixel variance {variance:.2f} below blank threshold "
                f"{blank_variance_threshold} — slide may be a flat-fill render "
                f"or a rendering failure (mode={mode}, mean RGB ~"
                f"({mean_r:.0f},{mean_g:.0f},{mean_b:.0f}))"
            ),
            "severity": "critical",
            "variance": round(variance, 4),
            "threshold": blank_variance_threshold,
        })

    return findings



_DEFAULT_FONT_SIZE_PT: float = 11.0
_PT_PER_EMU: float = 1.0 / 12700  # 1 EMU = 1/12700 pt


def _count_text_lines(text_frame: Any) -> int:
    """Count non-empty paragraph text lines in a text frame.

    Each paragraph whose stripped text is non-empty counts as one line.
    Intra-paragraph newlines (``\n``) are also split and counted individually.
    Returns at least 1 when the text frame has any non-empty content.
    """
    count = 0
    for para in text_frame.paragraphs:
        try:
            para_text = para.text
        except Exception:  # noqa: BLE001
            continue
        for line in para_text.split("\n"):
            if line.strip():
                count += 1
    return max(count, 1)


def _get_shape_font_size_pt(text_frame: Any) -> float:
    """Return the largest explicit run font size (in pt) found in *text_frame*.

    Falls back to *_DEFAULT_FONT_SIZE_PT* (11pt) when no run carries an
    explicit size.  pptx stores font size in EMU units (1 pt = 12700 EMU);
    sizes are converted to pt before comparison.
    """
    max_pt: float = 0.0
    for para in text_frame.paragraphs:
        for run in para.runs:
            try:
                size_emu = run.font.size
            except Exception:  # noqa: BLE001
                continue
            if size_emu is None:
                continue
            pt = float(size_emu) * _PT_PER_EMU
            if pt > max_pt:
                max_pt = pt
    return max_pt if max_pt > 0.0 else _DEFAULT_FONT_SIZE_PT


def _detect_text_overflow_heuristic(slide: Any, prs_obj: Any) -> list[dict]:
    """Detect likely text overflow using geometry when available, density otherwise.

    Geometry path (preferred — used when shape height > 0):
      Estimate text height as ``(num_lines * font_size_pt) / 72`` inches and
      compare against the shape bounding-box height.  If the estimate exceeds
      the shape height, a TEXT_OVERFLOW finding is raised (severity=high).
      The finding dict includes ``shape_id``, ``shape_height_in``,
      ``estimated_text_height_in``, and ``overflow_estimate_in`` for
      downstream consumers.

    Density fallback (used when shape height is 0 or geometry is absent):
      Legacy chars-per-sq-inch heuristic: HIGH > 200 chars/sq.in,
      MEDIUM > 100 chars/sq.in.
    """
    findings: list[dict] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if not text:
            continue
        if shape.width is None or shape.height is None:  # H-014: None guard before abs()
            continue
        width_in = abs(shape.width) / _EMU_PER_INCH   # BLK-D03: abs()
        height_in = abs(shape.height) / _EMU_PER_INCH  # BLK-D03: abs()

        # ── Geometry path ────────────────────────────────────────────────
        if height_in > 0:
            num_lines = _count_text_lines(shape.text_frame)
            font_pt = _get_shape_font_size_pt(shape.text_frame)
            estimated_text_height_in = (num_lines * font_pt) / 72.0
            if estimated_text_height_in > height_in:
                overflow_est = round(estimated_text_height_in - height_in, 4)
                findings.append({
                    "criterion": "TEXT_OVERFLOW",
                    "result": "FAIL",
                    "detail": (
                        f"Shape {shape.shape_id} {shape.name!r} "
                        f"({width_in:.2f}in × {height_in:.2f}in): "
                        f"estimated text height {estimated_text_height_in:.3f}in "
                        f"exceeds shape height {height_in:.3f}in "
                        f"(overflow ~{overflow_est:.3f}in, "
                        f"{num_lines} lines at {font_pt:.1f}pt)"
                    ),
                    "severity": "high",
                    "shape_id": shape.shape_id,
                    "shape_height_in": round(height_in, 4),
                    "estimated_text_height_in": round(estimated_text_height_in, 4),
                    "overflow_estimate_in": overflow_est,
                })
                _log.debug(
                    "TEXT_OVERFLOW (geometry): shape_id=%r lines=%d font_pt=%.1f "
                    "est_h=%.3fin shape_h=%.3fin",
                    shape.shape_id, num_lines, font_pt,
                    estimated_text_height_in, height_in,
                )
            continue  # geometry path handled; skip density fallback for this shape

        # ── Density fallback (no geometry) ──────────────────────────────
        area = max(width_in * height_in, 0.01)
        density = len(text) / area
        if density > 200:
            severity = "high"
        elif density > 100:
            severity = "medium"
        else:
            continue
        findings.append({
            "criterion": "TEXT_OVERFLOW",
            "result": "FAIL",
            "detail": (
                f"Shape {shape.shape_id} {shape.name!r} "
                f"({width_in:.2f}in × {height_in:.2f}in, {len(text)} chars) likely overflows"
            ),
            "severity": severity,
        })
        _log.debug(
            "TEXT_OVERFLOW (density): shape_id=%r density=%.1f severity=%s",
            shape.shape_id, density, severity,
        )
    return findings


def _detect_shapes_outside_slide(
    slide: Any, slide_width_emu: int, slide_height_emu: int
) -> list[dict]:
    """Detect shapes outside the slide boundary (RCI-US-16, severity=critical)."""
    findings: list[dict] = []
    slide_width_in = slide_width_emu / _EMU_PER_INCH
    slide_height_in = slide_height_emu / _EMU_PER_INCH
    for shape in slide.shapes:
        if any(v is None for v in [shape.left, shape.top, shape.width, shape.height]):
            continue
        left_in = shape.left / _EMU_PER_INCH
        top_in = shape.top / _EMU_PER_INCH
        width_in = shape.width / _EMU_PER_INCH
        height_in = shape.height / _EMU_PER_INCH
        right_in = left_in + width_in
        bottom_in = top_in + height_in

        overflow_in: dict[str, float] = {}
        violations: list[str] = []
        if left_in < 0:
            overflow_in["left"] = round(-left_in, 4)
            violations.append(f"left edge {left_in:.3f}in before slide left")
        if top_in < 0:
            overflow_in["top"] = round(-top_in, 4)
            violations.append(f"top edge {top_in:.3f}in above slide top")
        if right_in > slide_width_in:
            overflow_in["right"] = round(right_in - slide_width_in, 4)
            violations.append(
                f"right edge extends {right_in - slide_width_in:.3f}in beyond slide"
            )
        if bottom_in > slide_height_in:
            overflow_in["bottom"] = round(bottom_in - slide_height_in, 4)
            violations.append(
                f"bottom edge extends {bottom_in - slide_height_in:.3f}in beyond slide"
            )

        if violations:
            findings.append({
                "criterion": "SHAPE_OUTSIDE_SLIDE",
                "result": "FAIL",
                "detail": f"Shape {shape.shape_id} {shape.name!r}: "
                + "; ".join(violations),
                "severity": "critical",
                "shape_id": shape.shape_id,
                "shape_name": shape.name,
                "bounds_in": {
                    "left": round(left_in, 4),
                    "top": round(top_in, 4),
                    "width": round(width_in, 4),
                    "height": round(height_in, 4),
                },
                "overflow_in": overflow_in,
            })
    return findings


def _detect_table_content_clipping(slide: Any) -> list[dict]:
    """Check table shapes for content clipping (explicit rows over bounds, or >6 auto-rows)."""
    findings: list[dict] = []
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        if shape.height is None:
            continue
        table = shape.table
        shape_height: int = shape.height  # EMU
        known_height = 0
        auto_rows = 0
        for row in table.rows:
            try:
                h = row.height
            except InvalidXmlError:
                h = None  # missing h attr → treat as auto-height
            if h is None or h == 0:
                auto_rows += 1
            else:
                known_height += h

        if auto_rows == 0 and known_height > shape_height:
            findings.append({
                "criterion": "TABLE_CONTENT_CLIPPING",
                "result": "FAIL",
                "detail": (
                    f"Shape {shape.name!r} (id={shape.shape_id}): "
                    f"table row heights sum {known_height / _EMU_PER_INCH:.2f} in "
                    f"exceeds shape height {shape_height / _EMU_PER_INCH:.2f} in — "
                    "content will clip in export"
                ),
                "severity": "high",
            })
        elif auto_rows > 0 and known_height > shape_height:
            findings.append({
                "criterion": "TABLE_CONTENT_CLIPPING",
                "result": "FAIL",
                "detail": (
                    f"Shape {shape.name!r} (id={shape.shape_id}): "
                    f"table explicit row heights sum {known_height / _EMU_PER_INCH:.2f} in already exceeds "
                    f"shape height {shape_height / _EMU_PER_INCH:.2f} in — content WILL clip even ignoring "
                    f"{auto_rows} auto-height row(s)"
                ),
                "severity": "high",
            })
        elif auto_rows > 0 and len(table.rows) > 6:
            findings.append({
                "criterion": "TABLE_CONTENT_CLIPPING",
                "result": "FAIL",
                "detail": (
                    f"Shape {shape.name!r} (id={shape.shape_id}): "
                    f"table has {len(table.rows)} rows ({auto_rows} with auto-height) — "
                    "auto-height rows may expand beyond shape boundary at render time; "
                    "verify in export PNG"
                ),
                "severity": "medium",
            })
    return findings


def _detect_unfinished_placeholder(slide: Any) -> list[dict]:
    """Return findings for placeholders that still contain default/empty text."""
    findings: list[dict] = []
    for shape in slide.shapes:
        try:
            ph_fmt = shape.placeholder_format
        except (ValueError, AttributeError):
            continue
        if ph_fmt is None:
            continue
        try:
            text = shape.text
        except Exception as exc:
            _log.debug(
                "_detect_unfinished_placeholder: shape %r text access failed: %s",
                getattr(shape, "name", "<unknown>"),
                exc,
            )
            continue
        stripped = text.strip()
        if stripped == "" or stripped.lower() in DEFAULT_TEXT:
            findings.append({
                "severity": "high",
                "check": "UNFINISHED_PLACEHOLDER",
                "detail": f"Placeholder {shape.name!r}: text {repr(shape.text[:_TEXT_TRUNCATION_LIMIT])}",
            })
    return findings


def _detect_low_font_size(slide: Any) -> list[dict]:
    """Return findings for run font sizes below readable thresholds (EMU)."""
    findings: list[dict] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    size = run.font.size
                    if size is None:
                        continue
                    if size < _PT_HIGH:
                        severity = "high"
                    elif size < _PT_MEDIUM:
                        severity = "medium"
                    else:
                        continue
                    findings.append({
                        "severity": severity,
                        "check": "LOW_FONT_SIZE",
                        "detail": f"Shape {shape.name!r}: {size / 12700:.1f}pt",
                    })
        except Exception as exc:
            _log.debug(
                "_detect_low_font_size: shape %r iteration failed: %s",
                getattr(shape, "name", "<unknown>"),
                exc,
            )
            continue
    return findings


def _detect_sparse_slide(
    slide: Any, slide_width_emu: int, slide_height_emu: int
) -> list[dict]:
    """Return MEDIUM finding when slide content fill is below threshold (blank slides skipped)."""
    all_shapes = list(slide.shapes)
    if len(all_shapes) == 0:
        return []
    slide_area = max(abs(slide_width_emu) * abs(slide_height_emu), 1)
    content_shapes = [
        s for s in all_shapes if s.has_text_frame and s.text.strip()
    ]
    content_count = len(content_shapes)
    filled_area = sum(
        abs(getattr(s, "width", 0)) * abs(getattr(s, "height", 0))
        for s in content_shapes
    )
    filled_fraction = filled_area / slide_area
    findings: list[dict] = []
    if filled_fraction < _SPARSE_FILL_THRESHOLD and content_count < _SPARSE_COUNT_THRESHOLD:
        findings.append({
            "severity": "medium",
            "check": "SPARSE_SLIDE",
            "detail": (
                f"Content fill {filled_fraction:.1%} ({content_count} text shapes)"
            ),
        })
    return findings
