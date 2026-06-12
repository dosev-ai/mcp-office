"""Content extraction operations for pptmcp (tables, images, text).

Helpers imported from presentation_pptx; no circular dependency.
"""
from __future__ import annotations

import logging
import os
from datetime import datetime, timezone
from pathlib import Path

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptmcp.presentation_pptx import (
    ValidationError,
    _check_path,
    _load_prs,
)

_log = logging.getLogger(__name__)


def _table_rows(shape) -> list[list[str]]:  # noqa: ANN001
    return [[cell.text for cell in row.cells] for row in shape.table.rows]


def _table_row_texts(shape) -> list[str]:  # noqa: ANN001
    row_texts = []
    for row in _table_rows(shape):
        cells = [cell.strip() for cell in row]
        row_texts.append("\t".join(cells) if any(cells) else "")
    return row_texts


def extract_tables(path: str, slide_index: int | None = None) -> dict:
    """Extract all tables from a slide or the whole deck as row/cell arrays.

    Returns ``{"tables": [...]}`` — always a dict, even when empty.
    """
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    indices: list[int] = (
        [slide_index] if slide_index is not None else list(range(len(prs.slides)))
    )
    result = []
    for i in indices:
        if i < 0 or i >= len(prs.slides):
            raise ValidationError(f"slide_index {i} out of range")
        slide = prs.slides[i]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows = [[cell.text for cell in row.cells] for row in table.rows]
                result.append({
                    "slide_index": i,
                    "shape_name": shape.name,
                    "rows": rows,
                })
    return {"tables": result}


def extract_images(path: str) -> dict:
    """List all embedded images across the deck with dimensions and content type.

    Returns ``{"images": [...]}`` — always a dict, even when empty.
    """
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    result = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    content_type: str | None = img.content_type
                except Exception:
                    content_type = None
                result.append({
                    "slide_index": i,
                    "shape_name": shape.name,
                    "width": shape.width,
                    "height": shape.height,
                    "content_type": content_type,
                })
    return {"images": result}


def export_slide_as_text(path: str, slide_index: int | None = None) -> list[dict]:
    """Extract all text from shapes on one slide or the entire deck."""
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    indices: list[int] = (
        [slide_index] if slide_index is not None else list(range(len(prs.slides)))
    )
    result = []
    for i in indices:
        if i < 0 or i >= len(prs.slides):
            raise ValidationError(f"slide_index {i} out of range")
        slide = prs.slides[i]
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            elif shape.has_table:
                texts.extend(row_text for row_text in _table_row_texts(shape) if row_text)
        result.append({"slide_index": i, "texts": texts})
    return result


def extract_presentation_text(path: str) -> dict:
    """Export full deck text with per-slide structure for RAG indexing.

    Returns a richer structure than export_slide_as_text — all slides at once,
    with title, body_text, notes, and per-shape text in one response.
    """
    resolved = _check_path(path)
    prs = _load_prs(resolved)

    slides_out = []
    for idx, slide in enumerate(prs.slides):
        title_text = ""
        body_parts: list[str] = []
        shapes_out: list[dict] = []

        for shape in slide.shapes:
            table_rows: list[list[str]] | None = None
            if shape.has_table:
                table_rows = _table_rows(shape)
                shape_text = "\n".join(
                    row_text for row_text in _table_row_texts(shape) if row_text
                )
            elif shape.has_text_frame:
                shape_text_parts: list[str] = []
                for para in shape.text_frame.paragraphs:
                    para_text = "".join(run.text for run in para.runs).strip()
                    if para_text:
                        shape_text_parts.append(para_text)
                shape_text = "\n".join(shape_text_parts)
            else:
                continue

            # Detect title placeholders so body extraction does not duplicate them.
            is_title = (
                shape.is_placeholder
                and shape.placeholder_format is not None
                and shape.placeholder_format.type in (
                    PP_PLACEHOLDER.TITLE,
                    PP_PLACEHOLDER.CENTER_TITLE,
                    PP_PLACEHOLDER.VERTICAL_TITLE,
                )
            )
            if is_title and not title_text:
                title_text = shape_text
            elif shape_text:
                body_parts.append(shape_text)

            if shape_text:
                shape_entry = {
                    "shape_id": shape.shape_id,
                    "name": shape.name,
                    "text": shape_text,
                }
                if table_rows is not None:
                    shape_entry["shape_type"] = "table"
                    shape_entry["rows"] = table_rows
                shapes_out.append(shape_entry)

        # Extract notes
        notes_text = ""
        if slide.has_notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes_text = ""

        slides_out.append(
            {
                "slide_index": idx,
                "title": title_text,
                "body_text": "\n".join(body_parts),
                "notes": notes_text,
                "shapes": shapes_out,
            }
        )

    return {
        "path": str(resolved),
        "slide_count": len(slides_out),
        "slides": slides_out,
    }


# ---------------------------------------------------------------------------
# Evidence bundle (moved from presentation_pptx — read-only content aggregation)
# ---------------------------------------------------------------------------

def _safe_is_relative_to(path: Path, root: Path) -> bool:
    """Safe wrapper for Path.is_relative_to compatible with Python 3.9+."""
    try:
        return path.is_relative_to(root)
    except (ValueError, TypeError):
        return False


def produce_evidence_bundle(
    path: str,
    review_results: list,
    contract_path: str | None = None,
    exports_dir: str | None = None,
) -> dict:
    """Produce machine-readable evidence JSON aggregating review results.

    review_results: list of dicts returned by review_slide_export().
    Read-only — no write gate required.
    Returns evidence bundle dict with bundle_version, overall_passed, slides, total_findings.
    """
    # Validate path (allowlist + extension)
    resolved_path = _check_path(path)

    # Validate contract_path if provided (deferred import to avoid circular dep)
    resolved_contract: str | None = None
    if contract_path is not None:
        from pptmcp.contract_pptx import _check_json_path  # noqa: PLC0415
        resolved_contract = str(_check_json_path(contract_path))

    # Validate exports_dir if provided
    resolved_exports: str | None = None
    if exports_dir is not None:
        if "\x00" in str(exports_dir):
            raise ValidationError("Invalid exports_dir: embedded null character")
        roots_env = os.getenv("PPT_ALLOWLIST_ROOTS", "").strip()
        if not roots_env:
            raise ValidationError("No allowlist configured")
        exp_resolved = Path(exports_dir).resolve()
        roots = [r.strip() for r in roots_env.split(",") if r.strip()]
        in_allowlist = any(
            _safe_is_relative_to(exp_resolved, Path(r).resolve())
            for r in roots
        )
        if not in_allowlist:
            raise ValidationError("exports_dir is not within the configured allowlist")
        if not exp_resolved.is_dir():
            raise ValidationError("exports_dir does not exist")
        resolved_exports = str(exp_resolved)

    # Validate review_results type
    if not isinstance(review_results, list):
        raise ValidationError("review_results must be a list")

    # Build slides summary
    slides_summary = []
    total_findings = 0
    all_passed = True

    for i, entry in enumerate(review_results):
        if not isinstance(entry, dict):
            raise ValidationError(
                f"review_results entry at index {i} must be a dict, "
                f"got {type(entry).__name__}"
            )
        slide_idx = entry.get("slide_index", None)
        passed = bool(entry.get("passed", False))
        raw_fc = entry.get("findings_count", 0)
        if raw_fc is None:
            raw_fc = 0
        try:
            findings_count = int(raw_fc)
        except (TypeError, ValueError) as exc:
            raise ValidationError(
                f"review_results[{i}].findings_count must be an integer, got {raw_fc!r}"
            ) from exc
        slides_summary.append({
            "slide_index": slide_idx,
            "review_passed": passed,
            "findings_count": findings_count,
        })
        total_findings += findings_count
        if not passed:
            all_passed = False

    # overall_passed: True only if ALL entries passed AND list is non-empty
    overall_passed = all_passed and len(review_results) > 0

    return {
        "bundle_version": "1.0",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "presentation_path": str(resolved_path),
        "contract_path": resolved_contract,
        "exports_dir": resolved_exports,
        "overall_passed": overall_passed,
        "slides": slides_summary,
        "total_findings": total_findings,
    }
