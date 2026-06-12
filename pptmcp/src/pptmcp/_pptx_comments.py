"""
python-pptx backend for PowerPoint comment management (list / add / delete).

No FastMCP / MCP imports — pure python-pptx + lxml + stdlib only.

PPTX comment part paths:
  Per-slide:  /ppt/comments/comment{N}.xml   (N = 1-based slide number)
  Authors:    /ppt/commentAuthors.xml         (presentation-level)

Namespace: http://schemas.openxmlformats.org/presentationml/2006/main (p:)
"""
from __future__ import annotations

import logging
import os
import re
from datetime import datetime, timezone

from lxml import etree as ET

from pptmcp._presentation_governance import ValidationError
from pptmcp.presentation_pptx import (
    _atomic_save,
    _audit_log,
    _check_confirm,
    _check_path,
    _check_write,
    _evict_prs,
    _load_prs,
)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_MAX_COMMENTS_DEFAULT = 200
_AUTHOR_MAX = 64
_DEFAULT_X_EMU = 914400    # 1 inch
_DEFAULT_Y_EMU = 685800    # 0.75 inch
_CTRL_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")
_EMU_MAX = 2**31 - 1

_CM_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments"
_AUTHORS_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/commentAuthors"
_CM_CT = "application/vnd.openxmlformats-officedocument.presentationml.comments+xml"
_AUTHORS_CT = "application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml"

# Modern threaded comments (Microsoft 365 / PowerPoint 2018+)
_MODERN_CM_REL = "http://schemas.microsoft.com/office/2018/10/relationships/comments"
_MODERN_NS = "http://schemas.microsoft.com/office/powerpoint/2018/8/main"
_DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _max_comments() -> int:
    try:
        return max(0, int(os.getenv("PPT_MAX_COMMENTS", str(_MAX_COMMENTS_DEFAULT))))
    except (ValueError, TypeError):
        return _MAX_COMMENTS_DEFAULT


def _p(tag: str) -> str:
    return f"{{{_NS}}}{tag}"


def _p188(tag: str) -> str:
    return f"{{{_MODERN_NS}}}{tag}"


def _a(tag: str) -> str:
    return f"{{{_DML_NS}}}{tag}"


def _get_slide_modern_comments_xml(slide):
    """Return parsed p188:cmLst element for modern threaded comments, or None if absent."""
    try:
        part = slide.part.part_related_by(_MODERN_CM_REL)
        return ET.fromstring(part.blob)
    except KeyError:
        return None


def _extract_modern_text(elem) -> str:
    """Extract concatenated a:t text from a p188:cm or p188:reply txBody."""
    parts = []
    tx_body = elem.find(_p188("txBody"))
    if tx_body is None:
        return ""
    for t_elem in tx_body.iter(_a("t")):
        if t_elem.text:
            parts.append(t_elem.text)
    return "".join(parts)


def _sanitize_text(value: str) -> str:
    return _CTRL_RE.sub("", value)


def _sanitize_author(value: str) -> str:
    # Strip control chars first, then truncate — never truncate into a control char
    cleaned = _CTRL_RE.sub("", value)
    if len(cleaned) > _AUTHOR_MAX:
        logging.getLogger(__name__).warning("author name truncated to %d chars", _AUTHOR_MAX)
        cleaned = cleaned[:_AUTHOR_MAX]
    return cleaned


def _validate_emu(value: int, name: str) -> None:
    if not (0 <= value <= _EMU_MAX):
        raise ValidationError(f"{name} must be in range 0–{_EMU_MAX}, got {value}")


def _get_slide_number(prs, slide_index: int) -> int:
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValidationError(
            f"slide_index {slide_index} out of range (0–{len(prs.slides) - 1})"
        )
    return slide_index + 1


# ---------------------------------------------------------------------------
# Authors part helpers
# ---------------------------------------------------------------------------

def _get_authors_xml(prs):
    """Return parsed commentAuthors XML element, or None if part absent."""
    try:
        part = prs.part.part_related_by(_AUTHORS_REL)
        return ET.fromstring(part.blob)
    except KeyError:
        return None


def _get_or_create_authors_part(prs):
    """Return (xml_element, part). Creates the part if absent."""
    try:
        part = prs.part.part_related_by(_AUTHORS_REL)
        return ET.fromstring(part.blob), part
    except KeyError:
        pass

    from pptx.opc.package import Part, PackURI  # noqa: PLC0415

    root = ET.Element(_p("cmAuthorLst"), nsmap={None: _NS})
    blob = ET.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
    part = Part(PackURI("/ppt/commentAuthors.xml"), _AUTHORS_CT, prs.part.package, blob=blob)
    prs.part.relate_to(part, _AUTHORS_REL)
    return root, part


def _resolve_author_id(authors_xml, author_id: int):
    """Map authorId → (name, resolved_bool). Falls back to f'author_{id}'."""
    if authors_xml is None:
        return f"author_{author_id}", False
    for elem in authors_xml.findall(_p("cmAuthor")):
        try:
            if int(elem.get("id", -1)) == author_id:
                return elem.get("name", f"author_{author_id}"), True
        except (ValueError, TypeError):
            continue
    return f"author_{author_id}", False


def _get_or_add_author(authors_xml, author_name: str) -> int:
    """Return existing authorId for author_name or create a new entry."""
    for elem in authors_xml.findall(_p("cmAuthor")):
        if elem.get("name") == author_name:
            try:
                return int(elem.get("id", 0))
            except (ValueError, TypeError):
                pass  # malformed id on existing author — fall through to allocate new id
    existing_ids = []
    for elem in authors_xml.findall(_p("cmAuthor")):
        try:
            existing_ids.append(int(elem.get("id", 0)))
        except (ValueError, TypeError):
            pass
    new_id = max(existing_ids, default=-1) + 1
    new_elem = ET.SubElement(authors_xml, _p("cmAuthor"))
    new_elem.set("id", str(new_id))
    new_elem.set("name", author_name)
    new_elem.set("initials", (author_name[:2].upper() if len(author_name) >= 2 else author_name.upper()))
    new_elem.set("lastIdx", "0")
    new_elem.set("clrIdx", str(new_id % 8))
    return new_id


def _next_idx_for_author(authors_xml, author_id: int) -> int:
    """Return next comment idx for author and increment lastIdx."""
    for elem in authors_xml.findall(_p("cmAuthor")):
        try:
            if int(elem.get("id", -1)) == author_id:
                last = int(elem.get("lastIdx", "0"))
                elem.set("lastIdx", str(last + 1))
                return last + 1
        except (ValueError, TypeError):
            continue
    return 1


# ---------------------------------------------------------------------------
# Comments part helpers
# ---------------------------------------------------------------------------

def _get_slide_comments_xml(slide):
    """Return parsed cmLst XML element for this slide, or None if absent."""
    try:
        part = slide.part.part_related_by(_CM_REL)
        return ET.fromstring(part.blob)
    except KeyError:
        return None


def _get_or_create_slide_comments_part(slide, slide_number: int):
    """Return (xml_element, part). Creates the part if absent."""
    try:
        part = slide.part.part_related_by(_CM_REL)
        return ET.fromstring(part.blob), part
    except KeyError:
        pass

    from pptx.opc.package import Part, PackURI  # noqa: PLC0415

    root = ET.Element(_p("cmLst"), nsmap={None: _NS})
    blob = ET.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
    part_uri = f"/ppt/comments/comment{slide_number}.xml"
    part = Part(PackURI(part_uri), _CM_CT, slide.part.package, blob=blob)
    slide.part.relate_to(part, _CM_REL)
    return root, part


def _write_part_blob(part, xml_root) -> None:
    part._blob = ET.tostring(xml_root, xml_declaration=True, encoding="UTF-8", standalone=True)


def _extract_text(cm) -> str:
    parts = []
    for t_elem in cm.iter(_p("t")):
        if t_elem.text:
            parts.append(t_elem.text)
    return "".join(parts)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def list_comments(path: str, slide_index: int | None = None) -> dict:
    """List comments on one slide (or all slides if slide_index is None).

    Returns:
        {comments, total_found, truncated, max_comments}
    """
    resolved = _check_path(path)
    prs = _load_prs(resolved)
    limit = _max_comments()
    authors_xml = _get_authors_xml(prs)

    if slide_index is not None:
        _get_slide_number(prs, slide_index)
        slide_range = [slide_index]
    else:
        slide_range = list(range(len(prs.slides)))

    raw: list[dict] = []
    for si in slide_range:
        slide = prs.slides[si]

        # Legacy comments (OOXML 2006 format)
        cm_xml = _get_slide_comments_xml(slide)
        if cm_xml is not None:
            for cm in cm_xml.findall(_p("cm")):
                try:
                    author_id = int(cm.get("authorId", 0))
                except (ValueError, TypeError):
                    author_id = 0
                try:
                    cm_idx = int(cm.get("idx", 0))
                except (ValueError, TypeError):
                    cm_idx = 0
                author_name, author_resolved = _resolve_author_id(authors_xml, author_id)
                pos = cm.find(_p("pos"))
                raw.append({
                    "comment_id": cm_idx,
                    "author_id": author_id,
                    "author": author_name,
                    "author_resolved": author_resolved,
                    "slide_index": si,
                    "text": _extract_text(cm),
                    "dt": cm.get("dt", ""),
                    "pos_x_emu": int(pos.get("x", 0)) if pos is not None else None,
                    "pos_y_emu": int(pos.get("y", 0)) if pos is not None else None,
                    "format": "legacy",
                })

        # Modern threaded comments (Microsoft 365 / PowerPoint 2018+ format)
        modern_xml = _get_slide_modern_comments_xml(slide)
        if modern_xml is not None:
            for cm in modern_xml.findall(_p188("cm")):
                cm_guid = cm.get("id", "")
                author_guid = cm.get("authorId", "")
                dt = cm.get("created", "")
                text = _extract_modern_text(cm)
                raw.append({
                    "comment_id": cm_guid,
                    "author_id": author_guid,
                    "author": author_guid,
                    "author_resolved": False,
                    "slide_index": si,
                    "text": text,
                    "dt": dt,
                    "pos_x_emu": None,
                    "pos_y_emu": None,
                    "format": "modern",
                })
                # Replies are child comments in the same thread
                reply_lst = cm.find(_p188("replyLst"))
                if reply_lst is not None:
                    for reply in reply_lst.findall(_p188("reply")):
                        raw.append({
                            "comment_id": reply.get("id", ""),
                            "author_id": reply.get("authorId", ""),
                            "author": reply.get("authorId", ""),
                            "author_resolved": False,
                            "slide_index": si,
                            "text": _extract_modern_text(reply),
                            "dt": reply.get("created", ""),
                            "pos_x_emu": None,
                            "pos_y_emu": None,
                            "format": "modern_reply",
                            "parent_comment_id": cm_guid,
                        })

    total_found = len(raw)
    truncated = total_found > limit
    return {
        "comments": raw[:limit],
        "total_found": total_found,
        "truncated": truncated,
        "max_comments": limit,
    }


def add_comment(
    path: str,
    slide_index: int,
    text: str,
    author: str = "MCP",
    slide_x_emu: int = _DEFAULT_X_EMU,
    slide_y_emu: int = _DEFAULT_Y_EMU,
    confirm: bool = False,
) -> dict:
    """Add a comment to a slide. Requires PPT_ENABLE_WRITE=true + confirm=True."""
    # B5: gate order — write env first, confirm second, path third
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)

    # B3: EMU coordinate guard
    _validate_emu(slide_x_emu, "slide_x_emu")
    _validate_emu(slide_y_emu, "slide_y_emu")

    prs = _load_prs(resolved)
    slide_number = _get_slide_number(prs, slide_index)
    slide = prs.slides[slide_index]

    # B2: sanitise inputs before any XML write
    clean_text = _sanitize_text(text)
    clean_author = _sanitize_author(author)

    authors_xml, authors_part = _get_or_create_authors_part(prs)
    author_id = _get_or_add_author(authors_xml, clean_author)
    new_idx = _next_idx_for_author(authors_xml, author_id)

    cm_xml, cm_part = _get_or_create_slide_comments_part(slide, slide_number)

    dt_str = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")

    # B2: use lxml element API — never string templating for XML content
    cm_elem = ET.SubElement(cm_xml, _p("cm"))
    cm_elem.set("authorId", str(author_id))
    cm_elem.set("dt", dt_str)
    cm_elem.set("idx", str(new_idx))

    pos_elem = ET.SubElement(cm_elem, _p("pos"))
    pos_elem.set("x", str(slide_x_emu))
    pos_elem.set("y", str(slide_y_emu))

    r_elem = ET.SubElement(cm_elem, _p("r"))
    t_elem = ET.SubElement(r_elem, _p("t"))
    t_elem.text = clean_text  # lxml escapes XML metacharacters automatically

    _write_part_blob(cm_part, cm_xml)
    _write_part_blob(authors_part, authors_xml)

    _atomic_save(prs, resolved)
    _evict_prs(resolved)
    _audit_log("manage_comments_add", resolved, slide_idx=slide_index)

    return {
        "status": "added",
        "comment_id": new_idx,
        "author": clean_author,
        "slide_index": slide_index,
    }


def delete_comment(
    path: str,
    slide_index: int,
    comment_id: int,
    confirm: bool = False,
) -> dict:
    """Delete a comment by idx from a slide. Requires PPT_ENABLE_WRITE=true + confirm=True."""
    # B5: gate order
    _check_write()
    _check_confirm(confirm)
    resolved = _check_path(path)

    prs = _load_prs(resolved)
    _get_slide_number(prs, slide_index)
    slide = prs.slides[slide_index]

    cm_xml = _get_slide_comments_xml(slide)
    if cm_xml is None:
        raise ValidationError(f"comment_id {comment_id} not found on slide {slide_index}")

    target = None
    for cm in cm_xml.findall(_p("cm")):
        try:
            if int(cm.get("idx", -1)) == comment_id:
                target = cm
                break
        except (ValueError, TypeError):
            continue

    if target is None:
        raise ValidationError(f"comment_id {comment_id} not found on slide {slide_index}")

    cm_xml.remove(target)

    try:
        cm_part = slide.part.part_related_by(_CM_REL)
    except KeyError:
        raise ValidationError(f"comment_id {comment_id} not found on slide {slide_index}")

    _write_part_blob(cm_part, cm_xml)
    _atomic_save(prs, resolved)
    _evict_prs(resolved)
    _audit_log("manage_comments_delete", resolved, slide_idx=slide_index)

    return {
        "status": "deleted",
        "comment_id": comment_id,
        "slide_index": slide_index,
    }
