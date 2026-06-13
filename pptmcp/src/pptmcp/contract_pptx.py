"""
Contract validation layer for pptmcp Phase A: Output Contract.
No FastMCP / MCP imports — pure python-pptx + stdlib only.
"""
from __future__ import annotations

import hashlib
import json
import logging
import os
from pathlib import Path
from typing import Any

from pptx import Presentation

from pptmcp.presentation_pptx import (
    ValidationError,
    _check_confirm,
    _check_path,
    _check_write,
)

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

CONTRACT_MAX_BYTES: int = 512 * 1024  # 524_288 bytes — far exceeds any real contract

_SUPPORTED_CONTRACT_VERSIONS: frozenset[str] = frozenset({"1.0"})

_EMU_PER_INCH: int = 914_400            # English Metric Units per inch (python-pptx standard)
_DIM_TOLERANCE: float = 0.01            # inches — float roundtrip tolerance for dimensions

# W5 Amendment A1: contract tool names are now declared directly in
# presentation_pptx.capabilities() to keep the canonical tool registry in one place.
# This list is intentionally empty; server.py appends it (no-op) for backward compat.
CONTRACT_TOOLS: list[dict] = []

# JSON Schema v1.0 constant — pure dict, serialised on demand by get_contract_schema().
# DC-06: this is a module-level constant; get_contract_schema() performs NO file I/O.
_CONTRACT_SCHEMA_V1: dict[str, Any] = {
    "$schema": "http://json-schema.org/draft-07/schema#",
    "title": "pptmcp Output Contract v1.0",
    "type": "object",
    "required": ["contract_version", "slides"],
    "additionalProperties": True,
    "properties": {
        "contract_version": {
            "type": "string",
            "enum": ["1.0"],
            "description": "Contract format version. Currently only '1.0' is supported.",
        },
        "slide_width_inches": {
            "type": "number",
            "minimum": 0,
            "description": "Expected slide width in inches (e.g. 13.33 for 16:9 widescreen).",
        },
        "slide_height_inches": {
            "type": "number",
            "minimum": 0,
            "description": "Expected slide height in inches (e.g. 7.5 for 16:9 widescreen).",
        },
        "slides": {
            "type": "array",
            "description": (
                "Per-slide acceptance criteria. "
                "Empty array means no per-slide checks are performed."
            ),
            "items": {
                "type": "object",
                "required": ["slide_index"],
                "additionalProperties": True,
                "properties": {
                    "slide_index": {
                        "type": "integer",
                        "minimum": 0,
                        "description": "0-based slide index.",
                    },
                    "min_shapes": {
                        "type": "integer",
                        "minimum": 0,
                        "description": "Minimum number of shapes required on this slide.",
                    },
                    "max_shapes": {
                        "type": "integer",
                        "minimum": 0,
                        "description": "Maximum number of shapes allowed on this slide.",
                    },
                    "expected_title": {
                        "type": "string",
                        "description": (
                            "Expected slide title (case-insensitive substring match). "
                            "Fails if slide has no title shape or title does not contain this string."
                        ),
                    },
                },
            },
        },
    },
}


# ---------------------------------------------------------------------------
# Private path helper
# ---------------------------------------------------------------------------

def _check_json_path(path: str) -> Path:
    """Validate a .json path against PPT_ALLOWLIST_ROOTS.

    Guard ordering is security-critical (BLK-04). Do not reorder.
    """
    # Step 1 — empty allowlist guard (fail-closed; mirrors _check_path)
    roots_env = os.getenv("PPT_ALLOWLIST_ROOTS", "").strip()
    if not roots_env:
        raise ValidationError("No allowlist configured")

    # Step 2 — null-byte guard (BLK-01; must precede Path.resolve())
    if "\x00" in str(path):
        raise ValidationError("Invalid path: null byte")

    # Step 3 — resolve symlinks / relative components
    resolved = Path(path).resolve()

    # Step 4 — extension check
    ext = resolved.suffix.lower()
    if ext != ".json":
        raise ValidationError(f"Unsupported extension: {ext!r}")

    # Step 5 — existence check (UX quality: clear error before json.load; WARN-06)
    if not resolved.exists():
        raise ValidationError("Contract file not found")

    # Step 6 — file size cap (BLK-02; DoS guard; after exists() so stat() is safe)
    fstat = resolved.stat()
    if fstat.st_size > CONTRACT_MAX_BYTES:
        raise ValidationError(
            f"Contract file exceeds size limit "
            f"({fstat.st_size} bytes > {CONTRACT_MAX_BYTES})"
        )

    # Step 7 — allowlist loop (same pattern as _check_path)
    roots = [r.strip() for r in roots_env.split(",") if r.strip()]
    for root in roots:
        try:
            if resolved.is_relative_to(Path(root).resolve()):
                return resolved
        except (ValueError, TypeError):
            continue

    raise ValidationError("Path not in allowlist")


# ---------------------------------------------------------------------------
# Public API — schema
# ---------------------------------------------------------------------------

def get_contract_schema() -> str:
    """Return the JSON Schema for a pptmcp Output Contract v1.0 file.

    Pure constant serialisation — performs no file I/O.
    Safe to call without PPT_ALLOWLIST_ROOTS being configured (DC-06).
    """
    return json.dumps(_CONTRACT_SCHEMA_V1, indent=2)


# ---------------------------------------------------------------------------
# Public API — validation
# ---------------------------------------------------------------------------

def validate_contract(contract_path: str) -> dict:
    """Validate a JSON Output Contract file against the pptmcp v1.0 schema.

    Returns a dict with keys:
      valid (bool)            — True iff no schema errors found
      contract_version (str)  — value from file, or None if absent
      slide_count (int)       — number of slide spec objects in the contract
      errors (list[str])      — [] if valid; list of human-readable messages if not

    Raises ValidationError for path/IO errors:
      - allowlist denial, file not found, wrong extension, too large, non-JSON content,
        missing required field 'contract_version', bool in numeric field (BLK-03)
    Schema errors (wrong type excluding bool, wrong version, missing optional field)
      return {"valid": False, "errors": [...]}.
    """
    # Path and IO validation (raises on any failure)
    resolved = _check_json_path(contract_path)

    try:
        raw = resolved.read_text(encoding="utf-8")
        data = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise ValidationError(f"Invalid JSON: {exc.msg}") from exc
    except OSError as exc:
        raise ValidationError(f"Cannot read contract file: {exc.strerror}") from exc

    if not isinstance(data, dict):
        raise ValidationError("Contract must be a JSON object")

    errors: list[str] = []

    # --- contract_version — REQUIRED (ODQ-01: missing raises, not returns) ---
    cv = data.get("contract_version")
    if cv is None:
        raise ValidationError("Missing required field: contract_version")
    if isinstance(cv, bool):
        raise ValidationError(
            "Field 'contract_version' must be a string, got bool"
        )
    if not isinstance(cv, str):
        raise ValidationError(
            f"Field 'contract_version' must be a string, got {type(cv).__name__!r}"
        )
    # Version must be in the supported set; unrecognised → schema error (return)
    if cv not in _SUPPORTED_CONTRACT_VERSIONS:
        errors.append(
            f"Unsupported contract_version: {cv!r}. "
            f"Supported: {sorted(_SUPPORTED_CONTRACT_VERSIONS)}"
        )

    # --- slide_width_inches / slide_height_inches (optional floats) ---
    for dim_field in ("slide_width_inches", "slide_height_inches"):
        v = data.get(dim_field)
        if v is None:
            continue
        if isinstance(v, bool):
            # BLK-03: bool-as-float is a type confusion attack — raise
            raise ValidationError(
                f"Field '{dim_field}' must be a number, got bool (type confusion)"
            )
        if not isinstance(v, (int, float)):
            errors.append(
                f"Field '{dim_field}' must be a number, got {type(v).__name__!r}"
            )

    # --- slides (required array) ---
    raw_slides = data.get("slides")
    if raw_slides is None:
        errors.append("Missing required field: slides")
        raw_slides = []
    elif not isinstance(raw_slides, list):
        errors.append(
            f"Field 'slides' must be an array, got {type(raw_slides).__name__!r}"
        )
        raw_slides = []

    for i, spec in enumerate(raw_slides):
        if not isinstance(spec, dict):
            errors.append(f"slides[{i}] must be an object")
            continue

        # slide_index — required int
        si = spec.get("slide_index")
        if si is None:
            errors.append(f"slides[{i}] missing required field: slide_index")
        elif isinstance(si, bool):
            # BLK-03: raise
            raise ValidationError(
                f"slides[{i}].slide_index must be an integer, got bool (type confusion)"
            )
        elif not isinstance(si, int):
            errors.append(
                f"slides[{i}].slide_index must be an integer, got {type(si).__name__!r}"
            )

        # min_shapes — optional int
        if "min_shapes" in spec:
            v = spec["min_shapes"]
            if isinstance(v, bool):
                raise ValidationError(
                    f"slides[{i}].min_shapes must be an integer, got bool (type confusion)"
                )
            if not isinstance(v, int):
                errors.append(
                    f"slides[{i}].min_shapes must be an integer, got {type(v).__name__!r}"
                )

        # max_shapes — optional int
        if "max_shapes" in spec:
            v = spec["max_shapes"]
            if isinstance(v, bool):
                raise ValidationError(
                    f"slides[{i}].max_shapes must be an integer, got bool (type confusion)"
                )
            if not isinstance(v, int):
                errors.append(
                    f"slides[{i}].max_shapes must be an integer, got {type(v).__name__!r}"
                )

        # expected_title — optional str
        if "expected_title" in spec:
            v = spec["expected_title"]
            if not isinstance(v, str):
                errors.append(
                    f"slides[{i}].expected_title must be a string, got {type(v).__name__!r}"
                )

    if errors:
        return {
            "valid": False,
            "contract_version": cv,
            "slide_count": len(raw_slides),
            "errors": errors,
        }

    return {
        "valid": True,
        "contract_version": cv,
        "slide_count": len(raw_slides),
        "errors": [],
    }


def declare_slide_contract(
    slide_path: str,
    spec: dict[str, Any],
    confirm: bool = False,
) -> dict:
    """Persist a slide-export contract beside the PPTX (RCI-US-01).

    The provided ``spec`` is augmented with required-field defaults
    (``contract_version='1.0'``, ``slides=[]``) when missing, then written
    to ``contract.json`` in the same directory as ``slide_path``. The
    resulting file is validated against the v1.0 schema before the call
    returns — a schema mismatch rolls back the write so the disk is never
    left holding an invalid contract.

    Spec keys recognised by the validator: ``contract_version``,
    ``slide_width_inches``, ``slide_height_inches``, ``slides`` (each with
    ``slide_index`` + optional ``min_shapes``/``max_shapes``/
    ``expected_title``). Additional keys (e.g. ``element_inventory``,
    ``geometry_tokens``, ``color_palette``, ``export_plan``,
    ``acceptance_thresholds``) are accepted as-is via the schema's
    ``additionalProperties: true`` gate.

    Gates:
      1. ``_check_path(slide_path)``  — .pptx + PPT_ALLOWLIST_ROOTS
      2. ``_check_write()``           — PPT_ENABLE_WRITE=true
      3. ``_check_confirm(confirm)``  — confirm=True

    Returns:
        ``{"contract_path": str, "checksum_sha256": str, "contract_version": str}``

    Raises:
        NotAllowedError:  PPT_ENABLE_WRITE not set.
        ValidationError:  Path/spec/schema failure, confirm=False.
    """
    resolved_pptx = _check_path(slide_path)
    _check_write()
    _check_confirm(confirm)

    if not isinstance(spec, dict):
        raise ValidationError("spec must be a JSON object")

    payload: dict[str, Any] = dict(spec)
    payload.setdefault("contract_version", "1.0")
    payload.setdefault("slides", [])

    contract_path = resolved_pptx.parent / "contract.json"
    json_text = json.dumps(payload, indent=2, default=str, sort_keys=True)

    if len(json_text.encode("utf-8")) > CONTRACT_MAX_BYTES:
        raise ValidationError(
            f"Contract payload exceeds size limit "
            f"({CONTRACT_MAX_BYTES} bytes)"
        )

    # Write then validate; on failure remove the file so disk is never left
    # holding a contract that fails its own schema. Use write_bytes so the
    # checksum we return is exactly the on-disk bytes (Windows write_text
    # would translate LF→CRLF and invalidate the hash).
    json_bytes = json_text.encode("utf-8")
    contract_path.write_bytes(json_bytes)
    try:
        report = validate_contract(str(contract_path))
    except ValidationError:
        contract_path.unlink(missing_ok=True)
        raise
    if not report.get("valid"):
        contract_path.unlink(missing_ok=True)
        raise ValidationError(
            "spec failed v1.0 contract schema: " + "; ".join(report.get("errors", []))
        )

    checksum = hashlib.sha256(json_bytes).hexdigest()
    logger.info(
        "declare_slide_contract: wrote %s (sha256=%s)",
        contract_path, checksum,
    )
    return {
        "contract_path": str(contract_path),
        "checksum_sha256": checksum,
        "contract_version": payload["contract_version"],
    }


def check_presentation_against_contract(path: str, contract_path: str) -> dict:
    """Check a presentation against its Output Contract specification.

    Returns a dict with keys:
      passed (bool)           — True iff no error-severity findings
      slide_count (int)       — number of slides in the deck
      findings (list[dict])   — [] if fully passing; each entry has
                                 slide_index (int | None), criterion (str),
                                 detail (str), and optionally severity='warning'

    Raises ValidationError for any path/IO access error on either file.
    Per-slide contract violations are reported as findings, not exceptions.
    """
    # DC-02 Step 1+2: validate BOTH paths before any file I/O (fail-fast access control)
    resolved_pptx = _check_path(path)
    resolved_json = _check_json_path(contract_path)  # includes size cap (BLK-02)

    # DC-02 Step 3-4: load contract JSON (size already validated by _check_json_path)
    try:
        raw = resolved_json.read_text(encoding="utf-8")
        contract_data = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise ValidationError(f"Invalid JSON in contract: {exc.msg}") from exc
    except OSError as exc:
        raise ValidationError(f"Cannot read contract file: {exc.strerror}") from exc

    if not isinstance(contract_data, dict):
        raise ValidationError("Contract must be a JSON object")

    # DC-02 Step 5: open presentation
    try:
        presentation = Presentation(str(resolved_pptx))
    except Exception as exc:                         # python-pptx raises various types
        raise ValidationError(f"Cannot open presentation: {exc}") from exc

    # DC-02 Step 6: run checks, build findings
    findings: list[dict] = []
    slide_count = len(presentation.slides)

    # --- Deck-level: slide dimension checks ---
    expected_w = contract_data.get("slide_width_inches")
    if (
        expected_w is not None
        and not isinstance(expected_w, bool)
        and isinstance(expected_w, (int, float))
    ):
        actual_w = presentation.slide_width / _EMU_PER_INCH
        if abs(actual_w - expected_w) > _DIM_TOLERANCE:
            findings.append({
                "slide_index": None,
                "criterion": "slide_width_inches",
                "detail": f"Expected {expected_w:.3f}\", got {actual_w:.3f}\"",
            })

    expected_h = contract_data.get("slide_height_inches")
    if (
        expected_h is not None
        and not isinstance(expected_h, bool)
        and isinstance(expected_h, (int, float))
    ):
        actual_h = presentation.slide_height / _EMU_PER_INCH
        if abs(actual_h - expected_h) > _DIM_TOLERANCE:
            findings.append({
                "slide_index": None,
                "criterion": "slide_height_inches",
                "detail": f"Expected {expected_h:.3f}\", got {actual_h:.3f}\"",
            })

    # --- Per-slide checks ---
    contract_slides = contract_data.get("slides", [])
    if not isinstance(contract_slides, list):
        contract_slides = []

    # Collect all specified slide indices for DC-04 coverage check
    specified_indices: set[int] = set()

    for spec in contract_slides:
        if not isinstance(spec, dict):
            continue

        si = spec.get("slide_index")
        if si is None or isinstance(si, bool) or not isinstance(si, int):
            continue    # malformed spec — cannot check; validate_contract catches this

        specified_indices.add(si)

        # Slide beyond deck length
        if si >= slide_count:
            findings.append({
                "slide_index": si,
                "criterion": "slide_exists",
                "detail": (
                    f"Contract references slide_index={si} but deck has "
                    f"only {slide_count} slide(s)"
                ),
            })
            continue

        slide = presentation.slides[si]
        shape_count = len(slide.shapes)

        # min_shapes
        min_s = spec.get("min_shapes")
        if (
            min_s is not None
            and not isinstance(min_s, bool)
            and isinstance(min_s, int)
            and shape_count < min_s
        ):
            findings.append({
                "slide_index": si,
                "criterion": "min_shapes",
                "detail": f"Expected >= {min_s} shapes, found {shape_count}",
            })

        # max_shapes
        max_s = spec.get("max_shapes")
        if (
            max_s is not None
            and not isinstance(max_s, bool)
            and isinstance(max_s, int)
            and shape_count > max_s
        ):
            findings.append({
                "slide_index": si,
                "criterion": "max_shapes",
                "detail": f"Expected <= {max_s} shapes, found {shape_count}",
            })

        # expected_title (case-insensitive substring match)
        expected_title = spec.get("expected_title")
        if expected_title is not None and isinstance(expected_title, str):
            title_shape = slide.shapes.title
            if title_shape is None or not title_shape.has_text_frame:
                findings.append({
                    "slide_index": si,
                    "criterion": "expected_title",
                    "detail": (
                        f"No title shape found; expected_title={expected_title!r}"
                    ),
                })
            else:
                actual_title = title_shape.text_frame.text.strip()
                if expected_title.lower() not in actual_title.lower():
                    findings.append({
                        "slide_index": si,
                        "criterion": "expected_title",
                        "detail": (
                            f"Title does not contain expected_title={expected_title!r} "
                            f"(case-insensitive)"
                        ),
                    })

    # DC-04: WARN for slides present in deck but not covered by contract.
    # Only fires when contract specifies at least one slide (partial coverage scenario).
    # Empty slides list → pass silently (no warnings generated; see W4 Q4).
    if contract_slides:
        for deck_idx in range(slide_count):
            if deck_idx not in specified_indices:
                findings.append({
                    "slide_index": deck_idx,
                    "criterion": "unspecified_slide",
                    "detail": (
                        f"Slide {deck_idx} in deck is not covered by the contract"
                    ),
                    "severity": "warning",
                })

    # DC-02 Step 7: return result
    # passed=True iff all findings are severity=warning (no error-grade violations)
    error_findings = [f for f in findings if f.get("severity") != "warning"]
    passed = len(error_findings) == 0

    return {
        "passed": passed,
        "slide_count": slide_count,
        "findings": findings,
    }
