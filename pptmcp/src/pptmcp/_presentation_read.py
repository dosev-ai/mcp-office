"""
Read-only operations for pptmcp: read_presentation, read_slide, list_slides,
list_shapes, get_shape, get_presentation_metadata, list_layouts,
read_speaker_notes.

No FastMCP / MCP imports — pure python-pptx + stdlib only.
"""
from __future__ import annotations

import logging
from typing import Any

from pptx.util import Emu

from pptmcp._presentation_cache import _load_prs
from pptmcp._presentation_governance import ValidationError

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Internal slide helpers (used by read operations and re-exported by facade)
# ---------------------------------------------------------------------------


def _slide_title(slide: Any) -> str | None:
    try:
        text = slide.shapes.title.text if slide.shapes.title else None
        return text or None
    except Exception:
        return None


def _slide_notes(slide: Any) -> str | None:
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        text = tf.text.strip()
        return text if text else None
    except Exception:
        return None


def _emu_to_inches(emu: int | None) -> float | None:
    """Convert EMU to inches (4 dp); returns None if emu is None."""
    return None if emu is None else round(float(Emu(emu).inches), 4)


def _shape_fill_color_hex(shape: Any) -> str | None:
    """Return a shape's solid fill RGB hex when it is safely discoverable."""
    try:
        rgb = shape.fill.fore_color.rgb
    except Exception:
        return None
    return str(rgb).upper() if rgb is not None else None


# ---------------------------------------------------------------------------
# Public read functions
# Each function imports _check_path from presentation_pptx at call time so
# that tests patching pptmcp.presentation_pptx._check_path take effect.
# ---------------------------------------------------------------------------


def read_presentation(path: str) -> dict:
    """Return slide count and per-slide summary for the presentation at path."""
    from pptmcp.presentation_pptx import _check_path  # noqa: PLC0415

    resolved = _check_path(path)
    prs = _load_prs(resolved)
    slides = []
    for i, slide in enumerate(prs.slides):
        slides.append({
            "slide_index": i,
            "title": _slide_title(slide),
            "shapes_count": len(slide.shapes),
            "has_notes": _slide_notes(slide) is not None,
        })
    return {"slide_count": len(prs.slides), "slides": slides}


def get_presentation_metadata(path: str) -> dict:
    """Return title, author, slide count, created/modified dates."""
    from pptmcp.presentation_pptx import _check_path  # noqa: PLC0415

    resolved = _check_path(path)
    prs = _load_prs(resolved)
    cp = prs.core_properties

    def _iso(dt: Any) -> str | None:
        return dt.isoformat() if dt else None

    return {
        "title": cp.title,
        "author": cp.author,
        "subject": cp.subject,
        "keywords": cp.keywords,
        "slide_count": len(prs.slides),
        "created": _iso(cp.created),
        "modified": _iso(cp.modified),
    }


def list_slides(path: str) -> list[dict]:
    """Return index, title, layout and has_notes for every slide."""
    from pptmcp.presentation_pptx import _check_path  # noqa: PLC0415

    resolved = _check_path(path)
    prs = _load_prs(resolved)
    result = []
    for i, slide in enumerate(prs.slides):
        result.append({
            "slide_index": i,
            "title": _slide_title(slide),
            "layout": slide.slide_layout.name,
            "has_notes": _slide_notes(slide) is not None,
        })
    return result


def read_slide(path: str, slide_index: int) -> dict:
    """Return full shape list and notes_text for a single slide by index."""
    from pptmcp.presentation_pptx import _check_path  # noqa: PLC0415

    resolved = _check_path(path)
    prs = _load_prs(resolved)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValidationError(
            f"slide_index {slide_index} out of range (0-{len(prs.slides) - 1})"
        )
    slide = prs.slides[slide_index]
    shapes = []
    for shape in slide.shapes:
        ph_idx: int | None = None
        if shape.is_placeholder and shape.placeholder_format is not None:
            ph_idx = shape.placeholder_format.idx
        shapes.append({
            "shape_id": shape.shape_id,
            "name": shape.name,
            "shape_type": str(shape.shape_type),
            "has_text": shape.has_text_frame,
            "text": shape.text_frame.text if shape.has_text_frame else None,
            "placeholder_idx": ph_idx,
        })
    return {
        "slide_index": slide_index,
        "title": _slide_title(slide),
        "shapes": shapes,
        "notes_text": _slide_notes(slide),
    }


def read_speaker_notes(path: str, slide_index: int | None = None) -> list[dict]:
    """Return speaker notes for one slide or all slides (None means all).

    Every slide in the requested range is always included in the result.
    Slides with no notes have ``notes_text`` set to ``None``.
    """
    from pptmcp.presentation_pptx import _check_path  # noqa: PLC0415

    resolved = _check_path(path)
    prs = _load_prs(resolved)
    indices: list[int] = (
        [slide_index] if slide_index is not None else list(range(len(prs.slides)))
    )
    result = []
    for i in indices:
        if i < 0 or i >= len(prs.slides):
            raise ValidationError(f"slide_index {i} out of range")
        slide = prs.slides[i]
        notes = _slide_notes(slide)
        result.append({"slide_index": i, "notes_text": notes})
    return result


def list_shapes(path: str, slide_index: int) -> list[dict]:
    """Return all shapes on a slide with type, position, and size."""
    from pptmcp.presentation_pptx import _check_path  # noqa: PLC0415

    resolved = _check_path(path)
    prs = _load_prs(resolved)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValidationError(f"slide_index {slide_index} out of range")
    slide = prs.slides[slide_index]
    result = []
    for shape in slide.shapes:
        ph_idx: int | None = None
        if shape.is_placeholder and shape.placeholder_format is not None:
            ph_idx = shape.placeholder_format.idx
        result.append({
            "shape_id": shape.shape_id,
            "name": shape.name,
            "shape_type": str(shape.shape_type),
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "has_text_frame": shape.has_text_frame,
            "has_table": shape.has_table,
            "has_chart": shape.has_chart,
            "fill_color_hex": _shape_fill_color_hex(shape),
            "placeholder_idx": ph_idx,
            "bounds": {
                "left_in":   _emu_to_inches(shape.left),
                "top_in":    _emu_to_inches(shape.top),
                "width_in":  _emu_to_inches(shape.width),
                "height_in": _emu_to_inches(shape.height),
            },
        })
    return result


def get_shape(path: str, slide_index: int, shape_id: int) -> dict:
    """Return full detail for a specific shape by shape_id on a slide."""
    from pptmcp.presentation_pptx import _check_path  # noqa: PLC0415

    resolved = _check_path(path)
    prs = _load_prs(resolved)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValidationError(f"slide_index {slide_index} out of range")
    slide = prs.slides[slide_index]
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            result: dict[str, Any] = {
                "shape_id": shape.shape_id,
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "has_text_frame": shape.has_text_frame,
                "has_table": shape.has_table,
                "has_chart": shape.has_chart,
                "fill_color_hex": _shape_fill_color_hex(shape),
            }
            if shape.has_text_frame:
                result["text"] = shape.text_frame.text
            if shape.has_table:
                result["rows"] = [
                    {"height_emu": row.height}
                    for row in shape.table.rows
                ]
            return result
    raise ValidationError(f"shape_id {shape_id} not found on slide {slide_index}")


def list_layouts(path: str) -> list[dict]:
    """Return all slide layouts with index, name, and placeholder info."""
    from pptmcp.presentation_pptx import _check_path  # noqa: PLC0415

    resolved = _check_path(path)
    prs = _load_prs(resolved)
    result = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholder_types = [
            {
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
            }
            for ph in layout.placeholders
        ]
        result.append({
            "index": i,
            "name": layout.name,
            "placeholder_count": len(layout.placeholders),
            "placeholder_types": placeholder_types,
        })
    return result
