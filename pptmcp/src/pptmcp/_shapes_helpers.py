"""_shapes_helpers — shared helper functions for pptmcp shape operations."""
from __future__ import annotations

import logging
import math
import re
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt

from pptmcp.presentation_pptx import ValidationError, _hex_to_rgb

_log = logging.getLogger(__name__)

_EMPTY_PLACEHOLDER_PATTERNS: frozenset[str] = frozenset({
    "click to add title",
    "click to add text",
    "click to add subtitle",
    "click to edit master title style",
    "click to edit master text styles",
    "click to edit master subtitle style",
    "click to add notes",
})


def _calculate_overflow_risk(text: str, width_in: float, height_in: float) -> dict:
    """Estimate overflow risk (heuristic — NOT an accurate render prediction).

    Uses a character-density approximation. Results are estimates only —
    actual rendering by PowerPoint may differ due to font metrics, kerning,
    and line-break algorithm.
    """
    if not text:
        return {"overflow_risk": "low", "overflow_detail": "No text content"}
    chars_per_line = max(1, int(width_in * 13))
    line_height_in = 0.25
    estimated_lines = math.ceil(len(text) / chars_per_line)
    estimated_height_in = estimated_lines * line_height_in
    ratio = estimated_height_in / max(height_in, 0.01)
    if ratio < 0.5:
        return {
            "overflow_risk": "low",
            "overflow_detail": f"~{estimated_lines} lines estimated, {height_in:.2f}in capacity",
        }
    if ratio < 1.0:
        return {
            "overflow_risk": "medium",
            "overflow_detail": f"~{estimated_lines} lines estimated, {height_in:.2f}in capacity",
        }
    return {
        "overflow_risk": "high",
        "overflow_detail": f"~{estimated_lines} lines estimated, {height_in:.2f}in capacity (risk: text may be clipped)",
    }


def _first_overlap_for_shape(slide: Any, target_shape_id: int) -> dict | None:
    """Return first overlap pair involving target_shape_id, or None. None bounds skipped."""
    shapes_list = list(slide.shapes)
    target = next((s for s in shapes_list if s.shape_id == target_shape_id), None)
    if target is None:
        return None
    if any(v is None for v in [target.left, target.top, target.width, target.height]):
        return None
    for other in shapes_list:
        if other.shape_id == target_shape_id:
            continue
        if any(v is None for v in [other.left, other.top, other.width, other.height]):
            continue
        ox = max(0, min(target.left + target.width, other.left + other.width) - max(target.left, other.left))
        oy = max(0, min(target.top + target.height, other.top + other.height) - max(target.top, other.top))
        if ox * oy > 0:
            return {
                "shape_a": target_shape_id,
                "shape_b": other.shape_id,
                "overlap_area_in2": round(ox * oy / (914400 ** 2), 4),
            }
    return None


def _normalize_color_hex(color_hex: str) -> tuple[int, int, int]:
    """Validate a 6-digit RGB hex color with or without a leading #."""
    normalized = color_hex.strip()
    if normalized.startswith("#"):
        normalized = normalized[1:]
    if not re.match(r"^[0-9A-Fa-f]{6}$", normalized):
        raise ValidationError(
            f"Invalid hex color {color_hex!r}. Must be 6 hex chars, e.g. 'FF0000'."
        )
    return _hex_to_rgb("#" + normalized)


def _apply_cell_text_style(
    cell: Any,
    *,
    font_size_pt: float | None = None,
    bold: bool | None = None,
    color_hex: str | None = None,
) -> int:
    """Apply text style to every run in a table cell and return updated run count."""
    if font_size_pt is not None and font_size_pt <= 0:
        raise ValidationError("font_size_pt must be greater than 0")
    rgb = RGBColor(*_normalize_color_hex(color_hex)) if color_hex is not None else None
    runs_updated = 0
    for para in cell.text_frame.paragraphs:
        runs = list(para.runs)
        if not runs and para.text:
            runs = [para.add_run()]
            runs[0].text = para.text
        for run in runs:
            if font_size_pt is not None:
                run.font.size = Pt(font_size_pt)
            if bold is not None:
                run.font.bold = bold
            if rgb is not None:
                run.font.color.rgb = rgb
            runs_updated += 1
    return runs_updated


def _apply_cell_fill(cell: Any, fill_color_hex: str | None) -> None:
    if fill_color_hex is None:
        return
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(*_normalize_color_hex(fill_color_hex))


def _is_title_placeholder(shape: Any) -> bool:
    try:
        return shape.placeholder_format.type in {
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
            PP_PLACEHOLDER.VERTICAL_TITLE,
        }
    except Exception:
        return False


def _is_empty_placeholder(shape: Any) -> bool:
    if not shape.is_placeholder:
        return False
    if not shape.has_text_frame:
        return True
    text = shape.text_frame.text.strip().lower()
    return text == "" or text in _EMPTY_PLACEHOLDER_PATTERNS


def _shapes_overlap(shape_a: Any, shape_b: Any) -> bool:
    if any(v is None for v in [shape_a.left, shape_a.top, shape_a.width, shape_a.height,
                               shape_b.left, shape_b.top, shape_b.width, shape_b.height]):
        return False
    overlap_x = max(
        0,
        min(shape_a.left + shape_a.width, shape_b.left + shape_b.width)
        - max(shape_a.left, shape_b.left),
    )
    overlap_y = max(
        0,
        min(shape_a.top + shape_a.height, shape_b.top + shape_b.height)
        - max(shape_a.top, shape_b.top),
    )
    return overlap_x * overlap_y > 0


def _remove_empty_overlapping_content_placeholders(slide: Any, target_shape: Any) -> int:
    placeholders = [
        shape
        for shape in slide.shapes
        if shape.shape_id != target_shape.shape_id
        and _is_empty_placeholder(shape)
        and not _is_title_placeholder(shape)
        and _shapes_overlap(shape, target_shape)
    ]
    for placeholder in placeholders:
        placeholder._element.getparent().remove(placeholder._element)
    return len(placeholders)
