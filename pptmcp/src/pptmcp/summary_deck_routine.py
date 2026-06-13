"""Summary deck creation routine — python-pptx-based, COM-free.

Public entrypoint: create_summary_deck_v1

Layer: no COM, no MCP imports. Standard library + python-pptx only.
Write gate: PPT_ENABLE_WRITE env var must be "true" (case-insensitive).
Allowlist: PPT_ALLOWLIST_ROOTS (comma-separated directory prefixes).
"""
from __future__ import annotations

import os
from typing import Any

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

_MAX_TABLE_ROWS = 50
_TABLE_LEFT_IN = 0.5
_TABLE_TOP_IN = 1.5
_TABLE_WIDTH_IN = 9.0
_ROW_HEIGHT_IN = 0.4

# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _check_write_gate() -> dict | None:
    """Return error dict if PPT_ENABLE_WRITE is not 'true', else None."""
    val = os.environ.get("PPT_ENABLE_WRITE", "").strip().lower()
    if val != "true":
        return {"error": "PPT_ENABLE_WRITE not set to true", "error_kind": "gate_error"}
    return None


def _check_path_allowlist(presentation_path: str) -> dict | None:
    """Validate presentation_path against PPT_ALLOWLIST_ROOTS.

    Returns error dict if validation fails, None if OK.
    Uses os.path.realpath() to resolve symlinks before comparison.
    """
    roots_env = os.environ.get("PPT_ALLOWLIST_ROOTS", "").strip()
    if not roots_env:
        return {"error": "PPT_ALLOWLIST_ROOTS not configured", "error_kind": "config_error"}

    roots = [r.strip() for r in roots_env.split(",") if r.strip()]
    if not roots:
        return {"error": "PPT_ALLOWLIST_ROOTS not configured", "error_kind": "config_error"}

    real_path = os.path.normcase(os.path.realpath(presentation_path))
    for root in roots:
        real_root = os.path.normcase(os.path.realpath(root))
        if not real_root.endswith(os.sep):
            real_root = real_root + os.sep
        if real_path.startswith(real_root) or real_path == real_root.rstrip(os.sep):
            return None  # allowed

    return {"error": "path not in allowlist", "error_kind": "allowlist_violation"}


# ---------------------------------------------------------------------------
# Public function
# ---------------------------------------------------------------------------


def create_summary_deck_v1(
    presentation_path: str,
    title: str,
    rows: list[list] | None = None,
    headers: list[str] | None = None,
    source_metadata: dict | None = None,
    source_result: dict | None = None,
) -> dict[str, Any]:
    """Create a 2-slide PowerPoint summary deck from tabular data.

    Slide 1: title slide (layout 0) with title and optional subtitle.
    Slide 2: blank slide (layout 5) with a table of headers + rows.

    Uses python-pptx only — no COM, no win32com.

    Requires PPT_ENABLE_WRITE=true and path within PPT_ALLOWLIST_ROOTS.

    Args:
        presentation_path: Absolute path for the output .pptx file.
        title:             Title text for slide 1.
        rows:              Data rows (list of lists).
        headers:           Column headers (first row of table).
        source_metadata:   Optional dict; "description" key used as subtitle.

    Returns:
        On success:
            {
                "presentation_path": str,
                "slide_count": int,
                "evidence": {
                    "routine": "pptmcp.create_summary_deck_v1",
                    "row_count": int,
                    "headers": list[str],
                },
            }
        On error:
            {"error": str, "error_kind": str}
    """
    # Extract rows/headers from a prior governed_routine result dict if provided
    if source_result is not None:
        if rows is None:
            rows = source_result.get("rows", [])
        if headers is None:
            headers = source_result.get("headers", [])
    if rows is None:
        rows = []
    if headers is None:
        headers = []

    # --- write gate FIRST ---
    gate_err = _check_write_gate()
    if gate_err is not None:
        return gate_err

    try:
        # --- allowlist gate ---
        err = _check_path_allowlist(presentation_path)
        if err is not None:
            return err

        from pptx import Presentation  # noqa: PLC0415
        from pptx.util import Inches  # noqa: PLC0415

        prs = Presentation()

        # ---- Slide 1: title slide ----
        title_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(title_layout)

        # Set title placeholder
        if slide1.shapes.title is not None:
            slide1.shapes.title.text = title

        # Set subtitle placeholder (placeholder index 1)
        description = ""
        if source_metadata:
            description = source_metadata.get("description", "")
        placeholders = list(slide1.placeholders)
        if len(placeholders) > 1:
            try:
                placeholders[1].text = description
            except Exception:
                pass  # subtitle placeholder may not be writable on some layouts

        # ---- Slide 2: table slide (blank layout) ----
        blank_layout = prs.slide_layouts[5]
        slide2 = prs.slides.add_slide(blank_layout)

        # Truncate rows if needed
        truncated = False
        display_rows = rows
        if len(rows) > _MAX_TABLE_ROWS:
            display_rows = rows[:_MAX_TABLE_ROWS]
            truncated = True

        num_cols = max(len(headers), 1)
        num_data_rows = len(display_rows)
        total_table_rows = 1 + num_data_rows  # header + data

        if truncated:
            total_table_rows += 1  # extra row for truncation note

        left = Inches(_TABLE_LEFT_IN)
        top = Inches(_TABLE_TOP_IN)
        width = Inches(_TABLE_WIDTH_IN)
        height = Inches(_ROW_HEIGHT_IN * total_table_rows)

        table_shape = slide2.shapes.add_table(
            total_table_rows, num_cols, left, top, width, height
        )
        tbl = table_shape.table

        # Header row
        for col_idx, header_text in enumerate(headers):
            if col_idx < num_cols:
                tbl.cell(0, col_idx).text = str(header_text) if header_text is not None else ""

        # Data rows
        for row_idx, row_data in enumerate(display_rows, start=1):
            for col_idx in range(num_cols):
                val = row_data[col_idx] if col_idx < len(row_data) else ""
                tbl.cell(row_idx, col_idx).text = str(val) if val is not None else ""

        # Truncation notice row
        if truncated:
            notice_row = 1 + num_data_rows
            tbl.cell(notice_row, 0).text = (
                f"[Truncated: showing {_MAX_TABLE_ROWS} of {len(rows)} rows]"
            )

        prs.save(presentation_path)

        return {
            "presentation_path": presentation_path,
            "slide_count": len(prs.slides),
            "evidence": {
                "routine": "pptmcp.create_summary_deck_v1",
                "row_count": len(rows),
                "headers": headers,
            },
        }

    except Exception as exc:  # noqa: BLE001
        return {"error": str(exc), "error_kind": "unexpected_error"}


__all__ = ["create_summary_deck_v1"]
