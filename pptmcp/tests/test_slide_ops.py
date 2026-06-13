"""Slide operation tests — targets presentation_pptx.py only."""

import pytest
from pptx import Presentation as _Prs
from pptmcp.presentation_pptx import (
    ValidationError,
    NotAllowedError,
    create_presentation,
    add_slide,
    copy_slide,
)


# ── GROUP 21: create_presentation() ──────────────────────────────────────

@pytest.mark.unit
def test_create_presentation_success(tmp_path, monkeypatch):
    """Create a new presentation with a title slide; returns correct dict."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "new_deck.pptx")
    result = create_presentation(path, title="Test", confirm=True)
    assert "path" in result
    assert result["slide_count"] >= 1
    assert result["title"] == "Test"


@pytest.mark.unit
def test_create_presentation_blocked_write(tmp_path, monkeypatch):
    """create_presentation raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "should_not_create.pptx")
    with pytest.raises(NotAllowedError):
        create_presentation(path, title="X", confirm=True)


@pytest.mark.unit
def test_create_presentation_no_confirm(tmp_path, monkeypatch):
    """create_presentation raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "no_confirm.pptx")
    with pytest.raises(ValidationError):
        create_presentation(path, confirm=False)


@pytest.mark.unit
def test_create_presentation_path_not_in_allowlist(tmp_path, monkeypatch):
    """create_presentation raises ValidationError for path outside allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        create_presentation(path, confirm=True)


@pytest.mark.unit
def test_create_presentation_file_already_exists(tmp_path, monkeypatch):
    """create_presentation raises ValidationError if file already exists."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    existing = tmp_path / "existing.pptx"
    existing.write_bytes(b"dummy content")
    with pytest.raises(ValidationError, match="already exists"):
        create_presentation(str(existing), confirm=True)


# ── GROUP 29 fragment: create_presentation blank slide ──────────────────

@pytest.mark.unit
def test_create_presentation_no_title_blank(tmp_path, monkeypatch):
    """create_presentation with title=None creates a blank deck (0 slides)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "blank.pptx")
    result = create_presentation(path, title=None, confirm=True)
    assert result["slide_count"] == 0
    assert result["title"] is None


@pytest.mark.unit
def test_add_slide_suppress_content_placeholder(tmp_path, monkeypatch):
    """add_slide can remove empty non-title placeholders from content layouts."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "suppress_placeholder.pptx"
    p = _Prs()
    p.save(str(path))

    result = add_slide(
        str(path),
        layout_index=1,
        title="Table Slide",
        suppress_content_placeholder=True,
        confirm=True,
    )

    assert result["placeholders_removed"] >= 1
    loaded = _Prs(str(path))
    slide = loaded.slides[result["slide_index"]]
    non_title_placeholders = [
        shape
        for shape in slide.shapes
        if shape.is_placeholder
        and getattr(shape.placeholder_format, "idx", None) != 0
    ]
    assert non_title_placeholders == []


# ── GROUP 28: copy_slide() ────────────────────────────────────────────────

@pytest.mark.unit
def test_copy_slide_same_file_success(tmp_path, monkeypatch):
    """copy_slide within the same file returns new_slide_index >= 1 with source/target paths."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "copy_src.pptx"
    p = _Prs()
    layout = p.slide_layouts[1]
    slide = p.slides.add_slide(layout)
    slide.shapes.title.text = "Original Slide"
    p.save(str(path))
    result = copy_slide(str(path), 0, str(path), confirm=True)
    assert result["new_slide_index"] >= 1
    assert "source_path" in result
    assert "target_path" in result
    assert result["source_slide_index"] == 0


@pytest.mark.unit
def test_copy_slide_invalid_source_slide_index(tmp_path, monkeypatch):
    """copy_slide raises ValidationError when source_slide_index is out of range."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "copy_oob.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="source_slide"):
        copy_slide(str(path), 99, str(path), confirm=True)


@pytest.mark.unit
def test_copy_slide_source_path_not_allowed(tmp_path, monkeypatch):
    """copy_slide raises ValidationError when source_path is outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    evil_src = str(tmp_path.parent / "evil_src.pptx")
    valid_target = str(tmp_path / "target.pptx")
    with pytest.raises(ValidationError, match="[Nn]ot in allowlist|not in allowlist"):
        copy_slide(evil_src, 0, valid_target, confirm=True)


@pytest.mark.unit
def test_copy_slide_target_path_not_allowed(tmp_path, monkeypatch):
    """copy_slide raises ValidationError when target_path is outside the allowlist (OWASP A03)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    valid_src = str(tmp_path / "source.pptx")
    evil_target = str(tmp_path.parent / "evil_target.pptx")
    with pytest.raises(ValidationError, match="[Nn]ot in allowlist|not in allowlist"):
        copy_slide(valid_src, 0, evil_target, confirm=True)


@pytest.mark.unit
def test_copy_slide_blocked_write(tmp_path, monkeypatch):
    """copy_slide raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(NotAllowedError):
        copy_slide(path, 0, path, confirm=True)


@pytest.mark.unit
def test_copy_slide_no_confirm(tmp_path, monkeypatch):
    """copy_slide raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        copy_slide(path, 0, path, confirm=False)


# ── post-COM extras: copy_slide cross-file ───────────────────────────────

@pytest.mark.unit
def test_copy_slide_cross_file_success(tmp_path, monkeypatch):
    """copy_slide from source to a different target appends the slide and reports both paths."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    src = tmp_path / "cs_src.pptx"
    tgt = tmp_path / "cs_tgt.pptx"
    p_src = _Prs()
    slide = p_src.slides.add_slide(p_src.slide_layouts[1])
    slide.shapes.title.text = "Source Slide"
    p_src.save(str(src))
    p_tgt = _Prs()
    p_tgt.slides.add_slide(p_tgt.slide_layouts[5])
    p_tgt.save(str(tgt))
    result = copy_slide(str(src), 0, str(tgt), confirm=True)
    assert result["source_slide_index"] == 0
    assert result["new_slide_index"] >= 1
    assert "cs_src.pptx" in result["source_path"]
    assert "cs_tgt.pptx" in result["target_path"]


@pytest.mark.unit
def test_copy_slide_cross_file_target_not_found(tmp_path, monkeypatch):
    """copy_slide raises ValidationError when the cross-file target does not exist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    src = tmp_path / "cs_src2.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(src))
    nonexistent_tgt = str(tmp_path / "does_not_exist.pptx")
    with pytest.raises(ValidationError, match="does not exist"):
        copy_slide(str(src), 0, nonexistent_tgt, confirm=True)


@pytest.mark.unit
def test_copy_slide_invalid_target_slide_index(tmp_path, monkeypatch):
    """copy_slide raises ValidationError when target_slide_index is out of range."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "cs_tgt_oob.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="target_slide"):
        copy_slide(str(path), 0, str(path), target_slide_index=99, confirm=True)


@pytest.mark.unit
def test_copy_slide_target_slide_index_prepend(tmp_path, monkeypatch):
    """copy_slide with target_slide_index=0 positions the copy at slide 0."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "cs_prepend.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = copy_slide(str(path), 0, str(path), target_slide_index=0, confirm=True)
    assert result["new_slide_index"] == 0
