"""Tests for compound dispatcher tools: add_content, set_format, slide, shape, export."""
import sys
import os

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptmcp.server import add_content, set_format, slide, shape
from pptmcp._server_handlers_read import export
from fastmcp.exceptions import ToolError


class TestAddContentDispatcher:
    """Guard paths for add_content validation."""

    def test_content_type_shape_without_shape_type_raises_tool_error(self, tmp_path, write_env):
        pptx_path = str(tmp_path / "dummy.pptx")
        from pptx import Presentation
        Presentation().save(pptx_path)
        with pytest.raises(ToolError, match="shape_type"):
            add_content(
                content_type="shape",
                path=pptx_path,
                slide_index=0,
                left=1.0, top=1.0, width=2.0, height=1.0,
                shape_type="",
                confirm=True,
            )

    def test_content_type_table_with_zero_rows_raises_tool_error(self, tmp_path, write_env):
        pptx_path = str(tmp_path / "dummy.pptx")
        from pptx import Presentation
        Presentation().save(pptx_path)
        with pytest.raises(ToolError, match="rows and cols"):
            add_content(
                content_type="table",
                path=pptx_path,
                slide_index=0,
                left=1.0, top=1.0, width=4.0, height=2.0,
                rows=0, cols=3,
                confirm=True,
            )

    def test_unknown_content_type_raises_tool_error(self, tmp_path, write_env):
        pptx_path = str(tmp_path / "dummy.pptx")
        from pptx import Presentation
        Presentation().save(pptx_path)
        with pytest.raises(ToolError, match="Unknown content_type"):
            add_content(
                content_type="xyzzy",
                path=pptx_path,
                slide_index=0,
                left=1.0, top=1.0, width=2.0, height=1.0,
                confirm=True,
            )

class TestSetFormatDispatcher:
    """Guard paths for set_format validation."""

    def test_unknown_target_raises_tool_error(self, tmp_path, write_env):
        pptx_path = str(tmp_path / "dummy.pptx")
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide_obj = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide_obj.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = "test"
        prs.save(pptx_path)
        with pytest.raises(ToolError, match="Unknown target"):
            set_format(
                target="xyzzy",
                path=pptx_path,
                slide_index=0,
                shape_id=txBox.shape_id,
                confirm=True,
            )


class TestSlideDispatcher:
    """Tests for slide() unified dispatch tool (D1)."""

    def _make_pptx(self, tmp_path, monkeypatch, n_slides=1):
        from pptx import Presentation
        pptx_path = str(tmp_path / "t.pptx")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        prs = Presentation()
        for _ in range(n_slides):
            prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(pptx_path)
        return pptx_path

    def test_slide_unknown_operation_raises_tool_error(self, tmp_path, write_env, monkeypatch):
        """slide() with unknown operation raises ToolError with clean message."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        with pytest.raises(ToolError, match="Unknown operation"):
            slide(operation="xyzzy", path=pptx_path, confirm=True)

    def test_slide_delete_requires_slide_index(self, tmp_path, write_env, monkeypatch):
        """slide(operation=delete) without slide_index raises ToolError."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        with pytest.raises(ToolError):
            slide(operation="delete", path=pptx_path, confirm=True)

    def test_slide_reorder_requires_new_order(self, tmp_path, write_env, monkeypatch):
        """slide(operation=reorder) without new_order raises ToolError."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch, n_slides=2)
        with pytest.raises(ToolError):
            slide(operation="reorder", path=pptx_path, confirm=True)

    def test_slide_add_returns_slide_index(self, tmp_path, write_env, monkeypatch):
        """slide(operation=add) returns dict with slide_index key."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        result = slide(operation="add", path=pptx_path, confirm=True)
        assert isinstance(result, dict)
        assert "slide_index" in result

    def test_slide_add_suppresses_content_placeholder_by_default(self, tmp_path, write_env, monkeypatch):
        """slide(operation=add) defaults to suppressing empty content placeholders."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        result = slide(
            operation="add",
            path=pptx_path,
            layout_index=1,
            title="Content slide",
            confirm=True,
        )
        assert result["slide_index"] == 1
        assert result["placeholders_removed"] >= 1

    def test_slide_reorder_permutation(self, tmp_path, write_env, monkeypatch):
        """slide(operation=reorder, new_order=[1,0]) returns new_order in result."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch, n_slides=2)
        result = slide(operation="reorder", path=pptx_path, new_order=[1, 0], confirm=True)
        assert isinstance(result, dict)
        assert result.get("new_order") == [1, 0]

    def test_slide_delete_success(self, tmp_path, write_env, monkeypatch):
        """slide(operation=delete, slide_index=0) returns a result dict."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch, n_slides=2)
        result = slide(operation="delete", path=pptx_path, slide_index=0, confirm=True)
        assert result is not None

    # P8: slide() copy path tests
    def test_slide_copy_within_deck_appends_slide(self, tmp_path, write_env, monkeypatch):
        """slide(operation=copy) within the same deck appends a slide and returns new index."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch, n_slides=2)
        result = slide(
            operation="copy",
            path=pptx_path,
            source_slide_index=0,
            confirm=True,
        )
        assert isinstance(result, dict)
        assert "new_slide_index" in result
        assert result["new_slide_index"] >= 0
        # The deck now has 3 slides; new_slide_index must be 0-based within range.
        assert result["new_slide_index"] <= 2

    def test_slide_copy_missing_source_slide_index_raises_tool_error(
        self, tmp_path, write_env, monkeypatch
    ):
        """slide(operation=copy) without source_slide_index or slide_index raises ToolError."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch, n_slides=2)
        with pytest.raises(ToolError):
            slide(operation="copy", path=pptx_path, confirm=True)

class TestShapeDispatcher:
    """Tests for shape() unified dispatch tool (D2)."""

    def _make_pptx(self, tmp_path, monkeypatch):
        from pptx import Presentation
        pptx_path = str(tmp_path / "s.pptx")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(pptx_path)
        return pptx_path

    def test_shape_unknown_operation_raises_tool_error(self, tmp_path, write_env, monkeypatch):
        """shape() with unknown operation raises ToolError."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        with pytest.raises(ToolError, match="Unknown operation"):
            shape(operation="xyzzy", path=pptx_path, slide_index=0, confirm=True)

    def test_shape_delete_requires_shape_id(self, tmp_path, write_env, monkeypatch):
        """shape(operation=delete) without shape_id raises ToolError."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        with pytest.raises(ToolError):
            shape(operation="delete", path=pptx_path, slide_index=0, confirm=True)

    def test_shape_set_properties_requires_shape_id(self, tmp_path, write_env, monkeypatch):
        """shape(operation=set_properties) without shape_id raises ToolError."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        with pytest.raises(ToolError):
            shape(operation="set_properties", path=pptx_path, slide_index=0, confirm=True)

    def test_shape_add_text_box_success(self, tmp_path, write_env, monkeypatch):
        """shape(operation=add_text_box) returns dict with shape_id."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        result = shape(
            operation="add_text_box",
            path=pptx_path,
            slide_index=0,
            left=1.0, top=1.0, width=2.0, height=1.0,
            confirm=True,
        )
        assert isinstance(result, dict)
        assert "shape_id" in result

    def test_shape_add_autoshape_empty_shape_type_raises(self, tmp_path, write_env, monkeypatch):
        """shape(operation=add_autoshape, shape_type=) raises ToolError about shape_type."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        with pytest.raises(ToolError, match="shape_type"):
            shape(
                operation="add_autoshape",
                path=pptx_path,
                slide_index=0,
                left=1.0, top=1.0, width=2.0, height=1.0,
                shape_type="",
                confirm=True,
            )

    def test_shape_add_table_zero_rows_raises(self, tmp_path, write_env, monkeypatch):
        """shape(operation=add_table, rows=0) raises ToolError about rows and cols."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        with pytest.raises(ToolError, match="rows and cols"):
            shape(
                operation="add_table",
                path=pptx_path,
                slide_index=0,
                left=1.0, top=1.0, width=4.0, height=2.0,
                rows=0, cols=3,
                confirm=True,
            )

    def test_shape_set_properties_fill_is_applied_and_readable(self, tmp_path, write_env, monkeypatch):
        """shape(operation=set_properties, target=fill) applies a solid fill visible via get_shape."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        added = shape(
            operation="add_autoshape",
            path=pptx_path,
            slide_index=0,
            shape_type="RECTANGLE",
            left=1.0,
            top=1.0,
            width=2.0,
            height=1.0,
            confirm=True,
        )
        result = shape(
            operation="set_properties",
            target="fill",
            path=pptx_path,
            slide_index=0,
            shape_id=added["shape_id"],
            fill_color_hex="00AEEF",
            confirm=True,
        )
        assert result["fill_updated"] is True
        assert result["fill_color_hex"] == "00AEEF"

        from pptmcp.presentation_pptx import get_shape

        detail = get_shape(pptx_path, 0, added["shape_id"])
        assert detail["fill_color_hex"] == "00AEEF"

class TestExportDispatcherBackend:
    """CI-safe tests for export backend functions (no COM required)."""

    def _make_pptx(self, tmp_path, monkeypatch):
        from pptx import Presentation
        pptx_path = str(tmp_path / "e.pptx")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(pptx_path)
        return pptx_path

    def test_export_slide_text_backend_returns_list(self, tmp_path, monkeypatch):
        """content_pptx.export_slide_as_text returns a list of per-slide dicts."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        from pptmcp import content_pptx as content
        result = content.export_slide_as_text(pptx_path, 0)
        assert isinstance(result, list)

    def test_export_slide_images_backend_returns_dict(self, tmp_path, monkeypatch):
        """content_pptx.extract_images returns dict with images key."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        from pptmcp import content_pptx as content
        result = content.extract_images(pptx_path)
        assert isinstance(result, dict)
        assert "images" in result

    def test_export_slide_png_no_com_env_raises_or_returns_error(self, tmp_path, monkeypatch):
        """_server_handlers_com_tools.export_slide without COM raises ValueError or returns error dict."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        from pptmcp import _server_handlers_com_tools as _ct
        # Without PPT_ENABLE_COM=true the COM guard raises ValueError.
        # On non-Windows CI the import itself may raise ImportError.
        # Either outcome is acceptable; unexpected exception types propagate.
        # On Windows the COM guard raises; on non-Windows CI the function returns
        # an error dict (success=False) rather than raising — both are acceptable.
        try:
            result = _ct.export_slide(
                file_path=pptx_path,
                slide_index=0,
                fmt="PNG",
                output_path=str(tmp_path / "out.png"),
            )
            assert isinstance(result, dict) and not result.get("success", True), (
                f"Expected error dict from export_slide without COM; got {result!r}"
            )
        except (ValueError, ImportError, PermissionError, ToolError):
            pass  # COM guard raised as expected

    def test_export_slide_scope_inline_uses_shared_contract(self, tmp_path, monkeypatch):
        """export(scope=slide, return_inline=True) packages PNG via shared artifact contract."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        out = tmp_path / "slide.png"

        def fake_export_slide(**kwargs):
            out.write_bytes(b"\x89PNG\r\n\x1a\nmock")
            return {
                "exported": True,
                "format": kwargs["fmt"],
                "slide_number": kwargs["slide_index"] + 1,
                "output_path": kwargs["output_path"],
            }

        monkeypatch.setattr(
            "pptmcp._server_handlers_com_tools.export_slide",
            fake_export_slide,
        )
        result = export(
            scope="slide",
            path=pptx_path,
            fmt="png",
            output_path=str(out),
            return_inline=True,
            options={"slide_index": 0, "inline_dpi": 96},
            confirm=True,
        )

        assert result["exported"] is True
        assert result["artifact"]["artifact_type"] == "slide_render"
        assert result["artifact"]["mime_type"] == "image/png"
        assert result["artifact"]["inline"] is True
        assert result["artifact"]["delivery"] == "base64_json"
        assert result["artifact"]["metadata"]["slide_index"] == 0
        assert result["artifact"]["metadata"]["dpi"] == 96

    def test_export_slide_scope_inline_disabled_omits_payload(self, tmp_path, monkeypatch):
        """Default return_inline=False preserves export and omits payload."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        out = tmp_path / "slide.png"

        def fake_export_slide(**kwargs):
            out.write_bytes(b"\x89PNG\r\n\x1a\nmock")
            return {
                "exported": True,
                "format": kwargs["fmt"],
                "slide_number": kwargs["slide_index"] + 1,
                "output_path": kwargs["output_path"],
            }

        monkeypatch.setattr(
            "pptmcp._server_handlers_com_tools.export_slide",
            fake_export_slide,
        )
        result = export(
            scope="slide",
            path=pptx_path,
            fmt="png",
            output_path=str(out),
            options={"slide_index": 0},
            confirm=True,
        )

        assert result["exported"] is True
        assert result["artifact"]["inline"] is False
        assert result["artifact"]["omitted_reason"] == "inline disabled"
        assert "image_base64" not in result["artifact"]

    def test_export_slide_scope_oversize_omits_without_failing_export(self, tmp_path, monkeypatch):
        """Oversize inline payload returns omission metadata while export remains successful."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        out = tmp_path / "slide.png"

        def fake_export_slide(**kwargs):
            out.write_bytes(b"abcdef")
            return {
                "exported": True,
                "format": kwargs["fmt"],
                "slide_number": kwargs["slide_index"] + 1,
                "output_path": kwargs["output_path"],
            }

        monkeypatch.setattr(
            "pptmcp._server_handlers_com_tools.export_slide",
            fake_export_slide,
        )
        result = export(
            scope="slide",
            path=pptx_path,
            fmt="png",
            output_path=str(out),
            return_inline=True,
            max_inline_bytes=4,
            options={"slide_index": 0},
            confirm=True,
        )

        assert result["exported"] is True
        assert result["artifact"]["inline"] is False
        assert result["artifact"]["omitted_reason"] == "encoded image exceeds max_inline_bytes"

    def test_export_slide_scope_rejects_unsupported_options(self, tmp_path, monkeypatch):
        """Consolidated export options are tightly allowlisted."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        with pytest.raises(ToolError, match="Unsupported options"):
            export(
                scope="slide",
                path=pptx_path,
                output_path=str(tmp_path / "slide.png"),
                options={"slide_index": 0, "range": "A1:B2"},
                confirm=True,
            )

    def test_export_slide_scope_rejects_non_png_fmt(self, tmp_path, monkeypatch):
        """pptmcp inline slide render supports PNG only."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        with pytest.raises(ToolError, match="fmt='png'"):
            export(
                scope="slide",
                path=pptx_path,
                fmt="pdf",
                output_path=str(tmp_path / "slide.pdf"),
                options={"slide_index": 0},
                confirm=True,
            )

    def test_export_slide_png_scope_accepts_slide_index_in_options(
        self, tmp_path, monkeypatch
    ):
        """Regression for scope='slide_png' must accept
        slide_index inside the options dict for parity with scope='slide'.
        """
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        out = tmp_path / "slidepng.png"

        captured = {}

        def fake_export_slide(**kwargs):
            captured.update(kwargs)
            out.write_bytes(b"\x89PNG\r\n\x1a\nmock")
            return {
                "exported": True,
                "format": kwargs["fmt"],
                "slide_number": kwargs["slide_index"] + 1,
                "output_path": kwargs["output_path"],
            }

        monkeypatch.setattr(
            "pptmcp._server_handlers_com_tools.export_slide",
            fake_export_slide,
        )

        result = export(
            scope="slide_png",
            path=pptx_path,
            output_path=str(out),
            options={"slide_index": 2, "inline_dpi": 144},
            confirm=True,
        )

        assert result.get("exported") is True
        assert captured["slide_index"] == 2
        assert captured["dpi"] == 144
        assert captured["fmt"] == "png"

    def test_export_slide_png_scope_still_rejects_unknown_options(
        self, tmp_path, monkeypatch
    ):
        """Allowlist for scope='slide_png' is the same as scope='slide':
        only slide_index + inline_dpi are accepted; unknown keys still raise.
        """
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        with pytest.raises(ToolError, match="Unsupported options"):
            export(
                scope="slide_png",
                path=pptx_path,
                output_path=str(tmp_path / "slidepng.png"),
                options={"slide_index": 0, "range": "A1"},
                confirm=True,
            )

    # P9: export() unit tests for read-only scopes (no COM required)
    def test_export_scope_slide_text_returns_list_without_com(self, tmp_path, monkeypatch):
        """export(scope='slide_text') returns list of per-slide dicts — no COM needed."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        result = export(scope="slide_text", path=pptx_path)
        assert isinstance(result, list)
        assert len(result) >= 1
        assert "slide_index" in result[0]
        assert "texts" in result[0]

    def test_export_scope_slide_text_single_slide_index(self, tmp_path, monkeypatch):
        """export(scope='slide_text', slide_index=0) restricts result to that slide."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        result = export(scope="slide_text", path=pptx_path, slide_index=0)
        assert isinstance(result, list)
        assert result[0]["slide_index"] == 0

    def test_export_scope_slide_images_returns_dict_without_com(self, tmp_path, monkeypatch):
        """export(scope='slide_images') returns dict with images key — no COM needed."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        result = export(scope="slide_images", path=pptx_path)
        assert isinstance(result, dict)
        assert "images" in result
        assert isinstance(result["images"], list)

    def test_export_unknown_scope_raises_tool_error(self, tmp_path, monkeypatch):
        """export() with an unknown scope raises ToolError with descriptive message."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        with pytest.raises(ToolError, match="Unknown scope"):
            export(scope="not_a_scope", path=pptx_path)

    def test_export_scope_slide_text_rejects_extra_options(self, tmp_path, monkeypatch):
        """export(scope='slide_text') rejects any options dict entries."""
        pptx_path = self._make_pptx(tmp_path, monkeypatch)
        with pytest.raises(ToolError, match="Unsupported options"):
            export(scope="slide_text", path=pptx_path, options={"inline_dpi": 96})
