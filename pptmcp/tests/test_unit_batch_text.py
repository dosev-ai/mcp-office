"""
Unit tests for pptmcp batch_set_text backend and handler.
No COM required — CI-safe.
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptmcp.presentation_pptx import NotAllowedError, ValidationError


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_pptx_with_textboxes(path, num_boxes: int = 2) -> tuple[str, list[int]]:
    """Create a 1-slide pptx with `num_boxes` textboxes; return (path_str, shape_ids)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    shape_ids = []
    for i in range(num_boxes):
        tb = slide.shapes.add_textbox(
            Inches(i + 0.5), Inches(0.5), Inches(2.0), Inches(1.0)
        )
        tb.text_frame.text = f"original text {i}"
        shape_ids.append(tb.shape_id)
    prs.save(str(path))
    return str(path), shape_ids


# ---------------------------------------------------------------------------
# Class 1: Gate tests
# ---------------------------------------------------------------------------

class TestBatchSetTextGates:
    def test_write_gate_blocks_without_env(self, tmp_path, monkeypatch):
        """PPT_ENABLE_WRITE not set -> NotAllowedError."""
        monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        path = str(tmp_path / "x.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        with pytest.raises(NotAllowedError):
            batch_set_text(path, 0, [{"shape_id": 1, "text": "hi"}], confirm=True)

    def test_confirm_gate_blocks_without_confirm(self, tmp_path, monkeypatch):
        """confirm=False -> ValidationError even when write env is set."""
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        path = str(tmp_path / "x.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        with pytest.raises(ValidationError, match="confirm"):
            batch_set_text(path, 0, [{"shape_id": 1, "text": "hi"}], confirm=False)

    def test_path_outside_allowlist_raises(self, tmp_path, monkeypatch):
        """Path not in allowlist -> ValidationError."""
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        allowlist_dir = tmp_path / "allowed"
        allowlist_dir.mkdir()
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(allowlist_dir))
        # File is in tmp_path, not in allowlist_dir
        outside_path = str(tmp_path / "outside.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        with pytest.raises(ValidationError, match="not in allowlist"):
            batch_set_text(outside_path, 0, [{"shape_id": 1, "text": "hi"}], confirm=True)

    def test_no_allowlist_raises(self, tmp_path, monkeypatch):
        """No PPT_ALLOWLIST_ROOTS set -> ValidationError."""
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.delenv("PPT_ALLOWLIST_ROOTS", raising=False)
        path = str(tmp_path / "x.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        with pytest.raises(ValidationError, match="No allowlist"):
            batch_set_text(path, 0, [{"shape_id": 1, "text": "hi"}], confirm=True)

    def test_non_pptx_extension_raises(self, tmp_path, monkeypatch):
        """Non-.pptx extension -> ValidationError."""
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        path = str(tmp_path / "file.docx")
        from pptmcp._pptx_batch_text import batch_set_text
        with pytest.raises(ValidationError, match="[Uu]nsupported"):
            batch_set_text(path, 0, [{"shape_id": 1, "text": "hi"}], confirm=True)

    def test_null_byte_in_path_raises(self, tmp_path, monkeypatch):
        """Null byte in path -> ValidationError."""
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        from pptmcp._pptx_batch_text import batch_set_text
        with pytest.raises(ValidationError, match="null byte"):
            batch_set_text("\x00evil.pptx", 0, [{"shape_id": 1, "text": "hi"}], confirm=True)


# ---------------------------------------------------------------------------
# Class 2: Input validation tests
# ---------------------------------------------------------------------------

class TestBatchSetTextInputValidation:
    @pytest.fixture(autouse=True)
    def _gates(self, tmp_path, monkeypatch):
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        self._tmp = tmp_path

    def test_empty_updates_raises(self):
        from pptmcp._pptx_batch_text import batch_set_text
        path = str(self._tmp / "x.pptx")
        with pytest.raises(ValidationError, match="non-empty"):
            batch_set_text(path, 0, [], confirm=True)

    def test_slide_out_of_range_raises(self):
        path, _ = _make_pptx_with_textboxes(self._tmp / "oor.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        with pytest.raises(ValidationError, match="out of range"):
            batch_set_text(path, 99, [{"shape_id": 1, "text": "hi"}], confirm=True)

    def test_negative_slide_raises(self):
        path, _ = _make_pptx_with_textboxes(self._tmp / "neg.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        with pytest.raises(ValidationError, match="out of range"):
            batch_set_text(path, -1, [{"shape_id": 1, "text": "hi"}], confirm=True)

    def test_dict_missing_shape_id_raises(self):
        path, _ = _make_pptx_with_textboxes(self._tmp / "noshapeid.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        with pytest.raises(ValidationError, match="shape_id"):
            batch_set_text(path, 0, [{"text": "hi"}], confirm=True)

    def test_dict_missing_text_raises(self):
        path, _ = _make_pptx_with_textboxes(self._tmp / "notext.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        with pytest.raises(ValidationError, match="text"):
            batch_set_text(path, 0, [{"shape_id": 1}], confirm=True)


# ---------------------------------------------------------------------------
# Class 3: Happy path tests
# ---------------------------------------------------------------------------

class TestBatchSetTextHappyPath:
    @pytest.fixture(autouse=True)
    def _gates(self, tmp_path, monkeypatch):
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        self._tmp = tmp_path

    def test_single_shape_update(self):
        path, (sid1, sid2) = _make_pptx_with_textboxes(self._tmp / "single.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        result = batch_set_text(path, 0, [{"shape_id": sid1, "text": "new text"}], confirm=True)
        assert result["updated"] is True
        assert len(result["results"]) == 1
        assert result["results"][0]["status"] == "ok"
        assert result["results"][0]["new_text"] == "new text"

    def test_multi_shape_update(self):
        path, (sid1, sid2) = _make_pptx_with_textboxes(self._tmp / "multi.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        updates = [
            {"shape_id": sid1, "text": "text A"},
            {"shape_id": sid2, "text": "text B"},
        ]
        result = batch_set_text(path, 0, updates, confirm=True)
        assert len(result["results"]) == 2
        assert result["error_count"] == 0

    def test_previous_text_captured(self):
        path, (sid1, _) = _make_pptx_with_textboxes(self._tmp / "prev.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        result = batch_set_text(path, 0, [{"shape_id": sid1, "text": "replaced"}], confirm=True)
        assert result["results"][0]["previous_text"] == "original text 0"

    def test_path_and_slide_index_in_result(self):
        path, (sid1, _) = _make_pptx_with_textboxes(self._tmp / "meta.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        result = batch_set_text(path, 0, [{"shape_id": sid1, "text": "x"}], confirm=True)
        assert "path" in result
        assert result["slide_index"] == 0

    def test_changes_persisted_to_disk(self):
        path, (sid1, _) = _make_pptx_with_textboxes(self._tmp / "persist.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        from pptmcp.presentation_pptx import _prs_cache
        batch_set_text(path, 0, [{"shape_id": sid1, "text": "saved text"}], confirm=True)
        _prs_cache.clear()
        reloaded = Presentation(path)
        texts = [
            shape.text_frame.text
            for shape in reloaded.slides[0].shapes
            if shape.has_text_frame
        ]
        assert "saved text" in texts

    def test_cache_evicted_after_save(self):
        path, (sid1, _) = _make_pptx_with_textboxes(self._tmp / "evict.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        from pptmcp.presentation_pptx import _load_prs, _prs_cache
        from pathlib import Path
        resolved = Path(path).resolve()
        cache_key = str(resolved)
        # Pre-load into cache
        _load_prs(resolved)
        assert cache_key in _prs_cache
        batch_set_text(path, 0, [{"shape_id": sid1, "text": "evicted"}], confirm=True)
        assert cache_key not in _prs_cache

    def test_non_zero_slide_index(self):
        """batch_set_text works on slides other than index 0."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])  # slide 0
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # slide 1
        tb = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2.0), Inches(1.0))
        tb.text_frame.text = "before"
        sid = tb.shape_id
        path = str(self._tmp / "slide1.pptx")
        prs.save(path)

        from pptmcp._pptx_batch_text import batch_set_text
        result = batch_set_text(path, 1, [{"shape_id": sid, "text": "after"}], confirm=True)
        assert result["slide_index"] == 1
        assert result["results"][0]["new_text"] == "after"


# ---------------------------------------------------------------------------
# Class 4: Partial failure tests
# ---------------------------------------------------------------------------

class TestBatchSetTextPartialFailure:
    @pytest.fixture(autouse=True)
    def _gates(self, tmp_path, monkeypatch):
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        self._tmp = tmp_path

    def test_unknown_shape_id_goes_to_skipped(self):
        path, (sid1, _) = _make_pptx_with_textboxes(self._tmp / "unknown.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        result = batch_set_text(
            path, 0,
            [{"shape_id": 99999, "text": "nope"}],
            confirm=True,
        )
        assert len(result["skipped"]) == 1
        assert result["skipped"][0]["shape_id"] == 99999

    def test_no_text_frame_shape_goes_to_skipped(self):
        """A picture shape has no text_frame and should land in skipped."""
        import struct
        import zlib

        def _chunk(tag, data):
            c = struct.pack(">I", len(data)) + tag + data
            return c + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)

        png = (
            b"\x89PNG\r\n\x1a\n"
            + _chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
            + _chunk(b"IDAT", zlib.compress(b"\x00\xff\xff\xff"))
            + _chunk(b"IEND", b"")
        )
        png_path = self._tmp / "pixel.png"
        png_path.write_bytes(png)

        prs_obj = Presentation()
        slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[5])
        pic = slide.shapes.add_picture(
            str(png_path), Inches(0.5), Inches(0.5), Inches(1.0), Inches(1.0)
        )
        # picture shapes are MSO_SHAPE_TYPE.PICTURE (13) — no has_text_frame
        pic_sid = pic.shape_id
        prs_path = str(self._tmp / "with_pic.pptx")
        prs_obj.save(prs_path)

        from pptmcp._pptx_batch_text import batch_set_text
        result = batch_set_text(
            prs_path, 0,
            [{"shape_id": pic_sid, "text": "text on pic"}],
            confirm=True,
        )
        assert len(result["skipped"]) == 1
        assert result["skipped"][0]["shape_id"] == pic_sid

    def test_all_bad_shapes_still_saves_file(self):
        """Even when all shapes are skipped the file is still saved (no exception)."""
        path, _ = _make_pptx_with_textboxes(self._tmp / "allbad.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        result = batch_set_text(
            path, 0,
            [{"shape_id": 99998, "text": "x"}, {"shape_id": 99997, "text": "y"}],
            confirm=True,
        )
        assert result["updated"] is True
        assert result["error_count"] == 2
        assert len(result["results"]) == 0

    def test_skipped_entry_has_shape_id_and_reason(self):
        path, _ = _make_pptx_with_textboxes(self._tmp / "skipkeys.pptx")
        from pptmcp._pptx_batch_text import batch_set_text
        result = batch_set_text(
            path, 0,
            [{"shape_id": 99996, "text": "gone"}],
            confirm=True,
        )
        entry = result["skipped"][0]
        assert "shape_id" in entry
        assert "reason" in entry
        assert isinstance(entry["reason"], str)
        assert len(entry["reason"]) > 0


# ---------------------------------------------------------------------------
# Class 5: Layer separation tests
# ---------------------------------------------------------------------------

class TestBatchSetTextLayerSeparation:
    def test_backend_has_no_fastmcp_import(self):
        """_pptx_batch_text.py must not import fastmcp or mcp.server."""
        import ast
        import pathlib

        src = pathlib.Path(__file__).resolve().parents[1] / "src" / "pptmcp" / "_pptx_batch_text.py"
        tree = ast.parse(src.read_text(encoding="utf-8"))
        for node in ast.walk(tree):
            if isinstance(node, (ast.Import, ast.ImportFrom)):
                names = []
                if isinstance(node, ast.Import):
                    names = [alias.name for alias in node.names]
                elif isinstance(node, ast.ImportFrom):
                    names = [node.module or ""]
                for name in names:
                    assert "fastmcp" not in name, f"fastmcp import found in backend: {name}"
                    assert "mcp.server" not in name, f"mcp.server import found in backend: {name}"

    def test_handler_has_no_pptx_import(self):
        """_server_handlers_batch_text.py must not import pptx at module level."""
        import ast
        import pathlib

        src = pathlib.Path(__file__).resolve().parents[1] / "src" / "pptmcp" / "_server_handlers_batch_text.py"
        tree = ast.parse(src.read_text(encoding="utf-8"))
        for node in ast.walk(tree):
            # Only check module-level imports (depth 1 direct children of Module)
            if isinstance(node, (ast.Import, ast.ImportFrom)):
                names = []
                if isinstance(node, ast.Import):
                    names = [alias.name for alias in node.names]
                elif isinstance(node, ast.ImportFrom):
                    names = [node.module or ""]
                for name in names:
                    assert not name.startswith("pptx"), (
                        f"pptx import found at module level in handler: {name}"
                    )
                    assert name != "python_pptx", f"python_pptx import in handler: {name}"
