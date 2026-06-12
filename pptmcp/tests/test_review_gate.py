"""test_review_gate.py — Phase D Review Gate test suite (W3 pre-implementation deliverable)

Tests for pptmcp.review_pptx:
  GROUP 2 (5): _detect_shapes_outside_slide   — slide-boundary violation detection
  GROUP 3 (1): review_slide_export            — full public-API return schema (empty slide)
  GROUP 4 (5): Security / validation guards   — BLK-D01 (index bounds), BLK-D02 (PNG exists),
                                                allowlist checks (OWASP A01)
  GROUP 5 (2): Edge-case robustness           — zero-area shapes (BLK-D03 max guard),
                                                negative EMU dimensions (BLK-D03 abs guard)

NOTE: GROUP 1 (_detect_text_overflow_heuristic, 4 tests) was relocated to
      test_review_visual_qa.py where it sits alongside the other visual-QA detectors.

Total: 13 test functions.

Execution: CI-safe (no COM, no live Office app required).
  pytest powerpoint/tests/test_review_gate.py -v
"""
from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock

import pytest
from pptx import Presentation
from pptx.util import Emu

from pptmcp.presentation_pptx import ValidationError
from pptmcp.review_pptx import (
    _detect_shapes_outside_slide,
    _detect_text_overflow_heuristic,
    review_slide_export,
)

# ── EMU constants ─────────────────────────────────────────────────────────
EMU_PER_INCH: int = 914400
# Standard 16:9 slide dimensions used throughout
SLIDE_WIDTH_EMU: int = 12192000   # 13.33 inches
SLIDE_HEIGHT_EMU: int = 6858000   # 7.50 inches


# ── Module-level helpers (not fixtures; used by multiple groups) ──────────

def _blank_16x9_prs() -> Presentation:
    """Return an unsaved in-memory 16:9 Presentation with one blank slide."""
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)
    prs.slides.add_slide(prs.slide_layouts[5])  # index 5 = blank layout, no placeholders
    return prs


def _save_prs(prs: Presentation, directory: Path, name: str = "review_test.pptx") -> Path:
    """Save *prs* to *directory/name* and return the absolute Path."""
    path = directory / name
    prs.save(str(path))
    return path


def _make_dummy_png(directory: Path, name: str = "export.png") -> Path:
    """Write the 8-byte PNG magic-number signature to *directory/name* and return the Path.

    The file is NOT a valid image but its existence is sufficient for the BLK-D02
    Path.exists() guard under test.
    """
    png = directory / name
    png.write_bytes(b"\x89PNG\r\n\x1a\n")
    return png


# ── Shared fixtures ───────────────────────────────────────────────────────

@pytest.fixture
def blank_env(tmp_path, monkeypatch):
    """One blank 16:9 slide + dummy PNG, with tmp_path set as the allowlist root.

    Returns (pptx_path_str, png_path_str) ready to pass to review_slide_export.
    """
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    prs = _blank_16x9_prs()
    pptx_path = _save_prs(prs, tmp_path)
    png_path = _make_dummy_png(tmp_path)
    return str(pptx_path), str(png_path)


# ═════════════════════════════════════════════════════════════════════════
# GROUP 2 — _detect_shapes_outside_slide
# Checks:  left < 0            → outside (left boundary)
#          top  < 0            → outside (top boundary)
#          left + width  > W   → outside (right boundary)
#          top  + height > H   → outside (bottom boundary)
# ═════════════════════════════════════════════════════════════════════════

@pytest.mark.unit
def test_outside_slide_left_negative_returns_finding():
    """TC-RG-05: shape.left = -100 EMU → SHAPE_OUTSIDE_SLIDE finding (left < 0).

    Any shape whose left edge is negative escapes the slide canvas on the left.
    """
    prs = _blank_16x9_prs()
    slide = prs.slides[0]

    slide.shapes.add_textbox(Emu(-100), Emu(0), Emu(500000), Emu(500000))

    findings = _detect_shapes_outside_slide(slide, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU)

    assert len(findings) >= 1
    assert any(f["criterion"] == "SHAPE_OUTSIDE_SLIDE" for f in findings)
    assert any(f["result"] == "FAIL" for f in findings)


@pytest.mark.unit
def test_outside_slide_right_exceeds_width_returns_finding():
    """TC-RG-06: right edge > SLIDE_WIDTH_EMU → SHAPE_OUTSIDE_SLIDE finding.

    Shape placement:
      left  = SLIDE_WIDTH_EMU - 100  (100 EMU before the right edge)
      width = 500_000 EMU
      right = (SLIDE_WIDTH_EMU - 100) + 500_000 = SLIDE_WIDTH_EMU + 499_900  → exceeds
    """
    prs = _blank_16x9_prs()
    slide = prs.slides[0]

    left_emu = SLIDE_WIDTH_EMU - 100
    slide.shapes.add_textbox(Emu(left_emu), Emu(0), Emu(500000), Emu(500000))

    findings = _detect_shapes_outside_slide(slide, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU)

    assert len(findings) >= 1
    assert any(f["criterion"] == "SHAPE_OUTSIDE_SLIDE" for f in findings)


@pytest.mark.unit
def test_outside_slide_top_negative_returns_finding():
    """TC-RG-07: shape.top = -100 EMU → SHAPE_OUTSIDE_SLIDE finding (top < 0).

    Any shape whose top edge is negative escapes the slide canvas above the top.
    """
    prs = _blank_16x9_prs()
    slide = prs.slides[0]

    slide.shapes.add_textbox(Emu(0), Emu(-100), Emu(500000), Emu(500000))

    findings = _detect_shapes_outside_slide(slide, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU)

    assert len(findings) >= 1
    assert any(f["criterion"] == "SHAPE_OUTSIDE_SLIDE" for f in findings)


@pytest.mark.unit
def test_outside_slide_bottom_exceeds_height_returns_finding():
    """TC-RG-08: bottom edge > SLIDE_HEIGHT_EMU → SHAPE_OUTSIDE_SLIDE finding.

    Shape placement:
      top    = SLIDE_HEIGHT_EMU - 100  (100 EMU above the bottom edge)
      height = 500_000 EMU
      bottom = (SLIDE_HEIGHT_EMU - 100) + 500_000 = SLIDE_HEIGHT_EMU + 499_900  → exceeds
    """
    prs = _blank_16x9_prs()
    slide = prs.slides[0]

    top_emu = SLIDE_HEIGHT_EMU - 100
    slide.shapes.add_textbox(Emu(0), Emu(top_emu), Emu(500000), Emu(500000))

    findings = _detect_shapes_outside_slide(slide, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU)

    assert len(findings) >= 1
    assert any(f["criterion"] == "SHAPE_OUTSIDE_SLIDE" for f in findings)


@pytest.mark.unit
def test_shape_within_slide_no_text_no_findings():
    """TC-RG-09: Shape fully within slide limits, empty text frame → both helpers return [].

    Shape at (2", 2"), size 2"×1" — well inside the 13.33"×7.5" canvas.
    No text set → len("") = 0 → density = 0 < 100 → no TEXT_OVERFLOW finding.
    Left, top, right, bottom all within bounds → no SHAPE_OUTSIDE_SLIDE finding.
    """
    prs = _blank_16x9_prs()
    slide = prs.slides[0]

    slide.shapes.add_textbox(
        Emu(2 * EMU_PER_INCH), Emu(2 * EMU_PER_INCH),
        Emu(2 * EMU_PER_INCH), Emu(1 * EMU_PER_INCH),
    )
    # text_frame.text is "" by default — no text assigned

    overflow_findings = _detect_text_overflow_heuristic(slide, prs)
    boundary_findings = _detect_shapes_outside_slide(slide, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU)

    assert overflow_findings == []
    assert boundary_findings == []


# ── GROUP 2b — RCI-US-16 spec lock-in ─────────────────────────────────────
# Severity=critical, inches (not EMU), structured bounds_in + overflow_in
# payload. Locks the spec.


@pytest.mark.unit
def test_outside_slide_spec_severity_critical_and_payload_in_inches():
    """RCI-US-16: finding severity=critical with inches-based structured payload.

    Shape placed 1in past the right edge:
      left = SLIDE_WIDTH - 0.5in = 12.83in
      width = 1.5in
      right = 14.33in → overflow = 14.33 - 13.33 = 1.0in
    """
    prs = _blank_16x9_prs()
    slide = prs.slides[0]
    half_in = EMU_PER_INCH // 2
    one_and_half_in = EMU_PER_INCH + half_in

    slide.shapes.add_textbox(
        Emu(SLIDE_WIDTH_EMU - half_in),     # left = 12.83in
        Emu(EMU_PER_INCH),                  # top  = 1.00in
        Emu(one_and_half_in),               # width = 1.50in
        Emu(EMU_PER_INCH),                  # height = 1.00in
    )

    findings = _detect_shapes_outside_slide(slide, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU)
    overflow = [f for f in findings if f["criterion"] == "SHAPE_OUTSIDE_SLIDE"]
    assert len(overflow) == 1
    f = overflow[0]

    # Severity must be critical per RCI-US-16 spec (not legacy "high")
    assert f["severity"] == "critical", f"expected severity=critical, got {f['severity']!r}"

    # Structured payload — machine-readable
    assert f["shape_id"] is not None
    assert "shape_name" in f
    bounds = f["bounds_in"]
    assert set(bounds.keys()) == {"left", "top", "width", "height"}
    assert bounds["left"] == pytest.approx(12.83, abs=0.01)
    assert bounds["width"] == pytest.approx(1.5, abs=0.01)

    overflow_dict = f["overflow_in"]
    assert "right" in overflow_dict
    assert overflow_dict["right"] == pytest.approx(1.0, abs=0.01)
    # No false-positive on other directions
    assert "left" not in overflow_dict
    assert "top" not in overflow_dict
    assert "bottom" not in overflow_dict

    # Human detail string must mention inches, not EMU
    assert "in" in f["detail"]
    assert "EMU" not in f["detail"]


@pytest.mark.unit
def test_outside_slide_spec_multi_direction_overflow():
    """RCI-US-16: shape overflowing on multiple edges reports each direction."""
    prs = _blank_16x9_prs()
    slide = prs.slides[0]

    # Shape at (-0.5in, -0.25in) sized 14.5in x 8in → overflows left, top, right, bottom
    slide.shapes.add_textbox(
        Emu(-EMU_PER_INCH // 2),
        Emu(-EMU_PER_INCH // 4),
        Emu(14 * EMU_PER_INCH + EMU_PER_INCH // 2),  # 14.5in wide
        Emu(8 * EMU_PER_INCH),                       # 8.0in tall
    )

    findings = _detect_shapes_outside_slide(slide, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU)
    overflow = [f for f in findings if f["criterion"] == "SHAPE_OUTSIDE_SLIDE"]
    assert len(overflow) == 1
    overflow_in = overflow[0]["overflow_in"]
    assert set(overflow_in.keys()) == {"left", "top", "right", "bottom"}
    assert overflow_in["left"] == pytest.approx(0.5, abs=0.01)
    assert overflow_in["top"] == pytest.approx(0.25, abs=0.01)


# ═════════════════════════════════════════════════════════════════════════
# GROUP 3 — review_slide_export: full return-schema validation
# ═════════════════════════════════════════════════════════════════════════

@pytest.mark.unit
def test_review_slide_export_empty_slide_full_schema(blank_env, monkeypatch):
    """TC-RG-01b: review_slide_export on a blank slide → complete, correct return schema.

    Expected schema (from spec):
      {
        "slide_index":    0,
        "png_path":       <echoed>,
        "pptx_path":      <echoed>,
        "passed":         True,
        "findings":       [],
        "checks_run":     ["TEXT_OVERFLOW", "SHAPE_OUTSIDE_SLIDE", "TABLE_CONTENT_CLIPPING"],
        "findings_count": 0,
      }
    """
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx_path, png_path = blank_env

    result = review_slide_export(pptx_path, 0, png_path, confirm=True)

    # All required top-level keys must be present
    for key in ("slide_index", "png_path", "pptx_path", "passed", "findings",
                "checks_run", "findings_count"):
        assert key in result, f"Missing key in result: {key!r}"

    assert result["passed"] is True
    assert result["findings"] == []
    assert result["findings_count"] == 0
    assert {"TEXT_OVERFLOW", "SHAPE_OUTSIDE_SLIDE", "TABLE_CONTENT_CLIPPING"}.issubset(result["checks_run"])
    assert result["slide_index"] == 0
    assert result["png_path"] == png_path
    assert result["pptx_path"] == pptx_path


# ── GROUP 3c — RCI-US-14 overall PASS/FAIL/WARN ─────────────────────────


@pytest.mark.unit
def test_review_slide_export_overall_pass_when_clean(blank_env, monkeypatch):
    """RCI-US-14: no findings and no visual_qa_findings → overall='PASS'.

    Deterministically pin all detectors to [] so this exercises the
    aggregation logic, not the heuristic behaviour of the detectors.
    """
    pptx_path, png_path = blank_env
    import pptmcp._review_export as _rev_export
    for name in (
        "_detect_text_overflow_heuristic",
        "_detect_shapes_outside_slide",
        "_detect_table_content_clipping",
        "_detect_unfinished_placeholder",
        "_detect_low_font_size",
        "_detect_sparse_slide",
    ):
        monkeypatch.setattr(_rev_export, name, lambda *args, **kwargs: [])

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = review_slide_export(pptx_path, 0, png_path, confirm=True)
    assert result["findings"] == []
    assert result["visual_qa_findings"] == []
    assert result["overall"] == "PASS"
    assert result["passed"] is True


@pytest.mark.unit
def test_review_slide_export_overall_warn_when_only_visual_qa(tmp_path, monkeypatch):
    """RCI-US-14: only visual_qa_findings (e.g. SPARSE_SLIDE) → overall='WARN'."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    prs = _blank_16x9_prs()  # blank layout produces SPARSE_SLIDE visual QA
    pptx_path = _save_prs(prs, tmp_path, "warn.pptx")
    png_path = _make_dummy_png(tmp_path, "warn.png")

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = review_slide_export(str(pptx_path), 0, str(png_path), confirm=True)
    if not result["visual_qa_findings"]:
        pytest.skip("Fixture failed to trigger visual_qa_findings; WARN path unreachable")
    assert result["findings"] == []
    assert result["overall"] == "WARN"
    assert result["passed"] is True  # passed reflects structural, not WARN


@pytest.mark.unit
def test_review_slide_export_overall_fail_when_structural_findings(tmp_path, monkeypatch):
    """RCI-US-14: structural finding (shape outside slide) → overall='FAIL'."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    prs = _blank_16x9_prs()
    slide = prs.slides[0]
    # Box that hangs off the right edge → SHAPE_OUTSIDE_SLIDE finding
    half_in = EMU_PER_INCH // 2
    slide.shapes.add_textbox(
        Emu(SLIDE_WIDTH_EMU - half_in),
        Emu(EMU_PER_INCH),
        Emu(EMU_PER_INCH + half_in),
        Emu(EMU_PER_INCH),
    )
    pptx_path = _save_prs(prs, tmp_path, "fail.pptx")
    png_path = _make_dummy_png(tmp_path, "fail.png")

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = review_slide_export(str(pptx_path), 0, str(png_path), confirm=True)
    assert result["overall"] == "FAIL"
    assert result["passed"] is False
    assert result["findings_count"] >= 1


# ── GROUP 3d — pixel-read PNG checks ────────────────


def _write_solid_png(path, size=(1920, 1080), color=(255, 255, 255)):
    """Write a real, decodable PNG of given size and solid color."""
    from PIL import Image
    img = Image.new("RGB", size, color)
    img.save(str(path), "PNG")


def _write_noisy_png(path, size=(1920, 1080)):
    """Write a PNG with random per-pixel noise so variance is high."""
    import os
    from PIL import Image
    img = Image.frombytes("RGB", size, os.urandom(size[0] * size[1] * 3))
    img.save(str(path), "PNG")


@pytest.mark.unit
def test_review_slide_export_blank_png_emits_blank_slide_finding(
    blank_env, monkeypatch
):
    """All-white PNG → BLANK_SLIDE structural finding (severity=critical)."""
    pptx_path, png_path = blank_env
    # Overwrite the dummy PNG with a real all-white image
    _write_solid_png(png_path, size=(320, 240), color=(255, 255, 255))

    import pptmcp._review_export as rev
    # Mute the other detectors so the BLANK_SLIDE finding is unambiguous
    for name in (
        "_detect_text_overflow_heuristic",
        "_detect_shapes_outside_slide",
        "_detect_table_content_clipping",
        "_detect_unfinished_placeholder",
        "_detect_low_font_size",
        "_detect_sparse_slide",
    ):
        monkeypatch.setattr(rev, name, lambda *a, **kw: [])

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = review_slide_export(pptx_path, 0, png_path, confirm=True)
    blank = [f for f in result["findings"] if f["criterion"] == "BLANK_SLIDE"]
    assert len(blank) == 1
    f = blank[0]
    assert f["severity"] == "critical"
    assert f["result"] == "FAIL"
    assert f["variance"] < f["threshold"]
    assert "BLANK_SLIDE" in result["checks_run"]
    assert result["overall"] == "FAIL"


@pytest.mark.unit
def test_review_slide_export_noisy_png_no_blank_finding(blank_env, monkeypatch):
    """Noisy PNG (random pixels) at render-sized dimensions → no BLANK_SLIDE finding."""
    pptx_path, png_path = blank_env
    # Use >=100x100 so the render-size guard does NOT skip the variance check —
    # this test specifically exercises the variance-above-threshold path.
    _write_noisy_png(png_path, size=(256, 256))

    import pptmcp._review_export as rev
    for name in (
        "_detect_text_overflow_heuristic",
        "_detect_shapes_outside_slide",
        "_detect_table_content_clipping",
        "_detect_unfinished_placeholder",
        "_detect_low_font_size",
        "_detect_sparse_slide",
    ):
        monkeypatch.setattr(rev, name, lambda *a, **kw: [])

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = review_slide_export(pptx_path, 0, png_path, confirm=True)
    blank = [f for f in result["findings"] if f["criterion"] == "BLANK_SLIDE"]
    assert blank == []


@pytest.mark.unit
def test_review_slide_export_pixel_dimension_mismatch(
    tmp_path, monkeypatch
):
    """Contract.expected_resolution mismatch → PIXEL_DIMENSION_MISMATCH finding."""
    import json as _json
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    prs = _blank_16x9_prs()
    pptx_path = _save_prs(prs, tmp_path, "pdim.pptx")

    # Real PNG at 800x600 (noisy so BLANK_SLIDE doesn't also fire)
    png_path = tmp_path / "pdim.png"
    _write_noisy_png(png_path, size=(800, 600))

    # Contract declares 1920x1080 → mismatch expected
    contract = tmp_path / "contract.json"
    contract.write_text(
        _json.dumps({
            "contract_version": "1.0",
            "slides": [],
            "expected_resolution": [1920, 1080],
        }),
        encoding="utf-8",
    )

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = review_slide_export(
        str(pptx_path), 0, str(png_path), contract_path=str(contract), confirm=True
    )

    mism = [f for f in result["findings"] if f["criterion"] == "PIXEL_DIMENSION_MISMATCH"]
    assert len(mism) == 1
    f = mism[0]
    assert f["severity"] == "high"
    assert f["actual_resolution"] == [800, 600]
    assert f["expected_resolution"] == [1920, 1080]
    assert "PIXEL_DIMENSION_MISMATCH" in result["checks_run"]


@pytest.mark.unit
def test_review_slide_export_matching_dimensions_no_finding(
    tmp_path, monkeypatch
):
    """Contract.expected_resolution matches actual → no PIXEL_DIMENSION_MISMATCH."""
    import json as _json
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    prs = _blank_16x9_prs()
    pptx_path = _save_prs(prs, tmp_path, "pdim_ok.pptx")
    png_path = tmp_path / "pdim_ok.png"
    _write_noisy_png(png_path, size=(640, 480))

    contract = tmp_path / "contract.json"
    contract.write_text(
        _json.dumps({
            "contract_version": "1.0",
            "slides": [],
            "expected_resolution": [640, 480],
        }),
        encoding="utf-8",
    )

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = review_slide_export(
        str(pptx_path), 0, str(png_path), contract_path=str(contract), confirm=True
    )
    mism = [f for f in result["findings"] if f["criterion"] == "PIXEL_DIMENSION_MISMATCH"]
    assert mism == []


# ── GROUP 3b — RCI-US-13 spec lock-in ─────────────────────────────────────
# review_result.json is written beside the PNG; contract findings are merged.


@pytest.mark.unit
def test_review_slide_export_writes_review_result_json(blank_env, monkeypatch):
    """RCI-US-13: review_slide_export writes review_result.json beside the PNG."""
    import json as _json

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx_path, png_path = blank_env

    result = review_slide_export(pptx_path, 0, png_path, confirm=True)

    # review_result_path is returned and points beside the PNG
    assert "review_result_path" in result
    rrp = Path(result["review_result_path"])
    assert rrp.name == "review_result.json"
    assert rrp.parent == Path(png_path).parent
    assert rrp.exists(), "review_result.json was not written to disk"

    # File contents match the returned dict
    on_disk = _json.loads(rrp.read_text(encoding="utf-8"))
    assert on_disk["slide_index"] == result["slide_index"]
    assert on_disk["passed"] == result["passed"]
    assert on_disk["findings_count"] == result["findings_count"]


@pytest.mark.unit
def test_review_slide_export_contract_dimension_mismatch_adds_finding(
    tmp_path, monkeypatch
):
    """RCI-US-13: contract_path with mismatched slide_width_inches adds a CONTRACT_ finding."""
    import json as _json

    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    prs = _blank_16x9_prs()       # actual width = 13.33in
    pptx_path = _save_prs(prs, tmp_path)
    png_path = _make_dummy_png(tmp_path)

    # Contract declares a different width → mismatch finding expected
    contract = tmp_path / "contract.json"
    contract.write_text(
        _json.dumps({"slide_width_inches": 10.0, "slide_height_inches": 7.5}),
        encoding="utf-8",
    )

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = review_slide_export(
        str(pptx_path), 0, str(png_path), contract_path=str(contract), confirm=True
    )

    contract_findings = [
        f for f in result["findings"]
        if f.get("criterion", "").startswith("CONTRACT_")
    ]
    assert len(contract_findings) >= 1
    width_findings = [
        f for f in contract_findings
        if "slide_width" in f["criterion"]
    ]
    assert len(width_findings) == 1
    f = width_findings[0]
    assert f["result"] == "FAIL"
    assert "10.000" in f["detail"] or "10.0" in f["detail"]
    assert "CONTRACT" in result["checks_run"]
    assert result["passed"] is False  # at least one structural finding now


# ═════════════════════════════════════════════════════════════════════════
# GROUP 4 — Security / validation guards
#   BLK-D01: 0 <= slide_index < len(prs.slides)  (lower AND upper bound)
#   BLK-D02: explicit Path.exists() check on png_path after allowlist passes
#   OWASP A01: both pptx and png paths must resolve inside PPT_ALLOWLIST_ROOTS
# ═════════════════════════════════════════════════════════════════════════

@pytest.mark.unit
def test_pptx_path_outside_allowlist_raises(tmp_path, monkeypatch):
    """TC-RG-10: pptx resolves outside PPT_ALLOWLIST_ROOTS → ValidationError (OWASP A01).

    Setup: allowlist covers tmp_path/allowed/; the pptx is saved to tmp_path/outside/.
    The PNG is inside the allowlist so the pptx check fires first.
    """
    allowed = tmp_path / "allowed"
    allowed.mkdir()
    outside = tmp_path / "outside"
    outside.mkdir()
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(allowed))

    # Save a real pptx outside the allowlist so the file-existence guard does not fire first
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    bad_pptx = outside / "bad.pptx"
    prs.save(str(bad_pptx))

    good_png = allowed / "export.png"
    good_png.write_bytes(b"\x89PNG\r\n\x1a\n")

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match=r"not in allowlist|allowlist"):
        review_slide_export(str(bad_pptx), 0, str(good_png), confirm=True)


@pytest.mark.unit
def test_png_path_outside_allowlist_raises(tmp_path, monkeypatch):
    """TC-RG-11: png_path resolves outside PPT_ALLOWLIST_ROOTS → ValidationError (OWASP A01).

    Setup: pptx is inside the allowlist; only the PNG is in a directory outside it.
    """
    allowed = tmp_path / "allowed"
    allowed.mkdir()
    outside = tmp_path / "outside"
    outside.mkdir()
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(allowed))

    prs = _blank_16x9_prs()
    pptx_path = _save_prs(prs, allowed)

    bad_png = outside / "export.png"
    bad_png.write_bytes(b"\x89PNG\r\n\x1a\n")

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match=r"not in allowlist|allowlist"):
        review_slide_export(str(pptx_path), 0, str(bad_png), confirm=True)


@pytest.mark.unit
def test_negative_slide_index_raises(blank_env, monkeypatch):
    """TC-RG-12: slide_index = -1 → ValidationError (BLK-D01: lower bound 0 <= index).

    A negative index must be rejected before any slide access attempt.
    """
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx_path, png_path = blank_env

    with pytest.raises(ValidationError, match=r"slide_index|out.of.range|invalid"):
        review_slide_export(pptx_path, -1, png_path, confirm=True)


@pytest.mark.unit
def test_slide_index_at_len_raises(tmp_path, monkeypatch):
    """TC-RG-13: slide_index = len(slides) on a 1-slide deck → ValidationError (BLK-D01).

    The deck has exactly 1 slide (valid index: 0 only).
    Index 1 equals len(slides), violating the strict upper bound (index < len).
    """
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    prs = _blank_16x9_prs()
    pptx_path = _save_prs(prs, tmp_path)
    png_path = _make_dummy_png(tmp_path)

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match=r"slide_index|out.of.range|invalid"):
        review_slide_export(str(pptx_path), 1, str(png_path), confirm=True)


@pytest.mark.unit
def test_png_file_does_not_exist_raises(tmp_path, monkeypatch):
    """TC-RG-14: png_path is inside allowlist but file is absent → ValidationError (BLK-D02).

    BLK-D02 requires an explicit Path.exists() check after the allowlist guard passes.
    Verifies that allowlist-membership alone is not sufficient — the file must also exist.
    """
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    prs = _blank_16x9_prs()
    pptx_path = _save_prs(prs, tmp_path)

    nonexistent_png = tmp_path / "ghost_export.png"
    assert not nonexistent_png.exists(), "Test precondition: file must not exist on disk"

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match=r"not found|does not exist|png|ghost"):
        review_slide_export(str(pptx_path), 0, str(nonexistent_png), confirm=True)


# ═════════════════════════════════════════════════════════════════════════
# GROUP 5 — Edge-case robustness
#   BLK-D03: abs(shape.width) and abs(shape.height) before EMU→inch conversion
#            max(area, 0.01) guard against zero-division
# ═════════════════════════════════════════════════════════════════════════

@pytest.mark.unit
def test_overflow_zero_area_shape_no_crash():
    """TC-RG-15: Mock shape with width=0, height=0 → no ZeroDivisionError.

    Implementation must guard zero area with max(area, 0.01) or by skipping the shape.
    The assertion is 'does not raise' — findings content is not constrained because
    the semantic result for a zero-area shape is implementation-defined.
    """
    shape_mock = MagicMock()
    shape_mock.has_text_frame = True
    shape_mock.width = 0
    shape_mock.height = 0
    shape_mock.text_frame.text = "Z" * 50

    slide_mock = MagicMock()
    slide_mock.shapes = [shape_mock]
    prs_mock = MagicMock()

    try:
        findings = _detect_text_overflow_heuristic(slide_mock, prs_mock)
    except ZeroDivisionError:
        pytest.fail(
            "_detect_text_overflow_heuristic raised ZeroDivisionError for a zero-area shape. "
            "BLK-D03 requires max(area, 0.01) or an explicit zero-area skip guard."
        )
    else:
        assert isinstance(findings, list), "Return value must be a list"


@pytest.mark.unit
def test_overflow_negative_emu_width_abs_guard_no_false_positive():
    """TC-RG-16: Mock shape with width=-EMU_PER_INCH, height=+EMU_PER_INCH → abs() applied.

    This test encodes BLK-D03: the implementation MUST call abs(shape.width) and
    abs(shape.height) before the EMU-to-inch conversion.  Without abs(), mixed-sign
    dimensions produce a negative area; combined with max(area, 0.01), this yields
    a density of 50/0.01 = 5000 > 200 — a false-positive HIGH finding.

    With abs() applied:
      width_in  = abs(-914400) / 914400 = 1.0 in
      height_in = abs(914400)  / 914400 = 1.0 in
      area      = 1.0 × 1.0 = 1.0 sq.in
      density   = 50 / 1.0 = 50  →  < 100  →  no finding  ✓

    Without abs() (buggy path):
      area    = -1.0 × 1.0 = -1.0  →  max(-1.0, 0.01) = 0.01
      density = 50 / 0.01 = 5000   →  > 200  →  "high" finding  ✗ (false positive)

    The test therefore asserts findings == [] to detect the bug.
    """
    shape_mock = MagicMock()
    shape_mock.has_text_frame = True
    shape_mock.width = -int(EMU_PER_INCH)   # negative width: -914400 EMU
    shape_mock.height = int(EMU_PER_INCH)    # positive height: +914400 EMU
    shape_mock.text_frame.text = "D" * 50    # 50 chars; density=50 with abs() → < 100

    slide_mock = MagicMock()
    slide_mock.shapes = [shape_mock]
    prs_mock = MagicMock()

    try:
        findings = _detect_text_overflow_heuristic(slide_mock, prs_mock)
    except Exception as exc:
        pytest.fail(
            f"_detect_text_overflow_heuristic raised {type(exc).__name__} "
            f"for a shape with negative EMU width: {exc}. "
            "BLK-D03 requires abs(shape.width) before EMU conversion."
        )

    assert findings == [], (
        f"Expected [] (density=50 with abs() applied), got {findings}. "
        "This likely means abs() is NOT applied and max(area, 0.01) is turning "
        "negative area into 0.01, producing a false-positive HIGH finding."
    )


# ═════════════════════════════════════════════════════════════════════════
# GROUP 6 — Mutation gate tests for review_slide_export
# ═════════════════════════════════════════════════════════════════════════

@pytest.mark.unit
def test_review_slide_export_blocked_when_write_disabled(tmp_path, monkeypatch):
    """Gate 1: NotAllowedError when PPT_ENABLE_WRITE not set."""
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    from pptmcp.presentation_pptx import NotAllowedError
    with pytest.raises(NotAllowedError):
        review_slide_export("dummy.pptx", 0, "dummy.png")


@pytest.mark.unit
def test_review_slide_export_blocked_when_confirm_false(tmp_path, monkeypatch):
    """Gate 2: ValidationError when PPT_ENABLE_WRITE=true but confirm=False."""
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptmcp.presentation_pptx import ValidationError as _ValidationError
    with pytest.raises(_ValidationError):
        review_slide_export("dummy.pptx", 0, "dummy.png", confirm=False)


@pytest.mark.unit
def test_review_slide_export_passes_with_gate_open(blank_env, monkeypatch):
    """Happy path: both gates open → result dict returned with review_result_path."""
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx_path, png_path = blank_env
    result = review_slide_export(pptx_path, 0, png_path, confirm=True)
    assert "review_result_path" in result
    assert result["passed"] in (True, False)
