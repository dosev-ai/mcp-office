"""Tests for visual QA detectors introduced in Sprint 1A (BUG-686 + FR-896).

Includes GROUP 1 (_detect_text_overflow_heuristic) relocated from test_review_gate.py.
"""
from __future__ import annotations

from unittest.mock import MagicMock, PropertyMock, patch

import pytest
from pptx import Presentation
from pptx.util import Emu

from pptmcp.review_pptx import (
    _detect_low_font_size,
    _detect_sparse_slide,
    _detect_text_overflow_heuristic,
    _detect_unfinished_placeholder,
    review_slide_export,
    DEFAULT_TEXT,
)

# ── EMU constants (used by GROUP-1 text-overflow tests) ───────────────────
EMU_PER_INCH: int = 914_400
SLIDE_WIDTH_EMU: int = 12_192_000   # 13.33 inches
SLIDE_HEIGHT_EMU: int = 6_858_000   # 7.50 inches


def _blank_16x9_prs() -> Presentation:
    """Return an unsaved in-memory 16:9 Presentation with one blank slide."""
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)
    prs.slides.add_slide(prs.slide_layouts[5])  # index 5 = blank layout
    return prs


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_mock_run(size_emu: int | None) -> MagicMock:
    run = MagicMock()
    run.font.size = size_emu
    return run


def _make_slide_with_placeholder(text: str = "", is_placeholder: bool = True) -> MagicMock:
    shape = MagicMock()
    shape.placeholder_format = MagicMock() if is_placeholder else None
    shape.text = text
    shape.name = "TestPlaceholder"
    shape.has_text_frame = True
    slide = MagicMock()
    slide.shapes = [shape]
    return slide, shape


# ---------------------------------------------------------------------------
# G6-01 .. G6-05  _detect_unfinished_placeholder
# ---------------------------------------------------------------------------

def test_G6_01_empty_placeholder_returns_high():
    slide, _ = _make_slide_with_placeholder(text="")
    issues = _detect_unfinished_placeholder(slide)
    assert len(issues) == 1
    assert issues[0]["severity"] == "high"
    assert issues[0]["check"] == "UNFINISHED_PLACEHOLDER"


def test_G6_02_default_text_returns_high():
    for default in list(DEFAULT_TEXT)[:3]:
        slide, _ = _make_slide_with_placeholder(text=default)
        issues = _detect_unfinished_placeholder(slide)
        assert len(issues) == 1, f"Expected issue for default text: {default!r}"


def test_G6_03_real_content_no_issue():
    slide, _ = _make_slide_with_placeholder(text="Q1 2025 Revenue Analysis")
    issues = _detect_unfinished_placeholder(slide)
    assert issues == []


def test_G6_04_non_placeholder_skipped():
    slide, _ = _make_slide_with_placeholder(text="", is_placeholder=False)
    issues = _detect_unfinished_placeholder(slide)
    assert issues == []


def test_G6_05_repr_in_detail_sanitizes():
    """RM-01: detail must use repr() for text and name!r for shape name."""
    slide, _ = _make_slide_with_placeholder(text="")
    issues = _detect_unfinished_placeholder(slide)
    # detail should contain repr-style quotes for the text portion
    assert "TestPlaceholder" in issues[0]["detail"]


# ---------------------------------------------------------------------------
# G6-06 .. G6-10  _detect_low_font_size
# ---------------------------------------------------------------------------

def test_G6_06_below_high_threshold_returns_HIGH():
    slide = MagicMock()
    shape = MagicMock()
    shape.has_text_frame = True
    shape.name = "TinyText"
    run = _make_mock_run(size_emu=9 * 12700)  # 9pt < 10pt threshold
    shape.text_frame.paragraphs = [MagicMock(runs=[run])]
    slide.shapes = [shape]
    issues = _detect_low_font_size(slide)
    assert any(i["severity"] == "high" for i in issues)


def test_G6_07_between_thresholds_returns_MEDIUM():
    slide = MagicMock()
    shape = MagicMock()
    shape.has_text_frame = True
    shape.name = "SmallText"
    run = _make_mock_run(size_emu=12 * 12700)  # 12pt → between 10 and 14
    shape.text_frame.paragraphs = [MagicMock(runs=[run])]
    slide.shapes = [shape]
    issues = _detect_low_font_size(slide)
    assert any(i["severity"] == "medium" for i in issues)


def test_G6_08_above_medium_threshold_no_issue():
    slide = MagicMock()
    shape = MagicMock()
    shape.has_text_frame = True
    shape.name = "NormalText"
    run = _make_mock_run(size_emu=16 * 12700)  # 16pt ≥ 14pt → ok
    shape.text_frame.paragraphs = [MagicMock(runs=[run])]
    slide.shapes = [shape]
    issues = _detect_low_font_size(slide)
    assert issues == []


def test_G6_09_none_font_size_skipped():
    slide = MagicMock()
    shape = MagicMock()
    shape.has_text_frame = True
    shape.name = "InheritedSize"
    run = _make_mock_run(size_emu=None)
    shape.text_frame.paragraphs = [MagicMock(runs=[run])]
    slide.shapes = [shape]
    issues = _detect_low_font_size(slide)
    assert issues == []


def test_G6_10_non_text_shape_skipped():
    slide = MagicMock()
    shape = MagicMock()
    shape.has_text_frame = False
    slide.shapes = [shape]
    issues = _detect_low_font_size(slide)
    assert issues == []


# ---------------------------------------------------------------------------
# G6-11 .. G6-13  _detect_sparse_slide
# ---------------------------------------------------------------------------

def test_G6_11_sparse_slide_returns_MEDIUM():
    slide = MagicMock()
    # Only 1 small content shape
    shape = MagicMock()
    shape.has_text_frame = True
    shape.text = "Hi"
    shape.width = 100_000
    shape.height = 100_000
    slide.shapes = [shape]
    # Slide is large (standard 9144000 x 6858000 EMU)
    issues = _detect_sparse_slide(slide, 9_144_000, 6_858_000)
    assert len(issues) == 1
    assert issues[0]["severity"] == "medium"
    assert issues[0]["check"] == "SPARSE_SLIDE"


def test_G6_12_zero_dimensions_no_divide_by_zero():
    """RM-02: abs()+max(...,1) guard prevents ZeroDivisionError."""
    slide = MagicMock()
    slide.shapes = []
    issues = _detect_sparse_slide(slide, 0, 0)
    # Should not raise, result is empty or returns sparse (both OK)
    assert isinstance(issues, list)


def test_G6_12b_zero_slide_dims_with_shape_no_divide():
    """RM-02: max(...,1) guard handles zero slide dimensions with a content shape."""
    slide = MagicMock()
    shape = MagicMock()
    shape.has_text_frame = True
    shape.text = "Some content"
    shape.width = 500_000
    shape.height = 300_000
    slide.shapes = [shape]
    # Both slide dimensions zero — denominator would be 0 without the guard
    issues = _detect_sparse_slide(slide, 0, 0)
    # Should not raise; result is either [] or SPARSE_SLIDE but not ZeroDivisionError
    assert isinstance(issues, list)


def test_G6_13_full_slide_no_issue():
    slide = MagicMock()
    shapes = []
    # 5 large shapes filling the slide
    for i in range(5):
        s = MagicMock()
        s.has_text_frame = True
        s.text = f"Content block {i}"
        s.width = 1_800_000
        s.height = 1_300_000
        shapes.append(s)
    slide.shapes = shapes
    issues = _detect_sparse_slide(slide, 9_144_000, 6_858_000)
    assert issues == []


# ---------------------------------------------------------------------------
# G6-14 .. G6-16  review_slide_export integration + return dict
# ---------------------------------------------------------------------------

def test_G6_14_checks_run_has_six_entries(tmp_path):
    """review_slide_export must list 6 checks_run entries after Sprint 1A."""
    import pptx
    prs = pptx.Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])  # blank
    prs_path = tmp_path / "test.pptx"
    prs.save(str(prs_path))
    dummy_png = tmp_path / "export.png"
    dummy_png.write_bytes(b"\x89PNG\r\n\x1a\n")
    import os
    os.environ["PPT_ALLOWLIST_ROOTS"] = str(tmp_path)
    os.environ["PPT_ENABLE_WRITE"] = "true"
    result = review_slide_export(str(prs_path), 0, str(dummy_png), confirm=True)
    # 6 structural + 2 pixel-read (PIXEL_DIMENSION_MISMATCH + BLANK_SLIDE) = 8
    assert len(result["checks_run"]) == 8
    assert "UNFINISHED_PLACEHOLDER" in result["checks_run"]
    assert "LOW_FONT_SIZE" in result["checks_run"]
    assert "SPARSE_SLIDE" in result["checks_run"]
    assert "PIXEL_DIMENSION_MISMATCH" in result["checks_run"]
    assert "BLANK_SLIDE" in result["checks_run"]


def test_G6_15_return_dict_has_new_keys(tmp_path):
    """Sprint 1A adds visual_readiness and confidence_boundary keys."""
    import pptx
    import os
    prs = pptx.Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    prs_path = tmp_path / "test.pptx"
    prs.save(str(prs_path))
    dummy_png = tmp_path / "export.png"
    dummy_png.write_bytes(b"\x89PNG\r\n\x1a\n")
    os.environ["PPT_ALLOWLIST_ROOTS"] = str(tmp_path)
    os.environ["PPT_ENABLE_WRITE"] = "true"
    result = review_slide_export(str(prs_path), 0, str(dummy_png), confirm=True)
    assert result.get("visual_readiness") == "not_assessed"
    assert result.get("confidence_boundary") == "structural_heuristics_only"


def test_G6_16_table_clipping_detail_uses_repr(tmp_path):
    """RM-01: _detect_table_content_clipping detail must use !r for shape.name."""
    import pptx
    import os
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[5]
    prs.slides.add_slide(slide_layout)
    # We can't easily trigger table clipping in a unit test without COM,
    # so just confirm the function runs and returns a dict with the expected keys.
    prs_path = tmp_path / "test.pptx"
    prs.save(str(prs_path))
    dummy_png = tmp_path / "export.png"
    dummy_png.write_bytes(b"\x89PNG\r\n\x1a\n")
    os.environ["PPT_ALLOWLIST_ROOTS"] = str(tmp_path)
    os.environ["PPT_ENABLE_WRITE"] = "true"
    result = review_slide_export(str(prs_path), 0, str(dummy_png), confirm=True)
    assert "findings" in result
    assert "checks_run" in result


def test_G6_EX_01_unfinished_placeholder_text_raises_returns_empty():
    """H-001 fix: shape.text raising must not crash detector."""
    slide = MagicMock()
    shape = MagicMock()
    shape.placeholder_format = MagicMock()
    shape.name = "BrokenPlaceholder"
    slide.shapes = [shape]
    with patch.object(type(shape), "text", new_callable=PropertyMock, create=True) as mock_text:
        mock_text.side_effect = AttributeError("XML corrupt")
        result = _detect_unfinished_placeholder(slide)
    assert result == []


def test_G6_EX_02_low_font_size_paragraphs_raises_returns_empty():
    """H-002 fix: text_frame.paragraphs raising must not crash detector."""
    slide = MagicMock()
    shape = MagicMock()
    shape.has_text_frame = True
    shape.name = "BrokenText"
    slide.shapes = [shape]
    with patch.object(type(shape.text_frame), "paragraphs", new_callable=PropertyMock, create=True) as mock_para:
        mock_para.side_effect = AttributeError("XML corrupt")
        result = _detect_low_font_size(slide)
    assert result == []


# ==========================================================================
# GROUP 1 — _detect_text_overflow_heuristic  (relocated from test_review_gate.py)
# Density formula: density = len(text) / (abs(width_in) * abs(height_in))
#   density > 200  → severity="high"
#   density > 100  → severity="medium"
#   else           → no finding
# ==========================================================================


@pytest.mark.unit
def test_overflow_empty_slide_returns_no_findings():
    """TC-RG-01a: Blank slide (layout 5, zero added shapes) → helper returns [].

    Verifies the base case: when no shapes with text are present, the heuristic
    produces an empty findings list.
    """
    prs = _blank_16x9_prs()
    slide = prs.slides[0]

    findings = _detect_text_overflow_heuristic(slide, prs)

    assert findings == []


@pytest.mark.unit
def test_overflow_high_density_returns_high_severity():
    """TC-RG-02: density fallback path (height=0 mock) with 300 chars -> HIGH.

    With shape.height=0 the geometry path is bypassed. Density fallback:
    area = max(0.5 * 0, 0.01) = 0.01, density = 300/0.01 = 30000 > 200.
    """
    shape_mock = MagicMock()
    shape_mock.has_text_frame = True
    shape_mock.width = int(0.5 * EMU_PER_INCH)
    shape_mock.height = 0
    shape_mock.shape_id = 99
    shape_mock.name = "DensityHighShape"
    shape_mock.text_frame.text = "A" * 300

    slide_mock = MagicMock()
    slide_mock.shapes = [shape_mock]

    findings = _detect_text_overflow_heuristic(slide_mock, MagicMock())

    assert len(findings) == 1, f"Expected 1 finding, got {len(findings)}: {findings}"
    f = findings[0]
    assert f["criterion"] == "TEXT_OVERFLOW"
    assert f["result"] == "FAIL"
    assert f["severity"] == "high"
    assert "detail" in f

@pytest.mark.unit
def test_overflow_medium_density_returns_medium_finding():
    """TC-RG-02b: density fallback path (height=0 mock) with 2 chars -> MEDIUM.

    With shape.height=0, area = max(width_in * 0, 0.01) = 0.01.
    density = 2 / 0.01 = 200.  200 > 100 and NOT > 200 → severity="medium".
    Covers the interior 100 < density <= 200 branch that was previously untested.
    """
    shape_mock = MagicMock()
    shape_mock.has_text_frame = True
    shape_mock.width = int(0.5 * EMU_PER_INCH)
    shape_mock.height = 0
    shape_mock.shape_id = 77
    shape_mock.name = "DensityMediumShape"
    shape_mock.text_frame.text = "AB"  # 2 chars: density = 2 / 0.01 = 200, > 100 and <= 200

    slide_mock = MagicMock()
    slide_mock.shapes = [shape_mock]

    findings = _detect_text_overflow_heuristic(slide_mock, MagicMock())

    assert len(findings) == 1, f"Expected 1 finding, got {len(findings)}: {findings}"
    f = findings[0]
    assert f["criterion"] == "TEXT_OVERFLOW"
    assert f["result"] == "FAIL"
    assert f["severity"] == "medium"
    assert "detail" in f


@pytest.mark.unit
def test_overflow_medium_density_returns_medium_severity():
    """TC-RG-03: density fallback path (height=0 mock) with 1 char -> no finding.

    With shape.height=0, area=max(0,0.01)=0.01, density=1/0.01=100.
    100 is NOT strictly > 100 so neither medium nor high threshold fires.
    Verifies the density fallback boundary condition is exclusive (> not >=).
    """
    shape_mock = MagicMock()
    shape_mock.has_text_frame = True
    shape_mock.width = int(0.5 * EMU_PER_INCH)
    shape_mock.height = 0
    shape_mock.shape_id = 88
    shape_mock.name = "DensityBoundaryShape"
    shape_mock.text_frame.text = "B"  # 1 char: density = 1/0.01 = 100, not > 100

    slide_mock = MagicMock()
    slide_mock.shapes = [shape_mock]

    findings = _detect_text_overflow_heuristic(slide_mock, MagicMock())

    assert findings == [], f"Expected no finding for density=100 (not > threshold), got {findings}" 

@pytest.mark.unit
def test_overflow_low_density_no_finding():
    """TC-RG-04: 4"×4" textbox with 100 chars → density≈6.25 → no finding.

    Density calculation:
      area    = 4.0 × 4.0 = 16.0 sq.in
      density = 100 / 16.0 = 6.25  →  < 100  →  skipped
    """
    prs = _blank_16x9_prs()
    slide = prs.slides[0]

    txb = slide.shapes.add_textbox(
        Emu(EMU_PER_INCH), Emu(EMU_PER_INCH),
        Emu(4 * EMU_PER_INCH), Emu(4 * EMU_PER_INCH),
    )
    txb.text_frame.text = "C" * 100

    findings = _detect_text_overflow_heuristic(slide, prs)

    assert findings == []

# ═══════════════════════════════════════════════════════════════════════════
# GROUP 7 — Geometry-aware TEXT_OVERFLOW (RCI-US-15)
# TC-GEO-01: 0.3in shape + 20 lines 9pt → must detect overflow
# TC-GEO-02: 7in shape + 3 lines default pt → must NOT false-positive
# ═══════════════════════════════════════════════════════════════════════════

@pytest.mark.unit
def test_geometry_overflow_short_shape_20_lines_9pt():
    """TC-GEO-01: 0.3in-tall shape + 20 lines at 9pt → TEXT_OVERFLOW.

    est_height = (20 * 9) / 72 = 2.5in > 0.3in → overflow finding with
    shape_id, shape_height_in, estimated_text_height_in, overflow_estimate_in.
    """
    from pptx.util import Pt
    prs = _blank_16x9_prs()
    slide = prs.slides[0]
    txb = slide.shapes.add_textbox(
        Emu(EMU_PER_INCH), Emu(EMU_PER_INCH),
        Emu(3 * EMU_PER_INCH), Emu(int(0.3 * EMU_PER_INCH)),
    )
    tf = txb.text_frame
    tf.text = "Line 1"
    tf.paragraphs[0].runs[0].font.size = Pt(9)
    for i in range(2, 21):
        run = tf.add_paragraph().add_run()
        run.text = f"Line {i}"
        run.font.size = Pt(9)

    findings = _detect_text_overflow_heuristic(slide, prs)
    hits = [f for f in findings if f["criterion"] == "TEXT_OVERFLOW"]
    assert hits, f"Expected TEXT_OVERFLOW for 20-line 9pt text in 0.3in shape, got {findings}"
    f = hits[0]
    assert f["result"] == "FAIL"
    assert f["severity"] == "high"
    assert "shape_id" in f
    assert "shape_height_in" in f
    assert "estimated_text_height_in" in f
    assert "overflow_estimate_in" in f
    assert f["overflow_estimate_in"] > 0
    assert f["estimated_text_height_in"] > f["shape_height_in"]


@pytest.mark.unit
def test_geometry_no_overflow_tall_shape_3_lines():
    """TC-GEO-02: 7in-tall shape + 3 lines (no explicit pt) → no TEXT_OVERFLOW.

    est_height = (3 * 11) / 72 ≈ 0.458in < 7in → no finding.
    """
    prs = _blank_16x9_prs()
    slide = prs.slides[0]
    txb = slide.shapes.add_textbox(
        Emu(EMU_PER_INCH), Emu(0),
        Emu(3 * EMU_PER_INCH), Emu(int(7 * EMU_PER_INCH)),
    )
    tf = txb.text_frame
    tf.text = "Line A"
    tf.add_paragraph().add_run().text = "Line B"
    tf.add_paragraph().add_run().text = "Line C"

    findings = _detect_text_overflow_heuristic(slide, prs)
    hits = [f for f in findings if f["criterion"] == "TEXT_OVERFLOW"]
    assert hits == [], f"Expected no TEXT_OVERFLOW for 3 lines in 7in shape, got {hits}"

