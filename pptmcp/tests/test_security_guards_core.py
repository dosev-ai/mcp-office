"""
Security guard tests for pptmcp — H-001 through H-012 regression suite.
All tests are @pytest.mark.unit (no live server, no COM required).
"""
from __future__ import annotations

from unittest.mock import MagicMock, patch

import pytest
from fastmcp.exceptions import ToolError
from pptx import Presentation as _Prs

from pptmcp.presentation_pptx import (
    ValidationError,
    _check_path,
    _load_prs,
    _find_blank_layout_idx,
    _prs_cache,
)
from pptmcp.shapes_pptx import manage_hyperlinks
from pptmcp.review_pptx import (
    _detect_shapes_outside_slide,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_pptx(path) -> str:
    """Create a 1-slide .pptx with one text shape and return path as str."""
    prs = _Prs()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Test"
    slide.placeholders[1].text = "Body"
    prs.save(str(path))
    return str(path)


# ---------------------------------------------------------------------------
# H-001: PPTMCPError from COM tool wrappers propagates as ToolError
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH001ComToolErrorPropagation:
    """H-001: PPTMCPError raised inside COM wrappers must be re-raised as ToolError."""

    def test_export_slide_ppterror_becomes_tool_error(self, tmp_path, monkeypatch):
        """PPTMCPError from ppt_export_slide propagates as ValueError via export_slide."""
        from unittest.mock import MagicMock
        from pptmcp.presentation_pptx import PPTMCPError
        from pptmcp import server
        from pptmcp import _server_handlers_com as _cmod

        monkeypatch.setattr(_cmod, "_COM_AVAILABLE", True)
        monkeypatch.setattr(_cmod, "_pcom", MagicMock())
        with patch.object(_cmod._pcom, "ppt_export_slide", side_effect=PPTMCPError("COM failed")):
            with pytest.raises(ToolError, match="COM failed"):
                server.export_slide(
                    file_path="x.pptx",
                    slide_index=0,
                    fmt="png",
                    output_path="out.png",
                    confirm=True,
                )

    def test_export_deck_ppterror_becomes_tool_error(self, tmp_path, monkeypatch):
        """PPTMCPError from ppt_export_deck propagates as ValueError via export_deck."""
        from unittest.mock import MagicMock
        from pptmcp.presentation_pptx import PPTMCPError
        from pptmcp import server
        from pptmcp import _server_handlers_com as _cmod

        monkeypatch.setattr(_cmod, "_COM_AVAILABLE", True)
        monkeypatch.setattr(_cmod, "_pcom", MagicMock())
        with patch.object(_cmod._pcom, "ppt_export_deck", side_effect=PPTMCPError("deck error")):
            with pytest.raises(ToolError, match="deck error"):
                server.export_deck(
                    file_path="x.pptx",
                    fmt="pdf",
                    output_path="out.pdf",
                    confirm=True,
                )

    def test_export_slides_to_stamped_ppterror_becomes_tool_error(self, monkeypatch):
        """PPTMCPError from export_slides_to_stamped_dir propagates as ValueError."""
        from unittest.mock import MagicMock
        from pptmcp.presentation_pptx import PPTMCPError
        from pptmcp import server
        from pptmcp import _server_handlers_com as _cmod

        monkeypatch.setattr(_cmod, "_COM_AVAILABLE", True)
        monkeypatch.setattr(_cmod, "_pcom", MagicMock())
        with patch.object(
            _cmod._pcom, "export_slides_to_stamped_dir", side_effect=PPTMCPError("stamped error")
        ):
            with pytest.raises(ToolError, match="stamped error"):
                server.export_slides_to_stamped_dir(
                    path="x.pptx",
                    base_dir="/tmp",
                    confirm=True,
                )


# ---------------------------------------------------------------------------
# H-004: _check_path raises ValidationError when file size exceeds limit
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH004FileSizeLimit:
    """H-004: _check_path / _check_image_path enforce file size limits."""

    def test_check_path_rejects_oversized_file(self, tmp_path, monkeypatch):
        """_check_path raises ValidationError if file.stat().st_size > _MAX_PPTX_BYTES."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        monkeypatch.setenv("PPT_MAX_FILE_MB", "1")  # 1 MB limit

        path = tmp_path / "big.pptx"
        path.write_bytes(b"\x00" * 100)  # tiny file

        # Patch stat to simulate a file that exceeds the limit
        class _FakeStat:
            st_size = 2 * 1024 * 1024  # 2 MB

        with patch("pptmcp.presentation_pptx._MAX_PPTX_BYTES", 1 * 1024 * 1024):
            with patch.object(type(path.resolve()), "stat", return_value=_FakeStat()):
                with pytest.raises(ValidationError, match="size limit"):
                    _check_path(str(path))

    def test_check_path_allows_normally_sized_file(self, tmp_path, monkeypatch):
        """_check_path does not raise for files under the size limit."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        path = tmp_path / "ok.pptx"
        # Create a minimal valid pptx
        prs = _Prs()
        prs.save(str(path))

        # Should not raise (file is tiny, well under default 256 MB limit)
        result = _check_path(str(path))
        assert result == path.resolve()

    def test_check_path_missing_file_not_raises_during_size_check(self, tmp_path, monkeypatch):
        """_check_path does not raise FileNotFoundError during size check for non-existent files."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        # Non-existent file - size check should pass silently (existence checked later by pptx)
        path = tmp_path / "nonexistent.pptx"
        # Should raise ValidationError for wrong extension or not found, NOT FileNotFoundError
        try:
            _check_path(str(path))
        except ValidationError:
            pass  # Expected — allowlist or extension error
        except FileNotFoundError:
            pytest.fail("_check_path should not raise FileNotFoundError during size check")


# ---------------------------------------------------------------------------
# H-005: _detect_shapes_outside_slide skips shapes with None position attributes
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH005NonePositionGuard:
    """H-005: _detect_shapes_outside_slide must not raise TypeError on None coordinates."""

    def _make_shape(self, shape_id, left, top, width, height):
        s = MagicMock()
        s.shape_id = shape_id
        s.name = f"shape_{shape_id}"
        s.left = left
        s.top = top
        s.width = width
        s.height = height
        return s

    def test_none_left_is_skipped(self):
        """Shape with left=None must be skipped without TypeError."""
        slide = MagicMock()
        slide.shapes = [self._make_shape(1, None, 0, 100, 100)]
        # Should not raise
        result = _detect_shapes_outside_slide(slide, 9144000, 6858000)
        assert result == []

    def test_none_top_is_skipped(self):
        """Shape with top=None must be skipped without TypeError."""
        slide = MagicMock()
        slide.shapes = [self._make_shape(2, 0, None, 100, 100)]
        result = _detect_shapes_outside_slide(slide, 9144000, 6858000)
        assert result == []

    def test_none_width_is_skipped(self):
        """Shape with width=None must be skipped without TypeError."""
        slide = MagicMock()
        slide.shapes = [self._make_shape(3, 0, 0, None, 100)]
        result = _detect_shapes_outside_slide(slide, 9144000, 6858000)
        assert result == []

    def test_fully_none_coords_is_skipped(self):
        """Shape with all None coords must be skipped without TypeError."""
        slide = MagicMock()
        slide.shapes = [self._make_shape(4, None, None, None, None)]
        result = _detect_shapes_outside_slide(slide, 9144000, 6858000)
        assert result == []

    def test_valid_shape_outside_slide_is_detected(self):
        """A shape outside the slide boundary with valid coords must still be detected."""
        slide = MagicMock()
        # Shape extends 914400 EMU beyond the right edge
        slide.shapes = [self._make_shape(5, 9000000, 0, 1058400, 914400)]
        result = _detect_shapes_outside_slide(slide, 9144000, 6858000)
        assert len(result) == 1
        assert result[0]["criterion"] == "SHAPE_OUTSIDE_SLIDE"

    def test_mixed_none_and_valid_shapes(self):
        """None-coord shapes are skipped; valid out-of-bounds shapes are detected."""
        slide = MagicMock()
        slide.shapes = [
            self._make_shape(6, None, 0, 100, 100),   # skipped
            self._make_shape(7, -100, 0, 100, 100),   # detected (left < 0)
        ]
        result = _detect_shapes_outside_slide(slide, 9144000, 6858000)
        assert len(result) == 1
        assert result[0]["criterion"] == "SHAPE_OUTSIDE_SLIDE"


# ---------------------------------------------------------------------------
# H-006: _load_prs evicts oldest entry when cache reaches _PRS_CACHE_MAX
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH006CacheEviction:
    """H-006: _load_prs respects _PRS_CACHE_MAX and evicts oldest entry."""

    def test_cache_evicts_when_full(self, tmp_path, monkeypatch):
        """When cache reaches _PRS_CACHE_MAX, the oldest entry is evicted."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        _prs_cache.clear()

        fake_max = 3
        with patch("pptmcp.presentation_pptx._PRS_CACHE_MAX", fake_max):
            # Fill cache to max
            paths = []
            for i in range(fake_max):
                p = tmp_path / f"file{i}.pptx"
                prs = _Prs()
                prs.save(str(p))
                paths.append(p)
                _load_prs(p)

            assert len(_prs_cache) == fake_max
            first_key = str(paths[0].resolve())
            assert first_key in _prs_cache

            # Loading one more should evict the oldest
            p_new = tmp_path / "file_new.pptx"
            prs = _Prs()
            prs.save(str(p_new))
            _load_prs(p_new)

            assert len(_prs_cache) == fake_max
            assert first_key not in _prs_cache, "Oldest entry should have been evicted"
            assert str(p_new.resolve()) in _prs_cache

        _prs_cache.clear()

    def test_cache_under_max_no_eviction(self, tmp_path, monkeypatch):
        """Cache entries below _PRS_CACHE_MAX are never evicted."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        _prs_cache.clear()

        p = tmp_path / "single.pptx"
        prs = _Prs()
        prs.save(str(p))

        _load_prs(p)
        key = str(p.resolve())
        assert key in _prs_cache

        # Load a second time — still in cache, not evicted
        _load_prs(p)
        assert key in _prs_cache

        _prs_cache.clear()


# ---------------------------------------------------------------------------
# H-007: manage_hyperlinks rejects disallowed URL schemes via urlparse
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH007UrlParseSchemeValidation:
    """H-007: manage_hyperlinks uses urlparse for scheme extraction."""

    @pytest.fixture
    def pptx_with_text(self, tmp_path, monkeypatch):
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        path = tmp_path / "test.pptx"
        return _make_pptx(path)

    def test_javascript_scheme_blocked(self, pptx_with_text):
        """javascript: URL is rejected by urlparse-based scheme check."""
        with pytest.raises(ValidationError, match="scheme"):
            manage_hyperlinks(
                pptx_with_text, 0, "add",
                shape_id=2, target_url="javascript:alert(1)", confirm=True
            )

    def test_data_scheme_blocked(self, pptx_with_text):
        """data: URL is rejected by urlparse-based scheme check."""
        with pytest.raises(ValidationError, match="scheme"):
            manage_hyperlinks(
                pptx_with_text, 0, "add",
                shape_id=2, target_url="data:text/html,<b>x</b>", confirm=True
            )

    def test_ftp_scheme_blocked(self, pptx_with_text):
        """ftp: URL is rejected (not in allowed set)."""
        with pytest.raises(ValidationError, match="scheme"):
            manage_hyperlinks(
                pptx_with_text, 0, "add",
                shape_id=2, target_url="ftp://files.example.com/file", confirm=True
            )

    def test_double_slash_no_scheme_blocked(self, pptx_with_text):
        """Protocol-relative URL (//example.com) has empty scheme and is blocked."""
        with pytest.raises(ValidationError, match="scheme"):
            manage_hyperlinks(
                pptx_with_text, 0, "add",
                shape_id=2, target_url="//example.com/path", confirm=True
            )


# ---------------------------------------------------------------------------
# H-008: manage_hyperlinks strips whitespace from target_url before storing
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH008TargetUrlStrip:
    """H-008: manage_hyperlinks strips whitespace from target_url before storing in hyperlink."""

    def test_target_url_stripped_before_store(self, tmp_path, monkeypatch):
        """URL with leading/trailing whitespace is stripped before being stored."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")

        path = tmp_path / "hl_strip.pptx"
        prs = _Prs()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Test"
        slide.placeholders[1].text = "Link text"
        prs.save(str(path))

        result = manage_hyperlinks(
            str(path), 0, "add",
            shape_id=2,
            target_url="  https://example.com  ",
            confirm=True,
        )
        assert result["ok"] is True
        assert result["url"] == "https://example.com"  # H-013: returned url is stripped

        # Verify stored url is stripped by listing hyperlinks
        list_result = manage_hyperlinks(str(path), 0, "list")
        stored_urls = [lnk["url"] for lnk in list_result["links"]]
        assert any(u == "https://example.com" for u in stored_urls), (
            f"Expected stripped URL in stored links, got {stored_urls}"
        )


# ---------------------------------------------------------------------------
# H-009: _find_blank_layout_idx returns correct index for "Blank" named layout
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH009FindBlankLayoutIdx:
    """H-009: _find_blank_layout_idx locates blank layout by name."""

    @staticmethod
    def _layout(name: str):
        m = MagicMock()
        m.name = name
        return m

    def test_finds_blank_by_name(self):
        """Returns index of layout whose name contains 'blank' (case-insensitive)."""
        mock_prs = MagicMock()
        mock_prs.slide_layouts = [
            self._layout("Title Slide"),
            self._layout("Title and Content"),
            self._layout("Blank"),
            self._layout("Two Content"),
        ]
        assert _find_blank_layout_idx(mock_prs) == 2

    def test_finds_blank_case_insensitive(self):
        """Name match is case-insensitive ('BLANK', 'blank', 'Blank' all match)."""
        mock_prs = MagicMock()
        mock_prs.slide_layouts = [
            self._layout("Title Slide"),
            self._layout("BLANK LAYOUT"),
        ]
        assert _find_blank_layout_idx(mock_prs) == 1

    def test_falls_back_to_zero_when_no_blank(self):
        """Returns 0 when no layout contains 'blank' in its name."""
        mock_prs = MagicMock()
        mock_prs.slide_layouts = [
            self._layout("Title Slide"),
            self._layout("Two Content"),
        ]
        assert _find_blank_layout_idx(mock_prs) == 0

    def test_real_presentation_has_blank_layout(self, tmp_path):
        """A default python-pptx Presentation includes a Blank layout."""
        prs = _Prs()
        idx = _find_blank_layout_idx(prs)
        # python-pptx default theme includes a "Blank" layout
        assert prs.slide_layouts[idx].name.lower() == "blank"


# ---------------------------------------------------------------------------
# H-012: null-byte in base_dir raises ValidationError
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH012NullByteBaseDir:
    """H-012: ppt_export_slides_to_stamped rejects null-byte in base_dir."""

    def test_null_byte_in_check_path_raises(self, tmp_path, monkeypatch):
        """_check_path raises ValidationError for paths with null bytes."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        null_path = str(tmp_path / "file\x00.pptx")
        with pytest.raises(ValidationError, match="null"):
            _check_path(null_path)

    def test_null_byte_in_base_dir_raises(self, tmp_path, monkeypatch):
        """ppt_export_slides_to_stamped raises ValidationError for null byte in base_dir."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        monkeypatch.setenv("PPT_ENABLE_COM", "true")
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")

        import sys
        monkeypatch.setattr(sys, "platform", "win32")

        # Patch win32com import inside _check_com_enabled
        with patch.dict("sys.modules", {"win32com": MagicMock(), "win32com.client": MagicMock()}):
            from pptmcp.presentation_com import export_slides_to_stamped_dir

            path = tmp_path / "test.pptx"
            prs = _Prs()
            prs.save(str(path))

            with pytest.raises(ValidationError, match="null"):
                export_slides_to_stamped_dir(
                    path=str(path),
                    base_dir="/tmp/exports\x00injected",
                    confirm=True,
                )

