"""Unit tests for pptmcp table mutation tools.

No COM required — CI-safe. Tests use real python-pptx files.

Coverage:
  US-PPT-01: add_column appends new column; row data preserved
  US-PPT-02: edit_table_cell updates target cell; others unchanged
  US-PPT-03: set_column_width sets the column width in EMU
  US-PPT-04: Gate checks — PPT_ENABLE_WRITE=false and confirm=False both reject
  US-PPT-05: Out-of-range row/col indices return structured errors
"""
from __future__ import annotations

import os
import sys

# Ensure the local branch src is on the path before any installed version.
_SRC = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "src")
if os.path.isdir(_SRC) and _SRC not in sys.path:
    sys.path.insert(0, os.path.realpath(_SRC))

import pytest
from pptx import Presentation
from pptx.util import Inches

_EMU_PER_INCH = 914400


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_pptx_with_table(path, rows: int = 2, cols: int = 3) -> tuple[str, int]:
    """Create a .pptx with a blank slide containing one table; return (path, shape_id)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    left = top = Inches(1)
    width = Inches(6)
    height = Inches(2)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    # Pre-fill cells so we can verify preservation
    tbl = table_shape.table
    for r in range(rows):
        for c in range(cols):
            tbl.cell(r, c).text = f"R{r}C{c}"
    prs.save(str(path))
    return str(path), table_shape.shape_id


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def pptx_with_table(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path, shape_id = _make_pptx_with_table(tmp_path / "table.pptx", rows=2, cols=3)
    return path, shape_id


@pytest.fixture
def pptx_no_write(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path, shape_id = _make_pptx_with_table(tmp_path / "nowrite.pptx")
    return path, shape_id


# ---------------------------------------------------------------------------
# US-PPT-01: add_column
# ---------------------------------------------------------------------------


class TestAddColumn:
    def test_add_column_increases_column_count(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import add_column
        path, shape_id = pptx_with_table
        result = add_column(path, slide_index=0, shape_id=shape_id, width_inches=1.5, confirm=True)
        assert result["status"] == "column_added"
        assert result["column_count"] == 4
        assert result["width_inches"] == 1.5

    def test_add_column_preserves_existing_cell_data(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import add_column, _find_table_shape
        from pptmcp.presentation_pptx import _load_prs, _check_path
        path, shape_id = pptx_with_table
        add_column(path, slide_index=0, shape_id=shape_id, confirm=True)
        # Verify original cells untouched
        prs = _load_prs(_check_path(path))
        table = _find_table_shape(prs, 0, shape_id).table
        assert table.cell(0, 0).text == "R0C0"
        assert table.cell(1, 2).text == "R1C2"

    def test_add_column_new_cells_empty(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import add_column, _find_table_shape
        from pptmcp.presentation_pptx import _load_prs, _check_path
        path, shape_id = pptx_with_table
        add_column(path, slide_index=0, shape_id=shape_id, confirm=True)
        prs = _load_prs(_check_path(path))
        table = _find_table_shape(prs, 0, shape_id).table
        # New column (index 3) should have empty text
        assert table.cell(0, 3).text == ""
        assert table.cell(1, 3).text == ""

    def test_add_column_grid_width_in_emu(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import add_column, _find_table_shape
        from pptmcp.presentation_pptx import _load_prs, _check_path
        from pptx.oxml.ns import qn
        path, shape_id = pptx_with_table
        add_column(path, slide_index=0, shape_id=shape_id, width_inches=2.0, confirm=True)
        prs = _load_prs(_check_path(path))
        shape = _find_table_shape(prs, 0, shape_id)
        tbl = shape.table._tbl
        grid_cols = tbl.find(qn("a:tblGrid")).findall(qn("a:gridCol"))
        last_col_w = int(grid_cols[-1].get("w"))
        assert last_col_w == int(2.0 * _EMU_PER_INCH)


# ---------------------------------------------------------------------------
# US-PPT-02: edit_table_cell
# ---------------------------------------------------------------------------


class TestEditTableCell:
    def test_edit_cell_updates_target(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import edit_table_cell
        path, shape_id = pptx_with_table
        result = edit_table_cell(path, slide_index=0, shape_id=shape_id, row=1, col=2, text="Updated", confirm=True)
        assert result["status"] == "cell_updated"
        assert result["row"] == 1
        assert result["col"] == 2
        assert result["text"] == "Updated"

    def test_edit_cell_preserves_others(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import edit_table_cell, _find_table_shape
        from pptmcp.presentation_pptx import _load_prs, _check_path
        path, shape_id = pptx_with_table
        edit_table_cell(path, slide_index=0, shape_id=shape_id, row=0, col=0, text="Changed", confirm=True)
        prs = _load_prs(_check_path(path))
        table = _find_table_shape(prs, 0, shape_id).table
        assert table.cell(0, 0).text == "Changed"
        # Other cells untouched
        assert table.cell(0, 1).text == "R0C1"
        assert table.cell(1, 0).text == "R1C0"

    def test_edit_cell_row_out_of_range(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import edit_table_cell
        from pptmcp._presentation_governance import PPTMCPError
        path, shape_id = pptx_with_table
        with pytest.raises((ValueError, PPTMCPError)):
            edit_table_cell(path, slide_index=0, shape_id=shape_id, row=99, col=0, text="X", confirm=True)

    def test_edit_cell_col_out_of_range(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import edit_table_cell
        from pptmcp._presentation_governance import PPTMCPError
        path, shape_id = pptx_with_table
        with pytest.raises((ValueError, PPTMCPError)):
            edit_table_cell(path, slide_index=0, shape_id=shape_id, row=0, col=99, text="X", confirm=True)


# ---------------------------------------------------------------------------
# US-PPT-03: set_column_width
# ---------------------------------------------------------------------------


class TestSetColumnWidth:
    def test_set_column_width_updates_emu(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import set_column_width, _find_table_shape
        from pptmcp.presentation_pptx import _load_prs, _check_path
        path, shape_id = pptx_with_table
        result = set_column_width(path, slide_index=0, shape_id=shape_id, col=1, width_inches=2.5, confirm=True)
        assert result["status"] == "column_width_set"
        assert result["col"] == 1
        assert result["width_inches"] == 2.5
        prs = _load_prs(_check_path(path))
        table = _find_table_shape(prs, 0, shape_id).table
        assert table.columns[1].width == int(2.5 * _EMU_PER_INCH)

    def test_set_column_width_col_out_of_range(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import set_column_width
        path, shape_id = pptx_with_table
        with pytest.raises(ValueError, match="out of range"):
            set_column_width(path, slide_index=0, shape_id=shape_id, col=99, width_inches=1.0, confirm=True)


# ---------------------------------------------------------------------------
# US-PPT-04: gate checks
# ---------------------------------------------------------------------------


class TestGateChecks:
    def test_add_column_write_disabled_raises(self, pptx_no_write):
        from pptmcp._pptx_table_mutations import add_column
        from pptmcp._presentation_governance import NotAllowedError
        path, shape_id = pptx_no_write
        with pytest.raises(NotAllowedError):
            add_column(path, slide_index=0, shape_id=shape_id, confirm=True)

    def test_add_column_confirm_false_raises(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import add_column
        from pptmcp._presentation_governance import ValidationError
        path, shape_id = pptx_with_table
        with pytest.raises(ValidationError, match="confirm"):
            add_column(path, slide_index=0, shape_id=shape_id, confirm=False)

    def test_edit_cell_confirm_false_raises(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import edit_table_cell
        from pptmcp._presentation_governance import ValidationError
        path, shape_id = pptx_with_table
        with pytest.raises(ValidationError, match="confirm"):
            edit_table_cell(path, slide_index=0, shape_id=shape_id, row=0, col=0, text="X", confirm=False)

    def test_set_column_width_confirm_false_raises(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import set_column_width
        from pptmcp._presentation_governance import ValidationError
        path, shape_id = pptx_with_table
        with pytest.raises(ValidationError, match="confirm"):
            set_column_width(path, slide_index=0, shape_id=shape_id, col=0, width_inches=1.0, confirm=False)

    def test_edit_cell_write_disabled_raises(self, pptx_no_write):
        from pptmcp._pptx_table_mutations import edit_table_cell
        from pptmcp._presentation_governance import NotAllowedError
        path, shape_id = pptx_no_write
        with pytest.raises(NotAllowedError):
            edit_table_cell(path, slide_index=0, shape_id=shape_id, row=0, col=0, text="X", confirm=True)


# ---------------------------------------------------------------------------
# US-PPT-05: non-table shape and bad shape_id
# ---------------------------------------------------------------------------


class TestErrorCases:
    def test_shape_not_found_raises(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import edit_table_cell
        path, shape_id = pptx_with_table
        with pytest.raises(ValueError, match="No shape"):
            edit_table_cell(path, slide_index=0, shape_id=9999, row=0, col=0, text="X", confirm=True)

    def test_slide_out_of_range_raises(self, pptx_with_table):
        from pptmcp._pptx_table_mutations import add_column
        path, shape_id = pptx_with_table
        with pytest.raises(ValueError, match="out of range"):
            add_column(path, slide_index=99, shape_id=shape_id, confirm=True)


# ---------------------------------------------------------------------------
# Extension A: set_table_style row_heights_pt
# ---------------------------------------------------------------------------


@pytest.fixture
def pptx_with_table_3rows(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path, shape_id = _make_pptx_with_table(tmp_path / "table3r.pptx", rows=3, cols=2)
    return path, shape_id


class TestSetTableStyleRowHeightsPt:
    """US-PPT-06: set_table_style row_heights_pt sets per-row heights."""

    @staticmethod
    def _get_table(prs, slide_index: int, shape_id: int):
        """Find the table shape by shape_id, not by position index."""
        for s in prs.slides[slide_index].shapes:
            if s.shape_id == shape_id and s.has_table:
                return s.table
        raise AssertionError(f"No table shape with id {shape_id} on slide {slide_index}")

    def test_full_list_sets_all_row_heights(self, pptx_with_table_3rows):
        from pptmcp._shapes_geometry import set_table_style
        from pptmcp.presentation_pptx import _load_prs, _check_path
        from pptx.util import Pt
        path, shape_id = pptx_with_table_3rows
        result = set_table_style(
            path, slide_index=0, shape_id=shape_id,
            row_heights_pt=[20.0, 30.0, 40.0],
            confirm=True,
        )
        assert result["rows_resized"] == 3
        prs = _load_prs(_check_path(path))
        table = self._get_table(prs, 0, shape_id)
        assert table.rows[0].height == pytest.approx(int(Pt(20.0)), abs=1)
        assert table.rows[1].height == pytest.approx(int(Pt(30.0)), abs=1)
        assert table.rows[2].height == pytest.approx(int(Pt(40.0)), abs=1)

    def test_partial_list_leaves_tail_rows_unchanged(self, pptx_with_table_3rows):
        from pptmcp._shapes_geometry import set_table_style
        from pptmcp.presentation_pptx import _load_prs, _check_path
        from pptx.util import Pt
        path, shape_id = pptx_with_table_3rows
        # Read original row[2] height before mutation
        prs0 = _load_prs(_check_path(path))
        original_row2_height = self._get_table(prs0, 0, shape_id).rows[2].height
        # Apply row_heights_pt for only the first 2 rows
        result = set_table_style(
            path, slide_index=0, shape_id=shape_id,
            row_heights_pt=[18.0, 24.0],
            confirm=True,
        )
        assert result["rows_resized"] == 2
        prs = _load_prs(_check_path(path))
        table = self._get_table(prs, 0, shape_id)
        assert table.rows[0].height == pytest.approx(int(Pt(18.0)), abs=1)
        assert table.rows[1].height == pytest.approx(int(Pt(24.0)), abs=1)
        # Row 2 must be unchanged
        assert table.rows[2].height == original_row2_height

    def test_none_row_heights_pt_is_noop(self, pptx_with_table_3rows):
        from pptmcp._shapes_geometry import set_table_style
        from pptmcp.presentation_pptx import _load_prs, _check_path
        path, shape_id = pptx_with_table_3rows
        prs0 = _load_prs(_check_path(path))
        original_heights = [r.height for r in self._get_table(prs0, 0, shape_id).rows]
        result = set_table_style(
            path, slide_index=0, shape_id=shape_id,
            row_heights_pt=None,
            confirm=True,
        )
        assert result["rows_resized"] == 0
        prs = _load_prs(_check_path(path))
        heights_after = [r.height for r in self._get_table(prs, 0, shape_id).rows]
        assert heights_after == original_heights

    def test_empty_list_is_noop(self, pptx_with_table_3rows):
        from pptmcp._shapes_geometry import set_table_style
        from pptmcp.presentation_pptx import _load_prs, _check_path
        path, shape_id = pptx_with_table_3rows
        prs0 = _load_prs(_check_path(path))
        original_heights = [r.height for r in self._get_table(prs0, 0, shape_id).rows]
        result = set_table_style(
            path, slide_index=0, shape_id=shape_id,
            row_heights_pt=[],
            confirm=True,
        )
        assert result["rows_resized"] == 0
        prs = _load_prs(_check_path(path))
        heights_after = [r.height for r in self._get_table(prs, 0, shape_id).rows]
        assert heights_after == original_heights

    def test_row_heights_pt_zero_raises_validation_error(self, pptx_with_table_3rows):
        from pptmcp._shapes_geometry import set_table_style
        from pptmcp._presentation_governance import ValidationError
        path, shape_id = pptx_with_table_3rows
        with pytest.raises(ValidationError):
            set_table_style(
                path, slide_index=0, shape_id=shape_id,
                row_heights_pt=[20.0, 0.0, 30.0],
                confirm=True,
            )

    def test_row_heights_pt_negative_raises_validation_error(self, pptx_with_table_3rows):
        from pptmcp._shapes_geometry import set_table_style
        from pptmcp._presentation_governance import ValidationError
        path, shape_id = pptx_with_table_3rows
        with pytest.raises(ValidationError):
            set_table_style(
                path, slide_index=0, shape_id=shape_id,
                row_heights_pt=[-5.0],
                confirm=True,
            )


class TestGetShapeRowsHeightEmu:
    """US-PPT-07: get_shape returns rows[].height_emu for table shapes."""

    def test_get_shape_table_includes_rows_key(self, pptx_with_table_3rows):
        from pptmcp._presentation_read import get_shape
        path, shape_id = pptx_with_table_3rows
        result = get_shape(path, slide_index=0, shape_id=shape_id)
        assert "rows" in result
        assert len(result["rows"]) == 3

    def test_get_shape_table_rows_have_height_emu(self, pptx_with_table_3rows):
        from pptmcp._presentation_read import get_shape
        path, shape_id = pptx_with_table_3rows
        result = get_shape(path, slide_index=0, shape_id=shape_id)
        for row_entry in result["rows"]:
            assert "height_emu" in row_entry
            assert isinstance(row_entry["height_emu"], (int, type(None)))

    def test_get_shape_non_table_has_no_rows_key(self, tmp_path, monkeypatch):
        from pptx import Presentation
        from pptmcp._presentation_read import get_shape
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # First shape is a text placeholder
        pptx_path = tmp_path / "no_table.pptx"
        prs.save(str(pptx_path))
        # Reload to get live shape
        from pptmcp.presentation_pptx import _load_prs, _check_path
        prs2 = _load_prs(_check_path(str(pptx_path)))
        first_shape = prs2.slides[0].shapes[0]
        if not first_shape.has_table:
            result = get_shape(str(pptx_path), slide_index=0, shape_id=first_shape.shape_id)
            assert "rows" not in result

    def test_set_and_read_row_heights_roundtrip(self, pptx_with_table_3rows):
        """Set row heights then read them back via get_shape."""
        from pptmcp._shapes_geometry import set_table_style
        from pptmcp._presentation_read import get_shape
        from pptx.util import Pt
        path, shape_id = pptx_with_table_3rows
        heights_pt = [15.0, 25.0, 35.0]
        set_table_style(
            path, slide_index=0, shape_id=shape_id,
            row_heights_pt=heights_pt,
            confirm=True,
        )
        result = get_shape(path, slide_index=0, shape_id=shape_id)
        assert "rows" in result
        for i, (row_entry, expected_pt) in enumerate(zip(result["rows"], heights_pt)):
            expected_emu = int(Pt(expected_pt))
            assert row_entry["height_emu"] == pytest.approx(expected_emu, abs=1), (
                f"Row {i} height_emu mismatch: got {row_entry['height_emu']}, expected ~{expected_emu}"
            )
