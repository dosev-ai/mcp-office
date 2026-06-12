"""READ-ONLY tool tests — targets presentation_pptx.py only. No server.py imports."""

import pytest
from pptx import Presentation as _Prs
from pptmcp.presentation_pptx import (
    PPTMCPError,
    ValidationError,
    capabilities,
    read_presentation,
    get_presentation_metadata,
    list_slides,
    list_layouts,
    read_slide,
    read_speaker_notes,
    list_shapes,
    get_shape,
    extract_tables,
    extract_images,
    export_slide_as_text,
    extract_presentation_text,
)

# NOTE: No _prs_cache import — autouse fixture is in conftest.py


# ── GROUP 0: capabilities() ───────────────────────────────────────────────

@pytest.mark.unit
def test_capabilities_phase1():
    assert capabilities()["phase"] == "2.1"


@pytest.mark.unit
def test_capabilities_backend():
    assert capabilities()["backend"] == "python-pptx"


@pytest.mark.unit
def test_capabilities_tool_count():
    assert len(capabilities()["tools"]) == 51  # 51 base tools (COM-conditional excluded)


@pytest.mark.unit
def test_capabilities_all_tools():
    expected = {
        "capabilities", "read_presentation", "get_presentation_metadata",
        "list_slides", "list_layouts", "read_slide", "read_speaker_notes",
        "list_shapes", "get_shape", "extract_tables", "extract_images",
        "create_presentation", "add_slide", "edit_text_placeholder", "set_speaker_notes",
        "replace_slide_text", "delete_shape", "insert_image", "reorder_slides",
        "delete_slide", "export_slide_as_text", "save",
        "extract_presentation_text", "manage_hyperlinks",
        "add_content", "set_format", "set_table_style", "copy_slide", "add_hyperlink",
        "detect_overlapping_shapes", "clear_slide_content", "apply_slide_layout",
        "remove_empty_placeholders", "review_slide_export",
        "declare_slide_contract",
        "validate_contract", "check_presentation_against_contract", "produce_evidence_bundle",
        "export_slide", "export_deck", "export_changed_slides_only", "export_slides_to_stamped_dir",
        "get_presentation_context",
        "slide", "shape", "export",
        "manage_comments",
        "batch_set_text",
        "pptmcp_add_column",
        "pptmcp_edit_table_cell",
        "pptmcp_set_column_width",
    }
    actual = {e["tool"] for e in capabilities()["tools"]}
    assert expected == actual


# ── GROUP 1: _check_path / allowlist ──────────────────────────────────────

@pytest.mark.unit
def test_path_not_in_allowlist(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    outside = tmp_path.parent / "outside_dir" / "file.pptx"
    with pytest.raises(ValidationError, match="not in allowlist"):
        list_slides(str(outside))


@pytest.mark.unit
def test_no_allowlist_configured(monkeypatch, tmp_path):
    monkeypatch.delenv("PPT_ALLOWLIST_ROOTS", raising=False)
    with pytest.raises(ValidationError, match="No allowlist"):
        list_slides(str(tmp_path / "test.pptx"))


@pytest.mark.unit
def test_unsupported_extension(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    pdf_path = tmp_path / "data.pdf"
    pdf_path.write_bytes(b"fake")
    with pytest.raises(ValidationError, match="Unsupported extension"):
        list_slides(str(pdf_path))


@pytest.mark.unit
def test_path_in_allowlist_passes(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "test.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[0])
    p.save(str(path))
    result = list_slides(str(path))
    assert isinstance(result, list)


# ── GROUP 2: list_slides() ────────────────────────────────────────────────

@pytest.mark.unit
def test_list_slides_returns_list(sample_pptx):
    result = list_slides(sample_pptx)
    assert isinstance(result, list)
    assert len(result) >= 1


@pytest.mark.unit
def test_list_slides_first_slide_index(sample_pptx):
    result = list_slides(sample_pptx)
    assert result[0]["slide_index"] == 0


@pytest.mark.unit
def test_list_slides_outside_allowlist(tmp_path, monkeypatch):
    root = tmp_path / "allowed"
    root.mkdir()
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(root))
    outside = tmp_path / "outside" / "test.pptx"
    with pytest.raises(ValidationError, match="not in allowlist"):
        list_slides(str(outside))


# ── GROUP 3: read_presentation() ─────────────────────────────────────────

@pytest.mark.unit
def test_read_presentation_slide_count(sample_pptx):
    result = read_presentation(sample_pptx)
    assert result["slide_count"] >= 1


@pytest.mark.unit
def test_read_presentation_slides_is_list(sample_pptx):
    result = read_presentation(sample_pptx)
    assert isinstance(result["slides"], list)
    assert len(result["slides"]) == result["slide_count"]


# ── GROUP 4: get_presentation_metadata() ─────────────────────────────────

@pytest.mark.unit
def test_get_metadata_slide_count(sample_pptx):
    result = get_presentation_metadata(sample_pptx)
    assert result["slide_count"] >= 1


@pytest.mark.unit
def test_get_metadata_has_keys(sample_pptx):
    result = get_presentation_metadata(sample_pptx)
    for key in ("title", "author", "subject", "slide_count"):
        assert key in result


# ── GROUP 5: read_slide() ─────────────────────────────────────────────────

@pytest.mark.unit
def test_read_slide_shapes_list(sample_pptx):
    result = read_slide(sample_pptx, 0)
    assert "shapes" in result
    assert isinstance(result["shapes"], list)


@pytest.mark.unit
def test_read_slide_out_of_range(sample_pptx):
    with pytest.raises(ValidationError):
        read_slide(sample_pptx, 999)


# ── GROUP 6: read_speaker_notes() ────────────────────────────────────────

@pytest.mark.unit
def test_read_speaker_notes_single_slide(sample_pptx):
    result = read_speaker_notes(sample_pptx, slide_index=0)
    assert isinstance(result, list)


@pytest.mark.unit
def test_read_speaker_notes_all_slides(sample_pptx):
    result = read_speaker_notes(sample_pptx, slide_index=None)
    assert isinstance(result, list)


# ── GROUP 7: list_shapes() ───────────────────────────────────────────────

@pytest.mark.unit
def test_list_shapes_is_list(sample_pptx):
    result = list_shapes(sample_pptx, 0)
    assert isinstance(result, list)


@pytest.mark.unit
def test_list_shapes_shape_keys(sample_pptx):
    result = list_shapes(sample_pptx, 0)
    for shape in result:
        assert "shape_id" in shape
        assert "name" in shape
        assert "shape_type" in shape
        assert "placeholder_idx" in shape


# ── GROUP 8: extract_tables() ────────────────────────────────────────────

@pytest.mark.unit
def test_extract_tables_empty(sample_pptx):
    result = extract_tables(sample_pptx)
    assert isinstance(result, dict)
    assert result["tables"] == []


# ── GROUP 9: extract_images() ────────────────────────────────────────────

@pytest.mark.unit
def test_extract_images_empty(sample_pptx):
    result = extract_images(sample_pptx)
    assert isinstance(result, dict)
    assert result["images"] == []


# ── GROUP 10: export_slide_as_text() ─────────────────────────────────────

@pytest.mark.unit
def test_export_single_slide(sample_pptx):
    result = export_slide_as_text(sample_pptx, slide_index=0)
    assert isinstance(result, list)
    assert len(result) == 1
    assert result[0]["slide_index"] == 0
    assert isinstance(result[0]["texts"], list)


@pytest.mark.unit
def test_export_all_slides(sample_pptx):
    result = export_slide_as_text(sample_pptx, slide_index=None)
    assert isinstance(result, list)
    assert len(result) >= 1


@pytest.mark.unit
def test_export_slide_as_text_includes_table_rows(sample_pptx_with_table):
    result = export_slide_as_text(sample_pptx_with_table, slide_index=0)
    assert "R0C0\tR0C1\tR0C2" in result[0]["texts"]


# ── GROUP 14b: get_shape() ────────────────────────────────────────────────

@pytest.mark.unit
def test_get_shape_returns_detail(sample_pptx):
    shapes = list_shapes(sample_pptx, 0)
    assert shapes, "Expected at least one shape on slide 0"
    shape_id = shapes[0]["shape_id"]
    result = get_shape(sample_pptx, 0, shape_id)
    assert result["shape_id"] == shape_id
    assert "name" in result
    assert "shape_type" in result


@pytest.mark.unit
def test_get_shape_invalid_id(sample_pptx):
    with pytest.raises(ValidationError, match="not found"):
        get_shape(sample_pptx, 0, 99999)


# ── GROUP 18: extract_presentation_text() ───────────────────────────────

@pytest.mark.unit
class TestExtractPresentationText:
    @pytest.mark.unit
    def test_extract_returns_expected_keys(self, sample_pptx):
        result = extract_presentation_text(sample_pptx)
        assert "path" in result
        assert "slide_count" in result
        assert "slides" in result

    @pytest.mark.unit
    def test_extract_slide_count(self, sample_pptx):
        result = extract_presentation_text(sample_pptx)
        assert result["slide_count"] == 1
        assert len(result["slides"]) == 1

    @pytest.mark.unit
    def test_extract_title_captured(self, sample_pptx):
        result = extract_presentation_text(sample_pptx)
        slide = result["slides"][0]
        assert slide["slide_index"] == 0
        # sample_pptx has title "Test Slide" on layout 1 (Title and Content)
        assert slide["title"] == "Test Slide"

    @pytest.mark.unit
    def test_extract_shapes_list_populated(self, sample_pptx):
        result = extract_presentation_text(sample_pptx)
        shapes = result["slides"][0]["shapes"]
        assert isinstance(shapes, list)
        assert len(shapes) >= 1
        for s in shapes:
            assert "shape_id" in s
            assert "name" in s
            assert "text" in s

    @pytest.mark.unit
    def test_extract_presentation_text_includes_table_rows(self, sample_pptx_with_table):
        result = extract_presentation_text(sample_pptx_with_table)
        slide = result["slides"][0]
        table_shapes = [s for s in slide["shapes"] if s.get("shape_type") == "table"]
        assert table_shapes
        assert table_shapes[0]["rows"][0] == ["R0C0", "R0C1", "R0C2"]
        assert "R0C0\tR0C1\tR0C2" in slide["body_text"]

    @pytest.mark.unit
    def test_extract_notes_captured(self, tmp_path, monkeypatch):
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        path = tmp_path / "withnotes.pptx"
        p = _Prs()
        slide = p.slides.add_slide(p.slide_layouts[1])
        slide.shapes.title.text = "Notes Slide"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "My speaker notes"
        p.save(str(path))
        result = extract_presentation_text(str(path))
        assert result["slides"][0]["notes"] == "My speaker notes"

    @pytest.mark.unit
    def test_extract_file_not_found(self, tmp_path, monkeypatch):
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        missing = tmp_path / "nonexistent.pptx"
        with pytest.raises(PPTMCPError):
            extract_presentation_text(str(missing))


# ── GROUP 20: Negative test cases (file not found, corrupted, empty) ───────

@pytest.mark.unit
def test_list_slides_file_not_found(tmp_path, monkeypatch):
    """Test that list_slides raises PPTMCPError when file does not exist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    missing = tmp_path / "nonexistent.pptx"
    with pytest.raises(PPTMCPError, match="Cannot open presentation"):
        list_slides(str(missing))


@pytest.mark.unit
def test_corrupted_pptx_file(tmp_path, monkeypatch):
    """Test that a corrupted/malformed .pptx file raises PPTMCPError."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    # Create a file with .pptx extension but invalid ZIP/OOXML content
    corrupted = tmp_path / "corrupted.pptx"
    corrupted.write_bytes(b"This is not a valid PPTX file; it's just random bytes!")
    with pytest.raises(PPTMCPError, match="Cannot open presentation"):
        list_slides(str(corrupted))


@pytest.mark.unit
def test_empty_presentation(tmp_path, monkeypatch):
    """Test operations on a presentation with 0 slides."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "empty.pptx"
    p = _Prs()
    # Create a presentation but add NO slides
    p.save(str(path))

    # list_slides should return an empty list
    result = list_slides(str(path))
    assert isinstance(result, list)
    assert len(result) == 0

    # read_slide(0) should raise ValidationError because slide 0 doesn't exist
    with pytest.raises(ValidationError, match="out of range"):
        read_slide(str(path), 0)

    # read_speaker_notes(0) should raise ValidationError
    with pytest.raises(ValidationError, match="out of range"):
        read_speaker_notes(str(path), slide_index=0)


# ── GROUP 22: list_layouts() ──────────────────────────────────────────────

@pytest.mark.unit
def test_list_layouts_returns_list(tmp_path, monkeypatch):
    """list_layouts returns a non-empty list for a default Presentation."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "layouts_test.pptx"
    p = _Prs()
    p.save(str(path))
    result = list_layouts(str(path))
    assert isinstance(result, list)
    assert len(result) > 0


@pytest.mark.unit
def test_list_layouts_structure(tmp_path, monkeypatch):
    """Each layout dict has index, name, placeholder_count, placeholder_types keys."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "layouts_struct.pptx"
    p = _Prs()
    p.save(str(path))
    result = list_layouts(str(path))
    for item in result:
        assert "index" in item
        assert "name" in item
        assert "placeholder_count" in item
        assert "placeholder_types" in item


@pytest.mark.unit
def test_list_layouts_path_not_allowed(tmp_path, monkeypatch):
    """list_layouts raises ValidationError for a path outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        list_layouts(path)


# ── GROUP 29 fragment: list_layouts file not found ────────────────────────

@pytest.mark.unit
def test_list_layouts_file_not_found(tmp_path, monkeypatch):
    """list_layouts raises PPTMCPError when the .pptx file does not exist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    missing = str(tmp_path / "missing.pptx")
    with pytest.raises(PPTMCPError):
        list_layouts(missing)
