"""
Unit tests for pptmcp manage_comments (list / add / delete).
No COM required — CI-safe.
"""
from __future__ import annotations

import pytest
from pptx import Presentation


def _make_pptx(path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.save(str(path))
    return str(path)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def pptx_path(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    return _make_pptx(tmp_path / "test.pptx")


@pytest.fixture
def pptx_no_write(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    return _make_pptx(tmp_path / "nowrite.pptx")


@pytest.fixture
def pptx_with_comment(pptx_path):
    from pptmcp._pptx_comments import add_comment
    add_comment(pptx_path, slide_index=0, text="First comment", author="Alice", confirm=True)
    return pptx_path


# ---------------------------------------------------------------------------
# list tests
# ---------------------------------------------------------------------------

class TestListComments:
    def test_list_empty_returns_empty(self, pptx_path):
        from pptmcp._pptx_comments import list_comments
        r = list_comments(pptx_path)
        assert r["comments"] == []
        assert r["total_found"] == 0
        assert r["truncated"] is False

    def test_list_returns_required_keys(self, pptx_path):
        from pptmcp._pptx_comments import list_comments
        r = list_comments(pptx_path)
        for k in ("comments", "total_found", "truncated", "max_comments"):
            assert k in r

    def test_list_after_add_shows_comment(self, pptx_path):
        from pptmcp._pptx_comments import add_comment, list_comments
        add_comment(pptx_path, slide_index=0, text="visible", author="Bob", confirm=True)
        r = list_comments(pptx_path, slide_index=0)
        assert any(c["text"] == "visible" for c in r["comments"])

    def test_list_slide_scope_isolates_slide(self, pptx_path):
        from pptmcp._pptx_comments import add_comment, list_comments
        add_comment(pptx_path, slide_index=1, text="slide1 only", author="C", confirm=True)
        r = list_comments(pptx_path, slide_index=0)
        assert all(c["slide_index"] == 0 for c in r["comments"])

    def test_list_max_comments_cap(self, pptx_path, monkeypatch):
        from pptmcp._pptx_comments import add_comment, list_comments
        monkeypatch.setenv("PPT_MAX_COMMENTS", "2")
        for i in range(4):
            add_comment(pptx_path, slide_index=0, text=f"cm{i}", confirm=True)
        r = list_comments(pptx_path)
        assert len(r["comments"]) == 2
        assert r["total_found"] == 4
        assert r["truncated"] is True

    def test_list_oor_slide_raises(self, pptx_path):
        from pptmcp._pptx_comments import list_comments
        from pptmcp._presentation_governance import ValidationError
        with pytest.raises(ValidationError, match="out of range"):
            list_comments(pptx_path, slide_index=99)

    def test_list_outside_allowlist_raises(self, tmp_path, monkeypatch):
        from pptmcp._pptx_comments import list_comments
        from pptmcp._presentation_governance import ValidationError
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "other"))
        p = _make_pptx(tmp_path / "bad.pptx")
        with pytest.raises(ValidationError):
            list_comments(p)


# ---------------------------------------------------------------------------
# add tests
# ---------------------------------------------------------------------------

class TestAddComment:
    def test_add_returns_expected_shape(self, pptx_path):
        from pptmcp._pptx_comments import add_comment
        r = add_comment(pptx_path, slide_index=0, text="hi", confirm=True)
        assert r["status"] == "added"
        assert isinstance(r["comment_id"], int)
        assert r["slide_index"] == 0

    def test_add_persists_across_reload(self, pptx_path):
        from pptmcp._pptx_comments import add_comment, list_comments
        from pptmcp.presentation_pptx import _prs_cache
        add_comment(pptx_path, slide_index=0, text="persisted", author="Eve", confirm=True)
        _prs_cache.clear()
        r = list_comments(pptx_path, slide_index=0)
        assert any(c["text"] == "persisted" for c in r["comments"])

    def test_add_no_write_env_raises(self, pptx_no_write):
        from pptmcp._pptx_comments import add_comment
        from pptmcp._presentation_governance import NotAllowedError
        with pytest.raises(NotAllowedError):
            add_comment(pptx_no_write, slide_index=0, text="X", confirm=True)

    def test_add_confirm_false_raises(self, pptx_path):
        from pptmcp._pptx_comments import add_comment
        from pptmcp._presentation_governance import ValidationError
        with pytest.raises(ValidationError, match="confirm"):
            add_comment(pptx_path, slide_index=0, text="X", confirm=False)

    def test_add_negative_x_emu_raises(self, pptx_path):
        from pptmcp._pptx_comments import add_comment
        from pptmcp._presentation_governance import ValidationError
        with pytest.raises(ValidationError, match="slide_x_emu"):
            add_comment(pptx_path, slide_index=0, text="X", slide_x_emu=-1, confirm=True)

    def test_add_overflow_y_emu_raises(self, pptx_path):
        from pptmcp._pptx_comments import add_comment
        from pptmcp._presentation_governance import ValidationError
        with pytest.raises(ValidationError, match="slide_y_emu"):
            add_comment(pptx_path, slide_index=0, text="X", slide_y_emu=2**31, confirm=True)

    def test_add_oor_slide_raises(self, pptx_path):
        from pptmcp._pptx_comments import add_comment
        from pptmcp._presentation_governance import ValidationError
        with pytest.raises(ValidationError, match="out of range"):
            add_comment(pptx_path, slide_index=99, text="X", confirm=True)

    def test_add_control_chars_stripped(self, pptx_path):
        from pptmcp._pptx_comments import add_comment, list_comments
        add_comment(pptx_path, slide_index=0, text="Hi\x01\x07!", author="Al\x00ice", confirm=True)
        r = list_comments(pptx_path, slide_index=0)
        cm = r["comments"][-1]
        assert "\x01" not in cm["text"]
        assert "Hi!" in cm["text"]

    def test_add_author_truncated(self, pptx_path):
        from pptmcp._pptx_comments import add_comment, list_comments
        add_comment(pptx_path, slide_index=0, text="hi", author="A" * 80, confirm=True)
        r = list_comments(pptx_path, slide_index=0)
        assert len(r["comments"][-1]["author"]) <= 64

    def test_add_creates_authors_part(self, pptx_path):
        import zipfile
        from pptmcp._pptx_comments import add_comment
        add_comment(pptx_path, slide_index=0, text="hi", author="X", confirm=True)
        with zipfile.ZipFile(pptx_path, "r") as z:
            assert any("commentAuthors" in n for n in z.namelist())

    def test_add_creates_comments_part(self, pptx_path):
        import zipfile
        from pptmcp._pptx_comments import add_comment
        add_comment(pptx_path, slide_index=0, text="hi", confirm=True)
        with zipfile.ZipFile(pptx_path, "r") as z:
            assert any("comment" in n and "Authors" not in n for n in z.namelist())


# ---------------------------------------------------------------------------
# delete tests
# ---------------------------------------------------------------------------

class TestDeleteComment:
    def test_delete_removes_comment(self, pptx_with_comment):
        from pptmcp._pptx_comments import delete_comment, list_comments
        from pptmcp.presentation_pptx import _prs_cache
        cid = list_comments(pptx_with_comment, slide_index=0)["comments"][0]["comment_id"]
        r = delete_comment(pptx_with_comment, slide_index=0, comment_id=cid, confirm=True)
        assert r["status"] == "deleted"
        _prs_cache.clear()
        after = list_comments(pptx_with_comment, slide_index=0)
        assert not any(c["comment_id"] == cid for c in after["comments"])

    def test_delete_no_write_raises(self, pptx_no_write):
        from pptmcp._pptx_comments import delete_comment
        from pptmcp._presentation_governance import NotAllowedError
        with pytest.raises(NotAllowedError):
            delete_comment(pptx_no_write, slide_index=0, comment_id=0, confirm=True)

    def test_delete_confirm_false_raises(self, pptx_with_comment):
        from pptmcp._pptx_comments import delete_comment, list_comments
        from pptmcp._presentation_governance import ValidationError
        cid = list_comments(pptx_with_comment, slide_index=0)["comments"][0]["comment_id"]
        with pytest.raises(ValidationError, match="confirm"):
            delete_comment(pptx_with_comment, slide_index=0, comment_id=cid, confirm=False)

    def test_delete_not_found_raises(self, pptx_path):
        from pptmcp._pptx_comments import delete_comment
        from pptmcp._presentation_governance import ValidationError
        with pytest.raises(ValidationError, match="comment_id 9999 not found on slide 0"):
            delete_comment(pptx_path, slide_index=0, comment_id=9999, confirm=True)

    def test_delete_message_excludes_file_path(self, pptx_path):
        from pptmcp._pptx_comments import delete_comment
        from pptmcp._presentation_governance import ValidationError
        try:
            delete_comment(pptx_path, slide_index=0, comment_id=9999, confirm=True)
        except ValidationError as exc:
            assert pptx_path not in str(exc)
            assert ".pptx" not in str(exc)
        else:
            pytest.fail("Expected ValidationError not raised")


# ---------------------------------------------------------------------------
# Smoke tests
# ---------------------------------------------------------------------------

class TestManageCommentsSmoke:
    def test_tool_registered(self):
        from fastmcp import FastMCP
        from pptmcp._server_handlers_comments import register_comments_tools
        mcp = FastMCP("test")
        register_comments_tools(mcp)
        # Use the local provider's synchronous _components dict — no async/socket needed
        provider = mcp.providers[0]
        names = {v.name for v in provider._components.values()}
        assert "manage_comments" in names

    def test_compact_surface_has_manage_comments(self):
        from pptmcp._server_compact import COMPACT_TOOL_NAMES
        assert "manage_comments" in COMPACT_TOOL_NAMES
