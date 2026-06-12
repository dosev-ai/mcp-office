import os
import sys

# Ensure worktree src takes priority over installed editable package.
_wt_src = os.path.realpath(
    os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "src")
)
if _wt_src not in sys.path:
    sys.path.insert(0, _wt_src)

from pathlib import Path

import pytest
from pptx import Presentation as _Prs
from pptmcp.presentation_pptx import (
    ValidationError,
    _prs_cache,
    add_hyperlink,
    add_shape,
    add_slide,
    add_table_to_slide,
    add_textbox,
    create_presentation,
    edit_text_placeholder,
    extract_images,
    extract_presentation_text,
    insert_image,
    manage_hyperlinks,
    read_presentation,
    save,
    set_speaker_notes,
    set_text_format,
)


def build_test_deck(
    base_dir,
    filename: str,
    *,
    title: str,
    slide_titles: tuple[str, ...] = (),
    save_after: bool = False,
) -> str:
    """Create a deck in the allowlisted directory for UAT scenarios."""
    path = str(Path(base_dir) / filename)
    create_presentation(path, title=title, confirm=True)
    for slide_title in slide_titles:
        add_slide(path, layout_index=1, title=slide_title, confirm=True)
    if save_after:
        save(path, confirm=True)
    return path


def combined_slide_text(slide: dict) -> str:
    """Flatten a slide payload into one searchable text string."""
    return " ".join(
        part
        for part in [
            slide.get("title") or "",
            slide.get("body_text") or "",
            *(shape.get("text", "") for shape in slide.get("shapes", [])),
        ]
        if part
    )


def assert_file_size_at_least(path: str, min_bytes: int = 1000) -> None:
    """Assert an exported artifact exists and is not trivially empty."""
    output_path = Path(path)
    assert output_path.exists()
    assert output_path.stat().st_size > min_bytes


def build_image_roundtrip_deck(base_dir, image_path: str) -> tuple[str, dict]:
    """Create the US-14 deck and insert the sample image."""
    path = build_test_deck(base_dir, "image_test.pptx", title="Image Test")
    insert_result = insert_image(
        path,
        0,
        image_path,
        left=0.5,
        top=0.5,
        width=2.0,
        confirm=True,
    )
    return path, insert_result


def assert_inserted_image_roundtrip(path: str) -> dict:
    """Verify the inserted image is reported by pptmcp and persisted on disk."""
    images = extract_images(path)
    assert len(images["images"]) >= 1
    assert images["images"][0]["slide_index"] == 0
    assert images["images"][0]["content_type"] is not None

    save(path, confirm=True)

    reloaded = _Prs(path)
    picture_shapes = [shape for shape in reloaded.slides[0].shapes if shape.shape_type == 13]
    assert len(picture_shapes) >= 1
    return images


def assert_insert_image_invalid_paths(path: str, allowlisted_dir) -> None:
    """Verify US-14 rejects malformed, unsupported, and disallowed image paths."""
    with pytest.raises(ValidationError, match="null byte"):
        insert_image(path, 0, "\x00malicious.png", confirm=True)

    fake_docx = Path(allowlisted_dir) / "fake.docx"
    fake_docx.write_bytes(b"fake content")
    with pytest.raises(ValidationError, match="[Uu]nsupported|extension|not allowed"):
        insert_image(path, 0, str(fake_docx), confirm=True)

    with pytest.raises(
        ValidationError,
        match="not in allowlist|traversal|not allowed|No allowlist",
    ):
        insert_image(path, 0, "/etc/passwd.png", confirm=True)


def build_rag_hyperlink_deck(
    base_dir,
    *,
    body_text: str = "MCP enables AI tool use",
    display_text: str = "See python-pptx docs",
    url: str = "http://python-pptx.readthedocs.io",
) -> tuple[str, int, dict]:
    """Create the US-15 deck with body text and a hyperlink-bearing references slide."""
    path = build_test_deck(
        base_dir,
        "rag_test.pptx",
        title="RAG Introduction",
        slide_titles=("References",),
    )
    add_textbox(path, 0, 1.0, 2.0, 5.0, 1.5, text=body_text, confirm=True)
    references_box = add_textbox(
        path,
        1,
        1.0,
        2.0,
        5.0,
        1.5,
        text=display_text,
        confirm=True,
    )
    hyperlink = add_hyperlink(
        path,
        1,
        references_box["shape_id"],
        run_index=0,
        url=url,
        display_text=display_text,
        confirm=True,
    )
    save(path, confirm=True)
    return path, references_box["shape_id"], hyperlink


def assert_rag_text_and_links(
    path: str,
    *,
    expected_text: str,
    expected_url: str,
    slide_index: int = 1,
) -> tuple[dict, dict]:
    """Verify the US-15 deck exposes the expected extracted text and hyperlink list."""
    full_text = extract_presentation_text(path)
    assert full_text["slide_count"] == 2
    assert full_text["slides"][0]["slide_index"] == 0
    assert expected_text in combined_slide_text(full_text["slides"][0])

    link_list = manage_hyperlinks(path, slide_index, "list")
    assert link_list["hyperlink_count"] >= 1
    assert link_list["links"][0]["url"] == expected_url
    assert link_list["operation"] == "list"

    list_no_confirm = manage_hyperlinks(path, slide_index, "list")
    assert "links" in list_no_confirm
    return full_text, link_list


def build_lifecycle_com_deck(base_dir, *, filename: str, title: str) -> str:
    """Create the shared US-20b lifecycle deck before COM export verification."""
    path = build_test_deck(base_dir, filename, title=title, slide_titles=("Content", "Refs"))
    edit_text_placeholder(path, 0, 0, title, confirm=True)
    set_speaker_notes(path, 0, "Intro notes", confirm=True)

    textbox = add_textbox(path, 1, 1.0, 1.0, 5.0, 2.0, text="Key finding", confirm=True)
    set_text_format(path, 1, textbox["shape_id"], bold=True, confirm=True)

    reference_shape = add_shape(
        path,
        2,
        "RECTANGLE",
        0.5,
        0.5,
        3.0,
        1.5,
        "References",
        confirm=True,
    )
    add_hyperlink(
        path,
        2,
        reference_shape["shape_id"],
        run_index=0,
        url="https://modelcontextprotocol.io",
        confirm=True,
    )
    add_table_to_slide(
        path,
        2,
        2,
        2,
        0.5,
        2.0,
        5.0,
        2.0,
        data=[["Tool", "Status"], ["read_presentation", "OK"]],
        confirm=True,
    )
    save(path, confirm=True)
    return path


def assert_slide_count(path: str, expected_count: int) -> dict:
    """Verify the deck has the expected number of slides."""
    presentation = read_presentation(path)
    assert presentation["slide_count"] == expected_count
    return presentation


def assert_com_export_result(
    result: dict,
    output_path: str,
    *,
    expected_slide_index: int | None = None,
    expected_slide_count: int | None = None,
) -> None:
    """Verify a COM export reported success and produced a meaningful artifact."""
    assert result["exported"] is True
    assert result["output_path"] == output_path
    if expected_slide_index is not None:
        assert result["slide_index"] == expected_slide_index
    if expected_slide_count is not None:
        assert result["slide_count"] == expected_slide_count
    assert_file_size_at_least(output_path)


def pytest_configure(config):
    config.addinivalue_line("markers", "unit: unit tests")
    config.addinivalue_line("markers", "smoke: smoke tests")
    config.addinivalue_line("markers", "integration: integration tests requiring live Office app")


@pytest.fixture
def sample_pptx(tmp_path, monkeypatch):
    """Create a minimal test.pptx, set PPT_ALLOWLIST_ROOTS, and return the path as str."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "test.pptx"
    p = _Prs()
    layout = p.slide_layouts[1]  # Title and Content
    slide = p.slides.add_slide(layout)
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = "Body text"
    p.save(str(path))
    return str(path)


@pytest.fixture
def allowlisted_dir(tmp_path, monkeypatch):
    """Creates an allowlisted directory and sets PPT_ALLOWLIST_ROOTS."""
    allowed = tmp_path / "allowed"
    allowed.mkdir()
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(allowed))
    return allowed


@pytest.fixture
def write_env(monkeypatch):
    """Sets PPT_ENABLE_WRITE=true."""
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")


@pytest.fixture
def no_write_env(monkeypatch):
    """Ensures PPT_ENABLE_WRITE is absent/false."""
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)


@pytest.fixture
def sample_pptx_file(allowlisted_dir):
    """Creates a real 1-slide .pptx in the allowlisted dir and returns its path."""
    from pptx import Presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Test Slide"
    path = allowlisted_dir / "sample.pptx"
    prs.save(str(path))
    return str(path)


@pytest.fixture
def com_env(monkeypatch):
    """Sets PPT_ENABLE_COM=true and mocks sys.platform to win32."""
    import sys
    monkeypatch.setenv("PPT_ENABLE_COM", "true")
    monkeypatch.setattr(sys, "platform", "win32")


@pytest.fixture
def sample_pptx_with_table(allowlisted_dir, write_env):
    """1-slide .pptx with a 3×3 table."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    table = slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(6), Inches(3)).table
    for r in range(3):
        for c in range(3):
            table.cell(r, c).text = f"R{r}C{c}"
    path = allowlisted_dir / "table_sample.pptx"
    prs.save(str(path))
    return str(path)


@pytest.fixture
def sample_png_file(allowlisted_dir):
    """Minimal valid 1×1 white PNG."""
    import struct
    import zlib

    def _chunk(tag, data):
        c = struct.pack(">I", len(data)) + tag + data
        return c + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)

    png = (
        b"\x89PNG\r\n\x1a\n"
        + _chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
        + _chunk(b"IDAT", zlib.compress(b"\x00\xff\xff\xff"))
        + _chunk(b"IEND", b"")
    )
    path = allowlisted_dir / "pixel.png"
    path.write_bytes(png)
    return str(path)


@pytest.fixture
def three_slide_pptx(allowlisted_dir, write_env):
    """3-slide .pptx each with a title."""
    from pptx import Presentation

    prs = Presentation()
    for title_text in ("Slide One", "Slide Two", "Slide Three"):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title_text
    path = allowlisted_dir / "three_slides.pptx"
    prs.save(str(path))
    return str(path)


@pytest.fixture
def kill_powerpnt():
    """Yields; kills POWERPNT.EXE on teardown (shell=False, NC-06)."""
    yield
    import subprocess

    subprocess.run(
        ["taskkill", "/F", "/IM", "POWERPNT.EXE"],
        shell=False,
        check=False,
        capture_output=True,
    )


@pytest.fixture(autouse=True)
def _clear_prs_cache():
    _prs_cache.clear()
    yield
    _prs_cache.clear()


@pytest.fixture
def pptx_with_known_shape_ids(tmp_path, write_env, allowlisted_dir):
    """Creates a 1-slide pptx with 2 textboxes in the allowlisted dir.

    Returns (str_path, shape_id_1, shape_id_2).
    """
    from pptx import Presentation as _P
    from pptx.util import Inches

    prs = _P()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb1 = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2.0), Inches(1.0))
    tb1.text_frame.text = "shape one text"
    tb2 = slide.shapes.add_textbox(Inches(3.0), Inches(0.5), Inches(2.0), Inches(1.0))
    tb2.text_frame.text = "shape two text"
    path = allowlisted_dir / "known_shapes.pptx"
    prs.save(str(path))
    return str(path), tb1.shape_id, tb2.shape_id


@pytest.fixture
def pptx_with_mixed_shapes(tmp_path, write_env, allowlisted_dir):
    """Creates a 1-slide pptx with one textbox and one picture shape (no text_frame).

    Returns (str_path, text_shape_id, no_text_shape_id).
    """
    import struct
    import zlib
    from pptx import Presentation as _P
    from pptx.util import Inches

    def _chunk(tag, data):
        c = struct.pack(">I", len(data)) + tag + data
        return c + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)

    png = (
        b"\x89PNG\r\n\x1a\n"
        + _chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
        + _chunk(b"IDAT", zlib.compress(b"\x00\xff\xff\xff"))
        + _chunk(b"IEND", b"")
    )
    png_path = allowlisted_dir / "pixel.png"
    png_path.write_bytes(png)

    prs = _P()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2.0), Inches(1.0))
    tb.text_frame.text = "editable text"
    pic = slide.shapes.add_picture(
        str(png_path), Inches(3.0), Inches(0.5), Inches(1.0), Inches(1.0)
    )
    path = allowlisted_dir / "mixed_shapes.pptx"
    prs.save(str(path))
    return str(path), tb.shape_id, pic.shape_id
