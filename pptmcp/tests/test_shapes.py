"""Shape creation and formatting tests — targets presentation_pptx.py only."""

import pytest
from pptx import Presentation as _Prs
from pptmcp.presentation_pptx import (
    ValidationError,
    NotAllowedError,
    add_textbox,
    add_shape,
    add_table_to_slide,
)
from pptmcp.shapes_pptx import set_shape_geometry
# Inline `from pptx.util import Inches` inside individual test functions is OK;
# leave those inline imports exactly as they appear in test_unit.py.


# ── GROUP 23: add_textbox() ───────────────────────────────────────────────

@pytest.mark.unit
def test_add_textbox_success(tmp_path, monkeypatch):
    """add_textbox returns dict with shape_id."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "textbox.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = add_textbox(str(path), 0, 1.0, 1.0, 4.0, 1.0, text="Hello", confirm=True)
    assert "shape_id" in result
    assert result["slide_index"] == 0
    assert result["width_inches"] == 4.0
    assert result["height_inches"] == 1.0


@pytest.mark.unit
def test_add_textbox_blocked_write(tmp_path, monkeypatch):
    """add_textbox raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(NotAllowedError):
        add_textbox(path, 0, 1.0, 1.0, 4.0, 1.0, confirm=True)


@pytest.mark.unit
def test_add_textbox_no_confirm(tmp_path, monkeypatch):
    """add_textbox raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        add_textbox(path, 0, 1.0, 1.0, 4.0, 1.0, confirm=False)


@pytest.mark.unit
def test_add_textbox_invalid_slide_index(tmp_path, monkeypatch):
    """add_textbox raises ValidationError for an out-of-range slide index."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "one_slide.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError):
        add_textbox(str(path), 99, 1.0, 1.0, 4.0, 1.0, confirm=True)


@pytest.mark.unit
def test_add_textbox_path_not_allowed(tmp_path, monkeypatch):
    """add_textbox raises ValidationError for a path outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        add_textbox(path, 0, 1.0, 1.0, 4.0, 1.0, confirm=True)


# ── GROUP 24: add_shape() ─────────────────────────────────────────────────

@pytest.mark.unit
def test_add_shape_rectangle_success(tmp_path, monkeypatch):
    """add_shape with RECTANGLE returns dict with shape_id."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "shape.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = add_shape(str(path), 0, "RECTANGLE", 1.0, 1.0, 3.0, 2.0, confirm=True)
    assert "shape_id" in result
    assert result["slide_index"] == 0
    assert result["shape_type"] == "RECTANGLE"


@pytest.mark.unit
def test_add_shape_blocked_write(tmp_path, monkeypatch):
    """add_shape raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(NotAllowedError):
        add_shape(path, 0, "RECTANGLE", 1.0, 1.0, 3.0, 2.0, confirm=True)


@pytest.mark.unit
def test_add_shape_no_confirm(tmp_path, monkeypatch):
    """add_shape raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        add_shape(path, 0, "RECTANGLE", 1.0, 1.0, 3.0, 2.0, confirm=False)


@pytest.mark.unit
def test_add_shape_invalid_type(tmp_path, monkeypatch):
    """add_shape raises ValidationError for an unknown shape_type string."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "shape2.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="Unknown shape type"):
        add_shape(str(path), 0, "NOT_A_REAL_SHAPE_XYZ", 1.0, 1.0, 3.0, 2.0, confirm=True)


@pytest.mark.unit
def test_add_shape_invalid_slide_index(tmp_path, monkeypatch):
    """add_shape raises ValidationError for an out-of-range slide index."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "shape3.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError):
        add_shape(str(path), 99, "RECTANGLE", 1.0, 1.0, 3.0, 2.0, confirm=True)


@pytest.mark.unit
def test_add_shape_path_not_allowed(tmp_path, monkeypatch):
    """add_shape raises ValidationError for a path outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        add_shape(path, 0, "RECTANGLE", 1.0, 1.0, 2.0, 1.0, confirm=True)


# ── GROUP 25: add_table_to_slide() ────────────────────────────────────────

@pytest.mark.unit
def test_add_table_success(tmp_path, monkeypatch):
    """add_table_to_slide with 2x3 table returns rows=2 cols=3."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "table.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = add_table_to_slide(str(path), 0, 2, 3, 1.0, 1.0, 6.0, 2.0, confirm=True)
    assert result["rows"] == 2
    assert result["cols"] == 3
    assert "shape_id" in result


@pytest.mark.unit
def test_add_table_with_data(tmp_path, monkeypatch):
    """add_table_to_slide with data=[["A","B"],["C","D"]] succeeds and returns shape_id."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "table_data.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = add_table_to_slide(
        str(path), 0, 2, 2, 1.0, 1.0, 4.0, 2.0,
        data=[["A", "B"], ["C", "D"]], confirm=True,
    )
    assert "shape_id" in result
    assert result["rows"] == 2
    assert result["cols"] == 2


@pytest.mark.unit
def test_add_table_with_compact_table_style(tmp_path, monkeypatch):
    """add_table_to_slide accepts compact table font and row-height options."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "table_compact.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = add_table_to_slide(
        str(path),
        0,
        2,
        2,
        1.0,
        1.0,
        4.0,
        1.0,
        data=[["H1", "H2"], ["A", "B"]],
        font_size_pt=7,
        header_font_size_pt=8,
        header_bold=True,
        row_height_pt=16,
        confirm=True,
    )
    assert result["runs_updated"] == 4
    assert result["rows_resized"] == 2


@pytest.mark.unit
def test_add_table_removes_overlapping_empty_content_placeholder(tmp_path, monkeypatch):
    """Adding a table removes an empty overlapping content placeholder by default."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "table_placeholder.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(path))
    result = add_table_to_slide(
        str(path),
        0,
        2,
        2,
        1.0,
        1.5,
        6.0,
        2.0,
        data=[["H1", "H2"], ["A", "B"]],
        confirm=True,
    )
    assert result["placeholders_removed"] >= 1


@pytest.mark.unit
def test_add_table_data_shape_mismatch(tmp_path, monkeypatch):
    """add_table_to_slide raises ValidationError when data row count does not match rows."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "table_bad.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="rows"):
        add_table_to_slide(
            str(path), 0, 3, 2, 1.0, 1.0, 4.0, 2.0,
            data=[["A", "B"], ["C", "D"]],  # 2 rows, but rows=3
            confirm=True,
        )


@pytest.mark.unit
def test_add_table_blocked_write(tmp_path, monkeypatch):
    """add_table_to_slide raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(NotAllowedError):
        add_table_to_slide(path, 0, 2, 2, 1.0, 1.0, 4.0, 2.0, confirm=True)


@pytest.mark.unit
def test_add_table_no_confirm(tmp_path, monkeypatch):
    """add_table_to_slide raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        add_table_to_slide(path, 0, 2, 2, 1.0, 1.0, 4.0, 2.0, confirm=False)


@pytest.mark.unit
def test_add_table_invalid_slide_index(tmp_path, monkeypatch):
    """add_table_to_slide raises ValidationError for an out-of-range slide index."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "one_slide.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError):
        add_table_to_slide(str(path), 99, 2, 2, 0.0, 0.0, 4.0, 2.0, confirm=True)



# ── GROUP 29 fragments: shape/table/format gap-fill ───────────────────────

@pytest.mark.unit
def test_add_shape_alias_circle_oval(tmp_path, monkeypatch):
    """add_shape with 'CIRCLE' alias resolves to OVAL without ValueError."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "circle.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = add_shape(str(path), 0, "CIRCLE", 1.0, 1.0, 2.0, 2.0, confirm=True)
    assert "shape_id" in result


@pytest.mark.unit
def test_add_table_path_not_allowed(tmp_path, monkeypatch):
    """add_table_to_slide raises ValidationError for a path outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        add_table_to_slide(path, 0, 2, 2, 1.0, 1.0, 4.0, 2.0, confirm=True)


@pytest.mark.unit
def test_add_table_data_col_mismatch(tmp_path, monkeypatch):
    """add_table_to_slide raises ValidationError when data column count != cols."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "table_col_bad.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="cols"):
        add_table_to_slide(
            str(path), 0, 2, 3, 1.0, 1.0, 4.0, 2.0,
            data=[["A", "B"], ["C", "D"]],  # rows match (2) but 2 cols != cols=3
            confirm=True,
        )


@pytest.mark.unit
def test_add_table_invalid_dimensions_rows_zero(tmp_path, monkeypatch):
    """add_table_to_slide raises ValidationError when rows=0."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="rows"):
        add_table_to_slide(path, 0, 0, 2, 1.0, 1.0, 4.0, 2.0, confirm=True)


# ════════════════════════════════════════════════════════════════════════════
# PHASE B QUALITY GATES: US-04 – US-08  (18 tests)
# ════════════════════════════════════════════════════════════════════════════

from pptmcp.shapes_pptx import (  # noqa: E402
    _calculate_overflow_risk,
    detect_overlapping_shapes as _detect_overlapping_shapes,
)
from pptmcp.presentation_pptx import (  # noqa: E402
    _emu_to_inches,
    PPTMCPError,
    capabilities,
    list_shapes,
    save,
)


# ── GROUP 1: US-04 word_wrap ──────────────────────────────────────────────

@pytest.mark.unit
def test_add_textbox_word_wrap_enabled(tmp_path, monkeypatch):
    """add_textbox sets word_wrap=True; verified by saving and reloading from disk."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "ww_textbox.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[6])
    p.save(str(path))
    result = add_textbox(str(path), 0, 1.0, 1.0, 4.0, 1.0, text="Wrap me", confirm=True)
    assert "shape_id" in result
    save(str(path), confirm=True)
    loaded = _Prs(str(path))
    target = next(s for s in loaded.slides[0].shapes if s.shape_id == result["shape_id"])
    assert target.text_frame.word_wrap is True


@pytest.mark.unit
def test_add_shape_word_wrap_enabled(tmp_path, monkeypatch):
    """add_shape(RECTANGLE) sets word_wrap=True; verified by saving and reloading from disk."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "ww_shape.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[6])
    p.save(str(path))
    result = add_shape(str(path), 0, "RECTANGLE", 1.0, 1.0, 3.0, 2.0, confirm=True)
    assert "shape_id" in result
    save(str(path), confirm=True)
    loaded = _Prs(str(path))
    target = next(s for s in loaded.slides[0].shapes if s.shape_id == result["shape_id"])
    assert target.text_frame.word_wrap is True


# ── GROUP 2: US-05 overflow_risk ─────────────────────────────────────────

@pytest.mark.unit
def test_overflow_risk_low():
    """Short text in a large box returns overflow_risk='low'."""
    result = _calculate_overflow_risk("Hi", 8.0, 5.0)
    assert result["overflow_risk"] == "low"


@pytest.mark.unit
def test_overflow_risk_medium():
    """Medium-density text returns overflow_risk='medium'."""
    # 100 chars / (int(2.0*13)=26 chars/line) = ceil(3.84)=4 lines
    # 4 * 0.25in = 1.0in; ratio = 1.0/2.0 = 0.5 → medium (0.5 < 1.0)
    result = _calculate_overflow_risk("A" * 100, 2.0, 2.0)
    assert result["overflow_risk"] == "medium"


@pytest.mark.unit
def test_overflow_risk_high():
    """Very long text in a tiny box returns overflow_risk='high'."""
    result = _calculate_overflow_risk("A" * 500, 1.0, 0.5)
    assert result["overflow_risk"] == "high"


@pytest.mark.unit
def test_add_textbox_returns_overflow_risk(tmp_path, monkeypatch):
    """add_textbox result contains overflow_risk in the valid value set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "or_textbox.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[6])
    p.save(str(path))
    result = add_textbox(str(path), 0, 1.0, 1.0, 4.0, 1.0, text="Test overflow", confirm=True)
    assert result["overflow_risk"] in ("low", "medium", "high")


@pytest.mark.unit
def test_add_textbox_overflow_detail_is_string(tmp_path, monkeypatch):
    """add_textbox result contains overflow_detail as a str."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "od_textbox.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[6])
    p.save(str(path))
    result = add_textbox(str(path), 0, 1.0, 1.0, 4.0, 1.0, text="Detail check", confirm=True)
    assert isinstance(result["overflow_detail"], str)


# ── GROUP 3: US-06 bounds ────────────────────────────────────────────────

@pytest.mark.unit
def test_list_shapes_bounds_keys(tmp_path, monkeypatch):
    """Every shape returned by list_shapes has 'bounds' dict with exactly 4 keys."""
    from pptx.util import Inches as _In
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "bounds_check.pptx"
    prs = _Prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(1, _In(1), _In(1), _In(2), _In(2))
    prs.save(str(path))
    shapes = list_shapes(str(path), 0)
    assert len(shapes) >= 1
    for sh in shapes:
        assert "bounds" in sh
        assert set(sh["bounds"].keys()) == {"left_in", "top_in", "width_in", "height_in"}


@pytest.mark.unit
def test_list_shapes_bounds_none_allowed():
    """_emu_to_inches(None) returns None (supports shapes whose bounds may be None)."""
    assert _emu_to_inches(None) is None


@pytest.mark.unit
def test_emu_to_inches_rounds_to_4dp():
    """_emu_to_inches converts standard EMU values to correct inch measurements."""
    assert _emu_to_inches(914400) == 1.0
    assert _emu_to_inches(457200) == 0.5


# ── GROUP 4: US-07 detect_overlapping_shapes ─────────────────────────────

@pytest.mark.unit
def test_detect_overlapping_no_overlaps(tmp_path, monkeypatch):
    """detect_overlapping_shapes returns empty list when shapes are clearly separated."""
    from pptx.util import Inches as _In
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "no_overlap.pptx"
    prs = _Prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(1, _In(0), _In(0), _In(2), _In(2))   # 0–2in
    slide.shapes.add_shape(1, _In(3), _In(0), _In(2), _In(2))   # 3–5in (1in gap)
    prs.save(str(path))
    result = _detect_overlapping_shapes(str(path), 0)
    assert result["overlapping_pairs"] == []


@pytest.mark.unit
def test_detect_overlapping_with_overlaps(tmp_path, monkeypatch):
    """detect_overlapping_shapes detects two identical shapes at the same position."""
    from pptx.util import Inches as _In
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "with_overlap.pptx"
    prs = _Prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(1, _In(1), _In(1), _In(2), _In(2))   # shape A
    slide.shapes.add_shape(1, _In(1), _In(1), _In(2), _In(2))   # shape B — same position
    prs.save(str(path))
    result = _detect_overlapping_shapes(str(path), 0)
    assert result["pair_count"] == 1
    assert result["overlapping_pairs"][0]["overlap_area_in2"] > 0


@pytest.mark.unit
def test_detect_overlapping_result_keys(tmp_path, monkeypatch):
    """detect_overlapping_shapes result has slide_index, overlapping_pairs, pair_count."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "result_keys.pptx"
    prs = _Prs()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))
    result = _detect_overlapping_shapes(str(path), 0)
    assert {"slide_index", "overlapping_pairs", "pair_count"}.issubset(result.keys())


# ── GROUP 5: US-08 overlap_warning ───────────────────────────────────────

@pytest.mark.unit
def test_add_textbox_overlap_warning_key_present(tmp_path, monkeypatch):
    """add_textbox on a blank slide (no other shapes) returns overlap_warning=None."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "ow_none_tb.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[6])
    p.save(str(path))
    result = add_textbox(str(path), 0, 1.0, 1.0, 2.0, 2.0, text="Solo", confirm=True)
    assert "overlap_warning" in result
    assert result["overlap_warning"] is None


@pytest.mark.unit
def test_add_textbox_overlap_warning_detected(tmp_path, monkeypatch):
    """Second shape at same position as first shape returns overlap_warning with detail."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "ow_detected.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[6])
    p.save(str(path))
    # First shape — no other shapes, so overlap_warning must be None
    first = add_shape(str(path), 0, "RECTANGLE", 1.0, 1.0, 2.0, 2.0, confirm=True)
    assert first["overlap_warning"] is None
    # Second shape at same position — in-memory cache already has first shape
    second = add_textbox(str(path), 0, 1.0, 1.0, 2.0, 2.0, text="Overlap!", confirm=True)
    assert second["overlap_warning"] is not None
    warn = second["overlap_warning"]
    assert "shape_a" in warn
    assert "shape_b" in warn
    assert "overlap_area_in2" in warn


@pytest.mark.unit
def test_add_shape_overlap_warning_key_present(tmp_path, monkeypatch):
    """add_shape on a blank slide (no other shapes) returns overlap_warning=None."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "ow_none_sh.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[6])
    p.save(str(path))
    result = add_shape(str(path), 0, "RECTANGLE", 1.0, 1.0, 3.0, 2.0, confirm=True)
    assert "overlap_warning" in result
    assert result["overlap_warning"] is None


# ── GROUP 6: capabilities ─────────────────────────────────────────────────

@pytest.mark.unit
def test_capabilities_includes_detect_overlapping_shapes():
    """capabilities() lists detect_overlapping_shapes in its tools array."""
    result = capabilities()
    tool_names = {e["tool"] for e in result["tools"]}
    assert "detect_overlapping_shapes" in tool_names


# ── Additional: invalid slide index ───────────────────────────────────────

@pytest.mark.unit
def test_detect_overlapping_invalid_slide_index(tmp_path, monkeypatch):
    """detect_overlapping_shapes with slide_index=-1 raises PPTMCPError."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "inv_slide.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[6])
    p.save(str(path))
    with pytest.raises(PPTMCPError):
        _detect_overlapping_shapes(str(path), -1)


# ── GROUP 99: set_shape_geometry — RCI shape-dispatch set_geometry op ────


def _make_one_shape_pptx(tmp_path, monkeypatch):
    """Create a deck with one auto-shape and return (path, shape_id)."""
    from pptx.util import Inches
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx = tmp_path / "geom.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    p.save(str(pptx))
    return str(pptx), tb.shape_id


@pytest.mark.unit
def test_set_shape_geometry_partial_update_changes_only_specified_axes(tmp_path, monkeypatch):
    """Only axes provided change; omitted axes stay put."""
    pptx, sid = _make_one_shape_pptx(tmp_path, monkeypatch)

    # Update only left and width; top and height remain at 1.0in
    result = set_shape_geometry(
        pptx, 0, sid,
        left_in=4.0, width_in=5.0,
        confirm=True,
    )
    assert result["updated"] is True
    g = result["geometry_in"]
    assert g["left"] == pytest.approx(4.0, abs=0.001)
    assert g["width"] == pytest.approx(5.0, abs=0.001)
    assert g["top"] == pytest.approx(1.0, abs=0.001)
    assert g["height"] == pytest.approx(1.0, abs=0.001)
    assert result["warnings"] == []  # 4.0 + 5.0 = 9.0in inside the 13.33in slide


@pytest.mark.unit
def test_set_shape_geometry_off_slide_emits_warnings_but_does_not_block(tmp_path, monkeypatch):
    """Off-slide geometry warns, never raises."""
    pptx, sid = _make_one_shape_pptx(tmp_path, monkeypatch)

    # Pushed off the right and bottom edges of a 13.33in x 7.5in slide
    result = set_shape_geometry(
        pptx, 0, sid,
        left_in=12.0, top_in=7.0, width_in=4.0, height_in=2.0,
        confirm=True,
    )
    assert result["updated"] is True
    # 12.0 + 4.0 = 16.0in > 13.33in; 7.0 + 2.0 = 9.0in > 7.5in
    joined = " ".join(result["warnings"])
    assert "right edge" in joined
    assert "bottom edge" in joined


@pytest.mark.unit
def test_set_shape_geometry_no_axes_raises(tmp_path, monkeypatch):
    """Calling with all axes None raises ValidationError."""
    pptx, sid = _make_one_shape_pptx(tmp_path, monkeypatch)
    with pytest.raises(ValidationError, match=r"at least one"):
        set_shape_geometry(pptx, 0, sid, confirm=True)


@pytest.mark.unit
def test_set_shape_geometry_write_gate_blocks(tmp_path, monkeypatch):
    """PPT_ENABLE_WRITE missing → NotAllowedError."""
    pptx, sid = _make_one_shape_pptx(tmp_path, monkeypatch)
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    with pytest.raises(NotAllowedError):
        set_shape_geometry(pptx, 0, sid, left_in=2.0, confirm=True)


@pytest.mark.unit
def test_set_shape_geometry_confirm_required(tmp_path, monkeypatch):
    """confirm=False → ValidationError."""
    pptx, sid = _make_one_shape_pptx(tmp_path, monkeypatch)
    with pytest.raises(ValidationError, match=r"confirm=True"):
        set_shape_geometry(pptx, 0, sid, left_in=2.0)


@pytest.mark.unit
def test_set_shape_geometry_negative_left_emits_warning(tmp_path, monkeypatch):
    """Negative left flagged in warnings."""
    pptx, sid = _make_one_shape_pptx(tmp_path, monkeypatch)
    result = set_shape_geometry(pptx, 0, sid, left_in=-0.5, confirm=True)
    assert result["updated"] is True
    assert any("before slide left" in w for w in result["warnings"])

