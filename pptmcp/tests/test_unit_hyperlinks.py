"""Unit tests: hyperlinks — extract_presentation_text, manage_hyperlinks, add_hyperlink
(groups 18–19, 30)."""
import pytest
from pptx import Presentation as _Prs
from pptx.util import Inches

from pptmcp.presentation_pptx import (
    PPTMCPError,
    ValidationError,
    NotAllowedError,
    add_hyperlink,
    extract_presentation_text,
    manage_hyperlinks,
    save,
)


# ── GROUP 18: extract_presentation_text() ───────────────────────────────

@pytest.mark.unit
class TestExtractPresentationText:
    @pytest.mark.unit
    def test_extract_returns_expected_keys(self, sample_pptx):
        result = extract_presentation_text(sample_pptx)
        assert "path" in result
        assert "slide_count" in result
        assert "slides" in result

    @pytest.mark.unit
    def test_extract_slide_count(self, sample_pptx):
        result = extract_presentation_text(sample_pptx)
        assert result["slide_count"] == 1
        assert len(result["slides"]) == 1

    @pytest.mark.unit
    def test_extract_title_captured(self, sample_pptx):
        result = extract_presentation_text(sample_pptx)
        slide = result["slides"][0]
        assert slide["slide_index"] == 0
        # sample_pptx has title "Test Slide" on layout 1 (Title and Content)
        assert slide["title"] == "Test Slide"

    @pytest.mark.unit
    def test_extract_shapes_list_populated(self, sample_pptx):
        result = extract_presentation_text(sample_pptx)
        shapes = result["slides"][0]["shapes"]
        assert isinstance(shapes, list)
        assert len(shapes) >= 1
        for s in shapes:
            assert "shape_id" in s
            assert "name" in s
            assert "text" in s

    @pytest.mark.unit
    def test_extract_notes_captured(self, tmp_path, monkeypatch):
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        path = tmp_path / "withnotes.pptx"
        p = _Prs()
        slide = p.slides.add_slide(p.slide_layouts[1])
        slide.shapes.title.text = "Notes Slide"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "My speaker notes"
        p.save(str(path))
        result = extract_presentation_text(str(path))
        assert result["slides"][0]["notes"] == "My speaker notes"

    @pytest.mark.unit
    def test_extract_file_not_found(self, tmp_path, monkeypatch):
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        missing = tmp_path / "nonexistent.pptx"
        with pytest.raises(PPTMCPError):
            extract_presentation_text(str(missing))


# ── GROUP 19: manage_hyperlinks() ────────────────────────────────────────

@pytest.mark.unit
class TestManageHyperlinks:
    @pytest.mark.unit
    def test_list_no_hyperlinks(self, sample_pptx):
        result = manage_hyperlinks(sample_pptx, slide_index=0, operation="list")
        assert result["operation"] == "list"
        assert result["hyperlink_count"] == 0
        assert result["links"] == []

    @pytest.mark.unit
    def test_invalid_operation(self, sample_pptx):
        with pytest.raises(ValidationError, match="operation must be one of"):
            manage_hyperlinks(sample_pptx, slide_index=0, operation="bogus")

    @pytest.mark.unit
    def test_add_requires_write_gate(self, monkeypatch):
        monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
        # Error raised before path validation — path doesn't need to exist
        with pytest.raises(NotAllowedError):
            manage_hyperlinks(
                "dummy.pptx", slide_index=0, operation="add",
                shape_id=1, target_url="https://example.com", confirm=True,
            )

    @pytest.mark.unit
    def test_add_requires_confirm(self, monkeypatch):
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        # Error raised before path validation
        with pytest.raises(ValidationError, match="confirm"):
            manage_hyperlinks(
                "dummy.pptx", slide_index=0, operation="add",
                shape_id=1, target_url="https://example.com", confirm=False,
            )

    @pytest.mark.unit
    def test_add_no_target_url(self, monkeypatch):
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        with pytest.raises(ValidationError, match="target_url"):
            manage_hyperlinks(
                "dummy.pptx", slide_index=0, operation="add",
                shape_id=1, target_url=None, confirm=True,
            )

    @pytest.mark.unit
    def test_add_shape_not_found(self, sample_pptx, monkeypatch):
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        with pytest.raises(ValidationError, match="not found"):
            manage_hyperlinks(
                sample_pptx, slide_index=0, operation="add",
                shape_id=99999, target_url="https://example.com", confirm=True,
            )

    @pytest.mark.unit
    def test_invalid_slide_index(self, sample_pptx):
        with pytest.raises(ValidationError, match="out of range"):
            manage_hyperlinks(sample_pptx, slide_index=999, operation="list")

    @pytest.mark.unit
    def test_add_disallows_javascript_url(self, tmp_path, monkeypatch):
        """manage_hyperlinks add must reject javascript: URLs."""
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        from pptx import Presentation as _Prs
        from pptx.util import Inches
        prs_obj = _Prs()
        slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        txBox.text_frame.text = "Click me"
        path = str(tmp_path / "test.pptx")
        prs_obj.save(path)

        from pptmcp.presentation_pptx import manage_hyperlinks, ValidationError
        with pytest.raises(ValidationError, match="scheme"):
            manage_hyperlinks(
                path=path,
                slide_index=0,
                operation="add",
                shape_id=txBox.shape_id,
                target_url="javascript:alert(1)",
                confirm=True,
            )

    @pytest.mark.unit
    def test_add_hyperlink_appears_in_list(self, tmp_path, monkeypatch):
        """Happy path: add a hyperlink then list — url and shape_id must appear in result."""
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        from pptx.util import Inches
        prs_obj = _Prs()
        slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[5])  # blank layout
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        txBox.text_frame.text = "Click me"
        path = str(tmp_path / "hl_add.pptx")
        prs_obj.save(path)
        shape_id = txBox.shape_id

        add_result = manage_hyperlinks(
            path=path, slide_index=0, operation="add",
            shape_id=shape_id, target_url="https://example.com", confirm=True,
        )
        assert add_result["ok"] is True
        assert add_result["url"] == "https://example.com"
        assert add_result["runs_updated"] >= 1

        list_result = manage_hyperlinks(path=path, slide_index=0, operation="list")
        assert list_result["hyperlink_count"] >= 1
        urls = [h["url"] for h in list_result["links"]]
        assert "https://example.com" in urls
        ids = [h["shape_id"] for h in list_result["links"]]
        assert shape_id in ids

    @pytest.mark.unit
    def test_add_remove_hyperlink_round_trip(self, tmp_path, monkeypatch):
        """Round-trip: add → verify present → remove → verify gone."""
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        from pptx.util import Inches
        prs_obj = _Prs()
        slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        txBox.text_frame.text = "Round-trip link"
        path = str(tmp_path / "hl_roundtrip.pptx")
        prs_obj.save(path)
        shape_id = txBox.shape_id

        # Add
        manage_hyperlinks(
            path=path, slide_index=0, operation="add",
            shape_id=shape_id, target_url="https://roundtrip.test", confirm=True,
        )

        # Verify present
        after_add = manage_hyperlinks(path=path, slide_index=0, operation="list")
        assert any(h["url"] == "https://roundtrip.test" for h in after_add["links"]), \
            "Hyperlink not found after add"

        # Remove
        remove_result = manage_hyperlinks(
            path=path, slide_index=0, operation="remove",
            shape_id=shape_id, confirm=True,
        )
        assert remove_result["ok"] is True
        assert remove_result["runs_cleared"] >= 1

        # Verify gone
        after_remove = manage_hyperlinks(path=path, slide_index=0, operation="list")
        assert not any(
            h["url"] == "https://roundtrip.test" for h in after_remove["links"]
        ), "Hyperlink still present after remove"


# ── GROUP 30: add_hyperlink() ───────────────────────────────────────────────

@pytest.mark.unit
def test_add_hyperlink_success(tmp_path, monkeypatch):
    """add_hyperlink returns correct dict and the URL persists when re-opened."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "hyperlink.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = txBox.text_frame.paragraphs[0].add_run()
    run.text = "Click me"
    p.save(str(path))
    url = "https://example.com"
    result = add_hyperlink(str(path), 0, txBox.shape_id, run_index=0, url=url, confirm=True)
    assert result["url"] == url
    assert result["slide_index"] == 0
    assert result["shape_id"] == txBox.shape_id
    assert "path" in result
    # Save and re-open to verify hyperlink persists
    save(str(path), confirm=True)
    reopened = _Prs(str(path))
    reopened_shape = next(s for s in reopened.slides[0].shapes if s.shape_id == txBox.shape_id)
    assert reopened_shape.text_frame.paragraphs[0].runs[0].hyperlink.address == url


@pytest.mark.unit
def test_add_hyperlink_invalid_url_scheme(tmp_path, monkeypatch):
    """add_hyperlink raises ValidationError for a URL with a disallowed scheme (ftp://)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "hl_badscheme.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.paragraphs[0].add_run().text = "link"
    p.save(str(path))
    with pytest.raises(ValidationError, match="scheme"):
        add_hyperlink(str(path), 0, txBox.shape_id, url="ftp://evil.com", confirm=True)


@pytest.mark.unit
def test_add_hyperlink_javascript_url(tmp_path, monkeypatch):
    """add_hyperlink raises ValidationError for a javascript: URL (OWASP A03 injection)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "hl_js.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.paragraphs[0].add_run().text = "xss"
    p.save(str(path))
    with pytest.raises(ValidationError, match="scheme"):
        add_hyperlink(str(path), 0, txBox.shape_id, url="javascript:alert(1)", confirm=True)


@pytest.mark.unit
def test_add_hyperlink_run_index_out_of_range(tmp_path, monkeypatch):
    """add_hyperlink raises IndexError when run_index exceeds paragraph run count."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "hl_oob.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.paragraphs[0].add_run().text = "one run"
    p.save(str(path))
    with pytest.raises((IndexError, ValidationError), match="run_index"):
        add_hyperlink(str(path), 0, txBox.shape_id, run_index=99, url="https://x.com", confirm=True)


@pytest.mark.unit
def test_add_hyperlink_no_text_frame(tmp_path, monkeypatch):
    """add_hyperlink raises ValidationError for a shape with no text frame (e.g. table)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "hl_notf.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    tbl = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    p.save(str(path))
    with pytest.raises(ValidationError, match="no text frame"):
        add_hyperlink(str(path), 0, tbl.shape_id, url="https://x.com", confirm=True)


@pytest.mark.unit
def test_add_hyperlink_path_not_allowed(tmp_path, monkeypatch):
    """add_hyperlink raises ValidationError when path is outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    evil_path = str(tmp_path.parent / "evil.pptx")
    with pytest.raises(ValidationError, match="[Nn]ot in allowlist|not in allowlist"):
        add_hyperlink(evil_path, 0, 1, url="https://x.com", confirm=True)


@pytest.mark.unit
def test_add_hyperlink_no_confirm(tmp_path, monkeypatch):
    """add_hyperlink raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        add_hyperlink(path, 0, 1, url="https://x.com", confirm=False)


@pytest.mark.unit
def test_add_hyperlink_mailto_url(tmp_path, monkeypatch):
    """add_hyperlink accepts mailto: URLs and returns correct url in result."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "hl_mailto.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.paragraphs[0].add_run().text = "email"
    p.save(str(path))
    url = "mailto:user@example.com"
    result = add_hyperlink(str(path), 0, txBox.shape_id, run_index=0, url=url, confirm=True)
    assert result["url"] == url
    assert result["slide_index"] == 0


@pytest.mark.unit
def test_add_hyperlink_empty_url(tmp_path, monkeypatch):
    """add_hyperlink raises ValidationError for an empty url (no scheme)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "hl_emptyurl.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.paragraphs[0].add_run().text = "link"
    p.save(str(path))
    with pytest.raises(ValidationError, match="scheme"):
        add_hyperlink(str(path), 0, txBox.shape_id, url="", confirm=True)


@pytest.mark.unit
def test_add_hyperlink_shape_not_found(tmp_path, monkeypatch):
    """add_hyperlink raises ValidationError when shape_id does not exist on the slide."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "hl_noshape.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="not found"):
        add_hyperlink(str(path), 0, shape_id=99999, url="https://x.com", confirm=True)


@pytest.mark.unit
def test_add_hyperlink_run_index_zero_empty_paragraph(tmp_path, monkeypatch):
    """add_hyperlink raises IndexError for run_index=0 on a paragraph with no runs."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "hl_noruns.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    # Add a textbox but do NOT add any runs — paragraph.runs will be []
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    p.save(str(path))
    with pytest.raises((IndexError, ValidationError), match="run_index"):
        add_hyperlink(str(path), 0, txBox.shape_id, run_index=0, url="https://x.com", confirm=True)


@pytest.mark.unit
def test_add_hyperlink_display_text_roundtrip(tmp_path, monkeypatch):
    """add_hyperlink updates run.text when display_text is provided."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "hl_displaytext.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = txBox.text_frame.paragraphs[0].add_run()
    run.text = "original"
    p.save(str(path))
    url = "https://example.com/page"
    result = add_hyperlink(
        str(path), 0, txBox.shape_id, run_index=0, url=url,
        display_text="new label", confirm=True,
    )
    assert result["url"] == url
    # Save then reload from disk to verify both text and hyperlink persisted
    save(str(path), confirm=True)
    reloaded = _Prs(str(path))
    shape = next(s for s in reloaded.slides[0].shapes if s.shape_id == txBox.shape_id)
    run_reloaded = shape.text_frame.paragraphs[0].runs[0]
    assert run_reloaded.text == "new label"
    assert run_reloaded.hyperlink.address == url
