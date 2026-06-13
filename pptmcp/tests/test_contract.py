"""test_contract.py — Phase A Output Contract test suite

Tests for pptmcp.contract_pptx:
  GROUP 1 (4): _check_json_path      — path security guards
  GROUP 2 (5): validate_contract     — schema + IO validation
  GROUP 3 (6): check_presentation_against_contract — per-slide checks
  GROUP 4 (2): BLK-03 bool coercion guards
"""
from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation as _Prs
from pptx.util import Inches

from pptmcp.presentation_pptx import ValidationError
from pptmcp.presentation_pptx import NotAllowedError
from pptmcp.contract_pptx import (
    CONTRACT_MAX_BYTES,
    _check_json_path,
    declare_slide_contract,
    validate_contract,
    check_presentation_against_contract,
)


# ── GROUP 1: _check_json_path ─────────────────────────────────────────────


@pytest.mark.unit
def test_check_json_path_valid_returns_path(tmp_path, monkeypatch):
    """TC-CJP-01: Valid .json file inside allowlisted dir → returns resolved Path."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    contract_file = tmp_path / "contract.json"
    contract_file.write_text("{}")

    result = _check_json_path(str(contract_file))

    assert isinstance(result, Path)
    assert result == contract_file.resolve()


@pytest.mark.unit
def test_check_json_path_null_byte_raises(tmp_path, monkeypatch):
    """TC-CJP-02: Null byte in path → ValidationError (BLK-01)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    evil_path = str(tmp_path) + "/contract\x00evil.json"

    with pytest.raises(ValidationError, match=r"null byte"):
        _check_json_path(evil_path)


@pytest.mark.unit
def test_check_json_path_outside_allowlist_raises(tmp_path, monkeypatch):
    """TC-CJP-03: Path resolves outside allowlist root → ValidationError (BLK-04)."""
    allowed = tmp_path / "allowed"
    allowed.mkdir()
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(allowed))
    # File must exist so that the exists() check (step 5) passes and the allowlist
    # check (step 7) fires rather than the earlier "not found" error.
    outside_dir = tmp_path / "outside"
    outside_dir.mkdir()
    outside = outside_dir / "contract.json"
    outside.write_text("{}")

    with pytest.raises(ValidationError, match=r"not in allowlist"):
        _check_json_path(str(outside))


@pytest.mark.unit
def test_check_json_path_wrong_extension_raises(tmp_path, monkeypatch):
    """TC-CJP-04: .txt extension → ValidationError."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    txt_file = tmp_path / "contract.txt"
    txt_file.write_text("{}")

    with pytest.raises(ValidationError, match=r"[Uu]nsupported extension"):
        _check_json_path(str(txt_file))


# ── GROUP 2: validate_contract ────────────────────────────────────────────


@pytest.mark.unit
def test_validate_contract_valid_returns_true(tmp_path, monkeypatch):
    """TC-VC-01: Well-formed contract.json → {"valid": True, "slide_count": 1, ...}."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    contract_data = {
        "contract_version": "1.0",
        "slide_width_inches": 13.33,
        "slide_height_inches": 7.5,
        "slides": [{"slide_index": 0, "min_shapes": 1, "max_shapes": 8}],
    }
    contract_file = tmp_path / "contract.json"
    contract_file.write_text(json.dumps(contract_data))

    result = validate_contract(str(contract_file))

    assert result["valid"] is True
    assert result["contract_version"] == "1.0"
    assert result["slide_count"] == 1
    assert result["errors"] == []


@pytest.mark.unit
def test_validate_contract_missing_version_raises(tmp_path, monkeypatch):
    """TC-VC-02: Missing 'contract_version' field → ValidationError (ODQ-01: raises per spec)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    contract_data = {
        "slides": [{"slide_index": 0, "min_shapes": 1, "max_shapes": 8}]
    }
    contract_file = tmp_path / "contract.json"
    contract_file.write_text(json.dumps(contract_data))

    with pytest.raises(ValidationError, match=r"contract_version|missing|required"):
        validate_contract(str(contract_file))


@pytest.mark.unit
def test_validate_contract_malformed_json_raises(tmp_path, monkeypatch):
    """TC-VC-03: Non-JSON content in .json file → ValidationError."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    contract_file = tmp_path / "contract.json"
    contract_file.write_text("{ not valid JSON !!! }")

    with pytest.raises(ValidationError, match=r"[Ii]nvalid JSON|[Pp]arse|[Dd]ecode"):
        validate_contract(str(contract_file))


@pytest.mark.unit
def test_validate_contract_wrong_extension_raises(tmp_path, monkeypatch):
    """TC-VC-04: .txt extension → ValidationError (delegation chain check)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    contract_file = tmp_path / "contract.txt"
    contract_file.write_text(json.dumps({"contract_version": "1.0", "slides": []}))

    with pytest.raises(ValidationError, match=r"[Uu]nsupported extension"):
        validate_contract(str(contract_file))


@pytest.mark.unit
def test_validate_contract_file_too_large_raises(tmp_path, monkeypatch):
    """TC-VC-05: File is CONTRACT_MAX_BYTES+1 bytes → ValidationError (BLK-02 DoS guard)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    big_file = tmp_path / "big_contract.json"
    big_file.write_bytes(b"x" * (CONTRACT_MAX_BYTES + 1))

    with pytest.raises(ValidationError, match=r"[Ss]ize limit|exceed|524[_ ]?288|512"):
        validate_contract(str(big_file))


# ── GROUP 3: check_presentation_against_contract ──────────────────────────


@pytest.mark.unit
def test_check_presentation_pass_shapes_within_bounds(sample_pptx, tmp_path):
    """TC-CPC-01: 2 shapes on slide 0; contract min=1/max=5 → passed=True, findings=[]."""
    contract_data = {
        "contract_version": "1.0",
        "slides": [{"slide_index": 0, "min_shapes": 1, "max_shapes": 5}],
    }
    contract_file = tmp_path / "contract.json"
    contract_file.write_text(json.dumps(contract_data))

    result = check_presentation_against_contract(sample_pptx, str(contract_file))

    assert result["passed"] is True
    assert result["findings"] == []


@pytest.mark.unit
def test_check_presentation_fail_max_shapes_exceeded(tmp_path, monkeypatch):
    """TC-CPC-02: Slide has 3 shapes; max_shapes=1 → passed=False, max_shapes finding."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))

    pptx_path = tmp_path / "test.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[6])  # blank layout — 0 shapes
    for _ in range(3):
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
        tb.text_frame.text = "extra"
    p.save(str(pptx_path))

    contract_data = {
        "contract_version": "1.0",
        "slides": [{"slide_index": 0, "min_shapes": 0, "max_shapes": 1}],
    }
    contract_file = tmp_path / "contract.json"
    contract_file.write_text(json.dumps(contract_data))

    result = check_presentation_against_contract(str(pptx_path), str(contract_file))

    assert result["passed"] is False
    assert len(result["findings"]) >= 1
    assert any(
        "max_shapes" in str(f) or "max" in str(f).lower()
        for f in result["findings"]
    )
    assert result["findings"][0].get("slide_index") == 0


@pytest.mark.unit
def test_check_presentation_fail_min_shapes_not_met(tmp_path, monkeypatch):
    """TC-CPC-03: Blank slide (0 shapes); min_shapes=3 → passed=False, min_shapes finding."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))

    pptx_path = tmp_path / "test.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[6])  # blank layout — 0 shapes
    p.save(str(pptx_path))

    contract_data = {
        "contract_version": "1.0",
        "slides": [{"slide_index": 0, "min_shapes": 3, "max_shapes": 10}],
    }
    contract_file = tmp_path / "contract.json"
    contract_file.write_text(json.dumps(contract_data))

    result = check_presentation_against_contract(str(pptx_path), str(contract_file))

    assert result["passed"] is False
    assert len(result["findings"]) >= 1
    assert any(
        "min_shapes" in str(f) or "min" in str(f).lower()
        for f in result["findings"]
    )


@pytest.mark.unit
def test_check_presentation_fail_expected_title_mismatch(sample_pptx, tmp_path):
    """TC-CPC-04: Slide title 'Test Slide' ≠ expected 'Executive Summary' → passed=False."""
    contract_data = {
        "contract_version": "1.0",
        "slides": [
            {
                "slide_index": 0,
                "min_shapes": 1,
                "max_shapes": 5,
                "expected_title": "Executive Summary",
            }
        ],
    }
    contract_file = tmp_path / "contract.json"
    contract_file.write_text(json.dumps(contract_data))

    result = check_presentation_against_contract(sample_pptx, str(contract_file))

    assert result["passed"] is False
    assert len(result["findings"]) >= 1
    assert any(
        "title" in str(f).lower() or "expected_title" in str(f)
        for f in result["findings"]
    )


@pytest.mark.unit
def test_check_presentation_null_byte_in_pptx_path_raises(tmp_path, monkeypatch):
    """TC-CPC-05: Null byte in pptx path → ValidationError (BLK-01 on _check_path)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    contract_file = tmp_path / "contract.json"
    contract_file.write_text(json.dumps({"contract_version": "1.0", "slides": []}))
    evil_pptx = str(tmp_path) + "/file\x00evil.pptx"

    with pytest.raises(ValidationError, match=r"null byte"):
        check_presentation_against_contract(evil_pptx, str(contract_file))


@pytest.mark.unit
def test_check_presentation_contract_outside_allowlist_raises(sample_pptx, tmp_path):
    """TC-CPC-06: Contract path outside allowlist → ValidationError (BLK-04 + DC-02 step 2)."""
    # Place the contract file outside tmp_path (which is the configured allowlist).
    # It must physically exist so the existence check (step 5) passes and the
    # allowlist check (step 7) fires.
    outside_dir = tmp_path.parent / "contract_outside"
    outside_dir.mkdir(exist_ok=True)
    outside_contract = outside_dir / "contract.json"
    outside_contract.write_text(json.dumps({"contract_version": "1.0", "slides": []}))

    with pytest.raises(ValidationError, match=r"not in allowlist"):
        check_presentation_against_contract(sample_pptx, str(outside_contract))


# ── GROUP 4: BLK-03 bool coercion guards ──────────────────────────────────


@pytest.mark.unit
def test_validate_contract_min_shapes_bool_raises(tmp_path, monkeypatch):
    """TC-BC-01: min_shapes=True (JSON bool) → ValidationError (BLK-03 type confusion)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    contract_data = {
        "contract_version": "1.0",
        "slides": [{"slide_index": 0, "min_shapes": True, "max_shapes": 8}],
    }
    contract_file = tmp_path / "contract.json"
    contract_file.write_text(json.dumps(contract_data))

    with pytest.raises(ValidationError, match=r"min_shapes|bool|integer|int"):
        validate_contract(str(contract_file))


@pytest.mark.unit
def test_validate_contract_max_shapes_bool_raises(tmp_path, monkeypatch):
    """TC-BC-02: max_shapes=False (JSON bool) → ValidationError (BLK-03 type confusion)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    contract_data = {
        "contract_version": "1.0",
        "slides": [{"slide_index": 0, "min_shapes": 1, "max_shapes": False}],
    }
    contract_file = tmp_path / "contract.json"
    contract_file.write_text(json.dumps(contract_data))

    with pytest.raises(ValidationError, match=r"max_shapes|bool|integer|int"):
        validate_contract(str(contract_file))


# ── GROUP 5: declare_slide_contract — RCI-US-01 ───────────────────────────


def _save_blank_prs(tmp_path: Path, name: str = "deck.pptx") -> Path:
    prs = _Prs()
    prs.slides.add_slide(prs.slide_layouts[6])
    out = tmp_path / name
    prs.save(str(out))
    return out


@pytest.mark.unit
def test_declare_slide_contract_happy_path(tmp_path, monkeypatch):
    """RCI-US-01: writes contract.json beside the PPTX with sha256 checksum."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx = _save_blank_prs(tmp_path)

    spec = {
        "slide_width_inches": 13.33,
        "slide_height_inches": 7.5,
        "slides": [{"slide_index": 0, "min_shapes": 0}],
    }
    result = declare_slide_contract(str(pptx), spec, confirm=True)

    contract_path = Path(result["contract_path"])
    assert contract_path.name == "contract.json"
    assert contract_path.parent == pptx.parent
    assert contract_path.exists()

    on_disk = json.loads(contract_path.read_text(encoding="utf-8"))
    assert on_disk["contract_version"] == "1.0"  # auto-filled
    assert on_disk["slide_width_inches"] == 13.33

    # Checksum matches the written bytes (re-hash to verify)
    import hashlib
    expected = hashlib.sha256(contract_path.read_bytes()).hexdigest()
    assert result["checksum_sha256"] == expected
    assert result["contract_version"] == "1.0"


@pytest.mark.unit
def test_declare_slide_contract_write_gate_blocks(tmp_path, monkeypatch):
    """PPT_ENABLE_WRITE not set → NotAllowedError; contract.json must not appear."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    pptx = _save_blank_prs(tmp_path)

    with pytest.raises(NotAllowedError, match=r"PPT_ENABLE_WRITE"):
        declare_slide_contract(str(pptx), {"slides": []}, confirm=True)
    assert not (tmp_path / "contract.json").exists()


@pytest.mark.unit
def test_declare_slide_contract_confirm_required(tmp_path, monkeypatch):
    """confirm=False (default) → ValidationError; no file is written."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx = _save_blank_prs(tmp_path)

    with pytest.raises(ValidationError, match=r"confirm=True"):
        declare_slide_contract(str(pptx), {"slides": []})
    assert not (tmp_path / "contract.json").exists()


@pytest.mark.unit
def test_declare_slide_contract_outside_allowlist_raises(tmp_path, monkeypatch):
    """PPTX path outside PPT_ALLOWLIST_ROOTS → ValidationError (OWASP A01)."""
    allowed = tmp_path / "allowed"
    allowed.mkdir()
    outside = tmp_path / "outside"
    outside.mkdir()
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(allowed))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    bad_pptx = _save_blank_prs(outside)

    with pytest.raises(ValidationError, match=r"allowlist"):
        declare_slide_contract(str(bad_pptx), {"slides": []}, confirm=True)


@pytest.mark.unit
def test_declare_slide_contract_invalid_spec_rolls_back(tmp_path, monkeypatch):
    """Schema-invalid spec → ValidationError; no contract.json or .tmp left on disk."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx = _save_blank_prs(tmp_path)

    # slide_index must be int — string here forces schema failure
    bad_spec = {"slides": [{"slide_index": "zero"}]}
    with pytest.raises(ValidationError):
        declare_slide_contract(str(pptx), bad_spec, confirm=True)

    assert not (tmp_path / "contract.json").exists()
    assert not (tmp_path / "contract.json.tmp").exists()


@pytest.mark.unit
def test_declare_slide_contract_accepts_extra_keys(tmp_path, monkeypatch):
    """RCI-US-01 spec accepts element_inventory/geometry_tokens/color_palette/etc."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx = _save_blank_prs(tmp_path)

    spec = {
        "slide_width_inches": 13.33,
        "slide_height_inches": 7.5,
        "slides": [],
        "element_inventory": [{"name": "title", "type": "text"}],
        "geometry_tokens": {"col_w": 4.0},
        "color_palette": {"primary": "#0055AA"},
        "export_plan": ["pdf", "png"],
        "acceptance_thresholds": {"max_overflow_in": 0.0},
    }
    result = declare_slide_contract(str(pptx), spec, confirm=True)
    on_disk = json.loads(Path(result["contract_path"]).read_text(encoding="utf-8"))
    assert on_disk["color_palette"] == {"primary": "#0055AA"}
    assert on_disk["acceptance_thresholds"] == {"max_overflow_in": 0.0}
