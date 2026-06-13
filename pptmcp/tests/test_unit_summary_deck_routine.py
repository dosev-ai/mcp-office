"""Unit tests for summary_deck_routine.create_summary_deck_v1.

CI-safe: no COM, no win32com, no live Office required.
Uses python-pptx to verify the created presentation structure.
"""
from __future__ import annotations

import os

import pytest
from pptx import Presentation

from pptmcp.summary_deck_routine import create_summary_deck_v1


# ---------------------------------------------------------------------------
# Sample data
# ---------------------------------------------------------------------------

SAMPLE_HEADERS = ["Ticker", "Score", "Decision"]
SAMPLE_ROWS = [
    ["AAPL", 0.92, "BUY"],
    ["MSFT", 0.55, "HOLD"],
    ["TSLA", 0.30, "SELL"],
]


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


class TestWriteGate:
    """GAP-S2-02 T1+T2: PPT_ENABLE_WRITE gate checks."""

    def test_write_gate_not_set(
        self, monkeypatch: pytest.MonkeyPatch, tmp_path
    ) -> None:
        monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
        result = create_summary_deck_v1(
            presentation_path=str(tmp_path / "out.pptx"),
            title="Test",
            rows=SAMPLE_ROWS,
            headers=SAMPLE_HEADERS,
        )
        assert result["error_kind"] == "gate_error"
        assert "PPT_ENABLE_WRITE" in result["error"]

    def test_write_gate_false(
        self, monkeypatch: pytest.MonkeyPatch, tmp_path
    ) -> None:
        monkeypatch.setenv("PPT_ENABLE_WRITE", "false")
        result = create_summary_deck_v1(
            presentation_path=str(tmp_path / "out.pptx"),
            title="Test",
            rows=SAMPLE_ROWS,
            headers=SAMPLE_HEADERS,
        )
        assert result["error_kind"] == "gate_error"

    def test_write_gate_mixed_case_false(
        self, monkeypatch: pytest.MonkeyPatch, tmp_path
    ) -> None:
        monkeypatch.setenv("PPT_ENABLE_WRITE", "False")
        result = create_summary_deck_v1(
            presentation_path=str(tmp_path / "out.pptx"),
            title="Test",
            rows=SAMPLE_ROWS,
            headers=SAMPLE_HEADERS,
        )
        assert result["error_kind"] == "gate_error"

    def test_write_gate_empty_string(
        self, monkeypatch: pytest.MonkeyPatch, tmp_path
    ) -> None:
        monkeypatch.setenv("PPT_ENABLE_WRITE", "")
        result = create_summary_deck_v1(
            presentation_path=str(tmp_path / "out.pptx"),
            title="Test",
            rows=SAMPLE_ROWS,
            headers=SAMPLE_HEADERS,
        )
        assert result["error_kind"] == "gate_error"


class TestAllowlistViolation:
    """GAP-S2-02 T3: path outside PPT_ALLOWLIST_ROOTS → allowlist_violation."""

    def test_allowlist_violation(
        self, monkeypatch: pytest.MonkeyPatch, tmp_path
    ) -> None:
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")

        other_dir = tmp_path / "other"
        other_dir.mkdir()
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(other_dir))

        # Path outside allowed root
        result = create_summary_deck_v1(
            presentation_path=str(tmp_path / "out.pptx"),
            title="Test",
            rows=SAMPLE_ROWS,
            headers=SAMPLE_HEADERS,
        )
        assert result["error_kind"] == "allowlist_violation"
        assert "not in allowlist" in result["error"]

    def test_allowlist_not_configured(
        self, monkeypatch: pytest.MonkeyPatch, tmp_path
    ) -> None:
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.delenv("PPT_ALLOWLIST_ROOTS", raising=False)

        result = create_summary_deck_v1(
            presentation_path=str(tmp_path / "out.pptx"),
            title="Test",
            rows=SAMPLE_ROWS,
            headers=SAMPLE_HEADERS,
        )
        assert result["error_kind"] == "config_error"


class TestCreateDeckSuccess:
    """GAP-S2-02 T4: happy-path deck creation produces valid .pptx."""

    def test_create_deck_success(
        self, monkeypatch: pytest.MonkeyPatch, tmp_path
    ) -> None:
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))

        out_path = str(tmp_path / "summary.pptx")
        result = create_summary_deck_v1(
            presentation_path=out_path,
            title="Trader Summary",
            rows=SAMPLE_ROWS,
            headers=SAMPLE_HEADERS,
            source_metadata={"description": "Summary deck example"},
        )

        assert "error" not in result, f"Unexpected error: {result}"
        assert result["slide_count"] == 2
        assert result["presentation_path"] == out_path
        assert result["evidence"]["routine"] == "pptmcp.create_summary_deck_v1"
        assert result["evidence"]["row_count"] == len(SAMPLE_ROWS)
        assert result["evidence"]["headers"] == SAMPLE_HEADERS

        # Verify file exists and is a valid presentation
        assert os.path.isfile(out_path)
        prs = Presentation(out_path)
        assert len(prs.slides) == 2

    def test_create_deck_title_text_on_slide1(
        self, monkeypatch: pytest.MonkeyPatch, tmp_path
    ) -> None:
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))

        out_path = str(tmp_path / "titled.pptx")
        create_summary_deck_v1(
            presentation_path=out_path,
            title="My Custom Title",
            rows=SAMPLE_ROWS,
            headers=SAMPLE_HEADERS,
        )

        prs = Presentation(out_path)
        slide1 = prs.slides[0]
        assert slide1.shapes.title is not None
        assert slide1.shapes.title.text == "My Custom Title"

    def test_create_deck_table_has_headers(
        self, monkeypatch: pytest.MonkeyPatch, tmp_path
    ) -> None:
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))

        out_path = str(tmp_path / "table.pptx")
        create_summary_deck_v1(
            presentation_path=out_path,
            title="Table Test",
            rows=SAMPLE_ROWS,
            headers=SAMPLE_HEADERS,
        )

        prs = Presentation(out_path)
        slide2 = prs.slides[1]
        # Find the table shape
        tables = [s for s in slide2.shapes if s.has_table]
        assert len(tables) == 1
        tbl = tables[0].table
        # First row should be headers
        assert tbl.cell(0, 0).text == "Ticker"
        assert tbl.cell(0, 1).text == "Score"
        assert tbl.cell(0, 2).text == "Decision"

    def test_create_deck_truncates_at_50_rows(
        self, monkeypatch: pytest.MonkeyPatch, tmp_path
    ) -> None:
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))

        out_path = str(tmp_path / "big.pptx")
        big_rows = [[f"R{i}", i, "BUY"] for i in range(60)]

        result = create_summary_deck_v1(
            presentation_path=out_path,
            title="Big Table",
            rows=big_rows,
            headers=SAMPLE_HEADERS,
        )

        assert "error" not in result
        assert result["evidence"]["row_count"] == 60  # full count preserved

        prs = Presentation(out_path)
        slide2 = prs.slides[1]
        tables = [s for s in slide2.shapes if s.has_table]
        tbl = tables[0].table
        # header (1) + 50 data rows + 1 truncation notice = 52 rows
        assert len(tbl.rows) == 52

    def test_write_gate_true_case_insensitive(
        self, monkeypatch: pytest.MonkeyPatch, tmp_path
    ) -> None:
        monkeypatch.setenv("PPT_ENABLE_WRITE", "TRUE")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))

        out_path = str(tmp_path / "ci.pptx")
        result = create_summary_deck_v1(
            presentation_path=out_path,
            title="Case Test",
            rows=SAMPLE_ROWS,
            headers=SAMPLE_HEADERS,
        )
        assert "error" not in result
        assert result["slide_count"] == 2

    def test_no_source_metadata(
        self, monkeypatch: pytest.MonkeyPatch, tmp_path
    ) -> None:
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))

        out_path = str(tmp_path / "nometa.pptx")
        result = create_summary_deck_v1(
            presentation_path=out_path,
            title="No Meta",
            rows=SAMPLE_ROWS,
            headers=SAMPLE_HEADERS,
            source_metadata=None,
        )
        assert "error" not in result
        assert result["slide_count"] == 2
