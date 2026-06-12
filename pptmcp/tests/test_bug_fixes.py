"""Regression tests for the 5 UAT-reported bugs — pptmcp.

BUG 1 delete_shape tool
BUG 2 replace_slide_text fails on <a:br> elements
BUG 3 edit_text_placeholder shape_id + placeholder idx
BUG 4 copy_slide param name aliases
BUG 5 set_speaker_notes notes alias

New regressions (BF-01/02/03) — filed March 5 2026
BF-01 capabilities() omits review_slide_export + export_slides_to_stamped_dir (COM)
BF-02 create_presentation does not accept template_path parameter
BF-03 review_slide_export false-positive PASS on table content clipping
"""

from __future__ import annotations

import pytest
from pptx import Presentation as _Prs
from pptx.oxml.ns import qn
from lxml import etree

from pptmcp.presentation_pptx import ValidationError, NotAllowedError, copy_slide as _copy_slide, _prs_cache
from pptmcp.shapes_pptx import (
    delete_shape,
    replace_slide_text,
    edit_text_placeholder,
    set_speaker_notes,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_prs_with_textbox(tmp_path, text: str) -> tuple[str, int]:
    """Create a 1-slide pptx, add a textbox with text, return (path, shape_id)."""
    from pptx.util import Inches
    prs = _Prs()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txb.text_frame.text = text
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return str(path), txb.shape_id


def _make_prs_with_br(tmp_path, line1: str = "hello", line2: str = "world") -> tuple[str, int]:
    """Create a 1-slide pptx with a textbox containing an <a:br> line-break element."""
    from pptx.util import Inches
    prs = _Prs()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = txb.text_frame
    # Build para with run + <a:br> + run via XML to guarantee the <a:br> is present
    para_p = tf.paragraphs[0]._p
    # Clear existing children
    for child in list(para_p):
        para_p.remove(child)
    r1 = etree.SubElement(para_p, qn('a:r'))
    t1 = etree.SubElement(r1, qn('a:t'))
    t1.text = line1
    etree.SubElement(para_p, qn('a:br'))
    r2 = etree.SubElement(para_p, qn('a:r'))
    t2 = etree.SubElement(r2, qn('a:t'))
    t2.text = line2
    path = tmp_path / "br_test.pptx"
    prs.save(str(path))
    return str(path), txb.shape_id


# ===========================================================================
# BUG 1 — delete_shape
# ===========================================================================

@pytest.mark.unit
def test_delete_shape_success(tmp_path, monkeypatch):
    """delete_shape removes the shape and returns deleted=True."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path, shape_id = _make_prs_with_textbox(tmp_path, "to delete")

    result = delete_shape(path, 0, shape_id, confirm=True)

    assert result["deleted"] is True
    assert result["shape_id"] == shape_id
    assert result["slide_index"] == 0

    # Verify shape is gone from the in-memory cache (delete is in-memory until save())
    from pptmcp.presentation_pptx import _load_prs
    from pathlib import Path
    live = _load_prs(Path(path))
    assert shape_id not in [sh.shape_id for sh in live.slides[0].shapes]


@pytest.mark.unit
def test_delete_shape_not_found(tmp_path, monkeypatch):
    """delete_shape raises ValidationError if shape_id does not exist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path, _ = _make_prs_with_textbox(tmp_path, "some shape")
    with pytest.raises(ValidationError, match="shape_id 9999"):
        delete_shape(path, 0, 9999, confirm=True)


@pytest.mark.unit
def test_delete_shape_requires_confirm(tmp_path, monkeypatch):
    """delete_shape raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path, shape_id = _make_prs_with_textbox(tmp_path, "x")
    with pytest.raises(ValidationError, match="confirm"):
        delete_shape(path, 0, shape_id, confirm=False)


@pytest.mark.unit
def test_delete_shape_requires_write_gate(tmp_path, monkeypatch):
    """delete_shape raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path, shape_id = _make_prs_with_textbox(tmp_path, "x")
    with pytest.raises(NotAllowedError):
        delete_shape(path, 0, shape_id, confirm=True)


# ===========================================================================
# BUG 2 — replace_slide_text with <a:br> elements
# ===========================================================================

@pytest.mark.unit
def test_replace_slide_text_no_br(tmp_path, monkeypatch):
    """replace_slide_text still works on normal runs without <a:br>."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path, _ = _make_prs_with_textbox(tmp_path, "hello world")
    result = replace_slide_text(path, "world", "pptmcp", confirm=True)
    assert result["replacements_made"] == 1


@pytest.mark.unit
def test_replace_slide_text_with_br_no_newline_in_find(tmp_path, monkeypatch):
    """replace_slide_text works when find string appears in a run adjacent to <a:br>."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path, shape_id = _make_prs_with_br(tmp_path, "hello", "world")
    # "hello" should be found in the first run (before <a:br>)
    result = replace_slide_text(path, "hello", "greet", confirm=True)
    assert result["replacements_made"] == 1


@pytest.mark.unit
def test_replace_slide_text_spanning_br(tmp_path, monkeypatch):
    """replace_slide_text finds and replaces text that spans an <a:br> boundary."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path, shape_id = _make_prs_with_br(tmp_path, "line1", "line2")
    # The find string uses \n which maps to <a:br>
    result = replace_slide_text(path, "line1\nline2", "single line", confirm=True)
    assert result["replacements_made"] == 1


@pytest.mark.unit
def test_replace_slide_text_br_no_exception(tmp_path, monkeypatch):
    """replace_slide_text does NOT raise when paragraph contains <a:br> elements."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    # Create a pptx specifically with <a:br> where find isn't present
    path, _ = _make_prs_with_br(tmp_path, "alpha", "beta")
    # This should complete without exception regardless of br presence
    result = replace_slide_text(path, "NOTPRESENT", "anything", confirm=True)
    assert result["replacements_made"] == 0


# ===========================================================================
# BUG 3 — edit_text_placeholder accepts shape_id
# ===========================================================================

@pytest.mark.unit
def test_edit_text_placeholder_by_shape_id(tmp_path, monkeypatch):
    """edit_text_placeholder with shape_id targets any text-containing shape."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path, shape_id = _make_prs_with_textbox(tmp_path, "original text")
    result = edit_text_placeholder(
        path, 0, text="replaced", confirm=True, shape_id=shape_id
    )
    assert result["text"] == "replaced"
    assert result["shape_id"] == shape_id


@pytest.mark.unit
def test_edit_text_placeholder_shape_id_no_text_frame(tmp_path, monkeypatch):
    """edit_text_placeholder with shape_id raises if shape has no text frame."""
    from pptx.util import Inches
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    prs = _Prs()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Add a picture/image placeholder would be complex; add a table which has no text_frame
    tbl = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    path = tmp_path / "tbl.pptx"
    prs.save(str(path))
    with pytest.raises(ValidationError, match="no text frame"):
        edit_text_placeholder(
            str(path), 0, text="fail", confirm=True, shape_id=tbl.shape_id
        )


@pytest.mark.unit
def test_edit_text_placeholder_by_placeholder_idx(tmp_path, monkeypatch):
    """edit_text_placeholder without shape_id uses placeholder_format.idx matching."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    prs = _Prs()
    layout = prs.slide_layouts[1]  # Title and Content layout
    prs.slides.add_slide(layout)
    path = tmp_path / "ph.pptx"
    prs.save(str(path))
    # idx=0 is the title placeholder in standard layouts
    result = edit_text_placeholder(str(path), 0, placeholder_idx=0, text="My Title", confirm=True)
    assert result["text"] == "My Title"
    assert result["placeholder_idx"] == 0


@pytest.mark.unit
def test_edit_text_placeholder_invalid_idx(tmp_path, monkeypatch):
    """edit_text_placeholder raises ValidationError for a missing placeholder idx."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    prs = _Prs()
    prs.slides.add_slide(prs.slide_layouts[5])  # blank — no placeholders
    path = tmp_path / "blank.pptx"
    prs.save(str(path))
    with pytest.raises(ValidationError):
        edit_text_placeholder(str(path), 0, placeholder_idx=99, text="x", confirm=True)


# ===========================================================================
# BUG 4 — copy_slide param aliases in server.py
# ===========================================================================

@pytest.mark.unit
def test_copy_slide_alias_path_and_slide_index(tmp_path, monkeypatch):
    """copy_slide server tool accepts path/slide_index aliases for source params."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    prs = _Prs()
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.slides.add_slide(prs.slide_layouts[5])
    p = tmp_path / "deck.pptx"
    prs.save(str(p))

    # Use aliases path= and slide_index= (instead of source_path / source_slide_index)
    result = _copy_slide(
        path=str(p),
        slide_index=0,
        target_path=str(p),
        confirm=True,
    )
    assert "new_slide_index" in result


@pytest.mark.unit
def test_copy_slide_original_params_still_work(tmp_path, monkeypatch):
    """copy_slide server tool still works with original source_path / source_slide_index."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    prs = _Prs()
    prs.slides.add_slide(prs.slide_layouts[5])
    p = tmp_path / "deck2.pptx"
    prs.save(str(p))

    result = _copy_slide(
        source_path=str(p),
        source_slide_index=0,
        target_path=str(p),
        confirm=True,
    )
    assert "new_slide_index" in result


@pytest.mark.unit
def test_copy_slide_defaults_target_to_source(tmp_path, monkeypatch):
    """copy_slide defaults target_path to source path when target_path is omitted."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    prs = _Prs()
    prs.slides.add_slide(prs.slide_layouts[5])
    p = tmp_path / "deck3.pptx"
    prs.save(str(p))

    result = _copy_slide(
        path=str(p),
        slide_index=0,
        confirm=True,
        # target_path intentionally omitted — should default to source
    )
    assert result["source_path"] == result["target_path"]


@pytest.mark.unit
def test_copy_slide_missing_source_raises(tmp_path, monkeypatch):
    """copy_slide raises ToolError when neither source_path nor path is provided."""
    from fastmcp.exceptions import ToolError
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises((ToolError, Exception), match="source_path|path"):
        _copy_slide(slide_index=0, confirm=True)


# ===========================================================================
# BUG 5 — set_speaker_notes notes alias
# ===========================================================================

@pytest.mark.unit
def test_set_speaker_notes_with_notes_alias(tmp_path, monkeypatch):
    """set_speaker_notes accepts notes= alias instead of notes_text=."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    prs = _Prs()
    prs.slides.add_slide(prs.slide_layouts[5])
    path = tmp_path / "notes.pptx"
    prs.save(str(path))
    result = set_speaker_notes(str(path), 0, notes="my notes text", confirm=True)
    assert result["notes_text"] == "my notes text"


@pytest.mark.unit
def test_set_speaker_notes_notes_text_takes_precedence(tmp_path, monkeypatch):
    """When both notes_text and notes are passed, notes_text wins."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    prs = _Prs()
    prs.slides.add_slide(prs.slide_layouts[5])
    path = tmp_path / "notesprio.pptx"
    prs.save(str(path))
    result = set_speaker_notes(
        str(path), 0,
        notes_text="primary",
        notes="alias",
        confirm=True,
    )
    assert result["notes_text"] == "primary"


@pytest.mark.unit
def test_set_speaker_notes_original_param_still_works(tmp_path, monkeypatch):
    """set_speaker_notes still works with the original notes_text= param."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    prs = _Prs()
    prs.slides.add_slide(prs.slide_layouts[5])
    path = tmp_path / "notesold.pptx"
    prs.save(str(path))
    result = set_speaker_notes(str(path), 0, notes_text="original param", confirm=True)
    assert result["notes_text"] == "original param"


@pytest.mark.unit
def test_set_speaker_notes_neither_provided_raises(tmp_path, monkeypatch):
    """set_speaker_notes raises ValidationError when neither notes nor notes_text is given."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    prs = _Prs()
    prs.slides.add_slide(prs.slide_layouts[5])
    path = tmp_path / "notesnone.pptx"
    prs.save(str(path))
    with pytest.raises(ValidationError, match="notes"):
        set_speaker_notes(str(path), 0, confirm=True)


@pytest.mark.unit
def test_set_speaker_notes_server_notes_alias(tmp_path, monkeypatch):
    """server.set_speaker_notes passes notes alias through to shapes layer."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    prs = _Prs()
    prs.slides.add_slide(prs.slide_layouts[5])
    path = tmp_path / "srvnotes.pptx"
    prs.save(str(path))
    result = set_speaker_notes(str(path), 0, notes="via server", confirm=True)
    assert result["notes_text"] == "via server"

# ===========================================================================
# BF-01 — capabilities() missing review_slide_export + export_slides_to_stamped_dir
# ===========================================================================

@pytest.mark.unit
def test_bf01_capabilities_includes_review_slide_export(tmp_path, monkeypatch):
    """capabilities() must list review_slide_export in the tools list."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    from pptmcp.presentation_pptx import capabilities
    result = capabilities()
    tool_names = {e["tool"] for e in result["tools"]}
    assert "review_slide_export" in tool_names, (
        "BF-01: 'review_slide_export' missing from capabilities tools list"
    )


@pytest.mark.unit
def test_bf01_capabilities_includes_export_slides_to_stamped_dir(tmp_path, monkeypatch):
    """capabilities() must list export_slides_to_stamped_dir in tools (always-registered)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    from pptmcp.presentation_pptx import capabilities
    result = capabilities()
    # After FR-01, always-registered COM tools appear in "tools" (not com_tools)
    tool_names = {e["tool"] for e in result["tools"]}
    assert "export_slides_to_stamped_dir" in tool_names, (
        "BF-01: 'export_slides_to_stamped_dir' missing from capabilities tools"
    )


# ===========================================================================
# BF-02 — create_presentation template_path parameter
# ===========================================================================

@pytest.mark.unit
def test_bf02_create_presentation_no_template(tmp_path, monkeypatch):
    """create_presentation still works without template_path (no regression)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptmcp.presentation_pptx import create_presentation
    out = tmp_path / "new.pptx"
    result = create_presentation(str(out), title="Hello", template_path=None, confirm=True)
    assert result["path"] == str(out)
    assert result["template_path"] is None
    assert out.exists()
    from pptx import Presentation as _Prs2
    reloaded = _Prs2(str(out))
    assert len(reloaded.slides) == 1
    assert reloaded.slides[0].shapes.title is not None
    assert reloaded.slides[0].shapes.title.text == "Hello"


@pytest.mark.unit
def test_bf02_create_presentation_nonexistent_template_raises(tmp_path, monkeypatch):
    """create_presentation raises ValidationError when template_path does not exist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptmcp.presentation_pptx import create_presentation
    out = tmp_path / "new2.pptx"
    bad_template = tmp_path / "ghost.pptx"
    with pytest.raises(ValidationError, match="[Tt]emplate"):
        create_presentation(str(out), template_path=str(bad_template), confirm=True)


@pytest.mark.unit
def test_bf02_create_presentation_template_outside_allowlist_raises(tmp_path, monkeypatch):
    """create_presentation raises ValidationError when template_path is outside the allowlist."""
    allowed = tmp_path / "allowed"
    allowed.mkdir()
    outside = tmp_path / "outside"
    outside.mkdir()
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(allowed))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptmcp.presentation_pptx import create_presentation
    out = allowed / "new3.pptx"
    bad_template = outside / "template.pptx"
    # Create the file so the not-found check is skipped; allowlist check should fire
    _t = _Prs()
    _t.save(str(bad_template))
    with pytest.raises(ValidationError, match="allowlist|[Aa]llowlist"):
        create_presentation(str(out), template_path=str(bad_template), confirm=True)


# ===========================================================================
# BF-03 — TABLE_CONTENT_CLIPPING check in review_slide_export
# ===========================================================================

def _make_png(tmp_path):
    """Write a minimal 1x1 white PNG for use as the png_path argument."""
    import struct
    import zlib
    def _chunk(tag, data):
        c = zlib.crc32(tag + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", c)
    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = _chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
    raw = b"\x00\xFF\xFF\xFF"
    idat = _chunk(b"IDAT", zlib.compress(raw))
    iend = _chunk(b"IEND", b"")
    png_path = tmp_path / "slide.png"
    png_path.write_bytes(sig + ihdr + idat + iend)
    return str(png_path)


def _emu(inches: float) -> int:
    return int(inches * 914_400)


@pytest.mark.unit
def test_bf03_table_clipping_explicit_heights_overflow(tmp_path, monkeypatch):
    """Table with explicit row heights summing > shape height → error finding."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptmcp.review_pptx import _detect_table_content_clipping

    prs_obj = _Prs()
    slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[5])

    from pptx.util import Emu
    rows, cols = 8, 2
    # Shape height = 4 in; rows auto-distributed at 0.5 in each → total == shape height
    table_shape = slide.shapes.add_table(rows, cols, Emu(_emu(1)), Emu(_emu(1)), Emu(_emu(4)), Emu(_emu(4)))

    # Shrink the shape ext.cy to 1 in at the XML level — rows still sum to 4 in → clipping
    from pptx.oxml.ns import qn as _qn
    xfrm = table_shape.element.find(".//" + _qn("a:ext"))
    xfrm.set("cy", str(_emu(1)))  # shape box now says 1 in, rows sum to 4 in

    deck = tmp_path / "clip.pptx"
    prs_obj.save(str(deck))

    from pptmcp.presentation_pptx import _load_prs, _check_path
    loaded = _load_prs(_check_path(str(deck)))
    findings = _detect_table_content_clipping(loaded.slides[0])
    assert any(f["criterion"] == "TABLE_CONTENT_CLIPPING" and f["result"] == "FAIL" and f["severity"] == "high" for f in findings), (
        f"Expected TABLE_CONTENT_CLIPPING FAIL/high, got: {findings}"
    )


@pytest.mark.unit
def test_bf03_table_clipping_mixed_auto_explicit_overflow(tmp_path, monkeypatch):
    """Table with explicit rows summing > shape height AND auto-height rows → high severity (Branch 2)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    from pptmcp.review_pptx import _detect_table_content_clipping

    prs_obj = _Prs()
    slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[5])

    from pptx.util import Emu
    from pptx.oxml.ns import qn as _qn
    rows, cols = 4, 2
    # Shape height = 2 in; table created with that height
    table_shape = slide.shapes.add_table(rows, cols, Emu(_emu(1)), Emu(_emu(1)), Emu(_emu(4)), Emu(_emu(2)))

    # Set first 3 rows to explicit height of 1 in each → known_height = 3 in > shape_height = 2 in
    tr_elems = list(table_shape.element.iter(_qn("a:tr")))
    for tr in tr_elems[:3]:
        tr.set("h", str(_emu(1)))
    # Set last row to h="0" → auto-height row (auto_rows > 0 condition)
    tr_elems[3].set("h", "0")

    deck = tmp_path / "mixed_overflow.pptx"
    prs_obj.save(str(deck))

    from pptmcp.presentation_pptx import _load_prs, _check_path
    loaded = _load_prs(_check_path(str(deck)))
    findings = _detect_table_content_clipping(loaded.slides[0])
    assert any(
        f["criterion"] == "TABLE_CONTENT_CLIPPING"
        and f["result"] == "FAIL"
        and f["severity"] == "high"
        for f in findings
    ), f"Expected TABLE_CONTENT_CLIPPING high finding, got: {findings}"


@pytest.mark.unit
def test_bf03_table_clipping_auto_height_many_rows(tmp_path, monkeypatch):
    """Table with >6 rows and auto-height rows → warning finding."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    from pptmcp.review_pptx import _detect_table_content_clipping

    prs_obj = _Prs()
    slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[5])

    from pptx.util import Emu
    from pptx.oxml.ns import qn as _qn
    rows = 8
    table_shape = slide.shapes.add_table(rows, 2, Emu(_emu(1)), Emu(_emu(1)), Emu(_emu(4)), Emu(_emu(2)))
    # Set h=0 on each tr element — h=0 signals auto-height in OOXML
    for tr_elem in table_shape.element.iter(_qn("a:tr")):
        tr_elem.set("h", "0")

    deck = tmp_path / "autoheight.pptx"
    prs_obj.save(str(deck))

    from pptmcp.presentation_pptx import _load_prs, _check_path
    loaded = _load_prs(_check_path(str(deck)))
    findings = _detect_table_content_clipping(loaded.slides[0])
    assert any(f["criterion"] == "TABLE_CONTENT_CLIPPING" and f["result"] == "FAIL" and f["severity"] == "medium" for f in findings), (
        f"Expected TABLE_CONTENT_CLIPPING FAIL/medium, got: {findings}"
    )


@pytest.mark.unit
def test_bf03_table_clipping_small_table_no_finding(tmp_path, monkeypatch):
    """Small table (3 rows, explicit heights fitting shape) → no TABLE_CONTENT_CLIPPING finding."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    from pptmcp.review_pptx import _detect_table_content_clipping

    prs_obj = _Prs()
    slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[5])

    from pptx.util import Emu
    rows = 3
    # Shape height = 2 in; 3 rows auto-distributed at 2/3 in each → total == shape height (fits)
    slide.shapes.add_table(rows, 2, Emu(_emu(1)), Emu(_emu(1)), Emu(_emu(4)), Emu(_emu(2)))

    deck = tmp_path / "smalltable.pptx"
    prs_obj.save(str(deck))

    from pptmcp.presentation_pptx import _load_prs, _check_path
    loaded = _load_prs(_check_path(str(deck)))
    findings = _detect_table_content_clipping(loaded.slides[0])
    assert not any(f.get("criterion") == "TABLE_CONTENT_CLIPPING" for f in findings), (
        f"Expected no TABLE_CONTENT_CLIPPING finding for fitting table, got: {findings}"
    )


@pytest.mark.unit
def test_bf03_review_slide_export_no_regression_clean_slide(tmp_path, monkeypatch):
    """review_slide_export returns passed:True for a clean slide (no regression)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    from pptmcp.review_pptx import review_slide_export

    prs_obj = _Prs()
    prs_obj.slides.add_slide(prs_obj.slide_layouts[5])
    deck = tmp_path / "clean.pptx"
    prs_obj.save(str(deck))
    png = _make_png(tmp_path)

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = review_slide_export(str(deck), 0, png, confirm=True)
    assert result["passed"] is True
    assert "TABLE_CONTENT_CLIPPING" in result["checks_run"]
    assert result["findings"] == []


@pytest.mark.unit
def test_bf03_e2e_review_slide_export_with_clipped_table(tmp_path, monkeypatch):
    """review_slide_export returns passed:False with TABLE_CONTENT_CLIPPING FAIL for a clipped table (E2E)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    from pptmcp.review_pptx import review_slide_export

    prs_obj = _Prs()
    slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[5])

    from pptx.util import Emu
    from pptx.oxml.ns import qn as _qn
    rows, cols = 8, 2
    # Shape height = 4 in; rows auto-distributed at 0.5 in each → total == shape height
    table_shape = slide.shapes.add_table(
        rows, cols, Emu(_emu(1)), Emu(_emu(1)), Emu(_emu(4)), Emu(_emu(4))
    )
    # Shrink the shape ext.cy to 1 in — rows still sum to 4 in → definite clipping
    xfrm = table_shape.element.find(".//" + _qn("a:ext"))
    xfrm.set("cy", str(_emu(1)))

    tmp_pptx = tmp_path / "clipped_e2e.pptx"
    prs_obj.save(str(tmp_pptx))
    png = _make_png(tmp_path)

    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = review_slide_export(str(tmp_pptx), 0, png, confirm=True)
    assert result["passed"] is False
    assert len(result["findings"]) >= 1
    assert any(
        f["criterion"] == "TABLE_CONTENT_CLIPPING"
        and f["result"] == "FAIL"
        and f["severity"] == "high"
        for f in result["findings"]
    ), f"Expected TABLE_CONTENT_CLIPPING FAIL/high finding, got: {result['findings']}"


# ---------------------------------------------------------------------------
# BUG: DisplayAlerts COM property assignment fails in certain session modes
# ---------------------------------------------------------------------------

class TestDisplayAlertsGuard:
    """Verify that a COM error on app.DisplayAlerts = False is caught, not fatal."""

    def test_ppt_app_context_survives_display_alerts_failure(self, monkeypatch):
        """_ppt_app_context should yield even if DisplayAlerts assignment raises."""
        import importlib
        import sys
        from unittest.mock import patch

        monkeypatch.setenv("PPT_ENABLE_COM", "true")

        class _FakeApp:
            """Mock PowerPoint.Application that rejects DisplayAlerts."""

            Visible = False

            @property
            def DisplayAlerts(self):
                return True

            @DisplayAlerts.setter
            def DisplayAlerts(self, value):
                raise OSError("COM session mode does not support DisplayAlerts")

            def Quit(self):
                pass

        fake_app = _FakeApp()

        # Build fake modules so local imports inside _ppt_app_context resolve
        fake_pythoncom = type("module", (), {
            "CoInitializeEx": staticmethod(lambda x: None),
            "CoUninitialize": staticmethod(lambda: None),
            "COINIT_APARTMENTTHREADED": 0,
            "com_error": OSError,
        })()
        fake_client = type("module", (), {
            "Dispatch": staticmethod(lambda name: fake_app),
        })()
        fake_win32com = type("module", (), {"client": fake_client})()

        # Use monkeypatch for sys.modules so cleanup is automatic even on failure
        mod_keys = ("pythoncom", "win32com", "win32com.client")
        saved = {k: sys.modules.get(k) for k in mod_keys}
        for k, v in saved.items():
            if v is not None:
                monkeypatch.delitem(sys.modules, k)
        monkeypatch.setitem(sys.modules, "pythoncom", fake_pythoncom)
        monkeypatch.setitem(sys.modules, "win32com", fake_win32com)
        monkeypatch.setitem(sys.modules, "win32com.client", fake_client)

        # Force re-import so the module picks up fake modules
        import pptmcp._pptx_com_guards as guards_mod
        importlib.reload(guards_mod)

        with patch.object(guards_mod, "_check_com_enabled"):
            with guards_mod._ppt_app_context() as app:
                assert app is fake_app


# ---------------------------------------------------------------------------
# BUG: copy_slide drops slide-level background fills (<p:bg>)
# ---------------------------------------------------------------------------

class TestCopySlidePreservesBackground:
    """Regression: copy_slide must preserve <p:bg> element from source slide."""

    @staticmethod
    def _make_prs_with_bg(tmp_path, filename="bg_test.pptx"):
        """Create a 1-slide pptx with a solid red background via XML."""
        prs = _Prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                 'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        cSld = slide._element.find('.//p:cSld', nsmap)
        spTree = cSld.find('p:spTree', nsmap)
        # Build <p:bg> with solid red fill
        bg_xml = (
            '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<p:bgPr><a:solidFill><a:srgbClr val="FF0000"/></a:solidFill>'
            '<a:effectLst/></p:bgPr></p:bg>'
        )
        bg_elem = etree.fromstring(bg_xml)
        cSld.insert(list(cSld).index(spTree), bg_elem)
        path = tmp_path / filename
        prs.save(str(path))
        return str(path)

    def test_in_presentation_copy_preserves_bg(self, tmp_path, monkeypatch):
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        _prs_cache.clear()

        path = self._make_prs_with_bg(tmp_path)
        result = _copy_slide(
            source_path=path, source_slide_index=0, confirm=True,
        )
        assert result["new_slide_index"] >= 0

        # Verify the copied slide has <p:bg>
        from pptmcp.presentation_pptx import _load_prs
        prs = _load_prs(tmp_path / "bg_test.pptx")
        new_slide = prs.slides[result["new_slide_index"]]
        nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        bg = new_slide._element.find('.//p:cSld/p:bg', nsmap)
        assert bg is not None, "copy_slide must preserve <p:bg> element"
        # Verify it contains the red color
        a_ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        clr = bg.find('.//{%s}srgbClr' % a_ns['a'])
        assert clr is not None and clr.get('val') == 'FF0000'

    def test_cross_file_copy_preserves_bg(self, tmp_path, monkeypatch):
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        _prs_cache.clear()

        src_path = self._make_prs_with_bg(tmp_path, "source.pptx")
        # Create empty target
        tgt_prs = _Prs()
        tgt_prs.slides.add_slide(tgt_prs.slide_layouts[5])
        tgt_path = str(tmp_path / "target.pptx")
        tgt_prs.save(tgt_path)

        result = _copy_slide(
            source_path=src_path, source_slide_index=0,
            target_path=tgt_path, confirm=True,
        )

        from pptmcp.presentation_pptx import _load_prs
        prs = _load_prs(tmp_path / "target.pptx")
        new_slide = prs.slides[result["new_slide_index"]]
        nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        bg = new_slide._element.find('.//p:cSld/p:bg', nsmap)
        assert bg is not None, "cross-file copy_slide must preserve <p:bg> element"
