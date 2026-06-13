"""Unit tests: content creation — insert_image, create_presentation, list_layouts,
add_textbox, add_shape, add_table_to_slide, set_text_format, set_paragraph_format,
copy_slide, and phase-1.5 gap fills (groups 15–16, 21–29)."""
import pytest
from pptx import Presentation as _Prs

from pptmcp.presentation_pptx import (
    ValidationError,
    NotAllowedError,
    add_shape,
    add_table_to_slide,
    add_textbox,
    create_presentation,
    insert_image,
    list_layouts,
)


# ── GROUP 15: insert_image() ──────────────────────────────────────────────

@pytest.mark.unit
def test_insert_image_blocked_without_write(tmp_path, monkeypatch):
    """insert_image must raise NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    pptx_path = tmp_path / "deck.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(pptx_path))
    # image_path validation can be any image extension path — will fail write gate first
    img_path = tmp_path / "img.png"
    img_path.write_bytes(b"fake")
    with pytest.raises(NotAllowedError):
        from pptmcp.presentation_pptx import insert_image
        insert_image(str(pptx_path), 0, str(img_path))


@pytest.mark.unit
def test_insert_image_rejects_non_image_extension(tmp_path, monkeypatch):
    """insert_image must reject files with unsupported image extensions."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx_path = tmp_path / "deck.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(pptx_path))
    bad_img = tmp_path / "img.exe"
    bad_img.write_bytes(b"fake")
    with pytest.raises(ValidationError, match="Unsupported image extension"):
        from pptmcp.presentation_pptx import insert_image
        insert_image(str(pptx_path), 0, str(bad_img), confirm=True)


# ── GROUP 16: insert_image confirm gate ───────────────────────────────────

@pytest.mark.unit
def test_insert_image_requires_confirm(tmp_path, monkeypatch):
    """insert_image must raise ValidationError when confirm=False even if all other checks pass."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx_path = tmp_path / "deck.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(pptx_path))
    img_path = tmp_path / "img.png"
    img_path.write_bytes(b"fake-png")
    with pytest.raises(ValidationError, match="confirm"):
        insert_image(str(pptx_path), 0, str(img_path), 0, 0, 100, 100, confirm=False)


@pytest.mark.unit
def test_insert_image_blocked_path(tmp_path, monkeypatch):
    """insert_image must raise ValidationError when image_path is outside the allowlist."""
    allowed_dir = tmp_path / "allowed"
    allowed_dir.mkdir()
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(allowed_dir))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx_path = allowed_dir / "deck.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(pptx_path))
    outside_img = tmp_path / "img.png"
    outside_img.write_bytes(b"fake-png")
    with pytest.raises((ValidationError, NotAllowedError)):
        insert_image(str(pptx_path), 0, str(outside_img), 0, 0, 100, 100, confirm=True)


# ── GROUP 21: create_presentation() ──────────────────────────────────────

@pytest.mark.unit
def test_create_presentation_success(tmp_path, monkeypatch):
    """Create a new presentation with a title slide; returns correct dict."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "new_deck.pptx")
    result = create_presentation(path, title="Test", confirm=True)
    assert "path" in result
    assert result["slide_count"] >= 1
    assert result["title"] == "Test"


@pytest.mark.unit
def test_create_presentation_blocked_write(tmp_path, monkeypatch):
    """create_presentation raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "should_not_create.pptx")
    with pytest.raises(NotAllowedError):
        create_presentation(path, title="X", confirm=True)


@pytest.mark.unit
def test_create_presentation_no_confirm(tmp_path, monkeypatch):
    """create_presentation raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "no_confirm.pptx")
    with pytest.raises(ValidationError):
        create_presentation(path, confirm=False)


@pytest.mark.unit
def test_create_presentation_path_not_in_allowlist(tmp_path, monkeypatch):
    """create_presentation raises ValidationError for path outside allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        create_presentation(path, confirm=True)


@pytest.mark.unit
def test_create_presentation_file_already_exists(tmp_path, monkeypatch):
    """create_presentation raises ValidationError if file already exists."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    existing = tmp_path / "existing.pptx"
    existing.write_bytes(b"dummy content")
    with pytest.raises(ValidationError, match="already exists"):
        create_presentation(str(existing), confirm=True)


# ── GROUP 22: list_layouts() ──────────────────────────────────────────────

@pytest.mark.unit
def test_list_layouts_returns_list(tmp_path, monkeypatch):
    """list_layouts returns a non-empty list for a default Presentation."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "layouts_test.pptx"
    p = _Prs()
    p.save(str(path))
    result = list_layouts(str(path))
    assert isinstance(result, list)
    assert len(result) > 0


@pytest.mark.unit
def test_list_layouts_structure(tmp_path, monkeypatch):
    """Each layout dict has index, name, placeholder_count, placeholder_types keys."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "layouts_struct.pptx"
    p = _Prs()
    p.save(str(path))
    result = list_layouts(str(path))
    for item in result:
        assert "index" in item
        assert "name" in item
        assert "placeholder_count" in item
        assert "placeholder_types" in item


@pytest.mark.unit
def test_list_layouts_path_not_allowed(tmp_path, monkeypatch):
    """list_layouts raises ValidationError for a path outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        list_layouts(path)


# ── GROUP 23: add_textbox() ───────────────────────────────────────────────

@pytest.mark.unit
def test_add_textbox_success(tmp_path, monkeypatch):
    """add_textbox returns dict with shape_id."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "textbox.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = add_textbox(str(path), 0, 1.0, 1.0, 4.0, 1.0, text="Hello", confirm=True)
    assert "shape_id" in result
    assert result["slide_index"] == 0
    assert result["width_inches"] == 4.0
    assert result["height_inches"] == 1.0


@pytest.mark.unit
def test_add_textbox_blocked_write(tmp_path, monkeypatch):
    """add_textbox raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(NotAllowedError):
        add_textbox(path, 0, 1.0, 1.0, 4.0, 1.0, confirm=True)


@pytest.mark.unit
def test_add_textbox_no_confirm(tmp_path, monkeypatch):
    """add_textbox raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        add_textbox(path, 0, 1.0, 1.0, 4.0, 1.0, confirm=False)


@pytest.mark.unit
def test_add_textbox_invalid_slide_index(tmp_path, monkeypatch):
    """add_textbox raises ValidationError for an out-of-range slide index."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "one_slide.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError):
        add_textbox(str(path), 99, 1.0, 1.0, 4.0, 1.0, confirm=True)


@pytest.mark.unit
def test_add_textbox_path_not_allowed(tmp_path, monkeypatch):
    """add_textbox raises ValidationError for a path outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        add_textbox(path, 0, 1.0, 1.0, 4.0, 1.0, confirm=True)


# ── GROUP 24: add_shape() ─────────────────────────────────────────────────

@pytest.mark.unit
def test_add_shape_rectangle_success(tmp_path, monkeypatch):
    """add_shape with RECTANGLE returns dict with shape_id."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "shape.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = add_shape(str(path), 0, "RECTANGLE", 1.0, 1.0, 3.0, 2.0, confirm=True)
    assert "shape_id" in result
    assert result["slide_index"] == 0
    assert result["shape_type"] == "RECTANGLE"


@pytest.mark.unit
def test_add_shape_blocked_write(tmp_path, monkeypatch):
    """add_shape raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(NotAllowedError):
        add_shape(path, 0, "RECTANGLE", 1.0, 1.0, 3.0, 2.0, confirm=True)


@pytest.mark.unit
def test_add_shape_no_confirm(tmp_path, monkeypatch):
    """add_shape raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        add_shape(path, 0, "RECTANGLE", 1.0, 1.0, 3.0, 2.0, confirm=False)


@pytest.mark.unit
def test_add_shape_invalid_type(tmp_path, monkeypatch):
    """add_shape raises ValidationError for an unknown shape_type string."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "shape2.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="Unknown shape type"):
        add_shape(str(path), 0, "NOT_A_REAL_SHAPE_XYZ", 1.0, 1.0, 3.0, 2.0, confirm=True)


@pytest.mark.unit
def test_add_shape_invalid_slide_index(tmp_path, monkeypatch):
    """add_shape raises ValidationError for an out-of-range slide index."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "shape3.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError):
        add_shape(str(path), 99, "RECTANGLE", 1.0, 1.0, 3.0, 2.0, confirm=True)


@pytest.mark.unit
def test_add_shape_path_not_allowed(tmp_path, monkeypatch):
    """add_shape raises ValidationError for a path outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        add_shape(path, 0, "RECTANGLE", 1.0, 1.0, 2.0, 1.0, confirm=True)


# ── GROUP 25: add_table_to_slide() ────────────────────────────────────────

@pytest.mark.unit
def test_add_table_success(tmp_path, monkeypatch):
    """add_table_to_slide with 2x3 table returns rows=2 cols=3."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "table.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = add_table_to_slide(str(path), 0, 2, 3, 1.0, 1.0, 6.0, 2.0, confirm=True)
    assert result["rows"] == 2
    assert result["cols"] == 3
    assert "shape_id" in result


@pytest.mark.unit
def test_add_table_with_data(tmp_path, monkeypatch):
    """add_table_to_slide with data=[["A","B"],["C","D"]] succeeds and returns shape_id."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "table_data.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = add_table_to_slide(
        str(path), 0, 2, 2, 1.0, 1.0, 4.0, 2.0,
        data=[["A", "B"], ["C", "D"]], confirm=True,
    )
    assert "shape_id" in result
    assert result["rows"] == 2
    assert result["cols"] == 2


@pytest.mark.unit
def test_add_table_data_shape_mismatch(tmp_path, monkeypatch):
    """add_table_to_slide raises ValidationError when data row count does not match rows."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "table_bad.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="rows"):
        add_table_to_slide(
            str(path), 0, 3, 2, 1.0, 1.0, 4.0, 2.0,
            data=[["A", "B"], ["C", "D"]],  # 2 rows, but rows=3
            confirm=True,
        )


@pytest.mark.unit
def test_add_table_blocked_write(tmp_path, monkeypatch):
    """add_table_to_slide raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(NotAllowedError):
        add_table_to_slide(path, 0, 2, 2, 1.0, 1.0, 4.0, 2.0, confirm=True)


@pytest.mark.unit
def test_add_table_no_confirm(tmp_path, monkeypatch):
    """add_table_to_slide raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        add_table_to_slide(path, 0, 2, 2, 1.0, 1.0, 4.0, 2.0, confirm=False)


@pytest.mark.unit
def test_add_table_invalid_slide_index(tmp_path, monkeypatch):
    """add_table_to_slide raises ValidationError for an out-of-range slide index."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "one_slide.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError):
        add_table_to_slide(str(path), 99, 2, 2, 0.0, 0.0, 4.0, 2.0, confirm=True)


