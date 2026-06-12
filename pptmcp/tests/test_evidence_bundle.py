"""
Tests for US-17 (export_changed_slides_only) and US-18 (produce_evidence_bundle).
Phase E acceptance-criteria coverage.
"""
from __future__ import annotations

import contextlib
import sys
from pathlib import Path
from unittest.mock import MagicMock

import pytest
from pptx import Presentation as _Prs

from pptmcp.presentation_pptx import ValidationError, PPTMCPError


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_pptx(tmp_path: Path, slide_count: int = 2) -> Path:
    """Create a real minimal .pptx with blank slides for hashing tests."""
    p = _Prs()
    for _ in range(slide_count):
        p.slides.add_slide(p.slide_layouts[5])
    pptx_path = tmp_path / "test.pptx"
    p.save(str(pptx_path))
    return pptx_path


def _setup_env(monkeypatch, tmp_path: Path) -> None:
    """Configure env vars and mock win32/pythoncom for COM gate tests."""
    monkeypatch.setenv("PPT_ENABLE_COM", "true")
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setitem(sys.modules, "win32com.client", MagicMock())
    mock_pythoncom = MagicMock()
    mock_pythoncom.com_error = Exception
    monkeypatch.setitem(sys.modules, "pythoncom", mock_pythoncom)


# ---------------------------------------------------------------------------
# US-18: produce_evidence_bundle — pure python-pptx, no COM required
# ---------------------------------------------------------------------------

class TestProduceEvidenceBundle:
    """Tests for produce_evidence_bundle (no COM, no write gate)."""

    # ── T01: empty review_results ─────────────────────────────────────────

    @pytest.mark.unit
    def test_empty_review_results(self, tmp_path: Path, monkeypatch) -> None:
        """Empty review_results → overall_passed=False, slides=[], total_findings=0."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        pptx_path = _make_pptx(tmp_path)
        from pptmcp.presentation_pptx import produce_evidence_bundle

        result = produce_evidence_bundle(str(pptx_path), review_results=[])

        assert result["overall_passed"] is False
        assert result["slides"] == []
        assert result["total_findings"] == 0

    # ── T02: single passing slide ─────────────────────────────────────────

    @pytest.mark.unit
    def test_single_passing_slide(self, tmp_path: Path, monkeypatch) -> None:
        """Single passed=True entry → overall_passed=True, total_findings=0."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        pptx_path = _make_pptx(tmp_path)
        from pptmcp.presentation_pptx import produce_evidence_bundle

        review_results = [{"slide_index": 0, "passed": True, "findings_count": 0}]
        result = produce_evidence_bundle(str(pptx_path), review_results=review_results)

        assert result["overall_passed"] is True
        assert result["total_findings"] == 0
        assert len(result["slides"]) == 1
        assert result["slides"][0]["review_passed"] is True

    # ── T03: mixed passed/failed → overall_passed=False ───────────────────

    @pytest.mark.unit
    def test_mixed_results_overall_false(self, tmp_path: Path, monkeypatch) -> None:
        """Mixed passed/failed entries → overall_passed=False."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        pptx_path = _make_pptx(tmp_path)
        from pptmcp.presentation_pptx import produce_evidence_bundle

        review_results = [
            {"slide_index": 0, "passed": True, "findings_count": 0},
            {"slide_index": 1, "passed": False, "findings_count": 2},
        ]
        result = produce_evidence_bundle(str(pptx_path), review_results=review_results)

        assert result["overall_passed"] is False
        assert result["total_findings"] == 2

    # ── T04: bundle_version and generated_at format ───────────────────────

    @pytest.mark.unit
    def test_bundle_version_and_timestamp(self, tmp_path: Path, monkeypatch) -> None:
        """bundle_version=='1.0', generated_at is ISO-8601-like string."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        pptx_path = _make_pptx(tmp_path)
        from pptmcp.presentation_pptx import produce_evidence_bundle

        result = produce_evidence_bundle(str(pptx_path), review_results=[])

        assert result["bundle_version"] == "1.0"
        assert isinstance(result["generated_at"], str)
        # ISO-8601 UTC timestamps contain 'T' and '+' or 'Z'
        assert "T" in result["generated_at"]

    # ── T11: path outside allowlist ───────────────────────────────────────

    @pytest.mark.unit
    def test_path_outside_allowlist(self, tmp_path: Path, monkeypatch) -> None:
        """Path outside PPT_ALLOWLIST_ROOTS → ValidationError."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", "/nonexistent_XZY_root")
        from pptmcp.presentation_pptx import produce_evidence_bundle

        with pytest.raises(PPTMCPError):
            produce_evidence_bundle(str(tmp_path / "test.pptx"), review_results=[])


# ---------------------------------------------------------------------------
# US-17: export_changed_slides_only — COM gate + hash tests
# ---------------------------------------------------------------------------

@pytest.mark.skipif(sys.platform != "win32", reason="COM tests require Windows (win32com)")
@pytest.mark.unit
class TestExportChangedSlidesOnly:
    """Unit tests for export_changed_slides_only — COM gate chain and hash logic."""

    # ── T05: COM disabled ─────────────────────────────────────────────────

    def test_com_disabled(self, tmp_path: Path, monkeypatch) -> None:
        """PPT_ENABLE_COM not set → raises PPTMCPError mentioning PPT_ENABLE_COM."""
        monkeypatch.delenv("PPT_ENABLE_COM", raising=False)
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        from pptmcp.presentation_com import export_changed_slides_only

        with pytest.raises(PPTMCPError, match="PPT_ENABLE_COM"):
            export_changed_slides_only(
                str(tmp_path / "x.pptx"), str(tmp_path), {}, confirm=True
            )

    # ── T06: write gate disabled ──────────────────────────────────────────

    def test_write_disabled(self, tmp_path: Path, monkeypatch) -> None:
        """PPT_ENABLE_WRITE not set → raises PPTMCPError."""
        monkeypatch.setenv("PPT_ENABLE_COM", "true")
        monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        monkeypatch.setitem(sys.modules, "win32com.client", MagicMock())
        from pptmcp.presentation_com import export_changed_slides_only

        with pytest.raises(PPTMCPError):
            export_changed_slides_only(
                str(tmp_path / "x.pptx"), str(tmp_path), {}, confirm=True
            )

    # ── T07: confirm=False ────────────────────────────────────────────────

    def test_confirm_false(self, tmp_path: Path, monkeypatch) -> None:
        """confirm=False → raises ValidationError about confirm=True."""
        _setup_env(monkeypatch, tmp_path)
        from pptmcp.presentation_com import export_changed_slides_only

        with pytest.raises(ValidationError):
            export_changed_slides_only(
                str(tmp_path / "x.pptx"), str(tmp_path), {}, confirm=False
            )

    # ── T08: null byte in base_dir ────────────────────────────────────────

    def test_null_byte_in_base_dir(self, tmp_path: Path, monkeypatch) -> None:
        """Null byte in base_dir → raises ValidationError."""
        _setup_env(monkeypatch, tmp_path)
        pptx_path = _make_pptx(tmp_path)
        from pptmcp.presentation_com import export_changed_slides_only

        with pytest.raises(ValidationError, match="null"):
            export_changed_slides_only(
                str(pptx_path), str(tmp_path) + "/sub\x00dir", {}, confirm=True
            )

    # ── T09: malformed previous_hashes value ─────────────────────────────

    def test_malformed_previous_hashes(self, tmp_path: Path, monkeypatch) -> None:
        """previous_hashes with non-32-char hex value → raises ValidationError."""
        _setup_env(monkeypatch, tmp_path)
        pptx_path = _make_pptx(tmp_path)
        from pptmcp.presentation_com import export_changed_slides_only

        bad_hashes = {"0": "not_a_valid_md5_hash"}
        with pytest.raises(ValidationError, match="32-char MD5"):
            export_changed_slides_only(
                str(pptx_path), str(tmp_path), bad_hashes, confirm=True
            )

    # ── T10: all hashes match → no COM opened ────────────────────────────

    def test_all_hashes_match_no_com(self, tmp_path: Path, monkeypatch) -> None:
        """When previous_hashes match all slides, no COM call is made."""
        _setup_env(monkeypatch, tmp_path)

        # Track whether _ppt_app_context was entered
        com_opened = []

        @contextlib.contextmanager
        def _mock_app_ctx(visible=False):
            com_opened.append(True)
            yield MagicMock()

        monkeypatch.setattr(
            "pptmcp.presentation_com._ppt_app_context", _mock_app_ctx
        )

        pptx_path = _make_pptx(tmp_path, slide_count=2)
        from pptmcp.presentation_com import _compute_slide_hash, export_changed_slides_only

        # Compute actual hashes via python-pptx
        import pptx as _pptx
        prs = _pptx.Presentation(str(pptx_path))
        prev: dict[str, str] = {
            str(i): _compute_slide_hash(slide)
            for i, slide in enumerate(prs.slides)
        }

        result = export_changed_slides_only(
            str(pptx_path), str(tmp_path), prev, confirm=True
        )

        assert result["exported_slides"] == []
        assert set(result["unchanged_slides"]) == {0, 1}
        assert len(com_opened) == 0, "COM context should NOT have been opened"
        assert "run_dir" in result
        assert "new_hashes" in result

    # ── T12: oversized width_px ───────────────────────────────────────────

    def test_oversized_width_px(self, tmp_path: Path, monkeypatch) -> None:
        """width_px=7681 → raises ValidationError (exceeds 0-7680 limit)."""
        _setup_env(monkeypatch, tmp_path)
        pptx_path = _make_pptx(tmp_path)
        from pptmcp.presentation_com import export_changed_slides_only

        with pytest.raises(ValidationError, match="7680"):
            export_changed_slides_only(
                str(pptx_path), str(tmp_path), {}, width_px=7681, confirm=True
            )
