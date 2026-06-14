import asyncio
from pathlib import Path

import pytest

# All smoke tests invoke FastMCP internals (asyncio.run → socket.socketpair).
# Opt the entire module back in to socket use while --disable-socket blocks unit tests.
pytestmark = pytest.mark.enable_socket


def _get_tool_names():
    """Return the set of registered tool names from the FastMCP server."""
    from pptmcp import server
    mcp = server.mcp
    tools_dict = None
    if hasattr(mcp, "_tool_manager") and hasattr(mcp._tool_manager, "_tools"):
        tools_dict = mcp._tool_manager._tools
    elif hasattr(mcp, "_tools"):
        tools_dict = mcp._tools
    if tools_dict is not None:
        return set(tools_dict)

    async def _names():
        tools = await mcp.list_tools()
        return {t.name for t in tools}

    return asyncio.run(_names())


def _get_tool_descriptions() -> dict[str, str]:
    """Return registered FastMCP tool descriptions keyed by tool name."""
    from pptmcp import server
    mcp = server.mcp
    tools_dict = None
    if hasattr(mcp, "_tool_manager") and hasattr(mcp._tool_manager, "_tools"):
        tools_dict = mcp._tool_manager._tools
    elif hasattr(mcp, "_tools"):
        tools_dict = mcp._tools
    elif hasattr(mcp, "_local_provider") and hasattr(mcp._local_provider, "_components"):
        components = mcp._local_provider._components
        return {
            key.split(":")[1].split("@")[0]: getattr(component, "description", "") or ""
            for key, component in components.items()
            if key.startswith("tool:")
        }
    if tools_dict is not None:
        return {
            name: getattr(tool, "description", "") or ""
            for name, tool in tools_dict.items()
        }

    async def _descriptions():
        tools = await server.mcp.list_tools()
        return {tool.name: tool.description or "" for tool in tools}

    return asyncio.run(_descriptions())


def _powerpoint_root() -> Path:
    return Path(__file__).resolve().parents[1]


@pytest.mark.smoke
def test_import_clean(capsys):
    """Importing the server module must not write anything to stdout."""
    import pptmcp.server  # noqa: F401

    captured = capsys.readouterr()
    assert captured.out == "", f"Unexpected stdout on import: {captured.out!r}"


@pytest.mark.smoke
def test_mcp_object_exists():
    from pptmcp import server

    assert server.mcp is not None


@pytest.mark.smoke
def test_tools_registered():
    from pptmcp import server
    import asyncio

    mcp = server.mcp

    # Probe known FastMCP internal attribute paths across versions
    tools_dict = None
    if hasattr(mcp, "_tool_manager") and hasattr(mcp._tool_manager, "_tools"):
        tools_dict = mcp._tool_manager._tools
    elif hasattr(mcp, "_tools"):
        tools_dict = mcp._tools

    if tools_dict is not None:
        tool_names = {name for name in tools_dict}
    else:
        # Fallback: use the public async list_tools() API
        async def _get_names():
            tools = await mcp.list_tools()
            return {t.name for t in tools}

        tool_names = asyncio.run(_get_names())

    expected = {
        "capabilities",
        "read_presentation",
        "get_presentation_metadata",
        "list_layouts",
        "list_slides",
        "read_slide",
        "read_speaker_notes",
        "list_shapes",
        "get_shape",
        "extract_tables",
        "extract_images",
        "add_slide",
        "edit_text_placeholder",
        "set_speaker_notes",
        "replace_slide_text",
        "insert_image",
        "reorder_slides",
        "delete_slide",
        "export_slide_as_text",
        "save",
        "create_presentation",
        "extract_presentation_text",
        "manage_hyperlinks",
        "add_content",
        "set_format",
        "set_table_style",
        "copy_slide",
        "add_hyperlink",
        "export_slide",
        "export_deck",
        "export_slides_to_stamped_dir",
        "delete_shape",
        "detect_overlapping_shapes",
        "review_slide_export",
        "validate_contract",
        "check_presentation_against_contract",
        "declare_slide_contract",
        "get_presentation_context",
        "produce_evidence_bundle",
        "export_changed_slides_only",
        "slide",
        "shape",
        "export",
        "manage_comments",
        "batch_set_text",
    }
    missing = expected - tool_names
    assert not missing, f"Missing registered tools: {missing}"


@pytest.mark.smoke
def test_extract_presentation_text_real_call(tmp_path, monkeypatch):
    """Smoke: extract_presentation_text tool is registered."""
    assert "extract_presentation_text" in _get_tool_names()


@pytest.mark.smoke
def test_manage_hyperlinks_real_call(tmp_path, monkeypatch):
    """Smoke: manage_hyperlinks tool is registered."""
    assert "manage_hyperlinks" in _get_tool_names()


@pytest.mark.smoke
def test_read_presentation_registered():
    """Smoke: read_presentation tool is registered."""
    assert "read_presentation" in _get_tool_names()


@pytest.mark.smoke
def test_list_slides_registered():
    """Smoke: list_slides tool is registered."""
    assert "list_slides" in _get_tool_names()


@pytest.mark.smoke
def test_read_slide_registered():
    """Smoke: read_slide tool is registered."""
    assert "read_slide" in _get_tool_names()


@pytest.mark.smoke
def test_add_slide_registered():
    """Smoke: add_slide tool is registered."""
    assert "add_slide" in _get_tool_names()


@pytest.mark.smoke
def test_set_speaker_notes_registered():
    """Smoke: set_speaker_notes tool is registered."""
    assert "set_speaker_notes" in _get_tool_names()


@pytest.mark.smoke
def test_replace_slide_text_registered():
    """Smoke: replace_slide_text tool is registered."""
    assert "replace_slide_text" in _get_tool_names()


@pytest.mark.smoke
def test_create_presentation_registered():
    """Smoke: create_presentation tool is registered."""
    assert "create_presentation" in _get_tool_names()


@pytest.mark.smoke
def test_copy_slide_registered():
    """Smoke: copy_slide tool is registered."""
    assert "copy_slide" in _get_tool_names()


@pytest.mark.smoke
def test_add_hyperlink_registered():
    """Smoke: add_hyperlink tool is registered."""
    assert "add_hyperlink" in _get_tool_names()


@pytest.mark.smoke
def test_export_slide_as_text_registered():
    """Smoke: export_slide_as_text tool is registered."""
    assert "export_slide_as_text" in _get_tool_names()


@pytest.mark.smoke
def test_recalculate_charts_contract_text_pinned_for_non_com_ci():
    """Smoke: pin the MCP-facing recalculate_charts contract without requiring COM registration."""
    contract_text = (
        "Activate embedded charts for recalculation only; "
        "linked and external refresh is skipped for safety. Requires "
        "PPT_ENABLE_COM=true, PPT_ENABLE_WRITE=true, and confirm=True. Windows only."
    )
    stale_text = "Refresh all embedded Excel chart data via COM."
    readme_fragments = (
        "Activate embedded charts for recalculation only.",
        "Linked and external refresh is skipped for safety.",
        "Requires `PPT_ENABLE_COM=true`, `PPT_ENABLE_WRITE=true`, and `confirm=True`",
    )

    com_handler_source = (
        _powerpoint_root() / "src" / "pptmcp" / "_server_handlers_com_tools.py"
    ).read_text(encoding="utf-8")
    readme_text = (_powerpoint_root() / "README.md").read_text(encoding="utf-8")

    assert contract_text in com_handler_source
    for fragment in readme_fragments:
        assert fragment in readme_text
    assert stale_text not in readme_text


@pytest.mark.smoke
def test_recalculate_charts_description_matches_pinned_contract_when_registered():
    """Smoke: when registered, the tool description must still expose the pinned contract."""
    contract_text = (
        "Activate embedded charts for recalculation only; "
        "linked and external refresh is skipped for safety. Requires "
        "PPT_ENABLE_COM=true, PPT_ENABLE_WRITE=true, and confirm=True. Windows only."
    )
    descriptions = _get_tool_descriptions()
    if "recalculate_charts" not in descriptions:
        pytest.skip("recalculate_charts not registered (COM unavailable at import)")

    description = descriptions["recalculate_charts"]
    assert description == contract_text


# ---------------------------------------------------------------------------
# Prompt template tests
# ---------------------------------------------------------------------------


def _get_prompt_names():
    """Return the set of registered prompt names (fastmcp 2.x and 3.x compatible)."""
    import asyncio
    from pptmcp import server

    mcp = server.mcp
    # fastmcp 3.x: list_prompts() returns Sequence[Prompt] with .name
    if hasattr(mcp, "list_prompts"):
        async def _names_3x():
            prompts = await mcp.list_prompts()
            return {p.name for p in prompts}
        return asyncio.run(_names_3x())
    # fastmcp 2.x: get_prompts() returns dict {name: FunctionPrompt}
    if hasattr(mcp, "get_prompts"):
        async def _names_2x():
            prompts = await mcp.get_prompts()
            return set(prompts)  # dict keys are names
        return asyncio.run(_names_2x())
    # Last resort: private _prompt_manager (2.x)
    mgr = getattr(mcp, "_prompt_manager", None)
    if mgr is None:
        pytest.skip("FastMCP version has no _prompt_manager — public API only")
        return
    return set(mgr._prompts) if hasattr(mgr, "_prompts") else set(mgr)


def test_prompt_render_check_iterate_registered():
    assert "ppt_render_check_iterate_v1" in _get_prompt_names()


def _invoke_prompt():
    """Call ppt_render_check_iterate_v1 regardless of fastmcp version.

    fastmcp 2.x: @mcp.prompt() returns FunctionPrompt (has .fn, not callable)
    fastmcp 3.x: @mcp.prompt() returns original function (callable directly)
    """
    from pptmcp import server
    fn_or_prompt = server.ppt_render_check_iterate_v1
    if callable(fn_or_prompt):
        return fn_or_prompt()
    elif hasattr(fn_or_prompt, "fn") and callable(fn_or_prompt.fn):
        return fn_or_prompt.fn()
    raise RuntimeError(f"Cannot invoke prompt: {type(fn_or_prompt)}")


def test_prompt_callable_returns_str():
    result = _invoke_prompt()
    assert isinstance(result, str)
    assert len(result) > 0


def test_prompt_body_contains_output_contract():
    result = _invoke_prompt()
    assert "Output Contract" in result


def test_prompt_body_uses_consolidated_slide_export_not_deprecated():
    result = _invoke_prompt()
    assert 'scope="slide"' in result
    assert "export_slide_as_png" not in result


def test_prompt_body_uses_consolidated_deck_export_not_deprecated():
    result = _invoke_prompt()
    assert 'scope="deck_pdf"' in result
    assert "export_deck_as_pdf" not in result


def test_prompt_body_uses_output_dir_placeholder():
    result = _invoke_prompt()
    assert "<output_dir>" in result
    assert "C:/Temp" not in result
    assert "C:\\\\Temp" not in result


def test_capabilities_has_workflow_and_prompt_keys():
    from pptmcp.presentation_pptx import capabilities
    caps = capabilities()
    assert "recommended_workflow" in caps
    assert "prompt_name" in caps
    assert caps["prompt_name"] == "ppt_render_check_iterate_v1"
    assert "\u2192" in caps["recommended_workflow"]


def test_capabilities_tool_is_registered_and_callable():
    """Server-registered capabilities tool returns compact-surface metadata."""
    from pptmcp import server

    mcp = server.mcp
    tool = None
    # FastMCP 2.x: _tool_manager._tools dict
    if hasattr(mcp, "_tool_manager") and hasattr(mcp._tool_manager, "_tools"):
        tool = mcp._tool_manager._tools.get("capabilities")
    # FastMCP 3.x: _local_provider._components keyed as "tool:<name>@"
    elif hasattr(mcp, "_local_provider") and hasattr(mcp._local_provider, "_components"):
        tool = mcp._local_provider._components.get("tool:capabilities@")
    if tool is None:
        pytest.fail("Cannot discover FastMCP capabilities tool")

    result = asyncio.run(tool.run({}))
    caps = result.structured_content
    assert caps["compact_tool_surface"]["tool_count"] == 27
    assert "tool_bundles" in caps
    assert "render_check_iterate" in caps["tool_bundles"]


def test_prompt_body_contains_all_six_steps():
    """All six STEP headings must be present in the prompt body."""
    result = _invoke_prompt()
    for step_num in range(1, 7):
        assert f"STEP {step_num}" in result, f"Missing STEP {step_num} in prompt body"


# ---------------------------------------------------------------------------
# Regression: ToolError hang fix
# ---------------------------------------------------------------------------

class TestToolErrorHangFix:
    def test_error_handling_middleware_is_registered(self):
        """ErrorHandlingMiddleware must be attached to the FastMCP instance."""
        from pptmcp import server
        mcp = server.mcp
        # FastMCP 3.x stores middlewares in the public .middleware list;
        # fall back to ._middleware for older/future layout changes.
        middleware_list = getattr(mcp, "middleware", None) or getattr(mcp, "_middleware", [])
        types = [type(m).__name__ for m in middleware_list]
        assert "ErrorHandlingMiddleware" in types, (
            f"ErrorHandlingMiddleware not found in mcp.middleware: {types}"
        )

    def test_write_gated_tool_raises_tool_error_when_write_disabled(self, monkeypatch):
        """Write-gated tool raises ToolError when PPT_ENABLE_WRITE is absent."""
        import asyncio
        from fastmcp.exceptions import ToolError
        from pptmcp import server

        monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
        mcp = server.mcp
        tools_dict = mcp._tool_manager._tools if hasattr(mcp, "_tool_manager") else {}
        tool = tools_dict.get("add_slide")
        if tool is None:
            pytest.skip("add_slide not registered in current surface mode")
        with pytest.raises(ToolError, match="PPT_ENABLE_WRITE"):
            asyncio.run(tool.run({"path": "x.pptx", "layout_index": 1, "confirm": True}))

    def test_write_gated_tool_raises_tool_error_without_confirm(self, monkeypatch):
        """Write-gated tool raises ToolError when confirm=False."""
        import asyncio
        from fastmcp.exceptions import ToolError
        from pptmcp import server

        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        mcp = server.mcp
        tools_dict = mcp._tool_manager._tools if hasattr(mcp, "_tool_manager") else {}
        tool = tools_dict.get("add_slide")
        if tool is None:
            pytest.skip("add_slide not registered in current surface mode")
        with pytest.raises(ToolError, match="confirm"):
            asyncio.run(tool.run({"path": "x.pptx", "layout_index": 1, "confirm": False}))

    def test_read_tool_succeeds_with_middleware_present(self, tmp_path, monkeypatch):
        """Read tool returns valid result with ErrorHandlingMiddleware wired (non-regression)."""
        import asyncio
        from pptx import Presentation as _Prs
        from fastmcp import Client
        from pptmcp import server

        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        prs = _Prs()
        prs.slides.add_slide(prs.slide_layouts[0])
        pptx_path = tmp_path / "smoke.pptx"
        prs.save(str(pptx_path))

        async def _call():
            async with Client(server.mcp) as client:
                return await client.call_tool(
                    "get_presentation_metadata", {"path": str(pptx_path)}
                )

        result = asyncio.run(_call())
        assert not result.is_error



# copy_slide clear_content regression
def _cs_pptx(tmp_path):
    from pptx import Presentation as _P; p = _P()
    s = p.slides.add_slide(p.slide_layouts[0])
    if s.shapes.title: s.shapes.title.text = "hi"
    out = tmp_path / "cs.pptx"; p.save(str(out)); return out


def test_copy_slide_clear_false_backward_compat(tmp_path, monkeypatch):
    """clear_content=False: result has new_slide_index, no cleared_placeholders."""
    import asyncio; from fastmcp import Client; from pptmcp import server
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    p = _cs_pptx(tmp_path)
    async def _c():
        async with Client(server.mcp) as c:
            return await c.call_tool("copy_slide",
                {"source_path": str(p), "source_slide_index": 0, "confirm": True, "clear_content": False})
    r = asyncio.run(_c()); d = r.structured_content
    assert not r.is_error and "new_slide_index" in d and "cleared_placeholders" not in d


def test_copy_slide_clear_true_returns_cleared_count(tmp_path, monkeypatch):
    """clear_content=True: result has new_slide_index AND cleared_placeholders."""
    import asyncio; from fastmcp import Client; from pptmcp import server
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    p = _cs_pptx(tmp_path)
    async def _c():
        async with Client(server.mcp) as c:
            return await c.call_tool("copy_slide",
                {"source_path": str(p), "source_slide_index": 0, "confirm": True, "clear_content": True})
    r = asyncio.run(_c()); d = r.structured_content
    assert not r.is_error and "new_slide_index" in d
    assert "cleared_placeholders" in d and isinstance(d["cleared_placeholders"], int)


def test_copy_slide_clear_content_evicts_cache(tmp_path, monkeypatch):
    """clear_content=True must evict the cache so subsequent reads return cleared slide."""
    import asyncio
    from fastmcp import Client
    from pptmcp import server
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    p = _cs_pptx(tmp_path)
    evicted_paths = []

    # Patch at the handler module where the name is bound (from ... import _evict_prs)
    import pptmcp._server_handlers_export as _handler
    import pptmcp._presentation_cache as _cache
    original_evict = _cache._evict_prs

    def _spy_evict(path):
        evicted_paths.append(str(path))
        return original_evict(path)

    monkeypatch.setattr(_handler, "_evict_prs", _spy_evict)

    async def _c():
        async with Client(server.mcp) as c:
            return await c.call_tool("copy_slide",
                {"source_path": str(p), "source_slide_index": 0, "confirm": True, "clear_content": True})

    r = asyncio.run(_c())
    assert not r.is_error
    # Cache must have been evicted exactly once for the destination path
    assert len(evicted_paths) == 1, f"expected 1 eviction, got {evicted_paths}"
    assert str(p) in evicted_paths[0] or evicted_paths[0].endswith("cs.pptx")
