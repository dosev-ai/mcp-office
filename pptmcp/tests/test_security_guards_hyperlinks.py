"""
Security guard tests for pptmcp — H-013 through R-009 regression suite.
All tests are @pytest.mark.unit (no live server, no COM required).
Split from test_security_guards.py.
"""

from contextlib import contextmanager
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation as _Prs

from pptmcp.presentation_pptx import (
    ValidationError,
    _load_prs,
    _prs_cache,
)
from pptmcp.shapes_pptx import manage_hyperlinks
from pptmcp.presentation_pptx import OfficeCOMError
from pptmcp.review_pptx import (
    _detect_text_overflow_heuristic,
    review_slide_export,
)

# ---------------------------------------------------------------------------

def _make_pptx(path) -> str:
    """Create a 1-slide .pptx with one text shape and return path as str."""
    prs = _Prs()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Test"
    slide.placeholders[1].text = "Body"
    prs.save(str(path))
    return str(path)



# ---------------------------------------------------------------------------
# H-016: TestH001b — 4 additional COM wrapper tests (now covers all 7 wrappers)
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH001bComToolErrorPropagation:
    """H-016: PPTMCPError from remaining 4 COM wrappers propagates as ToolError."""

    def test_export_slide_as_png_ppterror_becomes_tool_error(self, monkeypatch):
        """PPTMCPError from ppt_export_slide propagates as ToolError via export_slide_as_png."""
        from fastmcp.exceptions import ToolError
        from pptmcp.presentation_pptx import PPTMCPError
        from pptmcp import server

        if not hasattr(server, "export_slide_as_png"):
            pytest.skip("export_slide_as_png not defined (COM unavailable at import)")

        monkeypatch.setattr(server, "_COM_AVAILABLE", True)
        with patch.object(server._pcom, "ppt_export_slide", side_effect=PPTMCPError("png com error")):
            with pytest.raises(ToolError, match="png com error"):
                server.export_slide_as_png(
                    file_path="x.pptx",
                    slide_index=0,
                    output_path="out.png",
                    confirm=True,
                )

    def test_export_deck_as_pdf_ppterror_becomes_tool_error(self, monkeypatch):
        """PPTMCPError from ppt_export_deck propagates as ToolError via export_deck_as_pdf."""
        from fastmcp.exceptions import ToolError
        from pptmcp.presentation_pptx import PPTMCPError
        from pptmcp import server

        if not hasattr(server, "export_deck_as_pdf"):
            pytest.skip("export_deck_as_pdf not defined (COM unavailable at import)")

        monkeypatch.setattr(server, "_COM_AVAILABLE", True)
        with patch.object(server._pcom, "ppt_export_deck", side_effect=PPTMCPError("pdf com error")):
            with pytest.raises(ToolError, match="pdf com error"):
                server.export_deck_as_pdf(
                    file_path="x.pptx",
                    output_path="out.pdf",
                    confirm=True,
                )

    def test_run_slide_show_ppterror_becomes_tool_error(self, monkeypatch):
        """PPTMCPError from run_slide_show propagates as ToolError via server.run_slide_show."""
        from fastmcp.exceptions import ToolError
        from pptmcp.presentation_pptx import PPTMCPError
        from pptmcp import server

        if not hasattr(server, "run_slide_show"):
            pytest.skip("run_slide_show not defined (COM unavailable at import)")

        monkeypatch.setattr(server, "_COM_AVAILABLE", True)
        with patch.object(server._pcom, "run_slide_show", side_effect=PPTMCPError("show error")):
            with pytest.raises(ToolError, match="show error"):
                server.run_slide_show(path="x.pptx", confirm=True)

    def test_recalculate_charts_ppterror_becomes_tool_error(self, monkeypatch):
        """PPTMCPError from recalculate_charts propagates as ToolError via server wrapper."""
        from fastmcp.exceptions import ToolError
        from pptmcp.presentation_pptx import PPTMCPError
        from pptmcp import server

        if not hasattr(server, "recalculate_charts"):
            pytest.skip("recalculate_charts not defined (COM unavailable at import)")

        monkeypatch.setattr(server, "_COM_AVAILABLE", True)
        with patch.object(server._pcom, "recalculate_charts", side_effect=PPTMCPError("calc error")):
            with pytest.raises(ToolError, match="calc error"):
                server.recalculate_charts(path="x.pptx", confirm=True)


# ---------------------------------------------------------------------------
# H-017: TestH002 — OfficeCOMError message strips HRESULT hex
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH002HresultStripped:
    """H-002: _ppt_app_context raises OfficeCOMError with message free of raw HRESULT hex."""

    def test_dispatch_com_error_hresult_not_in_message(self):
        """OfficeCOMError from _ppt_app_context must not expose raw HRESULT hex in message."""
        try:
            import pythoncom  # noqa: F401
            import win32com.client  # noqa: F401
        except ImportError:
            pytest.skip("pywin32 not available")

        from pptmcp.presentation_com import _ppt_app_context
        from pptmcp.presentation_pptx import OfficeCOMError

        fake_err = pythoncom.com_error(-2147221005, "OLE error 0x800401F3", None, None)
        with patch("win32com.client.Dispatch", side_effect=fake_err), \
             patch("pythoncom.CoInitializeEx"), \
             patch("pythoncom.CoUninitialize"):
            with pytest.raises(OfficeCOMError) as exc_info:
                with _ppt_app_context():
                    pass  # pragma: no cover — Dispatch raises before yield

        assert "0x" not in str(exc_info.value)


# ---------------------------------------------------------------------------
# H-017: TestH003 — com_error from Presentations.Open wrapped as OfficeCOMError
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH003StampedExportOpenWrapped:
    """H-003: com_error from Presentations.Open in export_slides_to_stamped_dir → OfficeCOMError."""

    def test_open_com_error_becomes_office_com_error(self):
        """When Presentations.Open raises com_error, OfficeCOMError is propagated."""
        try:
            import pythoncom  # noqa: F401
        except ImportError:
            pytest.skip("pywin32 not available")

        from pptmcp.presentation_com import export_slides_to_stamped_dir
        from pptmcp.presentation_pptx import OfficeCOMError

        fake_err = pythoncom.com_error(-2147352567, "Presentations.Open failed", None, None)
        mock_app = MagicMock()
        mock_app.Presentations.Open.side_effect = fake_err

        @contextmanager
        def _mock_app_ctx(*args, **kwargs):
            yield mock_app

        with patch("pptmcp._pptx_com_export._ppt_app_context", _mock_app_ctx), \
             patch("pptmcp._pptx_com_export._check_com_enabled"), \
             patch("pptmcp._pptx_com_export._check_write"), \
             patch("pptmcp._pptx_com_export._check_confirm"), \
             patch("pptmcp._pptx_com_export._check_path", return_value=Path("/fake/test.pptx")), \
             patch("pptmcp._pptx_com_export._check_in_allowlist"):
            with pytest.raises(OfficeCOMError):
                export_slides_to_stamped_dir(
                    path="/fake/test.pptx",
                    base_dir="/fake/base",
                    confirm=True,
                )


# ---------------------------------------------------------------------------
# H-017: TestH011 — _load_prs error message does not contain file path
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH011LoadPrsPathSanitized:
    """H-011: _load_prs OfficeCOMError message must not contain the file path."""

    def test_corrupt_pptx_error_strips_path(self, monkeypatch):
        """When Presentation() raises, the OfficeCOMError message omits the file path."""
        from pptmcp.presentation_pptx import OfficeCOMError

        test_path = Path("/some/secret/path.pptx")
        key = str(test_path)
        _prs_cache.pop(key, None)  # ensure not cached

        with patch("pptmcp.presentation_pptx.Presentation",
                   side_effect=Exception("fake/path/file.pptx: bad zip")):
            with pytest.raises(OfficeCOMError) as exc_info:
                _load_prs(test_path)

        assert str(test_path) not in str(exc_info.value)
        _prs_cache.pop(key, None)  # cleanup


# ---------------------------------------------------------------------------
# H-013: strip before urlparse prevents false-positive scheme rejection
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH013StripBeforeUrlparse:
    """H-013: target_url.strip() before urlparse prevents false-positive scheme rejection.

    Before the fix, urlparse(" https://example.com").scheme == '' (empty), which
    would cause the scheme guard to reject a perfectly valid HTTPS URL that happens
    to have leading whitespace.  The fix strips first, then validates.
    """

    @pytest.fixture
    def pptx_for_h013(self, tmp_path, monkeypatch):
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        return _make_pptx(tmp_path / "h013.pptx")

    def test_leading_space_https_not_falsely_rejected(self, pptx_for_h013):
        """URL with leading space (scheme appears empty to urlparse before strip) is accepted."""
        # Without H-013: urlparse(' https://example.com').scheme == '' → rejected
        # With H-013:    strip() → 'https://example.com' → scheme='https' → accepted
        result = manage_hyperlinks(
            pptx_for_h013, 0, "add",
            shape_id=2,
            target_url=" https://example.com/leading",
            confirm=True,
        )
        assert result["ok"] is True

    def test_trailing_space_https_not_falsely_rejected(self, pptx_for_h013):
        """URL with trailing space only is accepted after strip."""
        result = manage_hyperlinks(
            pptx_for_h013, 0, "add",
            shape_id=2,
            target_url="https://example.com/trailing ",
            confirm=True,
        )
        assert result["ok"] is True

    def test_strip_result_has_correct_scheme_via_urlparse(self):
        """After strip(), urlparse correctly extracts the scheme (defence-in-depth)."""
        from urllib.parse import urlparse
        # Stripping always produces a URL that urlparse handles deterministically:
        assert urlparse("  https://example.com  ".strip()).scheme == "https"
        assert urlparse("  http://example.com  ".strip()).scheme == "http"


# ---------------------------------------------------------------------------
# H-014: _detect_text_overflow_heuristic skips shapes with None width/height
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH014TextOverflowNoneGuard:
    """H-014: _detect_text_overflow_heuristic must not raise TypeError on None dimensions."""

    def _make_text_shape(self, shape_id, width, height, text="A" * 500):
        """Mock shape that has_text_frame=True with given dimensions."""
        tf = MagicMock()
        tf.text = text
        s = MagicMock()
        s.shape_id = shape_id
        s.name = f"shape_{shape_id}"
        s.has_text_frame = True
        s.text_frame = tf
        s.width = width
        s.height = height
        return s

    def _make_slide(self, shapes):
        slide = MagicMock()
        slide.shapes = shapes
        return slide

    def test_none_width_shape_skipped_no_finding(self):
        """Shape with text but width=None is skipped by the H-014 guard — no TypeError."""
        slide = self._make_slide([self._make_text_shape(1, None, 914400)])
        prs_mock = MagicMock()  # prs_obj not used by the function
        result = _detect_text_overflow_heuristic(slide, prs_mock)
        assert result == []

    def test_none_height_shape_skipped_no_finding(self):
        """Shape with text but height=None is skipped by the H-014 guard — no TypeError."""
        slide = self._make_slide([self._make_text_shape(2, 914400, None)])
        prs_mock = MagicMock()
        result = _detect_text_overflow_heuristic(slide, prs_mock)
        assert result == []

    def test_both_none_shape_skipped_no_finding(self):
        """Shape with both width=None and height=None is skipped — no TypeError."""
        slide = self._make_slide([self._make_text_shape(3, None, None)])
        prs_mock = MagicMock()
        result = _detect_text_overflow_heuristic(slide, prs_mock)
        assert result == []

    def test_none_width_skipped_valid_high_density_still_detected(self):
        """None-width shape is skipped; adjacent high-density shape is still reported."""
        # height=0 forces density fallback: area=max(0.5*0,0.01)=0.01, density=50000 -> high
        half_inch = 457200  # EMU
        valid_shape = self._make_text_shape(4, half_inch, 0, text="B" * 500)
        none_shape = self._make_text_shape(5, None, half_inch, text="C" * 500)
        slide = self._make_slide([none_shape, valid_shape])
        prs_mock = MagicMock()
        result = _detect_text_overflow_heuristic(slide, prs_mock)
        # Only the valid shape should appear
        assert len(result) == 1
        assert result[0]["criterion"] == "TEXT_OVERFLOW"
        assert result[0]["severity"] == "high"

    def test_no_text_frame_shape_still_skipped(self):
        """Shape without text frame is unaffected by H-014 (existing behaviour preserved)."""
        s = MagicMock()
        s.has_text_frame = False
        s.width = None
        s.height = None
        slide = self._make_slide([s])
        prs_mock = MagicMock()
        result = _detect_text_overflow_heuristic(slide, prs_mock)
        assert result == []


# ---------------------------------------------------------------------------
# H-015: review_slide_export uses _load_prs (corrupt pptx → OfficeCOMError)
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestH015ReviewUsesLoadPrs:
    """H-015: review_slide_export uses _load_prs instead of bare Presentation().

    The key consequence: a corrupt .pptx raises OfficeCOMError (from _load_prs)
    rather than a raw pptx PackageNotFoundError that could leak the file path.
    """

    def test_corrupt_pptx_raises_office_com_error(self, tmp_path, monkeypatch):
        """review_slide_export on a corrupt .pptx raises OfficeCOMError (via _load_prs)."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        _prs_cache.clear()

        # Write a file with .pptx extension but invalid zip content
        bad_pptx = tmp_path / "corrupt.pptx"
        bad_pptx.write_bytes(b"NOT A ZIP FILE")

        # Create a minimal dummy PNG so the png_path check passes
        dummy_png = tmp_path / "export.png"
        dummy_png.write_bytes(b"\x89PNG\r\n\x1a\n")

        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        with pytest.raises(OfficeCOMError):
            review_slide_export(str(bad_pptx), 0, str(dummy_png), confirm=True)

        _prs_cache.clear()

    def test_corrupt_pptx_error_message_omits_path(self, tmp_path, monkeypatch):
        """The OfficeCOMError message from _load_prs must not contain the file path."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        _prs_cache.clear()

        bad_pptx = tmp_path / "secret_path.pptx"
        bad_pptx.write_bytes(b"NOT A ZIP")

        dummy_png = tmp_path / "export.png"
        dummy_png.write_bytes(b"\x89PNG\r\n\x1a\n")

        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        with pytest.raises(OfficeCOMError) as exc_info:
            review_slide_export(str(bad_pptx), 0, str(dummy_png), confirm=True)

        assert str(bad_pptx) not in str(exc_info.value)

        _prs_cache.clear()


# ---------------------------------------------------------------------------
# R-001 / R-002: produce_evidence_bundle — findings_count + non-dict guards
# ---------------------------------------------------------------------------

@pytest.mark.unit
class TestR001R002EvidenceBundleValidation:
    """Regression tests for R-001 (findings_count coercion) and R-002 (non-dict entry)."""

    def _make_pptx(self, tmp_path: Path) -> str:
        from pptx import Presentation as _Prs
        prs = _Prs()
        prs.slides.add_slide(prs.slide_layouts[5])
        p = tmp_path / "eb.pptx"
        prs.save(str(p))
        return str(p)

    def test_r001_findings_count_none_treated_as_zero(self, tmp_path, monkeypatch):
        """R-001: findings_count=None must be treated as 0 (no TypeError/ValueError)."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        from pptmcp.presentation_pptx import produce_evidence_bundle

        pptx = self._make_pptx(tmp_path)
        review_results = [{"slide_index": 0, "passed": True, "findings_count": None}]
        result = produce_evidence_bundle(pptx, review_results=review_results)

        assert result["slides"][0]["findings_count"] == 0
        assert result["total_findings"] == 0

    def test_r001_findings_count_string_int_raises_validation_error(self, tmp_path, monkeypatch):
        """R-001: findings_count with non-numeric string raises ValidationError."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        from pptmcp.presentation_pptx import produce_evidence_bundle

        pptx = self._make_pptx(tmp_path)
        review_results = [{"slide_index": 0, "passed": True, "findings_count": "not_a_number"}]
        with pytest.raises(ValidationError, match="findings_count"):
            produce_evidence_bundle(pptx, review_results=review_results)

    def test_r002_non_dict_entry_raises_validation_error(self, tmp_path, monkeypatch):
        """R-002: non-dict entry in review_results raises ValidationError with index info."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        from pptmcp.presentation_pptx import produce_evidence_bundle

        pptx = self._make_pptx(tmp_path)
        review_results = [
            {"slide_index": 0, "passed": True, "findings_count": 0},
            "not_a_dict",  # R-002: string instead of dict
        ]
        with pytest.raises(ValidationError, match="index 1"):
            produce_evidence_bundle(pptx, review_results=review_results)

    def test_r002_none_entry_raises_validation_error(self, tmp_path, monkeypatch):
        """R-002: None entry in review_results raises ValidationError."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        from pptmcp.presentation_pptx import produce_evidence_bundle

        pptx = self._make_pptx(tmp_path)
        with pytest.raises(ValidationError, match="index 0"):
            produce_evidence_bundle(pptx, review_results=[None])

    def test_r008_nonexistent_exports_dir_raises_validation_error(self, tmp_path, monkeypatch):
        """R-008: exports_dir not existing on disk raises ValidationError."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        from pptmcp.presentation_pptx import produce_evidence_bundle

        pptx = self._make_pptx(tmp_path)
        nonexistent_dir = str(tmp_path / "ghost_dir")
        with pytest.raises(ValidationError, match="exports_dir does not exist"):
            produce_evidence_bundle(pptx, review_results=[], exports_dir=nonexistent_dir)


# ---------------------------------------------------------------------------
# R-003 / R-007 / R-009: export_changed_slides_only — validation guards
# ---------------------------------------------------------------------------

@pytest.mark.unit
@pytest.mark.skipif(
    __import__("sys").platform != "win32",
    reason="export_changed_slides_only COM tests require Windows",
)
class TestR003R007R009ExportChangedSlidesValidation:
    """Regression tests for R-003 (empty-note), R-007 (file-not-found), R-009 (base_dir)."""

    def _setup_env(self, monkeypatch, tmp_path: Path) -> None:
        import sys
        monkeypatch.setenv("PPT_ENABLE_COM", "true")
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
        monkeypatch.setitem(sys.modules, "win32com.client", MagicMock())
        mock_pc = MagicMock()
        mock_pc.com_error = Exception
        monkeypatch.setitem(sys.modules, "pythoncom", mock_pc)

    def _make_pptx(self, tmp_path: Path, slide_count: int = 1) -> Path:
        from pptx import Presentation as _Prs
        p = _Prs()
        for _ in range(slide_count):
            p.slides.add_slide(p.slide_layouts[5])
        pptx_path = tmp_path / "test_r.pptx"
        p.save(str(pptx_path))
        return pptx_path

    def test_r007_nonexistent_pptx_raises_validation_error(self, tmp_path, monkeypatch):
        """R-007: path that passes allowlist but does not exist raises ValidationError."""
        self._setup_env(monkeypatch, tmp_path)
        from pptmcp.presentation_com import export_changed_slides_only

        missing_path = str(tmp_path / "missing.pptx")  # never created
        with pytest.raises(ValidationError, match="not found"):
            export_changed_slides_only(missing_path, str(tmp_path), {}, confirm=True)

    def test_r009_nonexistent_base_dir_raises_validation_error(self, tmp_path, monkeypatch):
        """R-009: base_dir that does not exist raises ValidationError before COM is opened."""
        self._setup_env(monkeypatch, tmp_path)
        pptx_path = self._make_pptx(tmp_path)
        from pptmcp.presentation_com import export_changed_slides_only

        nonexistent_base = str(tmp_path / "no_such_dir")
        with pytest.raises(ValidationError, match="base_dir does not exist"):
            export_changed_slides_only(str(pptx_path), nonexistent_base, {}, confirm=True)

    def test_r003_empty_changed_list_returns_note(self, tmp_path, monkeypatch):
        """R-003: when no slides changed, result contains note and exported_count=0."""
        import contextlib
        import pptx as _pptx
        self._setup_env(monkeypatch, tmp_path)
        pptx_path = self._make_pptx(tmp_path, slide_count=1)
        from pptmcp.presentation_com import _compute_slide_hash, export_changed_slides_only

        # COM context must never open when nothing changed
        com_opened = []

        @contextlib.contextmanager
        def _mock_app_ctx(visible=False):
            com_opened.append(True)
            yield MagicMock()

        monkeypatch.setattr("pptmcp._pptx_com_export._ppt_app_context", _mock_app_ctx)
        monkeypatch.setattr("pptmcp.presentation_com._ppt_app_context", _mock_app_ctx)

        # Compute hashes to present as "already seen"
        prs = _pptx.Presentation(str(pptx_path))
        prev = {str(i): _compute_slide_hash(s) for i, s in enumerate(prs.slides)}

        result = export_changed_slides_only(
            str(pptx_path), str(tmp_path), prev, confirm=True
        )

        assert result["exported_slides"] == []
        assert result.get("exported_count") == 0
        assert result.get("note") == "no changed slides to export"
        assert len(com_opened) == 0
