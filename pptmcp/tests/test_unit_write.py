"""Unit tests: write guards, confirm guards, write operations, UX/GX fixes, negative cases
(groups 11–14, 17, 20)."""
import pytest
from pptx import Presentation as _Prs

from pptmcp.presentation_pptx import (
    PPTMCPError,
    ValidationError,
    NotAllowedError,
    add_slide,
    delete_slide,
    edit_text_placeholder,
    get_shape,
    list_shapes,
    list_slides,
    read_slide,
    read_speaker_notes,
    replace_slide_text,
    reorder_slides,
    save,
    set_speaker_notes,
    extract_presentation_text,
)
from pptmcp.shapes_pptx import set_shape_fill


# ── GROUP 11: Write guards (no PPT_ENABLE_WRITE) ─────────────────────────

@pytest.mark.unit
def test_add_slide_blocked_without_write(sample_pptx, monkeypatch):
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    with pytest.raises(NotAllowedError):
        add_slide(sample_pptx, confirm=True)


@pytest.mark.unit
def test_save_blocked_without_write(sample_pptx, monkeypatch):
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    with pytest.raises(NotAllowedError):
        save(sample_pptx, confirm=True)


@pytest.mark.unit
def test_delete_slide_blocked_without_write(sample_pptx, monkeypatch):
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    with pytest.raises(NotAllowedError):
        delete_slide(sample_pptx, 0, confirm=True)


@pytest.mark.unit
def test_set_speaker_notes_blocked_without_write(sample_pptx, monkeypatch):
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    with pytest.raises(NotAllowedError):
        set_speaker_notes(sample_pptx, 0, "notes")


# ── GROUP 12: confirm guards ──────────────────────────────────────────────

@pytest.mark.unit
def test_add_slide_blocked_without_confirm(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="confirm"):
        add_slide(sample_pptx, confirm=False)


@pytest.mark.unit
def test_delete_slide_blocked_without_confirm(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="confirm"):
        delete_slide(sample_pptx, 0, confirm=False)


@pytest.mark.unit
def test_save_blocked_without_confirm(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="confirm"):
        save(sample_pptx, confirm=False)


@pytest.mark.unit
def test_replace_slide_text_blocked_without_confirm(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="confirm"):
        replace_slide_text(sample_pptx, "Test", "New", confirm=False)


@pytest.mark.unit
def test_set_speaker_notes_requires_confirm(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="confirm"):
        set_speaker_notes(sample_pptx, 0, "notes", confirm=False)


@pytest.mark.unit
def test_edit_text_placeholder_requires_confirm(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="confirm"):
        edit_text_placeholder(sample_pptx, 0, 0, "new text", confirm=False)


# ── GROUP 13: write operations (PPT_ENABLE_WRITE=true, confirm=True) ────

@pytest.mark.unit
def test_set_speaker_notes_write(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = set_speaker_notes(sample_pptx, 0, "My speaker notes", confirm=True)
    assert result["notes_text"] == "My speaker notes"
    assert result["slide_index"] == 0


@pytest.mark.unit
def test_edit_text_placeholder_write(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = edit_text_placeholder(sample_pptx, 0, 0, "New Title", confirm=True)
    assert "slide_index" in result
    assert result["text"] == "New Title"


@pytest.mark.unit
def test_add_slide_write(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = add_slide(sample_pptx, layout_index=1, confirm=True)
    assert "slide_index" in result
    assert result["slide_index"] == 1  # second slide (0-indexed)
    assert "title" in result


@pytest.mark.unit
def test_add_slide_layout_name_resolves(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx import Presentation as _Prs
    path = tmp_path / "ln.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(path))
    layout_name = p.slide_layouts[0].name
    result = add_slide(str(path), layout_name=layout_name, confirm=True)
    assert "slide_index" in result


@pytest.mark.unit
def test_add_slide_layout_name_not_found(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx import Presentation as _Prs
    path = tmp_path / "ln2.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(path))
    with pytest.raises(ValidationError, match="not found"):
        add_slide(str(path), layout_name="No Such Layout", confirm=True)


@pytest.mark.unit
def test_add_slide_layout_name_takes_precedence(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx import Presentation as _Prs
    path = tmp_path / "ln3.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(path))
    layout_name = p.slide_layouts[2].name
    result = add_slide(str(path), layout_index=99, layout_name=layout_name, confirm=True)
    assert "slide_index" in result


@pytest.mark.unit
def test_replace_slide_text_write(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = replace_slide_text(sample_pptx, "Test", "Modified", confirm=True)
    assert "replacements_made" in result
    assert result["replacements_made"] >= 1


# ── GROUP 14: UX-01 behavioural guarantee ────────────────────────────────

@pytest.mark.unit
def test_read_speaker_notes_includes_empty_slides(sample_pptx):
    """UX-01: every slide is always included, even those with no notes."""
    result = read_speaker_notes(sample_pptx, slide_index=None)
    assert len(result) == 1  # fixture has exactly 1 slide
    notes_text = result[0]["notes_text"]
    assert notes_text is None or notes_text == ""


@pytest.mark.unit
def test_reorder_slides_write(sample_pptx, monkeypatch, tmp_path):
    """Add a second slide so we can reorder."""
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    # Re-create with 2 slides
    path = tmp_path / "two_slides.pptx"
    p = _Prs()
    for _ in range(2):
        p.slides.add_slide(p.slide_layouts[1])
    p.save(str(path))
    result = reorder_slides(str(path), [1, 0], confirm=True)
    assert result["new_order"] == [1, 0]


@pytest.mark.unit
def test_delete_slide_write(sample_pptx, monkeypatch, tmp_path):
    """Add a second slide first, then delete index 1."""
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "two_slides_del.pptx"
    p = _Prs()
    for _ in range(2):
        p.slides.add_slide(p.slide_layouts[1])
    p.save(str(path))
    result = delete_slide(str(path), 1, confirm=True)
    assert result["deleted_index"] == 1


@pytest.mark.unit
def test_save_write(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = save(sample_pptx, confirm=True)
    assert "saved" in result
    assert result["saved"] is True
    assert result["path"].endswith(".pptx")


# ── GROUP 14: get_shape() ─────────────────────────────────────────────────

@pytest.mark.unit
def test_get_shape_returns_detail(sample_pptx):
    shapes = list_shapes(sample_pptx, 0)
    assert shapes, "Expected at least one shape on slide 0"
    shape_id = shapes[0]["shape_id"]
    result = get_shape(sample_pptx, 0, shape_id)
    assert result["shape_id"] == shape_id
    assert "name" in result
    assert "shape_type" in result


@pytest.mark.unit
def test_get_shape_invalid_id(sample_pptx):
    with pytest.raises(ValidationError, match="not found"):
        get_shape(sample_pptx, 0, 99999)


# ── GROUP 17: GX UX-gap fixes ─────────────────────────────────────────────

@pytest.mark.unit
def test_edit_text_placeholder_response_key(sample_pptx, monkeypatch):
    """GX-01: response must use key 'text', not 'new_text'."""
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = edit_text_placeholder(sample_pptx, 0, 0, "GX01 Title", confirm=True)
    assert "text" in result
    assert "new_text" not in result
    assert result["text"] == "GX01 Title"


@pytest.mark.unit
def test_list_slides_no_title_is_null(tmp_path, monkeypatch):
    """GX-02: list_slides must return None (null) for slides with no/empty title."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "notitle.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(path))
    result = list_slides(str(path))
    assert result[0]["title"] is None, f"Expected None, got {result[0]['title']!r}"


@pytest.mark.unit
def test_read_slide_shapes_include_text(tmp_path, monkeypatch):
    """GX-03: read_slide shape dicts must include a 'text' field."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "withtext.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[0])  # layout 0 has title + content
    # Set the title placeholder text
    slide.shapes.title.text = "GX03 Text"
    p.save(str(path))
    result = read_slide(str(path), 0)
    shape_dicts = result["shapes"]
    # All shapes must have a 'text' key
    for s in shape_dicts:
        assert "text" in s, f"Shape {s.get('name')} missing 'text' field"
    # The title placeholder must have a non-empty text value
    title_shapes = [s for s in shape_dicts if s.get("placeholder_idx") == 0]
    assert title_shapes, "No title placeholder (idx=0) found"
    assert title_shapes[0]["text"] == "GX03 Text"


# ── GROUP 20: Negative test cases (file not found, corrupted, empty) ───────

@pytest.mark.unit
def test_list_slides_file_not_found(tmp_path, monkeypatch):
    """Test that list_slides raises PPTMCPError when file does not exist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    missing = tmp_path / "nonexistent.pptx"
    with pytest.raises(PPTMCPError, match="Cannot open presentation"):
        list_slides(str(missing))


@pytest.mark.unit
def test_edit_text_placeholder_invalid_index(sample_pptx, monkeypatch):
    """Test that edit_text_placeholder raises ValidationError for invalid placeholder index."""
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="placeholder_idx .* not found"):
        edit_text_placeholder(sample_pptx, 0, 999, "text", confirm=True)


@pytest.mark.unit
def test_corrupted_pptx_file(tmp_path, monkeypatch):
    """Test that a corrupted/malformed .pptx file raises PPTMCPError."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    # Create a file with .pptx extension but invalid ZIP/OOXML content
    corrupted = tmp_path / "corrupted.pptx"
    corrupted.write_bytes(b"This is not a valid PPTX file; it's just random bytes!")
    with pytest.raises(PPTMCPError, match="Cannot open presentation"):
        list_slides(str(corrupted))


@pytest.mark.unit
def test_empty_presentation(tmp_path, monkeypatch):
    """Test operations on a presentation with 0 slides."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "empty.pptx"
    p = _Prs()
    # Create a presentation but add NO slides
    p.save(str(path))

    # list_slides should return an empty list
    result = list_slides(str(path))
    assert isinstance(result, list)
    assert len(result) == 0

    # read_slide(0) should raise ValidationError because slide 0 doesn't exist
    with pytest.raises(ValidationError, match="out of range"):
        read_slide(str(path), 0)

    # read_speaker_notes(0) should raise ValidationError
    with pytest.raises(ValidationError, match="out of range"):
        read_speaker_notes(str(path), slide_index=0)


# ── GROUP 30: stale-cache regression ──────────────

@pytest.mark.unit
def test_extract_presentation_text_sees_updated_state_after_add_slide(
    tmp_path, monkeypatch
):
    """Regression: extract_presentation_text must reflect mutations made by
    add_slide without requiring an explicit save() call in between.

    Before the fix, add_slide mutated the in-memory cache but did not persist;
    subsequent reads from a reload (e.g. after cache eviction by another op)
    would return the pre-mutation disk state.  After the fix, add_slide
    atomically saves-and-evicts so every read is consistent with disk.
    """
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")

    path = tmp_path / "stale_test.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[0])
    slide.shapes.title.text = "Original"
    p.save(str(path))

    before = extract_presentation_text(str(path))
    assert before["slide_count"] == 1

    # Mutation — should atomically save+evict so disk is updated immediately
    add_slide(str(path), layout_index=0, title="New Slide", confirm=True)

    # extract_presentation_text must see the NEW slide without a separate save()
    after = extract_presentation_text(str(path))
    assert after["slide_count"] == 2, (
        f"Expected 2 slides after add_slide, got {after['slide_count']} "
        "(stale cache bug: cache not evicted after add_slide)"
    )
    # Verify list_slides also agrees — both read the same fresh disk state
    slides = list_slides(str(path))
    assert len(slides) == 2


@pytest.mark.unit
def test_extract_presentation_text_sees_updated_state_after_delete_slide(
    tmp_path, monkeypatch
):
    """Regression: extract_presentation_text must reflect delete_slide immediately."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")

    path = tmp_path / "stale_delete_test.pptx"
    p = _Prs()
    for title in ("Slide A", "Slide B"):
        slide = p.slides.add_slide(p.slide_layouts[0])
        slide.shapes.title.text = title
    p.save(str(path))

    before = extract_presentation_text(str(path))
    assert before["slide_count"] == 2

    delete_slide(str(path), slide_index=1, confirm=True)

    after = extract_presentation_text(str(path))
    assert after["slide_count"] == 1, (
        f"Expected 1 slide after delete_slide, got {after['slide_count']} "
        "(stale cache bug: cache not evicted after delete_slide)"
    )


@pytest.mark.unit
def test_list_slides_and_extract_text_agree_after_reorder(tmp_path, monkeypatch):
    """Regression: list_slides and extract_presentation_text must agree after reorder_slides."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")

    path = tmp_path / "stale_reorder_test.pptx"
    p = _Prs()
    for title in ("First", "Second", "Third"):
        slide = p.slides.add_slide(p.slide_layouts[0])
        slide.shapes.title.text = title
    p.save(str(path))

    reorder_slides(str(path), new_order=[2, 0, 1], confirm=True)

    list_result = list_slides(str(path))
    text_result = extract_presentation_text(str(path))

    assert len(list_result) == 3
    assert text_result["slide_count"] == 3
    # Both views must agree on slide count — previously list_slides could see
    # a different in-memory state than extract_presentation_text
    assert len(list_result) == text_result["slide_count"]


# ── GROUP: shape() target='fill' round-trip ──────────────────────────────────

@pytest.mark.unit
def test_set_shape_fill_round_trip(tmp_path, monkeypatch):
    """set_shape_fill applies a solid fill and reports the normalised hex back.

    Exercises the target='fill' path in shape(operation='set_properties').
    The return dict is the authoritative contract: fill_updated=True and
    fill_color_hex equals the normalised (uppercase, no '#') input hex.
    """
    from pptx.util import Inches

    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")

    path = tmp_path / "fill_test.pptx"
    prs = _Prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(1), Inches(1), Inches(2), Inches(1),
    )
    shape_id = shape.shape_id
    prs.save(str(path))

    result = set_shape_fill(str(path), slide_index=0, shape_id=shape_id,
                            fill_color_hex="FF0000", confirm=True)

    # Contract: the return dict carries the normalised hex and a True sentinel.
    assert result["fill_updated"] is True
    assert result["fill_color_hex"] == "FF0000"
    assert result["shape_id"] == shape_id
    assert result["slide_index"] == 0

    # Verify the in-memory presentation object has the fill applied
    # by saving then reloading fresh (bypassing the LRU cache).
    save(str(path), confirm=True)
    from pptx.dml.color import RGBColor
    reloaded = _Prs(str(path))
    reloaded_shape = next(
        s for s in reloaded.slides[0].shapes if s.shape_id == shape_id
    )
    assert reloaded_shape.fill.fore_color.rgb == RGBColor(0xFF, 0x00, 0x00)
