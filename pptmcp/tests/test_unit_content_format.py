"""Unit tests: content formatting — set_text_format, set_paragraph_format,
copy_slide, and phase-1.5 gap fills (advanced content ops).
Split from test_unit_content.py."""
import pytest
from pptx import Presentation as _Prs
from pptx.util import Inches

from pptmcp.presentation_pptx import (
    PPTMCPError,
    ValidationError,
    NotAllowedError,
    add_shape,
    add_table_to_slide,
    copy_slide,
    create_presentation,
    list_layouts,
    set_paragraph_format,
    set_text_format,
)

# ── GROUP 26: set_text_format() ────────────────────────────────────────────

@pytest.mark.unit
def test_set_text_format_bold_success(tmp_path, monkeypatch):
    """set_text_format with bold=True returns runs_updated >= 1."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "bold.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Bold me"
    p.save(str(path))
    result = set_text_format(str(path), 0, txBox.shape_id, bold=True, confirm=True)
    assert result["runs_updated"] >= 1
    assert result["slide_index"] == 0
    assert result["shape_id"] == txBox.shape_id


@pytest.mark.unit
def test_set_text_format_color_success(tmp_path, monkeypatch):
    """set_text_format with color_hex='FF0000' returns runs_updated >= 1."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "color.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Red text"
    p.save(str(path))
    result = set_text_format(str(path), 0, txBox.shape_id, color_hex="FF0000", confirm=True)
    assert result["runs_updated"] >= 1


@pytest.mark.unit
def test_set_text_format_invalid_color_hex(tmp_path, monkeypatch):
    """set_text_format raises ValidationError for an invalid color_hex."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "badcolor.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Bad"
    p.save(str(path))
    with pytest.raises(ValidationError, match="[Ii]nvalid hex"):
        set_text_format(str(path), 0, txBox.shape_id, color_hex="ZZZZZZ", confirm=True)


@pytest.mark.unit
def test_set_text_format_no_text_frame(tmp_path, monkeypatch):
    """set_text_format raises ValidationError for a shape with no text frame (table shape)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "notf.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    tbl = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    p.save(str(path))
    with pytest.raises(ValidationError, match="no text frame"):
        set_text_format(str(path), 0, tbl.shape_id, bold=True, confirm=True)


@pytest.mark.unit
def test_set_text_format_blocked_write(tmp_path, monkeypatch):
    """set_text_format raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(NotAllowedError):
        set_text_format(path, 0, 1, bold=True, confirm=True)


@pytest.mark.unit
def test_set_text_format_no_confirm(tmp_path, monkeypatch):
    """set_text_format raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        set_text_format(path, 0, 1, bold=True, confirm=False)


# ── GROUP 27: set_paragraph_format() ────────────────────────────────────────

@pytest.mark.unit
def test_set_paragraph_format_alignment_success(tmp_path, monkeypatch):
    """set_paragraph_format with alignment='CENTER' returns correct dict."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "align.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Center me"
    p.save(str(path))
    result = set_paragraph_format(
        str(path), 0, txBox.shape_id, paragraph_index=0, alignment="CENTER", confirm=True
    )
    assert result["paragraph_index"] == 0
    assert result["slide_index"] == 0
    assert result["shape_id"] == txBox.shape_id


@pytest.mark.unit
def test_set_paragraph_format_invalid_alignment(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError for an unknown alignment value."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "badalign.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Bad align"
    p.save(str(path))
    with pytest.raises(ValidationError, match="[Ii]nvalid alignment"):
        set_paragraph_format(str(path), 0, txBox.shape_id, alignment="BOGUS", confirm=True)


@pytest.mark.unit
def test_set_paragraph_format_invalid_para_index(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError for paragraph_index out of range."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "badpara.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "One paragraph"
    p.save(str(path))
    with pytest.raises(ValidationError, match="paragraph_index"):
        set_paragraph_format(str(path), 0, txBox.shape_id, paragraph_index=99, confirm=True)


@pytest.mark.unit
def test_set_paragraph_format_blocked_write(tmp_path, monkeypatch):
    """set_paragraph_format raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(NotAllowedError):
        set_paragraph_format(path, 0, 1, alignment="CENTER", confirm=True)


@pytest.mark.unit
def test_set_paragraph_format_no_confirm(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        set_paragraph_format(path, 0, 1, alignment="CENTER", confirm=False)


# ── GROUP 28: copy_slide() ──────────────────────────────────────────────────

@pytest.mark.unit
def test_copy_slide_same_file_success(tmp_path, monkeypatch):
    """copy_slide within the same file returns new_slide_index >= 1 with source/target paths."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "copy_src.pptx"
    p = _Prs()
    layout = p.slide_layouts[1]
    slide = p.slides.add_slide(layout)
    slide.shapes.title.text = "Original Slide"
    p.save(str(path))
    result = copy_slide(str(path), 0, str(path), confirm=True)
    assert result["new_slide_index"] >= 1
    assert "source_path" in result
    assert "target_path" in result
    assert result["source_slide_index"] == 0


@pytest.mark.unit
def test_copy_slide_invalid_source_slide_index(tmp_path, monkeypatch):
    """copy_slide raises ValidationError when source_slide_index is out of range."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "copy_oob.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="source_slide"):
        copy_slide(str(path), 99, str(path), confirm=True)


@pytest.mark.unit
def test_copy_slide_source_path_not_allowed(tmp_path, monkeypatch):
    """copy_slide raises ValidationError when source_path is outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    evil_src = str(tmp_path.parent / "evil_src.pptx")
    valid_target = str(tmp_path / "target.pptx")
    with pytest.raises(ValidationError, match="[Nn]ot in allowlist|not in allowlist"):
        copy_slide(evil_src, 0, valid_target, confirm=True)


@pytest.mark.unit
def test_copy_slide_target_path_not_allowed(tmp_path, monkeypatch):
    """copy_slide raises ValidationError when target_path is outside the allowlist (OWASP A03)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    valid_src = str(tmp_path / "source.pptx")
    evil_target = str(tmp_path.parent / "evil_target.pptx")
    with pytest.raises(ValidationError, match="[Nn]ot in allowlist|not in allowlist"):
        copy_slide(valid_src, 0, evil_target, confirm=True)


@pytest.mark.unit
def test_copy_slide_blocked_write(tmp_path, monkeypatch):
    """copy_slide raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(NotAllowedError):
        copy_slide(path, 0, path, confirm=True)


@pytest.mark.unit
def test_copy_slide_no_confirm(tmp_path, monkeypatch):
    """copy_slide raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        copy_slide(path, 0, path, confirm=False)


# ── GROUP 29: Phase 1.5 gap-fill tests ────────────────────────────────────

# create_presentation — no-title variant
@pytest.mark.unit
def test_create_presentation_no_title_blank(tmp_path, monkeypatch):
    """create_presentation with title=None creates a blank deck (0 slides)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "blank.pptx")
    result = create_presentation(path, title=None, confirm=True)
    assert result["slide_count"] == 0
    assert result["title"] is None


# list_layouts — file not found
@pytest.mark.unit
def test_list_layouts_file_not_found(tmp_path, monkeypatch):
    """list_layouts raises PPTMCPError when the .pptx file does not exist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    missing = str(tmp_path / "missing.pptx")
    with pytest.raises(PPTMCPError):
        list_layouts(missing)


# add_shape — CIRCLE alias resolves to OVAL
@pytest.mark.unit
def test_add_shape_alias_circle_oval(tmp_path, monkeypatch):
    """add_shape with 'CIRCLE' alias resolves to OVAL without ValueError."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "circle.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = add_shape(str(path), 0, "CIRCLE", 1.0, 1.0, 2.0, 2.0, confirm=True)
    assert "shape_id" in result


# add_table_to_slide — path not in allowlist
@pytest.mark.unit
def test_add_table_path_not_allowed(tmp_path, monkeypatch):
    """add_table_to_slide raises ValidationError for a path outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        add_table_to_slide(path, 0, 2, 2, 1.0, 1.0, 4.0, 2.0, confirm=True)


# add_table_to_slide — data column count mismatch
@pytest.mark.unit
def test_add_table_data_col_mismatch(tmp_path, monkeypatch):
    """add_table_to_slide raises ValidationError when data column count != cols."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "table_col_bad.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="cols"):
        add_table_to_slide(
            str(path), 0, 2, 3, 1.0, 1.0, 4.0, 2.0,
            data=[["A", "B"], ["C", "D"]],  # rows match (2) but 2 cols != cols=3
            confirm=True,
        )


# add_table_to_slide — invalid dimensions (rows=0)
@pytest.mark.unit
def test_add_table_invalid_dimensions_rows_zero(tmp_path, monkeypatch):
    """add_table_to_slide raises ValidationError when rows=0."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="rows"):
        add_table_to_slide(path, 0, 0, 2, 1.0, 1.0, 4.0, 2.0, confirm=True)


# set_text_format — path not in allowlist
@pytest.mark.unit
def test_set_text_format_path_not_allowed(tmp_path, monkeypatch):
    """set_text_format raises ValidationError for a path outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        set_text_format(path, 0, 1, bold=True, confirm=True)


# set_text_format — invalid slide index
@pytest.mark.unit
def test_set_text_format_invalid_slide_index(tmp_path, monkeypatch):
    """set_text_format raises ValidationError for an out-of-range slide index."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "stf_oob.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError):
        set_text_format(str(path), 99, 1, bold=True, confirm=True)


# set_text_format — invalid shape_id
@pytest.mark.unit
def test_set_text_format_invalid_shape_id(tmp_path, monkeypatch):
    """set_text_format raises ValidationError when shape_id does not exist on slide."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "stf_badshape.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="not found"):
        set_text_format(str(path), 0, 99999, bold=True, confirm=True)


# set_paragraph_format — path not in allowlist
@pytest.mark.unit
def test_set_paragraph_format_path_not_allowed(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError for a path outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        set_paragraph_format(path, 0, 1, alignment="CENTER", confirm=True)


# set_paragraph_format — invalid slide index
@pytest.mark.unit
def test_set_paragraph_format_invalid_slide_index(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError for an out-of-range slide index."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "spf_oob.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError):
        set_paragraph_format(str(path), 99, 1, alignment="CENTER", confirm=True)


# set_paragraph_format — invalid shape_id
@pytest.mark.unit
def test_set_paragraph_format_invalid_shape_id(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError when shape_id does not exist on slide."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "spf_badshape.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="not found"):
        set_paragraph_format(str(path), 0, 99999, alignment="CENTER", confirm=True)


# set_paragraph_format — shape with no text frame
@pytest.mark.unit
def test_set_paragraph_format_no_text_frame(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError for a shape with no text frame."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches as _In
    path = tmp_path / "spf_notf.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    tbl = slide.shapes.add_table(2, 2, _In(1), _In(1), _In(4), _In(2))
    p.save(str(path))
    with pytest.raises(ValidationError, match="no text frame"):
        set_paragraph_format(str(path), 0, tbl.shape_id, alignment="CENTER", confirm=True)


# set_paragraph_format — spacing happy path
@pytest.mark.unit
def test_set_paragraph_format_spacing_success(tmp_path, monkeypatch):
    """set_paragraph_format with line_spacing returns correct dict."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches as _In
    path = tmp_path / "spf_spacing.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(_In(1), _In(1), _In(4), _In(1))
    txBox.text_frame.text = "Spaced paragraph"
    p.save(str(path))
    result = set_paragraph_format(
        str(path), 0, txBox.shape_id, paragraph_index=0, line_spacing=1.5, confirm=True
    )
    assert result["paragraph_index"] == 0
    assert result["shape_id"] == txBox.shape_id


# copy_slide — cross-file happy path
@pytest.mark.unit
def test_copy_slide_cross_file_success(tmp_path, monkeypatch):
    """copy_slide from source to a different target appends the slide and reports both paths."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    src = tmp_path / "cs_src.pptx"
    tgt = tmp_path / "cs_tgt.pptx"
    p_src = _Prs()
    slide = p_src.slides.add_slide(p_src.slide_layouts[1])
    slide.shapes.title.text = "Source Slide"
    p_src.save(str(src))
    p_tgt = _Prs()
    p_tgt.slides.add_slide(p_tgt.slide_layouts[5])
    p_tgt.save(str(tgt))
    result = copy_slide(str(src), 0, str(tgt), confirm=True)
    assert result["source_slide_index"] == 0
    assert result["new_slide_index"] >= 1
    assert "cs_src.pptx" in result["source_path"]
    assert "cs_tgt.pptx" in result["target_path"]


# copy_slide — cross-file target file does not exist
@pytest.mark.unit
def test_copy_slide_cross_file_target_not_found(tmp_path, monkeypatch):
    """copy_slide raises ValidationError when the cross-file target does not exist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    src = tmp_path / "cs_src2.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(src))
    nonexistent_tgt = str(tmp_path / "does_not_exist.pptx")
    with pytest.raises(ValidationError, match="does not exist"):
        copy_slide(str(src), 0, nonexistent_tgt, confirm=True)


# copy_slide — invalid target_slide_index
@pytest.mark.unit
def test_copy_slide_invalid_target_slide_index(tmp_path, monkeypatch):
    """copy_slide raises ValidationError when target_slide_index is out of range."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "cs_tgt_oob.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="target_slide"):
        copy_slide(str(path), 0, str(path), target_slide_index=99, confirm=True)


# copy_slide — target_slide_index=0 inserts the copy at position 0
@pytest.mark.unit
def test_copy_slide_target_slide_index_prepend(tmp_path, monkeypatch):
    """copy_slide with target_slide_index=0 positions the copy at slide 0."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "cs_prepend.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    result = copy_slide(str(path), 0, str(path), target_slide_index=0, confirm=True)
    assert result["new_slide_index"] == 0
