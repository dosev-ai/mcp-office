"""Unit tests: capabilities, path guards, and read-only operations (groups 0–10)."""
import pytest
from pptx import Presentation as _Prs

from pptmcp.presentation_pptx import (
    ValidationError,
    capabilities,
    export_slide_as_text,
    extract_images,
    extract_tables,
    get_presentation_metadata,
    list_shapes,
    list_slides,
    read_presentation,
    read_slide,
    read_speaker_notes,
)


# ── GROUP 0: capabilities() ───────────────────────────────────────────────

@pytest.mark.unit
def test_capabilities_phase1():
    assert capabilities()["phase"] == "2.1"


@pytest.mark.unit
def test_capabilities_backend():
    assert capabilities()["backend"] == "python-pptx"


@pytest.mark.unit
def test_capabilities_tool_count():
    assert len(capabilities()["tools"]) == 51  # 51 base tools (COM-conditional excluded)


@pytest.mark.unit
def test_capabilities_all_tools():
    expected = {
        "capabilities", "read_presentation", "get_presentation_metadata",
        "list_slides", "list_layouts", "read_slide", "read_speaker_notes",
        "list_shapes", "get_shape", "extract_tables", "extract_images",
        "create_presentation", "add_slide", "edit_text_placeholder", "set_speaker_notes",
        "replace_slide_text", "insert_image", "reorder_slides",
        "delete_slide", "export_slide_as_text", "save",
        "extract_presentation_text", "manage_hyperlinks",
        "add_content", "set_format", "set_table_style", "copy_slide", "add_hyperlink",
        "detect_overlapping_shapes", "clear_slide_content", "apply_slide_layout",
        "remove_empty_placeholders", "review_slide_export",
        "delete_shape", "declare_slide_contract",
        "validate_contract", "check_presentation_against_contract",
        "produce_evidence_bundle", "export_slide", "export_deck",
        "export_changed_slides_only", "export_slides_to_stamped_dir",
        "get_presentation_context",
        "slide", "shape", "export",
        "manage_comments",
        "batch_set_text",
        "pptmcp_add_column",
        "pptmcp_edit_table_cell",
        "pptmcp_set_column_width",
    }
    assert expected == {e["tool"] for e in capabilities()["tools"]}


# ── GROUP 1: _check_path / allowlist ──────────────────────────────────────

@pytest.mark.unit
def test_path_not_in_allowlist(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    outside = tmp_path.parent / "outside_dir" / "file.pptx"
    with pytest.raises(ValidationError, match="not in allowlist"):
        list_slides(str(outside))


@pytest.mark.unit
def test_no_allowlist_configured(monkeypatch, tmp_path):
    monkeypatch.delenv("PPT_ALLOWLIST_ROOTS", raising=False)
    with pytest.raises(ValidationError, match="No allowlist"):
        list_slides(str(tmp_path / "test.pptx"))


@pytest.mark.unit
def test_unsupported_extension(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    pdf_path = tmp_path / "data.pdf"
    pdf_path.write_bytes(b"fake")
    with pytest.raises(ValidationError, match="Unsupported extension"):
        list_slides(str(pdf_path))


@pytest.mark.unit
def test_path_in_allowlist_passes(tmp_path, monkeypatch):
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "test.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[0])
    p.save(str(path))
    result = list_slides(str(path))
    assert isinstance(result, list)


# ── GROUP 2: list_slides() ────────────────────────────────────────────────

@pytest.mark.unit
def test_list_slides_returns_list(sample_pptx):
    result = list_slides(sample_pptx)
    assert isinstance(result, list)
    assert len(result) >= 1


@pytest.mark.unit
def test_list_slides_first_slide_index(sample_pptx):
    result = list_slides(sample_pptx)
    assert result[0]["slide_index"] == 0


@pytest.mark.unit
def test_list_slides_outside_allowlist(tmp_path, monkeypatch):
    root = tmp_path / "allowed"
    root.mkdir()
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(root))
    outside = tmp_path / "outside" / "test.pptx"
    with pytest.raises(ValidationError, match="not in allowlist"):
        list_slides(str(outside))


# ── GROUP 3: read_presentation() ─────────────────────────────────────────

@pytest.mark.unit
def test_read_presentation_slide_count(sample_pptx):
    result = read_presentation(sample_pptx)
    assert result["slide_count"] >= 1


@pytest.mark.unit
def test_read_presentation_slides_is_list(sample_pptx):
    result = read_presentation(sample_pptx)
    assert isinstance(result["slides"], list)
    assert len(result["slides"]) == result["slide_count"]


# ── GROUP 4: get_presentation_metadata() ─────────────────────────────────

@pytest.mark.unit
def test_get_metadata_slide_count(sample_pptx):
    result = get_presentation_metadata(sample_pptx)
    assert result["slide_count"] >= 1


@pytest.mark.unit
def test_get_metadata_has_keys(sample_pptx):
    result = get_presentation_metadata(sample_pptx)
    for key in ("title", "author", "subject", "slide_count"):
        assert key in result


# ── GROUP 5: read_slide() ─────────────────────────────────────────────────

@pytest.mark.unit
def test_read_slide_shapes_list(sample_pptx):
    result = read_slide(sample_pptx, 0)
    assert "shapes" in result
    assert isinstance(result["shapes"], list)


@pytest.mark.unit
def test_read_slide_out_of_range(sample_pptx):
    with pytest.raises(ValidationError):
        read_slide(sample_pptx, 999)


# ── GROUP 6: read_speaker_notes() ────────────────────────────────────────

@pytest.mark.unit
def test_read_speaker_notes_single_slide(sample_pptx):
    result = read_speaker_notes(sample_pptx, slide_index=0)
    assert isinstance(result, list)


@pytest.mark.unit
def test_read_speaker_notes_all_slides(sample_pptx):
    result = read_speaker_notes(sample_pptx, slide_index=None)
    assert isinstance(result, list)


# ── GROUP 7: list_shapes() ───────────────────────────────────────────────

@pytest.mark.unit
def test_list_shapes_is_list(sample_pptx):
    result = list_shapes(sample_pptx, 0)
    assert isinstance(result, list)


@pytest.mark.unit
def test_list_shapes_shape_keys(sample_pptx):
    result = list_shapes(sample_pptx, 0)
    for shape in result:
        assert "shape_id" in shape
        assert "name" in shape
        assert "shape_type" in shape
        assert "placeholder_idx" in shape


# ── GROUP 8: extract_tables() ────────────────────────────────────────────

@pytest.mark.unit
def test_extract_tables_empty(sample_pptx):
    result = extract_tables(sample_pptx)
    assert isinstance(result, dict)
    assert result["tables"] == []


# ── GROUP 9: extract_images() ────────────────────────────────────────────

@pytest.mark.unit
def test_extract_images_empty(sample_pptx):
    result = extract_images(sample_pptx)
    assert isinstance(result, dict)
    assert result["images"] == []


# ── GROUP 10: export_slide_as_text() ─────────────────────────────────────

@pytest.mark.unit
def test_export_single_slide(sample_pptx):
    result = export_slide_as_text(sample_pptx, slide_index=0)
    assert isinstance(result, list)
    assert len(result) == 1
    assert result[0]["slide_index"] == 0
    assert isinstance(result[0]["texts"], list)


@pytest.mark.unit
def test_export_all_slides(sample_pptx):
    result = export_slide_as_text(sample_pptx, slide_index=None)
    assert isinstance(result, list)
    assert len(result) >= 1
