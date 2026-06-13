"""
Integration tests for pptmcp — real .pptx file I/O, no mocks, no COM.
All tests are @pytest.mark.integration.
"""
import pytest
from pptx import Presentation as RawPresentation
from pptmcp.presentation_pptx import (
    read_presentation,
    read_slide,
    list_slides,
    extract_presentation_text,
    list_shapes,
    read_speaker_notes,
    extract_tables,
    extract_images,
    get_presentation_metadata,
)


@pytest.mark.integration
class TestReadOperations:
    """Real .pptx file reads — no write gate needed."""

    def test_read_presentation_returns_slide_count(self, sample_pptx_file):
        result = read_presentation(sample_pptx_file)
        assert isinstance(result, dict)
        assert result.get("slide_count", 0) >= 1

    def test_list_slides_returns_list(self, sample_pptx_file):
        result = list_slides(sample_pptx_file)
        assert isinstance(result, list)
        assert len(result) >= 1

    def test_read_slide_returns_dict(self, sample_pptx_file):
        result = read_slide(sample_pptx_file, 0)
        assert isinstance(result, dict)
        assert "slide_index" in result or "shapes" in result

    def test_extract_presentation_text_0based(self, sample_pptx_file):
        """NC-1 fix: slide_index in output must be 0-based."""
        result = extract_presentation_text(sample_pptx_file)
        slides = result if isinstance(result, list) else result.get("slides", [])
        if slides:
            first_slide_index = slides[0].get("slide_index")
            assert first_slide_index == 0, f"Expected 0-based slide_index=0, got {first_slide_index}"

    def test_list_shapes_returns_list(self, sample_pptx_file):
        result = list_shapes(sample_pptx_file, 0)
        assert isinstance(result, list)

    def test_read_speaker_notes_no_error(self, sample_pptx_file):
        result = read_speaker_notes(sample_pptx_file, 0)
        assert result is not None

    def test_extract_tables_no_error(self, sample_pptx_file):
        result = extract_tables(sample_pptx_file)
        assert result is not None

    def test_extract_images_no_error(self, sample_pptx_file):
        result = extract_images(sample_pptx_file)
        assert result is not None

    def test_get_presentation_metadata_returns_dict(self, sample_pptx_file):
        result = get_presentation_metadata(sample_pptx_file)
        assert isinstance(result, dict)


@pytest.mark.integration
class TestWriteOperations:
    """Integration write tests — require write_env fixture."""

    def test_add_slide_writes_to_disk(self, write_env, sample_pptx_file):
        from pptmcp.presentation_pptx import add_slide, save
        original_count = len(RawPresentation(sample_pptx_file).slides)
        add_slide(sample_pptx_file, confirm=True)
        save(sample_pptx_file, confirm=True)
        new_count = len(RawPresentation(sample_pptx_file).slides)
        assert new_count == original_count + 1

    def test_set_speaker_notes_persists(self, write_env, sample_pptx_file):
        from pptmcp.presentation_pptx import set_speaker_notes
        set_speaker_notes(sample_pptx_file, 0, "integration test note", confirm=True)
        result = read_speaker_notes(sample_pptx_file, 0)
        assert "integration test note" in str(result)

    def test_replace_slide_text_persists(self, write_env, sample_pptx_file):
        from pptmcp.presentation_pptx import replace_slide_text, save
        replace_slide_text(sample_pptx_file, "Test Slide", "Replaced", confirm=True)
        save(sample_pptx_file, confirm=True)
        text_result = extract_presentation_text(sample_pptx_file)
        slides = text_result if isinstance(text_result, list) else text_result.get("slides", [])
        all_text = " ".join(str(s) for s in slides)
        assert "Replaced" in all_text
