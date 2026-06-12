"""Write-gate and mutation guard tests — targets presentation_pptx.py only."""

import pytest
from pptx import Presentation as _Prs
from pptmcp.presentation_pptx import (
    ValidationError,
    NotAllowedError,
    add_slide,
    edit_text_placeholder,
    set_speaker_notes,
    replace_slide_text,
    insert_image,
    reorder_slides,
    delete_slide,
    save,
    list_slides,
    read_slide,
    read_speaker_notes,
)


# ── GROUP 11: Write guards (no PPT_ENABLE_WRITE) ─────────────────────────

@pytest.mark.unit
def test_add_slide_blocked_without_write(sample_pptx, monkeypatch):
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    with pytest.raises(NotAllowedError):
        add_slide(sample_pptx, confirm=True)


@pytest.mark.unit
def test_save_blocked_without_write(sample_pptx, monkeypatch):
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    with pytest.raises(NotAllowedError):
        save(sample_pptx, confirm=True)


@pytest.mark.unit
def test_delete_slide_blocked_without_write(sample_pptx, monkeypatch):
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    with pytest.raises(NotAllowedError):
        delete_slide(sample_pptx, 0, confirm=True)


@pytest.mark.unit
def test_set_speaker_notes_blocked_without_write(sample_pptx, monkeypatch):
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    with pytest.raises(NotAllowedError):
        set_speaker_notes(sample_pptx, 0, "notes")


# ── GROUP 12: confirm guards ──────────────────────────────────────────────

@pytest.mark.unit
def test_add_slide_blocked_without_confirm(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="confirm"):
        add_slide(sample_pptx, confirm=False)


@pytest.mark.unit
def test_delete_slide_blocked_without_confirm(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="confirm"):
        delete_slide(sample_pptx, 0, confirm=False)


@pytest.mark.unit
def test_save_blocked_without_confirm(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="confirm"):
        save(sample_pptx, confirm=False)


@pytest.mark.unit
def test_replace_slide_text_blocked_without_confirm(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="confirm"):
        replace_slide_text(sample_pptx, "Test", "New", confirm=False)


@pytest.mark.unit
def test_set_speaker_notes_requires_confirm(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="confirm"):
        set_speaker_notes(sample_pptx, 0, "notes", confirm=False)


@pytest.mark.unit
def test_edit_text_placeholder_requires_confirm(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="confirm"):
        edit_text_placeholder(sample_pptx, 0, 0, "new text", confirm=False)


# ── GROUP 13: write operations (PPT_ENABLE_WRITE=true, confirm=True) ────

@pytest.mark.unit
def test_set_speaker_notes_write(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = set_speaker_notes(sample_pptx, 0, "My speaker notes", confirm=True)
    assert result["notes_text"] == "My speaker notes"
    assert result["slide_index"] == 0


@pytest.mark.unit
def test_edit_text_placeholder_write(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = edit_text_placeholder(sample_pptx, 0, 0, "New Title", confirm=True)
    assert "slide_index" in result
    assert result["text"] == "New Title"


@pytest.mark.unit
def test_add_slide_write(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = add_slide(sample_pptx, layout_index=1, confirm=True)
    assert "slide_index" in result
    assert result["slide_index"] == 1  # second slide (0-indexed)
    assert "title" in result


@pytest.mark.unit
def test_replace_slide_text_write(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = replace_slide_text(sample_pptx, "Test", "Modified", confirm=True)
    assert "replacements_made" in result
    assert result["replacements_made"] >= 1


# ── GROUP 14a: UX-01 behavioural guarantee ────────────────────────────────

@pytest.mark.unit
def test_read_speaker_notes_includes_empty_slides(sample_pptx):
    """UX-01: every slide is always included, even those with no notes."""
    result = read_speaker_notes(sample_pptx, slide_index=None)
    assert len(result) == 1  # fixture has exactly 1 slide
    notes_text = result[0]["notes_text"]
    assert notes_text is None or notes_text == ""


@pytest.mark.unit
def test_reorder_slides_write(sample_pptx, monkeypatch, tmp_path):
    """Add a second slide so we can reorder."""
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    # Re-create with 2 slides
    path = tmp_path / "two_slides.pptx"
    p = _Prs()
    for _ in range(2):
        p.slides.add_slide(p.slide_layouts[1])
    p.save(str(path))
    result = reorder_slides(str(path), [1, 0], confirm=True)
    assert result["new_order"] == [1, 0]


@pytest.mark.unit
def test_delete_slide_write(sample_pptx, monkeypatch, tmp_path):
    """Add a second slide first, then delete index 1."""
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    path = tmp_path / "two_slides_del.pptx"
    p = _Prs()
    for _ in range(2):
        p.slides.add_slide(p.slide_layouts[1])
    p.save(str(path))
    result = delete_slide(str(path), 1, confirm=True)
    assert result["deleted_index"] == 1


@pytest.mark.unit
def test_save_write(sample_pptx, monkeypatch):
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = save(sample_pptx, confirm=True)
    assert "saved" in result
    assert result["saved"] is True
    assert result["path"].endswith(".pptx")


# ── GROUP 15: insert_image() ──────────────────────────────────────────────

@pytest.mark.unit
def test_insert_image_blocked_without_write(tmp_path, monkeypatch):
    """insert_image must raise NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    pptx_path = tmp_path / "deck.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(pptx_path))
    # image_path validation can be any image extension path — will fail write gate first
    img_path = tmp_path / "img.png"
    img_path.write_bytes(b"fake")
    with pytest.raises(NotAllowedError):
        from pptmcp.presentation_pptx import insert_image
        insert_image(str(pptx_path), 0, str(img_path))


@pytest.mark.unit
def test_insert_image_rejects_non_image_extension(tmp_path, monkeypatch):
    """insert_image must reject files with unsupported image extensions."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx_path = tmp_path / "deck.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(pptx_path))
    bad_img = tmp_path / "img.exe"
    bad_img.write_bytes(b"fake")
    with pytest.raises(ValidationError, match="Unsupported image extension"):
        from pptmcp.presentation_pptx import insert_image
        insert_image(str(pptx_path), 0, str(bad_img), confirm=True)


# ── GROUP 16: insert_image confirm gate ───────────────────────────────────

@pytest.mark.unit
def test_insert_image_requires_confirm(tmp_path, monkeypatch):
    """insert_image must raise ValidationError when confirm=False even if all other checks pass."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx_path = tmp_path / "deck.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(pptx_path))
    img_path = tmp_path / "img.png"
    img_path.write_bytes(b"fake-png")
    with pytest.raises(ValidationError, match="confirm"):
        insert_image(str(pptx_path), 0, str(img_path), 0, 0, 100, 100, confirm=False)


@pytest.mark.unit
def test_insert_image_blocked_path(tmp_path, monkeypatch):
    """insert_image must raise ValidationError when image_path is outside the allowlist."""
    allowed_dir = tmp_path / "allowed"
    allowed_dir.mkdir()
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(allowed_dir))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    pptx_path = allowed_dir / "deck.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(pptx_path))
    outside_img = tmp_path / "img.png"
    outside_img.write_bytes(b"fake-png")
    with pytest.raises((ValidationError, NotAllowedError)):
        insert_image(str(pptx_path), 0, str(outside_img), 0, 0, 100, 100, confirm=True)


# ── GROUP 17: GX UX-gap fixes ─────────────────────────────────────────────

@pytest.mark.unit
def test_edit_text_placeholder_response_key(sample_pptx, monkeypatch):
    """GX-01: response must use key 'text', not 'new_text'."""
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    result = edit_text_placeholder(sample_pptx, 0, 0, "GX01 Title", confirm=True)
    assert "text" in result
    assert "new_text" not in result
    assert result["text"] == "GX01 Title"


@pytest.mark.unit
def test_list_slides_no_title_is_null(tmp_path, monkeypatch):
    """GX-02: list_slides must return None (null) for slides with no/empty title."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "notitle.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[1])
    p.save(str(path))
    result = list_slides(str(path))
    assert result[0]["title"] is None, f"Expected None, got {result[0]['title']!r}"


@pytest.mark.unit
def test_read_slide_shapes_include_text(tmp_path, monkeypatch):
    """GX-03: read_slide shape dicts must include a 'text' field."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "withtext.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[0])  # layout 0 has title + content
    # Set the title placeholder text
    slide.shapes.title.text = "GX03 Text"
    p.save(str(path))
    result = read_slide(str(path), 0)
    shape_dicts = result["shapes"]
    # All shapes must have a 'text' key
    for s in shape_dicts:
        assert "text" in s, f"Shape {s.get('name')} missing 'text' field"
    # The title placeholder must have a non-empty text value
    title_shapes = [s for s in shape_dicts if s.get("placeholder_idx") == 0]
    assert title_shapes, "No title placeholder (idx=0) found"
    assert title_shapes[0]["text"] == "GX03 Text"


# ── W5 CHANGE 1 MODIFICATION: moved from GROUP 20 ────────────────────────

@pytest.mark.unit
def test_edit_text_placeholder_invalid_index(sample_pptx, monkeypatch):
    """Test that edit_text_placeholder raises ValidationError for invalid placeholder index."""
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    with pytest.raises(ValidationError, match="placeholder_idx .* not found"):
        edit_text_placeholder(sample_pptx, 0, 999, "text", confirm=True)
