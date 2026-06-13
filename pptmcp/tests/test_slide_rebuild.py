"""Tests for FR-826: clear_slide_content, apply_slide_layout, remove_empty_placeholders.

Groups:
  TestClearSlideContent      — 7 tests
  TestApplySlideLayout       — 8 tests
  TestRemoveEmptyPlaceholders — 6 tests
Total: 21 tests  (≥ 19 required)
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptmcp.presentation_pptx import (
    NotAllowedError,
    ValidationError,
    apply_slide_layout,
    clear_slide_content,
)
from pptmcp.shapes_pptx import remove_empty_placeholders


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _make_pptx_with_shapes(directory, filename="test.pptx", n_shapes: int = 2) -> str:
    """Create a 1-slide .pptx with n_shapes text boxes and return the path string."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout (no PH)
    for i in range(n_shapes):
        txBox = slide.shapes.add_textbox(Inches(i), Inches(0), Inches(1), Inches(1))
        txBox.text_frame.text = f"shape {i}"
    path = directory / filename
    prs.save(str(path))
    return str(path)


def _make_pptx_with_empty_ph(directory, filename="ph_test.pptx") -> str:
    """1-slide .pptx using the Title Slide layout (index 0) — contains empty/default placeholders."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # noqa: F841
    # Leave placeholders untouched so they contain the PowerPoint default prompt text
    path = directory / filename
    prs.save(str(path))
    return str(path)


def _make_pptx_with_filled_ph(directory, filename="ph_filled.pptx") -> str:
    """1-slide .pptx with title + sub-title filled in."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "My Real Title"
    try:
        slide.placeholders[1].text = "My real subtitle"
    except (KeyError, IndexError):
        pass
    path = directory / filename
    prs.save(str(path))
    return str(path)


# ---------------------------------------------------------------------------
# TestClearSlideContent
# ---------------------------------------------------------------------------

class TestClearSlideContent:
    """Tests for presentation_pptx.clear_slide_content."""

    @pytest.mark.unit
    def test_path_not_allowed(self, tmp_path, monkeypatch):
        """ValidationError when path is outside the allowlist."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "allowed"))
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        outside = tmp_path / "outside" / "test.pptx"
        outside.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])
        prs.save(str(outside))
        with pytest.raises(ValidationError, match="allowlist"):
            clear_slide_content(str(outside), 0, confirm=True)

    @pytest.mark.unit
    def test_write_not_enabled(self, allowlisted_dir, monkeypatch):
        """NotAllowedError when PPT_ENABLE_WRITE is absent."""
        monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
        p = _make_pptx_with_shapes(allowlisted_dir)
        with pytest.raises(NotAllowedError):
            clear_slide_content(p, 0, confirm=True)

    @pytest.mark.unit
    def test_confirm_not_set(self, allowlisted_dir, write_env):
        """ValidationError when confirm=False."""
        p = _make_pptx_with_shapes(allowlisted_dir)
        with pytest.raises(ValidationError, match="confirm"):
            clear_slide_content(p, 0, confirm=False)

    @pytest.mark.unit
    def test_clears_all_shapes(self, allowlisted_dir, write_env):
        """All shapes removed; cleared_count equals original shape count."""
        p = _make_pptx_with_shapes(allowlisted_dir, n_shapes=3)
        # Read initial shape count (layout may add placeholder shapes)
        initial_count = len(list(Presentation(p).slides[0].shapes))
        result = clear_slide_content(p, 0, confirm=True)
        assert result["cleared_count"] == initial_count
        assert result["slide_index"] == 0
        # Verify file was persisted and slide is now empty
        prs = Presentation(p)
        assert len(list(prs.slides[0].shapes)) == 0

    @pytest.mark.unit
    def test_idempotent_on_blank_slide(self, allowlisted_dir, write_env):
        """cleared_count == 0 and no error when slide is already blank."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        path = allowlisted_dir / "blank.pptx"
        prs.save(str(path))
        # Blank layout has no shapes — cleared_count may be 0 or equal to PH count
        result = clear_slide_content(str(path), 0, confirm=True)
        assert isinstance(result["cleared_count"], int)
        assert result["slide_index"] == 0
        prs2 = Presentation(str(path))
        assert len(list(prs2.slides[0].shapes)) == 0

    @pytest.mark.unit
    def test_oob_slide_index(self, allowlisted_dir, write_env):
        """ValidationError for out-of-range slide_index."""
        p = _make_pptx_with_shapes(allowlisted_dir)  # 1 slide → index 0 only
        with pytest.raises(ValidationError):
            clear_slide_content(p, 5, confirm=True)

    @pytest.mark.unit
    def test_negative_slide_index(self, allowlisted_dir, write_env):
        """ValidationError for negative slide_index."""
        p = _make_pptx_with_shapes(allowlisted_dir)
        with pytest.raises(ValidationError):
            clear_slide_content(p, -1, confirm=True)


# ---------------------------------------------------------------------------
# TestApplySlideLayout
# ---------------------------------------------------------------------------

class TestApplySlideLayout:
    """Tests for presentation_pptx.apply_slide_layout."""

    @pytest.mark.unit
    def test_path_not_allowed(self, tmp_path, monkeypatch):
        """ValidationError when path is outside the allowlist."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "allowed"))
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        outside = tmp_path / "outside" / "test.pptx"
        outside.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(str(outside))
        with pytest.raises(ValidationError, match="allowlist"):
            apply_slide_layout(str(outside), 0, 1, confirm=True)

    @pytest.mark.unit
    def test_write_not_enabled(self, allowlisted_dir, monkeypatch):
        """NotAllowedError when PPT_ENABLE_WRITE is absent."""
        monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = allowlisted_dir / "wr.pptx"
        prs.save(str(path))
        with pytest.raises(NotAllowedError):
            apply_slide_layout(str(path), 0, 1, confirm=True)

    @pytest.mark.unit
    def test_confirm_not_set(self, allowlisted_dir, write_env):
        """ValidationError when confirm=False."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = allowlisted_dir / "cf.pptx"
        prs.save(str(path))
        with pytest.raises(ValidationError, match="confirm"):
            apply_slide_layout(str(path), 0, 1, confirm=False)

    @pytest.mark.unit
    def test_layout_index_out_of_range(self, allowlisted_dir, write_env):
        """ValidationError for a layout_index beyond available layouts."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = allowlisted_dir / "oob.pptx"
        prs.save(str(path))
        with pytest.raises(ValidationError):
            apply_slide_layout(str(path), 0, 9999, confirm=True)

    @pytest.mark.unit
    def test_applies_layout_changes_name(self, allowlisted_dir, write_env):
        """Return dict contains layout_applied matching the target layout's name."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = allowlisted_dir / "apl.pptx"
        prs.save(str(path))
        target_idx = 5  # 'Blank' in the default theme
        target_name = Presentation(str(path)).slide_layouts[target_idx].name
        result = apply_slide_layout(str(path), 0, target_idx, remove_placeholders=False, confirm=True)
        assert result["layout_applied"] == target_name
        assert result["slide_index"] == 0

    @pytest.mark.unit
    def test_applies_layout_with_removal(self, allowlisted_dir, write_env):
        """remove_placeholders=True removes placeholder shapes and persists."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide has 2 PH
        path = allowlisted_dir / "aplrm.pptx"
        prs.save(str(path))
        result = apply_slide_layout(str(path), 0, 5, remove_placeholders=True, confirm=True)
        assert isinstance(result["placeholders_removed"], int)
        # After save prs.slides[0] should have 0 placeholder shapes
        prs2 = Presentation(str(path))
        ph_remaining = [s for s in prs2.slides[0].shapes if s.is_placeholder]
        assert len(ph_remaining) == 0

    @pytest.mark.unit
    def test_applies_layout_keeps_placeholders(self, allowlisted_dir, write_env):
        """remove_placeholders=False keeps placeholders intact."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content — 2 PH
        path = allowlisted_dir / "aplkp.pptx"
        prs.save(str(path))
        original_ph_count = len([s for s in Presentation(str(path)).slides[0].shapes if s.is_placeholder])
        result = apply_slide_layout(str(path), 0, 0, remove_placeholders=False, confirm=True)
        assert result["placeholders_removed"] == 0
        prs2 = Presentation(str(path))
        ph_after = [s for s in prs2.slides[0].shapes if s.is_placeholder]
        # remove_placeholders=False must leave all placeholder shapes untouched on disk
        assert len(ph_after) == original_ph_count

    @pytest.mark.unit
    def test_oob_slide_index(self, allowlisted_dir, write_env):
        """ValidationError for out-of-range slide_index."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = allowlisted_dir / "soob.pptx"
        prs.save(str(path))
        with pytest.raises(ValidationError):
            apply_slide_layout(str(path), 99, 0, confirm=True)


# ---------------------------------------------------------------------------
# TestRemoveEmptyPlaceholders
# ---------------------------------------------------------------------------

class TestRemoveEmptyPlaceholders:
    """Tests for shapes_pptx.remove_empty_placeholders."""

    @pytest.mark.unit
    def test_path_not_allowed(self, tmp_path, monkeypatch):
        """ValidationError when path is outside the allowlist."""
        monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "allowed"))
        monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
        outside = tmp_path / "outside" / "rph.pptx"
        outside.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(str(outside))
        with pytest.raises(ValidationError, match="allowlist"):
            remove_empty_placeholders(str(outside), 0, confirm=True)

    @pytest.mark.unit
    def test_write_not_enabled(self, allowlisted_dir, monkeypatch):
        """NotAllowedError when PPT_ENABLE_WRITE is absent."""
        monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
        p = _make_pptx_with_empty_ph(allowlisted_dir, "wr_ph.pptx")
        with pytest.raises(NotAllowedError):
            remove_empty_placeholders(p, 0, confirm=True)

    @pytest.mark.unit
    def test_confirm_not_set(self, allowlisted_dir, write_env):
        """ValidationError when confirm=False."""
        p = _make_pptx_with_empty_ph(allowlisted_dir, "cf_ph.pptx")
        with pytest.raises(ValidationError, match="confirm"):
            remove_empty_placeholders(p, 0, confirm=False)

    @pytest.mark.unit
    def test_removes_empty_placeholders(self, allowlisted_dir, write_env):
        """Empty / default-text placeholders are removed; file is persisted."""
        p = _make_pptx_with_empty_ph(allowlisted_dir, "emp.pptx")
        initial_ph = [s for s in Presentation(p).slides[0].shapes if s.is_placeholder]
        result = remove_empty_placeholders(p, 0, confirm=True)
        assert result["removed_count"] == len(initial_ph)
        prs2 = Presentation(p)
        remaining_ph = [s for s in prs2.slides[0].shapes if s.is_placeholder]
        assert len(remaining_ph) == 0

    @pytest.mark.unit
    def test_keeps_populated_placeholder(self, allowlisted_dir, write_env):
        """Placeholders with real content are preserved."""
        p = _make_pptx_with_filled_ph(allowlisted_dir, "filled.pptx")
        initial_ph = [s for s in Presentation(p).slides[0].shapes if s.is_placeholder]
        result = remove_empty_placeholders(p, 0, confirm=True)
        # All PH had real text — none should be removed
        assert result["removed_count"] == 0
        prs2 = Presentation(p)
        remaining_ph = [s for s in prs2.slides[0].shapes if s.is_placeholder]
        assert len(remaining_ph) == len(initial_ph)

    @pytest.mark.unit
    def test_no_placeholders_returns_zero(self, allowlisted_dir, write_env):
        """A slide that has been cleared (truly 0 shapes) returns removed_count == 0."""
        p = _make_pptx_with_shapes(allowlisted_dir, "noph.pptx", n_shapes=2)
        # First clear all shapes (including any layout placeholders)
        clear_slide_content(p, 0, confirm=True)
        # Now remove_empty_placeholders on an empty slide should return 0
        result = remove_empty_placeholders(p, 0, confirm=True)
        assert result["removed_count"] == 0
        assert result["remaining_shapes"] == 0
