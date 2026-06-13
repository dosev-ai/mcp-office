"""Unit tests for set_text_format paragraph_range and run_index params,
and set_paragraph_format paragraph_range param.

No COM required — CI-safe. Tests use real python-pptx files.

Coverage:
  US-PPT-SF-01: paragraph_range limits which paragraphs are formatted
  US-PPT-SF-02: run_index limits which run within each paragraph is formatted
  US-PPT-SF-03: both params omitted == backward-compat (all paragraphs, all runs)
  US-PPT-SF-04: paragraph_range on set_paragraph_format processes only the range
  US-PPT-SF-05: paragraph_range and run_index combined — only run 0 in paras 1-2
  US-PPT-SF-06: paragraph_range with out-of-range indices skips silently (no error)
  US-PPT-SF-07: invalid paragraph_range (wrong length) raises ValidationError
"""
from __future__ import annotations

import os
import sys

_SRC = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "src")
if os.path.isdir(_SRC) and _SRC not in sys.path:
    sys.path.insert(0, os.path.realpath(_SRC))

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_pptx_with_multirun_paragraphs(path, n_paragraphs: int = 3, n_runs: int = 2) -> tuple[str, int]:
    """Create a 1-slide .pptx with a text box containing n_paragraphs paragraphs,
    each with n_runs distinct runs. Returns (path_str, shape_id)."""
    from pptx.oxml.ns import qn
    from lxml import etree

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(4))
    tf = tb.text_frame
    tf.word_wrap = True

    # Clear any auto-created paragraph and rebuild
    tf.clear()

    for para_idx in range(n_paragraphs):
        if para_idx == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        for run_idx in range(n_runs):
            run = para.add_run()
            run.text = f"P{para_idx}R{run_idx}"

    prs.save(str(path))
    return str(path), tb.shape_id


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def pptx_multirun(tmp_path, monkeypatch):
    """3-paragraph, 2-run-each textbox in an allowlisted dir."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path, shape_id = _make_pptx_with_multirun_paragraphs(
        tmp_path / "multirun.pptx", n_paragraphs=3, n_runs=2
    )
    return path, shape_id


# ---------------------------------------------------------------------------
# US-PPT-SF-01: paragraph_range filters paragraphs
# ---------------------------------------------------------------------------


class TestSetFormatParagraphRange:
    def test_paragraph_range_applies_only_to_specified_paragraphs(self, pptx_multirun):
        """paragraph_range=[1,2] formats only paragraphs 1 and 2, not paragraph 0."""
        from pptmcp._shapes_text import set_text_format
        from pptmcp.presentation_pptx import _load_prs, _check_path

        path, shape_id = pptx_multirun
        result = set_text_format(
            path, 0, shape_id,
            bold=True,
            paragraph_range=[1, 2],
            confirm=True,
        )
        # Runs in paras 1 and 2 were formatted (2 paras * 2 runs = 4)
        assert result["runs_updated"] == 4

        # Verify via reload: para 0 runs must NOT be bold, paras 1-2 must be bold
        prs = _load_prs(_check_path(path))
        shape = next(s for s in prs.slides[0].shapes if s.shape_id == shape_id)
        paragraphs = shape.text_frame.paragraphs
        for run in paragraphs[0].runs:
            assert run.font.bold is not True, "Para 0 run should not have been made bold"
        for para_idx in [1, 2]:
            for run in paragraphs[para_idx].runs:
                assert run.font.bold is True, f"Para {para_idx} run should be bold"

    def test_paragraph_range_single_paragraph(self, pptx_multirun):
        """paragraph_range=[0,0] formats only paragraph 0 (single paragraph range)."""
        from pptmcp._shapes_text import set_text_format
        from pptmcp.presentation_pptx import _load_prs, _check_path

        path, shape_id = pptx_multirun
        result = set_text_format(
            path, 0, shape_id,
            italic=True,
            paragraph_range=[0, 0],
            confirm=True,
        )
        assert result["runs_updated"] == 2  # 1 para * 2 runs

        prs = _load_prs(_check_path(path))
        shape = next(s for s in prs.slides[0].shapes if s.shape_id == shape_id)
        paragraphs = shape.text_frame.paragraphs
        for run in paragraphs[0].runs:
            assert run.font.italic is True
        for para_idx in [1, 2]:
            for run in paragraphs[para_idx].runs:
                assert run.font.italic is not True


# ---------------------------------------------------------------------------
# US-PPT-SF-02: run_index limits the run formatted
# ---------------------------------------------------------------------------


class TestSetFormatRunIndex:
    def test_run_index_applies_only_to_specified_run(self, pptx_multirun):
        """run_index=0 formats only run 0 in every paragraph."""
        from pptmcp._shapes_text import set_text_format
        from pptmcp.presentation_pptx import _load_prs, _check_path

        path, shape_id = pptx_multirun
        result = set_text_format(
            path, 0, shape_id,
            bold=True,
            run_index=0,
            confirm=True,
        )
        # 3 paragraphs * 1 run each = 3
        assert result["runs_updated"] == 3

        prs = _load_prs(_check_path(path))
        shape = next(s for s in prs.slides[0].shapes if s.shape_id == shape_id)
        for para in shape.text_frame.paragraphs:
            runs = para.runs
            if len(runs) >= 1:
                assert runs[0].font.bold is True, "Run 0 should be bold"
            if len(runs) >= 2:
                assert runs[1].font.bold is not True, "Run 1 should NOT be bold"

    def test_run_index_out_of_bounds_skips_silently(self, pptx_multirun):
        """run_index=99 skips all runs silently — runs_updated == 0."""
        from pptmcp._shapes_text import set_text_format

        path, shape_id = pptx_multirun
        result = set_text_format(
            path, 0, shape_id,
            bold=True,
            run_index=99,
            confirm=True,
        )
        assert result["runs_updated"] == 0


# ---------------------------------------------------------------------------
# US-PPT-SF-03: backward compat — no params == format all
# ---------------------------------------------------------------------------


class TestSetFormatNoParams:
    def test_no_params_formats_all_paragraphs_and_runs(self, pptx_multirun):
        """Omitting paragraph_range and run_index formats all runs in all paragraphs."""
        from pptmcp._shapes_text import set_text_format

        path, shape_id = pptx_multirun
        result = set_text_format(
            path, 0, shape_id,
            bold=True,
            confirm=True,
        )
        # 3 paragraphs * 2 runs = 6
        assert result["runs_updated"] == 6


# ---------------------------------------------------------------------------
# US-PPT-SF-04: set_paragraph_format paragraph_range
# ---------------------------------------------------------------------------


class TestSetParagraphFormatRange:
    def test_paragraph_range_processes_multiple_paragraphs(self, pptx_multirun):
        """paragraph_range=[1,2] applies alignment to paragraphs 1 and 2."""
        from pptmcp._shapes_text import set_paragraph_format
        from pptmcp.presentation_pptx import _load_prs, _check_path
        from pptx.enum.text import PP_ALIGN

        path, shape_id = pptx_multirun
        result = set_paragraph_format(
            path, 0, shape_id,
            alignment="CENTER",
            paragraph_range=[1, 2],
            confirm=True,
        )
        # Result should still report paragraph_index=0 (default), range is the effective filter
        assert result["shape_id"] == shape_id

        prs = _load_prs(_check_path(path))
        shape = next(s for s in prs.slides[0].shapes if s.shape_id == shape_id)
        paragraphs = shape.text_frame.paragraphs
        assert paragraphs[0].alignment != PP_ALIGN.CENTER or paragraphs[0].alignment is None
        assert paragraphs[1].alignment == PP_ALIGN.CENTER
        assert paragraphs[2].alignment == PP_ALIGN.CENTER

    def test_paragraph_range_none_falls_back_to_paragraph_index(self, pptx_multirun):
        """When paragraph_range=None, only paragraph_index paragraph is formatted."""
        from pptmcp._shapes_text import set_paragraph_format
        from pptmcp.presentation_pptx import _load_prs, _check_path
        from pptx.enum.text import PP_ALIGN

        path, shape_id = pptx_multirun
        set_paragraph_format(
            path, 0, shape_id,
            paragraph_index=2,
            alignment="RIGHT",
            paragraph_range=None,
            confirm=True,
        )

        prs = _load_prs(_check_path(path))
        shape = next(s for s in prs.slides[0].shapes if s.shape_id == shape_id)
        paragraphs = shape.text_frame.paragraphs
        assert paragraphs[2].alignment == PP_ALIGN.RIGHT


# ---------------------------------------------------------------------------
# US-PPT-SF-05: Combined paragraph_range + run_index
# ---------------------------------------------------------------------------


class TestSetFormatCombined:
    def test_paragraph_range_and_run_index_combined(self, pptx_multirun):
        """paragraph_range=[1,2] + run_index=0 formats only run 0 in paras 1 and 2."""
        from pptmcp._shapes_text import set_text_format
        from pptmcp.presentation_pptx import _load_prs, _check_path

        path, shape_id = pptx_multirun
        result = set_text_format(
            path, 0, shape_id,
            bold=True,
            paragraph_range=[1, 2],
            run_index=0,
            confirm=True,
        )
        # 2 paragraphs * 1 run each = 2
        assert result["runs_updated"] == 2

        prs = _load_prs(_check_path(path))
        shape = next(s for s in prs.slides[0].shapes if s.shape_id == shape_id)
        paragraphs = shape.text_frame.paragraphs
        # Para 0: untouched
        for run in paragraphs[0].runs:
            assert run.font.bold is not True
        # Paras 1-2: only run 0 bold
        for para_idx in [1, 2]:
            runs = paragraphs[para_idx].runs
            assert runs[0].font.bold is True
            assert runs[1].font.bold is not True


# ---------------------------------------------------------------------------
# US-PPT-SF-06: paragraph_range with out-of-range indices skips silently
# ---------------------------------------------------------------------------


class TestSetFormatOutOfRangeParagraphRange:
    def test_paragraph_range_partially_out_of_range_skips_bad_indices(self, pptx_multirun):
        """paragraph_range=[2,5] with only 3 paras (0-2) formats only para 2 silently."""
        from pptmcp._shapes_text import set_text_format

        path, shape_id = pptx_multirun
        result = set_text_format(
            path, 0, shape_id,
            bold=True,
            paragraph_range=[2, 5],
            confirm=True,
        )
        # Only para 2 (index 2) exists in range [2,5]; 2 runs
        assert result["runs_updated"] == 2

    def test_paragraph_range_entirely_out_of_range_formats_nothing(self, pptx_multirun):
        """paragraph_range=[10,20] formats nothing when no such paragraphs exist."""
        from pptmcp._shapes_text import set_text_format

        path, shape_id = pptx_multirun
        result = set_text_format(
            path, 0, shape_id,
            bold=True,
            paragraph_range=[10, 20],
            confirm=True,
        )
        assert result["runs_updated"] == 0


# ---------------------------------------------------------------------------
# US-PPT-SF-07: invalid paragraph_range raises ValidationError
# ---------------------------------------------------------------------------


class TestSetFormatInvalidParagraphRange:
    def test_paragraph_range_wrong_length_raises(self, pptx_multirun):
        """paragraph_range with != 2 elements raises ValidationError."""
        from pptmcp._shapes_text import set_text_format
        from pptmcp.presentation_pptx import ValidationError

        path, shape_id = pptx_multirun
        with pytest.raises(ValidationError, match="paragraph_range"):
            set_text_format(
                path, 0, shape_id,
                bold=True,
                paragraph_range=[1],
                confirm=True,
            )

    def test_set_paragraph_format_paragraph_range_wrong_length_raises(self, pptx_multirun):
        """set_paragraph_format with paragraph_range != 2 elements raises ValidationError."""
        from pptmcp._shapes_text import set_paragraph_format
        from pptmcp.presentation_pptx import ValidationError

        path, shape_id = pptx_multirun
        with pytest.raises(ValidationError, match="paragraph_range"):
            set_paragraph_format(
                path, 0, shape_id,
                alignment="CENTER",
                paragraph_range=[0, 1, 2],
                confirm=True,
            )
