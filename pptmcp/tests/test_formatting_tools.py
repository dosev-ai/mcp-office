"""Text and paragraph formatting tool tests — targets presentation_pptx.py only."""

import pytest
from pptx import Presentation as _Prs
from pptmcp.presentation_pptx import (
    ValidationError,
    NotAllowedError,
    set_text_format,
    set_paragraph_format,
    set_table_style,
)


# ── GROUP 26: set_text_format() ────────────────────────────────────────────

@pytest.mark.unit
def test_set_text_format_bold_success(tmp_path, monkeypatch):
    """set_text_format with bold=True returns runs_updated >= 1."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches
    path = tmp_path / "bold.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Bold me"
    p.save(str(path))
    result = set_text_format(str(path), 0, txBox.shape_id, bold=True, confirm=True)
    assert result["runs_updated"] >= 1
    assert result["slide_index"] == 0
    assert result["shape_id"] == txBox.shape_id


@pytest.mark.unit
def test_set_text_format_color_success(tmp_path, monkeypatch):
    """set_text_format with color_hex='FF0000' returns runs_updated >= 1."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches
    path = tmp_path / "color.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Red text"
    p.save(str(path))
    result = set_text_format(str(path), 0, txBox.shape_id, color_hex="FF0000", confirm=True)
    assert result["runs_updated"] >= 1


@pytest.mark.unit
def test_set_text_format_invalid_color_hex(tmp_path, monkeypatch):
    """set_text_format raises ValidationError for an invalid color_hex."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches
    path = tmp_path / "badcolor.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Bad"
    p.save(str(path))
    with pytest.raises(ValidationError, match="[Ii]nvalid hex"):
        set_text_format(str(path), 0, txBox.shape_id, color_hex="ZZZZZZ", confirm=True)


@pytest.mark.unit
def test_set_text_format_no_text_frame(tmp_path, monkeypatch):
    """set_text_format raises ValidationError for a shape with no text frame (table shape)."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches
    path = tmp_path / "notf.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    tbl = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    p.save(str(path))
    with pytest.raises(ValidationError, match="no text frame"):
        set_text_format(str(path), 0, tbl.shape_id, bold=True, confirm=True)


@pytest.mark.unit
def test_set_text_format_blocked_write(tmp_path, monkeypatch):
    """set_text_format raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(NotAllowedError):
        set_text_format(path, 0, 1, bold=True, confirm=True)


@pytest.mark.unit
def test_set_text_format_no_confirm(tmp_path, monkeypatch):
    """set_text_format raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        set_text_format(path, 0, 1, bold=True, confirm=False)


# ── GROUP 27: set_paragraph_format() ────────────────────────────────────────

@pytest.mark.unit
def test_set_paragraph_format_alignment_success(tmp_path, monkeypatch):
    """set_paragraph_format with alignment='CENTER' returns correct dict."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches
    path = tmp_path / "align.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Center me"
    p.save(str(path))
    result = set_paragraph_format(
        str(path), 0, txBox.shape_id, paragraph_index=0, alignment="CENTER", confirm=True
    )
    assert result["paragraph_index"] == 0
    assert result["slide_index"] == 0
    assert result["shape_id"] == txBox.shape_id


@pytest.mark.unit
def test_set_paragraph_format_invalid_alignment(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError for an unknown alignment value."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches
    path = tmp_path / "badalign.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Bad align"
    p.save(str(path))
    with pytest.raises(ValidationError, match="[Ii]nvalid alignment"):
        set_paragraph_format(str(path), 0, txBox.shape_id, alignment="BOGUS", confirm=True)


@pytest.mark.unit
def test_set_paragraph_format_invalid_para_index(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError for paragraph_index out of range."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches
    path = tmp_path / "badpara.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "One paragraph"
    p.save(str(path))
    with pytest.raises(ValidationError, match="paragraph_index"):
        set_paragraph_format(str(path), 0, txBox.shape_id, paragraph_index=99, confirm=True)


@pytest.mark.unit
def test_set_paragraph_format_blocked_write(tmp_path, monkeypatch):
    """set_paragraph_format raises NotAllowedError when PPT_ENABLE_WRITE is not set."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.delenv("PPT_ENABLE_WRITE", raising=False)
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(NotAllowedError):
        set_paragraph_format(path, 0, 1, alignment="CENTER", confirm=True)


@pytest.mark.unit
def test_set_paragraph_format_no_confirm(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError when confirm=False."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "noop.pptx")
    with pytest.raises(ValidationError, match="confirm"):
        set_paragraph_format(path, 0, 1, alignment="CENTER", confirm=False)


# ── Gap-fill: set_text_format ─────────────────────────────────────────────

@pytest.mark.unit
def test_set_text_format_path_not_allowed(tmp_path, monkeypatch):
    """set_text_format raises ValidationError for a path outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        set_text_format(path, 0, 1, bold=True, confirm=True)


@pytest.mark.unit
def test_set_text_format_invalid_slide_index(tmp_path, monkeypatch):
    """set_text_format raises ValidationError for an out-of-range slide index."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "stf_oob.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError):
        set_text_format(str(path), 99, 1, bold=True, confirm=True)


@pytest.mark.unit
def test_set_text_format_invalid_shape_id(tmp_path, monkeypatch):
    """set_text_format raises ValidationError when shape_id does not exist on slide."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "stf_badshape.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="not found"):
        set_text_format(str(path), 0, 99999, bold=True, confirm=True)


# ── Gap-fill: set_paragraph_format ───────────────────────────────────────

@pytest.mark.unit
def test_set_paragraph_format_path_not_allowed(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError for a path outside the allowlist."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path / "subdir"))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = str(tmp_path / "outside.pptx")
    with pytest.raises(ValidationError):
        set_paragraph_format(path, 0, 1, alignment="CENTER", confirm=True)


@pytest.mark.unit
def test_set_paragraph_format_invalid_slide_index(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError for an out-of-range slide index."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "spf_oob.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError):
        set_paragraph_format(str(path), 99, 1, alignment="CENTER", confirm=True)


# ── post-COM extras: set_paragraph_format ─────────────────────────────────

@pytest.mark.unit
def test_set_paragraph_format_invalid_shape_id(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError when shape_id does not exist on slide."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    path = tmp_path / "spf_badshape.pptx"
    p = _Prs()
    p.slides.add_slide(p.slide_layouts[5])
    p.save(str(path))
    with pytest.raises(ValidationError, match="not found"):
        set_paragraph_format(str(path), 0, 99999, alignment="CENTER", confirm=True)


@pytest.mark.unit
def test_set_paragraph_format_no_text_frame(tmp_path, monkeypatch):
    """set_paragraph_format raises ValidationError for a shape with no text frame."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches as _In
    path = tmp_path / "spf_notf.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    tbl = slide.shapes.add_table(2, 2, _In(1), _In(1), _In(4), _In(2))
    p.save(str(path))
    with pytest.raises(ValidationError, match="no text frame"):
        set_paragraph_format(str(path), 0, tbl.shape_id, alignment="CENTER", confirm=True)


@pytest.mark.unit
def test_set_paragraph_format_spacing_success(tmp_path, monkeypatch):
    """set_paragraph_format with line_spacing returns correct dict."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches as _In
    path = tmp_path / "spf_spacing.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    txBox = slide.shapes.add_textbox(_In(1), _In(1), _In(4), _In(1))
    txBox.text_frame.text = "Spaced paragraph"
    p.save(str(path))
    result = set_paragraph_format(
        str(path), 0, txBox.shape_id, paragraph_index=0, line_spacing=1.5, confirm=True
    )
    assert result["paragraph_index"] == 0
    assert result["shape_id"] == txBox.shape_id


# ── Table styling ─────────────────────────────────────────────────────────

@pytest.mark.unit
def test_set_table_style_font_header_and_row_height(tmp_path, monkeypatch):
    """set_table_style applies compact body/header fonts and row heights."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches as _In
    path = tmp_path / "table_style.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    tbl = slide.shapes.add_table(2, 2, _In(1), _In(1), _In(4), _In(2))
    tbl.table.cell(0, 0).text = "Header"
    tbl.table.cell(1, 0).text = "Body"
    p.save(str(path))

    result = set_table_style(
        str(path),
        0,
        tbl.shape_id,
        font_size_pt=8,
        header_font_size_pt=9,
        header_bold=True,
        row_height_pt=18,
        confirm=True,
    )

    assert result["cells_updated"] == 4
    assert result["rows_resized"] == 2
    from pptmcp.presentation_pptx import save
    save(str(path), confirm=True)
    loaded = _Prs(str(path))
    loaded_table = next(s for s in loaded.slides[0].shapes if s.shape_id == tbl.shape_id).table
    header_run = loaded_table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    body_run = loaded_table.cell(1, 0).text_frame.paragraphs[0].runs[0]
    assert round(header_run.font.size.pt) == 9
    assert header_run.font.bold is True
    assert round(body_run.font.size.pt) == 8
    assert round(loaded_table.rows[0].height.pt) == 18


@pytest.mark.unit
def test_set_table_style_rejects_non_table_shape(tmp_path, monkeypatch):
    """set_table_style raises ValidationError when shape_id is not a table."""
    monkeypatch.setenv("PPT_ALLOWLIST_ROOTS", str(tmp_path))
    monkeypatch.setenv("PPT_ENABLE_WRITE", "true")
    from pptx.util import Inches as _In
    path = tmp_path / "not_table_style.pptx"
    p = _Prs()
    slide = p.slides.add_slide(p.slide_layouts[5])
    tx_box = slide.shapes.add_textbox(_In(1), _In(1), _In(4), _In(1))
    tx_box.text_frame.text = "Not a table"
    p.save(str(path))

    with pytest.raises(ValidationError, match="not a table"):
        set_table_style(str(path), 0, tx_box.shape_id, font_size_pt=8, confirm=True)
